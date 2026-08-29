import os
from fastapi.testclient import TestClient
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from apps.admin_backend.main import app
from apps.admin_backend.generators.deliverables import deliverable_generator, OUTPUT_DIR

def test_deliverables_and_download_live():
    print("=" * 70)
    print("  MRPL SOVEREIGN WORKBENCH - DELIVERABLES & DOWNLOAD LIVE TEST")
    print("=" * 70)

    client = TestClient(app)

    # 1. Generate & Validate Real Word Document (.docx)
    print("\n--- [1] GENERATE & RE-OPEN MICROSOFT WORD (.DOCX) DELIVERABLE ---")
    docx_path = deliverable_generator.generate_approval_note_docx()
    print(f"Generated DOCX Path: {docx_path} ({os.path.getsize(docx_path)} bytes)")
    assert os.path.exists(docx_path)
    
    # Reopen with python-docx to verify document structure
    reopened_doc = Document(docx_path)
    headings = [p.text for p in reopened_doc.paragraphs if p.style.name.startswith('Heading')]
    print(f"  • Verified Heading Count: {len(headings)}")
    print(f"  • Table Count in Docx: {len(reopened_doc.tables)}")
    assert len(reopened_doc.tables) >= 3
    assert any("Executive Anomaly Summary" in h for h in headings)

    # 2. Generate & Validate Real Excel Workbooks (.xlsx)
    print("\n--- [2] GENERATE & RE-OPEN EXCEL WORKBOOKS (.XLSX) ---")
    xlsx_calc_path = deliverable_generator.generate_hydraulic_register_xlsx()
    print(f"Generated Hydraulic XLSX Path: {xlsx_calc_path} ({os.path.getsize(xlsx_calc_path)} bytes)")
    assert os.path.exists(xlsx_calc_path)

    wb1 = load_workbook(xlsx_calc_path)
    ws1 = wb1.active
    print(f"  • Sheet Title: {ws1.title}")
    print(f"  • Efficiency Value in Cell C14: {ws1['C14'].value}%")
    assert ws1['C14'].value == 81.39

    xlsx_asset_path = deliverable_generator.generate_asset_register_xlsx()
    print(f"Generated Asset Register XLSX Path: {xlsx_asset_path} ({os.path.getsize(xlsx_asset_path)} bytes)")
    wb2 = load_workbook(xlsx_asset_path)
    ws2 = wb2.active
    print(f"  • Asset Sheet Rows: {ws2.max_row} (9 ISA 5.1 Tags + Header)")
    assert ws2.max_row == 10

    # 3. Generate & Validate Real PowerPoint Deck (.pptx)
    print("\n--- [3] GENERATE & RE-OPEN POWERPOINT DECK (.PPTX) ---")
    pptx_path = deliverable_generator.generate_turnaround_briefing_pptx()
    print(f"Generated PPTX Path: {pptx_path} ({os.path.getsize(pptx_path)} bytes)")
    assert os.path.exists(pptx_path)

    prs = Presentation(pptx_path)
    print(f"  • Slide Count: {len(prs.slides)}")
    print(f"  • Slide 1 Title: '{prs.slides[0].shapes.title.text}'")
    assert len(prs.slides) == 3

    # 4. Test GET /api/files/download/{filename} Valid Download
    print("\n--- [4] TEST SECURE FILE DOWNLOAD ENDPOINT ---")
    res_dl = client.get("/api/files/download/MRPL_Furnace_Inspection_Approval_Note.docx")
    print(f"GET /api/files/download/MRPL_Furnace_Inspection_Approval_Note.docx -> Status: {res_dl.status_code}")
    print(f"  • Content-Type: {res_dl.headers.get('content-type')}")
    print(f"  • Content-Length: {len(res_dl.content)} bytes")
    assert res_dl.status_code == 200
    assert len(res_dl.content) > 1000

    # 5. Security Test: Path Traversal Attack Rejection (../../etc/passwd)
    print("\n--- [5] TEST SECURITY: PATH TRAVERSAL ATTACK REJECTION ---")
    traversal_filenames = [
        "../../etc/passwd",
        "..\\..\\Windows\\System32\\cmd.exe",
        "../../apps/admin_backend/main.py",
        "%2e%2e%2f%2e%2e%2fetc%2fpasswd"
    ]
    for attack in traversal_filenames:
        res_attack = client.get(f"/api/files/download/{attack}")
        print(f"  • Attack Payload '{attack}' -> Status: {res_attack.status_code} (Blocked)")
        assert res_attack.status_code in (403, 404), f"Traversal '{attack}' was not blocked!"

    print("\n" + "=" * 70)
    print("  ALL FILE GENERATION & DOWNLOAD TESTS PASSED (100% VERIFIED)")
    print("=" * 70)

if __name__ == "__main__":
    test_deliverables_and_download_live()
