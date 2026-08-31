from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches
import os
import time

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))


class ExtendedGenerators_FinanceEsg:

    def __init__(self):
        os.makedirs(OUTPUT_DIR, exist_ok=True)

    def _get_output_path(self, filename):
        return os.path.join(OUTPUT_DIR, filename)

    # =========================================================================
    # FINANCE METHODS
    # =========================================================================

    def generate_annual_financial_stmt_docx(self, filename="Annual_Financial_Statements.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Annual Financial Statements", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Standalone Financial Statements", level=1)
        doc.add_paragraph("Balance Sheet as at 31st March 2025", style="List Bullet")
        doc.add_paragraph("Statement of Profit and Loss for the year ended 31st March 2025", style="List Bullet")
        doc.add_paragraph("Cash Flow Statement for the year ended 31st March 2025", style="List Bullet")
        doc.add_paragraph("Statement of Changes in Equity", style="List Bullet")
        doc.add_paragraph("Notes to Accounts (Schedules 1-45)", style="List Bullet")
        doc.add_heading("Consolidated Financial Statements", level=1)
        doc.add_paragraph("Consolidated Balance Sheet", style="List Bullet")
        doc.add_paragraph("Consolidated Statement of Profit and Loss", style="List Bullet")
        doc.add_paragraph("Consolidated Cash Flow Statement", style="List Bullet")
        doc.add_paragraph("Consolidated Statement of Changes in Equity", style="List Bullet")
        doc.add_paragraph("Consolidated Notes to Accounts", style="List Bullet")
        doc.add_heading("Key Financial Highlights", level=2)
        table = doc.add_table(rows=7, cols=3, style="Light Grid Accent 1")
        headers = ["Particulars", "Standalone (Rs Cr)", "Consolidated (Rs Cr)"]
        data = [
            ["Total Revenue", "45,230.50", "72,890.75"],
            ["EBITDA", "12,450.80", "19,870.25"],
            ["Profit Before Tax", "8,920.30", "14,560.60"],
            ["Profit After Tax", "6,690.25", "10,920.45"],
            ["Total Assets", "1,85,400.00", "2,98,650.30"],
            ["Net Worth", "62,300.15", "98,750.40"],
        ]
        for i, h in enumerate(headers):
            cell = table.rows[0].cells[i]
            cell.text = h
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Auditor Opinion", level=2)
        doc.add_paragraph(
            "The audited financial statements have been prepared in accordance with Indian Accounting "
            "Standards (Ind AS) notified under the Companies (Indian Accounting Standards) Rules, 2015 "
            "and comply with Section 129 and Section 134 of the Companies Act, 2013."
        )
        doc.add_heading("Significant Accounting Policies", level=2)
        doc.add_paragraph("Revenue Recognition (Ind AS 115)", style="List Number")
        doc.add_paragraph("Lease Accounting (Ind AS 116)", style="List Number")
        doc.add_paragraph("Employee Benefits (Ind AS 19)", style="List Number")
        doc.add_paragraph("Impairment of Assets (Ind AS 36)", style="List Number")
        doc.add_paragraph("Income Taxes (Ind AS 12)", style="List Number")
        doc.add_paragraph("Provisions and Contingent Liabilities (Ind AS 37)", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_statutory_audit_report_docx(self, filename="Statutory_Audit_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Statutory Audit Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Independent Auditor Report", style="List Bullet")
        doc.add_paragraph("To the Members of [Company Name] Limited", style="List Bullet")
        doc.add_heading("Report on the Standalone Financial Statements", level=1)
        doc.add_paragraph(
            "We have audited the accompanying standalone financial statements of [Company Name] Limited, "
            "which comprise the Balance Sheet as at 31st March 2025, the Statement of Profit and Loss, "
            "the Cash Flow Statement and the Statement of Changes in Equity for the year then ended, "
            "and notes to the financial statements."
        )
        doc.add_heading("Opinion", level=2)
        doc.add_paragraph(
            "In our opinion and to the best of our information and according to the explanations given "
            "to us, the aforesaid standalone financial statements give the information required by the "
            "Companies Act, 2013 in the manner so required and give a true and fair view in conformity "
            "with the Indian Accounting Standards prescribed under Section 133 of the Act."
        )
        doc.add_heading("Basis for Opinion", level=2)
        doc.add_paragraph(
            "We conducted our audit in accordance with the Standards on Auditing specified under "
            "Section 143(10) of the Act."
        )
        doc.add_heading("Key Audit Matters", level=2)
        doc.add_paragraph("Revenue Recognition and Contract Assets", style="List Number")
        doc.add_paragraph("Impairment of Goodwill and Intangible Assets", style="List Number")
        doc.add_paragraph("Provision for Litigations and Claims", style="List Number")
        doc.add_paragraph("Going Concern Assessment", style="List Number")
        doc.add_heading("Report on Other Legal and Regulatory Requirements", level=2)
        doc.add_paragraph(
            "As required by Section 143(3) of the Act, we report that: (a) We have sought and obtained "
            "all the information and explanations which to the best of our knowledge and belief were "
            "necessary for the purposes of our audit."
        )
        doc.add_heading("Report on the Internal Financial Controls", level=2)
        doc.add_paragraph(
            "We have audited the internal financial controls over financial reporting of [Company Name] "
            "Limited as of 31st March 2025 in conjunction with our audit of the standalone financial statements."
        )
        doc.save(filepath)
        return filepath

    def generate_internal_audit_report_docx(self, filename="Internal_Audit_Report_Finance.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Internal Audit Report - Finance", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Audit Scope and Objective", level=1)
        doc.add_paragraph(
            "The internal audit function is conducted in accordance with the Standards on Internal Audit "
            "issued by ICAI and the Internal Audit Policy approved by the Audit Committee."
        )
        doc.add_heading("Finance & Accounts Audit Findings", level=1)
        table = doc.add_table(rows=6, cols=4, style="Light Grid Accent 1")
        headers = ["Finding No.", "Description", "Risk Rating", "Status"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        findings = [
            ["FIN-001", "Reconciliation delays in inter-company transactions", "Medium", "Resolved"],
            ["FIN-002", "Manual adjustments in general ledger without proper authorization", "High", "In Progress"],
            ["FIN-003", "Inventory valuation methodology inconsistency across units", "Medium", "In Progress"],
            ["FIN-004", "Pending confirmation of trade receivables > 180 days", "Low", "Resolved"],
            ["FIN-005", "Depreciation calculation errors in fixed asset register", "Medium", "Closed"],
        ]
        for r_idx, row_data in enumerate(findings):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Recommendations", level=1)
        doc.add_paragraph("Implement automated inter-company reconciliation module", style="List Number")
        doc.add_paragraph("Strengthen approval workflows for manual journal entries", style="List Number")
        doc.add_paragraph("Standardize inventory valuation policy across all business units", style="List Number")
        doc.add_paragraph("Establish monthly aging review process for trade receivables", style="List Number")
        doc.add_paragraph("Conduct periodic fixed asset physical verification drives", style="List Number")
        doc.add_heading("Conclusion", level=1)
        doc.add_paragraph(
            "The overall internal control environment in the Finance & Accounts function is satisfactory. "
            "A follow-up audit is scheduled for Q2 FY 2025-26."
        )
        doc.save(filepath)
        return filepath

    def generate_cost_audit_report_docx(self, filename="Cost_Audit_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Cost Audit Report - Section 148", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Cost Audit Report under Section 148 of the Companies Act, 2013")
        doc.add_heading("Cost Records Maintained", level=1)
        doc.add_paragraph("The Company has maintained cost records as specified under Companies (Cost Records and Audit) Rules, 2014.")
        doc.add_heading("Cost Statement Summary", level=2)
        table = doc.add_table(rows=9, cols=3, style="Light Grid Accent 1")
        headers = ["Cost Component", "Amount (Rs Cr)", "% of Total"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["Raw Materials", "18,500.00", "42.5%"],
            ["Employee Benefits", "6,200.00", "14.3%"],
            ["Manufacturing Expenses", "8,750.00", "20.1%"],
            ["Administration Expenses", "3,200.00", "7.4%"],
            ["Selling & Distribution", "2,850.00", "6.6%"],
            ["Finance Costs", "2,100.00", "4.8%"],
            ["Depreciation & Amortization", "1,900.00", "4.3%"],
            ["Total Cost of Production", "43,500.00", "100.0%"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Compliance with Cost Accounting Standards", level=1)
        doc.add_paragraph("CAS-1: Classification of Costs", style="List Number")
        doc.add_paragraph("CAS-2: Capacity Determination", style="List Number")
        doc.add_paragraph("CAS-3: Production, Planning and Control", style="List Number")
        doc.add_paragraph("CAS-5: Measurement of Cost of Materials", style="List Number")
        doc.add_paragraph("CAS-7: Measurement of Cost of Machine Capacity", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_income_tax_return_docx(self, filename="Income_Tax_Return_ITR6.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Income Tax Returns - ITR-6", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Assessment Year 2025-26", level=1)
        doc.add_paragraph("Previous Year: 2024-25 (April 2024 to March 2025)")
        doc.add_heading("Return Filing Details", level=2)
        table = doc.add_table(rows=9, cols=2, style="Light Grid Accent 1")
        details = [
            ["Form Type", "ITR-6"], ["PAN", "AABCM1234N"],
            ["Assessment Year", "2025-26"], ["Due Date of Filing", "30th November 2025"],
            ["Mode of Filing", "Electronic (e-Verification)"], ["Chartered Accountant", "[CA Name], FRN: 123456"],
            ["Audit Report Section", "44AB (Tax Audit)"], ["Transfer Pricing", "Yes (Section 92E)"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Tax Computation Summary", level=2)
        doc.add_paragraph("Income from Business & Profession: Rs 5,595.35 Cr", style="List Bullet")
        doc.add_paragraph("Short-term Capital Gains: Rs 125.20 Cr", style="List Bullet")
        doc.add_paragraph("Income from Other Sources: Rs 45.80 Cr", style="List Bullet")
        doc.add_paragraph("Gross Total Income: Rs 5,766.35 Cr", style="List Bullet")
        doc.add_paragraph("Deductions under Chapter VI-A: Rs 85.50 Cr", style="List Bullet")
        doc.add_paragraph("Taxable Income: Rs 5,680.85 Cr", style="List Bullet")
        doc.add_paragraph("Income Tax Payable: Rs 1,477.02 Cr", style="List Bullet")
        doc.add_paragraph("TDS Already Paid: Rs 1,320.45 Cr", style="List Bullet")
        doc.add_paragraph("Net Tax Payable / (Refund): Rs 156.57 Cr", style="List Bullet")
        doc.add_heading("Advance Tax Schedule", level=2)
        adv_table = doc.add_table(rows=5, cols=3, style="Light Grid Accent 1")
        adv_headers = ["Due Date", "Due (%)", "Amount (Rs Cr)"]
        for i, h in enumerate(adv_headers):
            adv_table.rows[0].cells[i].text = h
        adv_data = [
            ["15th June 2024", "15%", "221.55"], ["15th September 2024", "45%", "664.65"],
            ["15th December 2024", "75%", "1,107.75"], ["15th March 2025", "100%", "1,477.02"],
        ]
        for r_idx, row_data in enumerate(adv_data):
            for c_idx, val in enumerate(row_data):
                adv_table.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_board_report_docx(self, filename="Board_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Board Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("To the Members of [Company Name] Limited", level=1)
        doc.add_heading("1. Financial Performance", level=2)
        doc.add_paragraph("The Board of Directors is pleased to present the Annual Report along with Audited Financial Statements for the financial year ended 31st March 2025. The Company has posted a robust performance with total income of Rs 45,230.50 Crore and PAT of Rs 6,690.25 Crore.")
        doc.add_heading("2. State of Company Affairs", level=2)
        doc.add_paragraph("The operations of the Company remained strong during the year with production of XX MMTPA of crude oil and XX MMTPA of natural gas.")
        doc.add_heading("3. Directors Responsibility Statement", level=2)
        doc.add_paragraph("The Directors confirm that:", style="List Bullet")
        doc.add_paragraph("Applicable Indian Accounting Standards have been followed", style="List Bullet 2")
        doc.add_paragraph("Selected accounting policies are applied consistently", style="List Bullet 2")
        doc.add_paragraph("Judgments and estimates are reasonable and prudent", style="List Bullet 2")
        doc.add_paragraph("Financial statements prepared on a going concern basis", style="List Bullet 2")
        doc.add_heading("4. Board Composition and Meetings", level=2)
        doc.add_paragraph("Total Board Meetings held during FY 2024-25: 8", style="List Bullet")
        doc.add_heading("5. Directors and Key Managerial Personnel", level=2)
        doc.add_paragraph("Shri [Name] - Chairman & Managing Director", style="List Bullet")
        doc.add_paragraph("Shri [Name] - Director (Finance) & CFO", style="List Bullet")
        doc.add_paragraph("Shri [Name] - Company Secretary & Compliance Officer", style="List Bullet")
        doc.add_heading("6. Corporate Social Responsibility (CSR)", level=2)
        doc.add_paragraph("During the year, the Company has spent Rs XXX Crore on various CSR activities in accordance with Section 135 of the Companies Act, 2013.")
        doc.save(filepath)
        return filepath

    def generate_annual_return_mgt7_docx(self, filename="Annual_Return_MGT7.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Annual Return - Form MGT-7", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Particulars of the Company", level=1)
        table = doc.add_table(rows=8, cols=2, style="Light Grid Accent 1")
        details = [
            ["Company Name", "[Company Name] Limited"], ["CIN", "L99999MH19XXPLCXXXXX"],
            ["Registered Office", "[Address], Mumbai - 400001"], ["Date of Incorporation", "15th June 19XX"],
            ["Category of Company", "Public Company Limited by Shares"], ["Sub-category", "Central Government Company"],
            ["Listed Company", "Yes (BSE/NSE)"], ["Annual Return Filed On", "30th November 2025"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
        doc.add_heading("Share Capital Details", level=2)
        doc.add_paragraph("Authorized Share Capital: Rs XX,000.00 Crore", style="List Bullet")
        doc.add_paragraph("Paid-up Share Capital: Rs XX,000.00 Crore", style="List Bullet")
        doc.add_heading("Meetings of Members", level=2)
        doc.add_paragraph("Annual General Meeting held on [Date] at [Venue]", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_register_contracts_docx(self, filename="Register_of_Contracts_S189.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Register of Contracts - Section 189", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Register maintained under Section 189 of the Companies Act, 2013")
        doc.add_heading("Particulars of Contracts", level=1)
        table = doc.add_table(rows=6, cols=5, style="Light Grid Accent 1")
        headers = ["Date", "Contract Details", "Party Name", "Value (Rs Cr)", "Board Approval"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["15/04/2024", "Crude Oil Supply Agreement", "[Party A]", "12,500.00", "Yes"],
            ["01/06/2024", "Technology License Agreement", "[Party B]", "850.25", "Yes"],
            ["12/08/2024", "Capital Equipment Procurement", "[Party C]", "2,300.50", "Yes"],
            ["20/10/2024", "Pipeline Construction Contract", "[Party D]", "5,600.75", "Yes"],
            ["15/01/2025", "LNG Import Agreement", "[Party E]", "8,900.00", "Yes"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_register_charges_docx(self, filename="Register_of_Carges_CHG1_CHG4.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Register of Charges - CHG-1/CHG-4", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Particulars of Charges Created/Modified", level=1)
        table = doc.add_table(rows=4, cols=6, style="Light Grid Accent 1")
        headers = ["Form", "Date", "Asset Charged", "Lender", "Amount (Rs Cr)", "Status"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["CHG-1", "01/04/2024", "Immovable Property", "[Bank A]", "15,000.00", "Satisfied"],
            ["CHG-1", "15/07/2024", "Plant & Machinery", "[Bank B]", "8,500.00", "Active"],
            ["CHG-4", "20/12/2024", "Floating Assets", "[Bank C]", "5,200.00", "Modified"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_dividend_distribution_docx(self, filename="Dividend_Distribution_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Dividend Distribution Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Pursuant to Regulation 43A of SEBI LODR Regulations, 2015")
        doc.add_heading("1. Objective", level=2)
        doc.add_paragraph("This policy sets out the parameters and factors that will guide the Board of Directors in declaring dividend on equity shares of the Company.")
        doc.add_heading("2. Factors for Determination of Dividend", level=2)
        doc.add_paragraph("Historical dividend track record", style="List Number")
        doc.add_paragraph("Current year profitability and free cash flows", style="List Number")
        doc.add_paragraph("Future capital expenditure requirements", style="List Number")
        doc.add_paragraph("Debt servicing obligations", style="List Number")
        doc.add_heading("3. Dividend Distribution Framework", level=2)
        doc.add_paragraph("The Company targets a minimum dividend payout ratio of 30% of net profit after tax.")
        doc.add_heading("4. Dividend History", level=2)
        table = doc.add_table(rows=4, cols=3, style="Light Grid Accent 1")
        headers = ["Financial Year", "Dividend Per Share (Rs)", "Payout Ratio (%)"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [["FY 2022-23", "12.50", "32.5%"], ["FY 2023-24", "14.00", "35.2%"], ["FY 2024-25", "15.50", "38.8%"]]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_finance_review_pptx(self, filename="Finance_Review_Presentation.pptx"):
        filepath = self._get_output_path(filename)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
        title.text = "Finance Review - FY 2024-25"
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = 1
            for run in paragraph.runs:
                run.font.size = PptxInches(0.7)
                run.bold = True
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Financial Highlights"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        table = slide.shapes.add_table(5, 3, PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(3)).table
        headers = ["Metric", "FY 2024-25", "FY 2023-24"]
        for i, h in enumerate(headers):
            table.cell(0, i).text = h
        data = [["Revenue", "Rs 45,230 Cr", "Rs 42,100 Cr"], ["EBITDA", "Rs 12,451 Cr", "Rs 11,200 Cr"],
                ["PAT", "Rs 6,690 Cr", "Rs 5,890 Cr"], ["EPS (Rs)", "45.50", "39.80"]]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.cell(r_idx + 1, c_idx).text = val
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Financial Outlook & Guidance"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "FY 2025-26 Outlook:\n- Revenue guidance: 8-10% growth\n- EBITDA margin: 28-30%\n- Capex plan: Rs 50,000 Cr\n- Dividend payout ratio: >30%"
        prs.save(filepath)
        return filepath

    # =========================================================================
    # ESG METHODS
    # =========================================================================

    def generate_crz_clearance_docx(self, filename="CRZ_Clearance.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("CRZ Clearance - Coastal Regulation Zone", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Coastal Regulation Zone Clearance under CRZ Notification, 2019")
        doc.add_heading("Project Details", level=1)
        table = doc.add_table(rows=7, cols=2, style="Light Grid Accent 1")
        details = [
            ["Project Name", "[Refinery Expansion Project]"], ["Location", "[Site], Gujarat Coastline"],
            ["CRZ Zone", "CRZ-II (Existing Development)"], ["CRZ Area (Ha)", "25.50"],
            ["Sea Frontage (m)", "1,250"], ["MoEFCC Reference No.", "CRZ/2024/XXXXX"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("CRZ Compliance Requirements", level=2)
        doc.add_paragraph("No permanent construction within 200m of HTL", style="List Bullet")
        doc.add_paragraph("No new land reclamation", style="List Bullet")
        doc.add_paragraph("Sewage treatment and disposal arrangements", style="List Bullet")
        doc.add_paragraph("Oil spill contingency plan", style="List Bullet")
        doc.add_paragraph("Monthly monitoring of coastal water quality", style="List Bullet")
        doc.add_heading("Environmental Safeguards", level=2)
        doc.add_paragraph("The project has been designed to minimize coastal zone impact with provisions for oil spill prevention, zero liquid discharge (ZLD) system, and construction of coastal protection bunds as per CRZ norms.")
        doc.save(filepath)
        return filepath

    def generate_hazwaste_authorization_docx(self, filename="Hazardous_Waste_Authorization.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Hazardous Waste Authorization", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Hazardous and Other Wastes (Management and Transboundary Movement) Rules, 2016")
        doc.add_heading("Authorization Details", level=1)
        table = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
        details = [
            ["Consent Number", "[Consent/Authorization No.]"], ["Issuing Authority", "Gujarat Pollution Control Board (GPCB)"],
            ["Date of Issue", "1st April 2024"], ["Validity", "5 Years (2024-2029)"],
            ["Authorization Type", "Generation, Storage, Transport & Disposal"], ["Schedule", "Schedule I (Category A)"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Hazardous Waste Categories", level=2)
        doc.add_paragraph("Used/Spent Oil & Solvents", style="List Bullet")
        doc.add_paragraph("Spent Catalyst", style="List Bullet")
        doc.add_paragraph("Acid Sludge", style="List Bullet")
        doc.add_paragraph("Contaminated Soil", style="List Bullet")
        doc.add_heading("Annual Waste Generation Summary", level=2)
        table2 = doc.add_table(rows=5, cols=3, style="Light Grid Accent 1")
        table2.rows[0].cells[0].text = "Waste Type"
        table2.rows[0].cells[1].text = "Quantity (MT/Year)"
        table2.rows[0].cells[2].text = "Authorized Limit (MT/Year)"
        data = [["Spent Oil", "12,500", "15,000"], ["Spent Catalyst", "2,800", "3,500"],
                ["Acid Sludge", "5,600", "7,000"], ["Other HW", "3,200", "4,000"]]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table2.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_biomedical_waste_auth_docx(self, filename="Biomedical_Waste_Authorization.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Biomedical Waste Authorization", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Biomedical Waste Management Rules, 2016")
        doc.add_heading("Facility Details", level=1)
        table = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
        details = [
            ["Facility Name", "[Refinery Township Hospital]"], ["Category", "Category A (Bed Strength > 30)"],
            ["BMW Collection Centre", "Central Waste Management Unit"],
            ["Treatment Facility", "On-site autoclave + third-party incineration"], ["Annual Generation (Kg/Day)", "45"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("BMW Categories", level=2)
        doc.add_paragraph("Yellow - Human anatomical waste, infectious waste", style="List Bullet")
        doc.add_paragraph("Red - Contaminated recyclable waste", style="List Bullet")
        doc.add_paragraph("Blue - Glassware, metallic waste", style="List Bullet")
        doc.add_paragraph("White - Sharps (needles, blades)", style="List Bullet")
        doc.add_heading("Compliance Requirements", level=2)
        doc.add_paragraph("Form IV filing with SPCB", style="List Bullet")
        doc.add_paragraph("Monthly BMW management committee meetings", style="List Bullet")
        doc.add_paragraph("Annual BMW audit by authorized agency", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_ewaste_return_docx(self, filename="E_Waste_Management_Return.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("E-Waste Management Return", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under E-Waste (Management) Rules, 2022")
        doc.add_heading("Producer/OEM Details", level=1)
        table = doc.add_table(rows=5, cols=2, style="Light Grid Accent 1")
        details = [
            ["Category", "Producer (Category B - Import > 1000 MT)"], ["EPR Registration No.", "EPR/XXXX/XXXXXX"],
            ["Reporting Period", "FY 2024-25"], ["E-Waste Generated (MT)", "850"],
            ["Target EPR Credits", "750 (88% compliance)"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("E-Waste Recycled/Processed", level=2)
        table2 = doc.add_table(rows=5, cols=3, style="Light Grid Accent 1")
        table2.rows[0].cells[0].text = "Category"
        table2.rows[0].cells[1].text = "Qty (MT)"
        table2.rows[0].cells[2].text = "Recycler"
        data = [["IT Equipment", "320", "[Authorized Recycler 1]"], ["Consumer Electronics", "180", "[Authorized Recycler 2]"],
                ["Batteries", "150", "[Authorized Recycler 3]"], ["Cables & Wires", "200", "[Authorized Recycler 4]"]]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table2.rows[r_idx + 1].cells[c_idx].text = val
        doc.save(filepath)
        return filepath

    def generate_energy_audit_report_docx(self, filename="Energy_Audit_Report_PAT.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Energy Audit Report - PAT Scheme", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Perform, Achieve and Trade (PAT) Scheme - Designated Consumer Report")
        doc.add_heading("Designated Consumer Details", level=1)
        table = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
        details = [
            ["DC Number", "DC/XXXX/XXXX"], ["Sector", "Oil Refining"],
            ["Baseline Year", "FY 2018-19"], ["Target Year", "FY 2024-25 (Cycle 3)"],
            ["Specific Energy Consumption (SEC)", "68.5 GJ/MT of crude"],
            ["Target SEC", "65.2 GJ/MT of crude (4.8% reduction)"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Energy Saving Measures Implemented", level=2)
        doc.add_paragraph("Waste heat recovery system installation", style="List Number")
        doc.add_paragraph("Variable Frequency Drive (VFD) on compressors", style="List Number")
        doc.add_paragraph("LED lighting replacement across facility", style="List Number")
        doc.add_paragraph("Solar power plant (25 MW) commissioning", style="List Number")
        doc.add_paragraph("Heat integration optimization in distillation units", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_biodiversity_register_docx(self, filename="Biodiversity_Register.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Biodiversity Register", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Biological Diversity Act, 2002 and Biological Diversity Rules, 2004")
        doc.add_heading("Site Information", level=1)
        table = doc.add_table(rows=5, cols=2, style="Light Grid Accent 1")
        details = [
            ["Project Site", "[Refinery Complex], Gujarat"], ["Area (Ha)", "500"],
            ["Biodiversity Zone", "Within 10 km of Protected Area"],
            ["State Biodiversity Board", "Gujarat State Biodiversity Board (GSBB)"],
            ["NBA Registration No.", "NBA/XXXX/XXXXXX"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Flora Documentation", level=2)
        doc.add_paragraph("Total plant species recorded: 125", style="List Bullet")
        doc.add_paragraph("Threatened species: 3", style="List Bullet")
        doc.add_heading("Fauna Documentation", level=2)
        doc.add_paragraph("Total fauna species recorded: 85", style="List Bullet")
        doc.add_paragraph("Mammals: 15 species", style="List Bullet")
        doc.add_paragraph("Birds: 45 species", style="List Bullet")
        doc.add_heading("Conservation Measures", level=2)
        doc.add_paragraph("Tree transplantation program (2,500 trees)", style="List Number")
        doc.add_paragraph("Mangrove restoration (15 Ha)", style="List Number")
        doc.add_paragraph("Native species plantation drive", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_eia_report_docx(self, filename="EIA_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Environmental Impact Assessment (EIA) Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Project Description", level=1)
        doc.add_paragraph("This EIA report has been prepared for the [Refinery Expansion] project in accordance with the Environmental Impact Assessment Notification, 2006 (amended 2024) issued by MoEFCC.")
        doc.add_heading("Baseline Environmental Conditions", level=2)
        table = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
        details = [
            ["Air Quality", "PM10: 45 ug/m3, PM2.5: 28 ug/m3 (Within NAAQS)"],
            ["Water Quality", "pH: 7.2, BOD: 3.2 mg/L (Within SWQS)"],
            ["Noise Level", "Day: 58 dB(A), Night: 45 dB(A) (Within limits)"],
            ["Soil Quality", "No significant contamination detected"],
            ["Ecology", "Moderate biodiversity, no protected species within site"],
            ["Socio-economic", "Population within 10 km: 2,50,000"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Environmental Management Plan", level=2)
        doc.add_paragraph("Zero Liquid Discharge (ZLD) system", style="List Number")
        doc.add_paragraph("Continuous Emission Monitoring System (CEMS)", style="List Number")
        doc.add_paragraph("Compensatory afforestation (2:1 ratio)", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_esg_review_pptx(self, filename="ESG_Review_Presentation.pptx"):
        filepath = self._get_output_path(filename)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
        title.text = "ESG Review - FY 2024-25"
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = 1
            for run in paragraph.runs:
                run.font.size = PptxInches(0.7)
                run.bold = True
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "ESG Performance Scorecard"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Environmental:\n- Carbon Emissions: -12% YoY\n- Water Recycling: 92%\n- Renewable Energy: 25%\n\nSocial:\n- Employee Safety: LTIFR 0.15\n- Diversity: 18% women\n- CSR Spend: Rs 125 Cr\n\nGovernance:\n- Board Independence: 50%\n- ESG Committee: 6 meetings\n- Cyber Security: ISO 27001"
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Environmental Performance"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Carbon Footprint:\n- Scope 1: 2.8 MT CO2e (down 8%)\n- Scope 2: 1.2 MT CO2e (down 15%)\n\nWater Management:\n- Consumption: 18.5 MCM\n- Recycling Rate: 92%\n\nBiodiversity:\n- Mangrove Restoration: 15 Ha\n- Tree Plantation: 50,000 trees"
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Social Impact & Governance"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Social:\n- LTIFR: 0.15 (Industry Avg: 0.45)\n- Healthcare Centers: 12\n- Schools Supported: 8\n- Women in Leadership: 12%\n\nGovernance:\n- Anti-Bribery: ISO 37001\n- Ethics Hotline: 24/7\n- GRI Standards: Comprehensive\n- DJSI Score: 72/100"
        prs.save(filepath)
        return filepath

    # =========================================================================
    # AUDIT METHODS
    # =========================================================================

    def generate_governance_report_docx(self, filename="Corporate_Governance_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Corporate Governance Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("In compliance with Regulation 34 of SEBI LODR Regulations, 2015")
        doc.add_heading("Board of Directors", level=1)
        table = doc.add_table(rows=6, cols=4, style="Light Grid Accent 1")
        headers = ["Name", "Designation", "Category", "Date of Joining"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["Shri [Name]", "Chairman", "Non-Executive", "01/04/2018"],
            ["Shri [Name]", "MD & CEO", "Whole-time", "01/07/2020"],
            ["Shri [Name]", "Director (Finance)", "Whole-time", "15/01/2021"],
            ["Smt. [Name]", "Independent Director", "Independent", "01/04/2019"],
            ["Shri [Name]", "Independent Director", "Independent", "01/04/2022"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Board Meetings", level=2)
        doc.add_paragraph("Total Board Meetings held: 8", style="List Bullet")
        doc.add_paragraph("Average attendance: 92%", style="List Bullet")
        doc.add_heading("Committees of the Board", level=2)
        doc.add_paragraph("Audit Committee: 6 meetings, 100% attendance", style="List Number")
        doc.add_paragraph("Nomination & Remuneration Committee: 4 meetings", style="List Number")
        doc.add_paragraph("Stakeholders Relationship Committee: 4 meetings", style="List Number")
        doc.add_paragraph("CSR Committee: 4 meetings", style="List Number")
        doc.add_paragraph("Risk Management Committee: 4 meetings", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_rpt_policy_docx(self, filename="Related_Party_Transaction_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Related Party Transaction Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Regulation 23 of SEBI LODR Regulations, 2015 and Section 188 of the Companies Act, 2013")
        doc.add_heading("1. Objective", level=2)
        doc.add_paragraph("This policy governs all Related Party Transactions (RPTs) to ensure that such transactions are conducted at arm's length basis and in the best interest of the Company.")
        doc.add_heading("2. Definition of Related Party", level=2)
        doc.add_paragraph("Directors and Key Managerial Personnel", style="List Bullet")
        doc.add_paragraph("Companies under same management", style="List Bullet")
        doc.add_paragraph("Associates and Joint Ventures", style="List Bullet")
        doc.add_paragraph("Promoters and their relatives", style="List Bullet")
        doc.add_heading("3. Approval Process", level=2)
        doc.add_paragraph("All RPTs exceeding Rs 10 Crore require Audit Committee prior approval. Transactions exceeding Rs 100 Crore require shareholder approval.")
        doc.add_heading("4. Materiality Threshold", level=2)
        doc.add_paragraph("Individual Transaction: 2% of annual revenue or Rs 10 Crore", style="List Bullet")
        doc.add_paragraph("Aggregate Transactions: 10% of annual revenue or Rs 50 Crore", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_vigil_mechanism_docx(self, filename="Vigil_Mechanism_Whistleblower.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Vigil Mechanism / Whistle Blower Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Section 177(10) of Companies Act, 2013 and Regulation 22 of SEBI LODR")
        doc.add_heading("1. Purpose", level=2)
        doc.add_paragraph("To provide a mechanism for employees and directors to report genuine concerns about unethical behavior, fraud, or violations without fear of retaliation.")
        doc.add_heading("2. Coverage", level=2)
        doc.add_paragraph("Financial fraud or misrepresentation", style="List Bullet")
        doc.add_paragraph("Bribery or corruption", style="List Bullet")
        doc.add_paragraph("Violation of laws or regulations", style="List Bullet")
        doc.add_paragraph("Health and safety violations", style="List Bullet")
        doc.add_paragraph("Environmental violations", style="List Bullet")
        doc.add_heading("3. Reporting Channels", level=2)
        doc.add_paragraph("Internal Audit Department", style="List Number")
        doc.add_paragraph("Chairman of Audit Committee", style="List Number")
        doc.add_paragraph("Ethics Hotline: [Toll-free Number]", style="List Number")
        doc.add_paragraph("Online Portal: [URL]", style="List Number")
        doc.add_heading("4. Investigation Process", level=2)
        doc.add_paragraph("Preliminary assessment within 48 hours", style="List Bullet")
        doc.add_paragraph("Formal investigation within 30 days", style="List Bullet")
        doc.add_paragraph("Closure report to Audit Committee", style="List Bullet")
        doc.add_heading("5. Protection to Whistle Blower", level=2)
        doc.add_paragraph("The identity of the whistle blower shall be kept confidential. No retaliation shall be taken against the whistle blower in any form.")
        doc.save(filepath)
        return filepath

    def generate_insider_trading_code_docx(self, filename="Insider_Trading_Code.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Code of Conduct for Prevention of Insider Trading", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under SEBI (Prohibition of Insider Trading) Regulations, 2015")
        doc.add_heading("1. Definitions", level=2)
        doc.add_paragraph("Insider: Person with Unpublished Price Sensitive Information (UPSI)", style="List Bullet")
        doc.add_paragraph("Connected Person: Director, officer, employee, or person with access to UPSI", style="List Bullet")
        doc.add_paragraph("Designated Persons: Senior management and their immediate relatives", style="List Bullet")
        doc.add_heading("2. Trading Restrictions", level=2)
        doc.add_paragraph("No trading during trading window closure (before financial results)", style="List Bullet")
        doc.add_paragraph("Pre-clearance required for all trades by Designated Persons", style="List Bullet")
        doc.add_paragraph("Minimum holding period: 6 months from acquisition", style="List Bullet")
        doc.add_paragraph("Maximum trade limit: 25% of total holdings per quarter", style="List Bullet")
        doc.add_heading("3. Disclosure Requirements", level=2)
        doc.add_paragraph("Trade disclosure within 2 working days to Compliance Officer", style="List Bullet")
        doc.add_paragraph("Annual disclosure of holdings by all Designated Persons", style="List Bullet")
        doc.add_heading("4. Penalties", level=2)
        doc.add_paragraph("Violation of this code may result in disciplinary action including termination, suspension, or cancellation of stock options.")
        doc.save(filepath)
        return filepath

    def generate_csr_policy_docx(self, filename="CSR_Policy_S135.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("CSR Policy - Section 135", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Corporate Social Responsibility Policy under Section 135 of Companies Act, 2013")
        doc.add_heading("1. CSR Objective", level=2)
        doc.add_paragraph("To integrate social, environmental, and economic concerns into business operations and contribute to sustainable development of communities.")
        doc.add_heading("2. CSR Committee Composition", level=2)
        doc.add_paragraph("Chairperson: Independent Director", style="List Bullet")
        doc.add_paragraph("Members: 2 Non-Executive Directors", style="List Bullet")
        doc.add_heading("3. CSR Activities (Schedule VII)", level=2)
        doc.add_paragraph("Education: Schools, colleges, skill development centers", style="List Number")
        doc.add_paragraph("Healthcare: Hospitals, mobile health units, sanitation", style="List Number")
        doc.add_paragraph("Environment: Afforestation, water conservation, renewable energy", style="List Number")
        doc.add_paragraph("Rural Development: Infrastructure, drinking water, roads", style="List Number")
        doc.add_heading("4. CSR Spend Target", level=2)
        doc.add_paragraph("CSR Spend Requirement (2% of average net profit): Rs XX Crore", style="List Bullet")
        doc.add_paragraph("Actual CSR Spend: Rs XX Crore", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_secretarial_compliance_docx(self, filename="Secretarial_Compliance_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Secretarial Compliance Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Regulation 24A of SEBI LODR Regulations, 2015")
        doc.add_heading("Compliance Status Summary", level=1)
        table = doc.add_table(rows=6, cols=3, style="Light Grid Accent 1")
        headers = ["Category", "Complied", "Non-Compliant"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [["SEBI LODR", "42", "0"], ["SEBI PIT", "8", "0"], ["SEBI SAST", "3", "0"],
                ["Companies Act", "35", "0"], ["Total", "88", "0"]]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Certification", level=2)
        doc.add_paragraph("I hereby certify that [Company Name] Limited has complied with all applicable provisions of SEBI LODR Regulations, 2015 and the Companies Act, 2013 during the financial year 2024-25.")
        doc.save(filepath)
        return filepath

    def generate_material_subsidiary_policy_docx(self, filename="Material_Subsidiary_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Material Subsidiary Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Regulation 16(1)(c) of SEBI LODR Regulations, 2015")
        doc.add_heading("1. Definition of Material Subsidiary", level=2)
        doc.add_paragraph("A subsidiary shall be considered material if its income or net worth exceeds 10% of the consolidated income or net worth of the listed entity, respectively.")
        doc.add_heading("2. Governance Framework", level=2)
        doc.add_paragraph("Independent Director on the Board of material subsidiary", style="List Bullet")
        doc.add_paragraph("Audit Committee review of material subsidiary accounts", style="List Bullet")
        doc.add_paragraph("Quarterly reporting to parent company Board", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_risk_management_docx(self, filename="Risk_Management_Framework.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Risk Management Framework", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Enterprise Risk Management (ERM) Approach", level=2)
        doc.add_paragraph("The Company follows an integrated Enterprise Risk Management framework aligned with COSO ERM and ISO 31000.")
        doc.add_heading("2. Risk Governance Structure", level=2)
        doc.add_paragraph("Board Level: Risk Management Committee", style="List Bullet")
        doc.add_paragraph("Management Level: Chief Risk Officer (CRO)", style="List Bullet")
        doc.add_heading("3. Risk Categories", level=2)
        table = doc.add_table(rows=8, cols=4, style="Light Grid Accent 1")
        headers = ["Risk Category", "Key Risks", "Likelihood", "Impact"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["Strategic", "Commodity price volatility", "High", "High"],
            ["Operational", "Plant shutdown/failure", "Medium", "High"],
            ["Financial", "Foreign exchange exposure", "High", "Medium"],
            ["Compliance", "Regulatory changes", "Medium", "High"],
            ["Environmental", "Climate change impact", "Medium", "High"],
            ["Cyber", "Data breach/cyber attack", "Low", "High"],
            ["Geopolitical", "Sanctions/trade restrictions", "Medium", "Medium"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("4. Risk Mitigation Measures", level=2)
        doc.add_paragraph("Hedging of 60-70% of forex exposure", style="List Bullet")
        doc.add_paragraph("Business continuity planning and disaster recovery", style="List Bullet")
        doc.add_paragraph("Insurance coverage for major assets and liabilities", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_ifc_report_docx(self, filename="Internal_Financial_Controls_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Internal Financial Controls (IFC) Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Section 143(3)(f) of the Companies Act, 2013")
        doc.add_heading("1. Scope of Review", level=2)
        doc.add_paragraph("We have audited the internal financial controls over financial reporting of [Company Name] Limited as of 31st March 2025.")
        doc.add_heading("2. Control Environment", level=2)
        table = doc.add_table(rows=6, cols=3, style="Light Grid Accent 1")
        headers = ["Control Area", "Design Effectiveness", "Operating Effectiveness"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["Revenue Recognition", "Adequate", "Effective"],
            ["Procurement & Payable", "Adequate", "Effective"],
            ["Fixed Assets", "Adequate", "Effective"],
            ["Payroll", "Adequate", "Effective"],
            ["Financial Reporting", "Adequate", "Effective"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("4. Opinion", level=2)
        doc.add_paragraph("In our opinion, [Company Name] Limited has, in all material respects, an adequate internal financial controls system over financial reporting and such internal financial controls are operating effectively as of 31st March 2025.")
        doc.save(filepath)
        return filepath

    def generate_statutory_returns_docx(self, filename="Statutory_Returns_ROC.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Statutory Returns - ROC Filings", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Returns filed with the Registrar of Companies (ROC) during FY 2024-25")
        doc.add_heading("Annual Returns", level=1)
        table = doc.add_table(rows=3, cols=4, style="Light Grid Accent 1")
        headers = ["Form", "Description", "Due Date", "Filing Date"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["MGT-7", "Annual Return", "29/11/2025", "30/11/2025"],
            ["MGT-14", "Board Resolution (Section 117)", "30 days", "Within time"],
            ["ADT-1", "Appointment of Auditor", "15 days from AGM", "Within time"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Filing Compliance Summary", level=2)
        doc.add_paragraph("Total forms filed during the year: 28", style="List Bullet")
        doc.add_paragraph("Forms filed within due date: 27 (96.4%)", style="List Bullet")
        doc.add_paragraph("Forms filed with delay: 1 (MGT-7 - 1 day delay)", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_audit_review_pptx(self, filename="Audit_Committee_Review.pptx"):
        filepath = self._get_output_path(filename)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
        title.text = "Audit Committee Review - FY 2024-25"
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = 1
            for run in paragraph.runs:
                run.font.size = PptxInches(0.7)
                run.bold = True
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Audit Committee Composition"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Committee Members:\n- Chairman: Smt. [Name] - Independent Director\n- Member: Shri [Name] - Independent Director\n\nMeetings Held: 6\nAttendance: 100%\n\nKey Reviews:\n- Quarterly financial results\n- Statutory & internal audit reports\n- Related party transactions\n- Risk management framework"
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Key Audit Findings & Compliance"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Statutory Audit:\n- Unqualified opinion\n\nInternal Audit:\n- 5 findings (3 resolved)\n\nIFC Review:\n- No material weaknesses\n\nCompliance:\n- SEBI LODR: 42/42\n- Companies Act: 35/35"
        prs.save(filepath)
        return filepath

    # =========================================================================
    # HR METHODS
    # =========================================================================

    def generate_hr_policy_docx(self, filename="HR_Welfare_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("HR Welfare Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Employment Terms & Conditions", level=2)
        doc.add_paragraph("Appointment letter issued with terms of employment", style="List Bullet")
        doc.add_paragraph("Probation period: 6 months (extendable by 6 months)", style="List Bullet")
        doc.add_paragraph("Notice period: 3 months for management, 1 month for workmen", style="List Bullet")
        doc.add_paragraph("Working hours: 8 hours/day, 48 hours/week", style="List Bullet")
        doc.add_heading("2. Compensation & Benefits", level=2)
        doc.add_paragraph("Pay structure: Basic + HRA + Special Allowance + Performance Pay", style="List Bullet")
        doc.add_paragraph("Annual increment: 8-12% based on performance", style="List Bullet")
        doc.add_paragraph("Performance bonus: Up to 20% of basic", style="List Bullet")
        doc.add_heading("3. Leave Entitlements", level=2)
        table = doc.add_table(rows=8, cols=2, style="Light Grid Accent 1")
        leaves = [
            ["Earned Leave", "30 days/year"], ["Casual Leave", "12 days/year"],
            ["Sick Leave", "12 days/year"], ["Privilege Leave", "10 days/year"],
            ["Maternity Leave", "26 weeks"], ["Paternity Leave", "15 days"],
            ["Compassionate Leave", "5 days"],
        ]
        for r_idx, (leave, days) in enumerate(leaves):
            table.rows[r_idx + 1].cells[0].text = leave
            table.rows[r_idx + 1].cells[1].text = days
        doc.add_heading("4. Welfare Measures", level=2)
        doc.add_paragraph("Group health insurance for employee and family", style="List Bullet")
        doc.add_paragraph("Life insurance coverage (5x annual CTC)", style="List Bullet")
        doc.add_paragraph("Provident fund contribution: 12% of basic", style="List Bullet")
        doc.add_paragraph("Gratuity as per Payment of Gratuity Act", style="List Bullet")
        doc.add_paragraph("Employee stock purchase plan", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_transfer_policy_docx(self, filename="Transfer_Benefit_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Transfer & Posting Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Transfer Policy Objectives", level=2)
        doc.add_paragraph("Career development through exposure to different roles/locations", style="List Bullet")
        doc.add_paragraph("Succession planning and talent management", style="List Bullet")
        doc.add_heading("2. Transfer Categories", level=2)
        doc.add_paragraph("Administrative Transfer: Based on organizational requirements", style="List Bullet")
        doc.add_paragraph("Development Transfer: For career enhancement", style="List Bullet")
        doc.add_paragraph("Voluntary Transfer: Employee-initiated", style="List Bullet")
        doc.add_heading("3. Transfer Benefits", level=2)
        doc.add_paragraph("House shifting allowance: Up to Rs 2,00,000", style="List Bullet")
        doc.add_paragraph("Temporary accommodation: Up to 30 days", style="List Bullet")
        doc.add_paragraph("Travel allowance for family", style="List Bullet")
        doc.add_paragraph("School admission assistance for children", style="List Bullet")
        doc.add_heading("4. Transfer Cycle", level=2)
        doc.add_paragraph("Minimum posting period: 3 years", style="List Bullet")
        doc.add_paragraph("Maximum posting period: 5 years", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_promotion_policy_docx(self, filename="Promotion_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Promotion Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Promotion Philosophy", level=2)
        doc.add_paragraph("Promotion is based on merit, performance, potential, and organizational needs.")
        doc.add_heading("2. Eligibility Criteria", level=2)
        doc.add_paragraph("Minimum 2 years in current grade", style="List Bullet")
        doc.add_paragraph("Rating of Above Expectations or higher in last 2 appraisals", style="List Bullet")
        doc.add_paragraph("Completion of required training/certifications", style="List Bullet")
        doc.add_heading("3. Career Ladder", level=2)
        table = doc.add_table(rows=8, cols=3, style="Light Grid Accent 1")
        headers = ["Level", "Designation", "Typical Experience"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["L1", "Officer", "0-3 years"], ["L2", "Senior Officer", "3-6 years"],
            ["L3", "Assistant Manager", "6-10 years"], ["L4", "Manager", "10-15 years"],
            ["L5", "Senior Manager", "15-20 years"], ["L6", "DGM", "20-25 years"],
            ["L7", "GM & above", "25+ years"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("4. Promotion Process", level=2)
        doc.add_paragraph("Self-nomination or manager recommendation", style="List Number")
        doc.add_paragraph("Departmental promotion committee screening", style="List Number")
        doc.add_paragraph("HR review and validation", style="List Number")
        doc.add_paragraph("Communication of decision within 30 days", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_medical_benefit_policy_docx(self, filename="Post_Retirement_Medical_Benefit.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Post-Retirement Medical Benefit Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Objective", level=2)
        doc.add_paragraph("To provide medical insurance coverage to retired employees and their dependents.")
        doc.add_heading("2. Eligibility", level=2)
        doc.add_paragraph("Employees retiring after 10 years of service", style="List Bullet")
        doc.add_paragraph("Voluntary retirement under VRS after 5 years of service", style="List Bullet")
        doc.add_heading("3. Coverage", level=2)
        table = doc.add_table(rows=4, cols=3, style="Light Grid Accent 1")
        headers = ["Service Category", "Coverage (Rs)", "Floater"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["20+ years", "10,00,000", "Employee + Spouse"],
            ["15-20 years", "7,50,000", "Employee + Spouse"],
            ["10-15 years", "5,00,000", "Employee + Spouse"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("4. Benefits Covered", level=2)
        doc.add_paragraph("In-patient hospitalization", style="List Bullet")
        doc.add_paragraph("Pre and post hospitalization (60 days each)", style="List Bullet")
        doc.add_paragraph("Day care procedures", style="List Bullet")
        doc.add_paragraph("Annual health check-up", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_standing_orders_docx(self, filename="Standing_Orders.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Standing Orders", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Under Industrial Employment (Standing Orders) Act, 1946")
        doc.add_heading("1. Classification of Workmen", level=2)
        doc.add_paragraph("Permanent workmen: Employed continuously for 240 days or more", style="List Bullet")
        doc.add_paragraph("Temporary workmen: Employed for less than 240 days", style="List Bullet")
        doc.add_paragraph("Probationers: Undergoing probation of 6 months", style="List Bullet")
        doc.add_heading("2. Shift Working", level=2)
        doc.add_paragraph("Three shifts: General (6 AM - 2 PM), Second (2 PM - 10 PM), Third (10 PM - 6 AM)", style="List Bullet")
        doc.add_paragraph("Night shift allowance: 25% of basic wages", style="List Bullet")
        doc.add_paragraph("Overtime: Double wages for work beyond 8 hours/day", style="List Bullet")
        doc.add_heading("3. Attendance & Leave", level=2)
        doc.add_paragraph("Roll call at beginning of each shift", style="List Bullet")
        doc.add_paragraph("Leave as per Factory Act, 1948", style="List Bullet")
        doc.add_heading("4. Termination & Dismissal", level=2)
        doc.add_paragraph("Notice period: 1 month for permanent workmen", style="List Bullet")
        doc.add_paragraph("Domestic enquiry mandatory before dismissal", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_human_rights_policy_docx(self, filename="Human_Rights_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Human Rights Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Aligned with UN Global Compact Principles")
        doc.add_heading("1. Policy Statement", level=2)
        doc.add_paragraph("[Company Name] Limited is committed to respecting and promoting human rights across its operations and value chain.")
        doc.add_heading("2. Core Principles", level=2)
        doc.add_paragraph("Non-discrimination and equal opportunity", style="List Bullet")
        doc.add_paragraph("Freedom of association and collective bargaining", style="List Bullet")
        doc.add_paragraph("Prohibition of child labor and forced labor", style="List Bullet")
        doc.add_paragraph("Right to safe and healthy working conditions", style="List Bullet")
        doc.add_paragraph("Fair wages and working hours", style="List Bullet")
        doc.add_heading("3. Due Diligence Process", level=2)
        doc.add_paragraph("Human rights impact assessment for new projects", style="List Number")
        doc.add_paragraph("Annual human rights audit", style="List Number")
        doc.add_paragraph("Grievance mechanism for affected stakeholders", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_diversity_report_docx(self, filename="Diversity_Inclusion_Report.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Diversity & Inclusion Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("For the Financial Year 2024-25")
        doc.add_heading("1. Workforce Composition", level=1)
        table = doc.add_table(rows=7, cols=2, style="Light Grid Accent 1")
        data = [
            ["Total Employees", "28,500"], ["Women Employees", "5,130 (18%)"],
            ["Persons with Disabilities", "855 (3%)"], ["Employees from SC/ST", "4,275 (15%)"],
            ["Fresh Graduates", "1,425 (5%)"], ["Experienced Hires", "27,075 (95%)"],
        ]
        for r_idx, (label, value) in enumerate(data):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("2. Gender Diversity Initiatives", level=2)
        doc.add_paragraph("Women Leadership Program: 50 women trained", style="List Bullet")
        doc.add_paragraph("Mentorship program for high-potential women", style="List Bullet")
        doc.add_paragraph("Flexible work arrangements for women", style="List Bullet")
        doc.add_heading("3. Training & Development", level=2)
        doc.add_paragraph("Average training hours per employee: 45 hours", style="List Bullet")
        doc.add_paragraph("Leadership development programs: 120 participants", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_training_register_docx(self, filename="Training_Competency_Register.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Training & Competency Register", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("For the Financial Year 2024-25")
        doc.add_heading("1. Training Summary", level=1)
        table = doc.add_table(rows=6, cols=2, style="Light Grid Accent 1")
        data = [
            ["Total Training Hours", "12,82,500 hours"], ["Average per Employee", "45 hours/year"],
            ["Number of Programs", "2,850"], ["Internal Training", "60%"], ["External Training", "40%"],
        ]
        for r_idx, (label, value) in enumerate(data):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("2. Competency Matrix", level=2)
        table2 = doc.add_table(rows=6, cols=3, style="Light Grid Accent 1")
        headers = ["Competency", "Target Level", "Achieved Level"]
        for i, h in enumerate(headers):
            table2.rows[0].cells[i].text = h
            for paragraph in table2.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data2 = [
            ["Technical Skills", "80%", "78%"], ["Safety Awareness", "100%", "98%"],
            ["Leadership", "70%", "72%"], ["Digital Literacy", "60%", "65%"],
            ["Compliance", "100%", "100%"],
        ]
        for r_idx, row_data in enumerate(data2):
            for c_idx, val in enumerate(row_data):
                table2.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("3. Key Training Programs", level=2)
        doc.add_paragraph("Safety induction (mandatory for all): 8,500 employees", style="List Bullet")
        doc.add_paragraph("Technical competency development: 4,200 employees", style="List Bullet")
        doc.add_paragraph("Leadership development: 120 high-potential employees", style="List Bullet")
        doc.add_paragraph("Compliance training: 100% coverage", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_hr_review_pptx(self, filename="HR_Review_Presentation.pptx"):
        filepath = self._get_output_path(filename)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
        title.text = "HR Review - FY 2024-25"
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = 1
            for run in paragraph.runs:
                run.font.size = PptxInches(0.7)
                run.bold = True
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Workforce Overview"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Total Workforce: 28,500\n\nHeadcount by Function:\n- Upstream: 8,500 (30%)\n- Refining: 6,200 (22%)\n- Marketing: 5,700 (20%)\n- Corporate: 4,275 (15%)\n\nAttrition Rate: 8.5%\nAverage Age: 38 years\nAverage Tenure: 12 years"
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Diversity & Engagement"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Diversity:\n- Women in Workforce: 18%\n- Women in Leadership: 12%\n- PwD Employees: 3%\n\nTraining:\n- Hours per employee: 45\n- Investment: Rs 85 Cr\n\nEngagement Score: 4.2/5.0"
        prs.save(filepath)
        return filepath

    # =========================================================================
    # VIGILANCE METHODS
    # =========================================================================

    def generate_anti_bribery_policy_docx(self, filename="Anti_Bribery_Anti_Corruption_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Anti-Bribery & Anti-Corruption Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("ISO 37001 Compliant")
        doc.add_heading("1. Policy Statement", level=2)
        doc.add_paragraph("[Company Name] Limited has zero tolerance for bribery and corruption. The Company is committed to conducting business ethically and in compliance with all applicable anti-bribery laws.")
        doc.add_heading("2. Prohibited Conduct", level=2)
        doc.add_paragraph("Offering, promising, or giving bribes to any person", style="List Bullet")
        doc.add_paragraph("Requesting, agreeing to receive, or accepting bribes", style="List Bullet")
        doc.add_paragraph("Facilitation payments", style="List Bullet")
        doc.add_heading("3. Gift & Hospitality Guidelines", level=2)
        doc.add_paragraph("Maximum value of gift: Rs 5,000 per instance", style="List Bullet")
        doc.add_paragraph("Maximum value of hospitality: Rs 10,000 per instance", style="List Bullet")
        doc.add_paragraph("Annual aggregate limit: Rs 25,000", style="List Bullet")
        doc.add_heading("4. Due Diligence", level=2)
        doc.add_paragraph("Third-party due diligence for agents and consultants", style="List Number")
        doc.add_paragraph("Enhanced due diligence for government-facing roles", style="List Number")
        doc.add_paragraph("Annual declaration of interests by all employees", style="List Number")
        doc.save(filepath)
        return filepath

    def generate_gift_policy_docx(self, filename="Gift_Policy.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Gift Policy", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("1. Scope", level=2)
        doc.add_paragraph("This policy applies to all employees, directors, and contractors of [Company Name] Limited.")
        doc.add_heading("2. Permissible Gifts", level=2)
        doc.add_paragraph("Branded promotional items (value up to Rs 1,000)", style="List Bullet")
        doc.add_paragraph("Festival greetings (value up to Rs 5,000)", style="List Bullet")
        doc.add_paragraph("Business meals (value up to Rs 2,000 per person)", style="List Bullet")
        doc.add_heading("3. Restricted Gifts", level=2)
        table = doc.add_table(rows=5, cols=2, style="Light Grid Accent 1")
        data = [
            ["Gift Type", "Action Required"],
            ["Cash or cash equivalents", "Prohibited"],
            ["Gifts > Rs 5,000", "Prior approval from Compliance Officer"],
            ["Gifts from/to government officials", "Enhanced due diligence"],
            ["Gifts to/from suppliers", "Gift register entry"],
        ]
        for r_idx, (label, value) in enumerate(data):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
        doc.add_heading("4. Approval Matrix", level=2)
        doc.add_paragraph("Up to Rs 5,000: Department Head approval", style="List Bullet")
        doc.add_paragraph("Rs 5,000 - Rs 25,000: Compliance Officer approval", style="List Bullet")
        doc.add_paragraph("Above Rs 25,000: Ethics Committee approval", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_vigilance_clearance_docx(self, filename="Vigilance_Clearance_Certificate.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Vigilance Clearance Certificate", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_heading("Certificate", level=2)
        doc.add_paragraph("This is to certify that Shri/Smt. [Name], [Designation], Employee No. [XXXXXX], has been granted Vigilance Clearance by the Competent Authority for the purpose of [Promotion/Transfer/Foreign Visit/Key Position] as on [Date].")
        doc.add_heading("Clearance Details", level=1)
        table = doc.add_table(rows=8, cols=2, style="Light Grid Accent 1")
        details = [
            ["Employee Name", "[Name]"], ["Employee ID", "[XXXXXX]"],
            ["Designation", "[Current Designation]"], ["Department", "[Department]"],
            ["Location", "[Location]"], ["Clearance For", "[Purpose]"],
            ["Clearance Date", "[DD/MM/YYYY]"], ["Valid Until", "[DD/MM/YYYY]"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Verification", level=2)
        doc.add_paragraph("The following checks have been conducted and found satisfactory:")
        doc.add_paragraph("Annual Property Return verification", style="List Bullet")
        doc.add_paragraph("Integrity pact compliance", style="List Bullet")
        doc.add_paragraph("No pending vigilance cases", style="List Bullet")
        doc.add_paragraph("No adverse findings in background verification", style="List Bullet")
        doc.add_heading("Authorizing Authority", level=2)
        doc.add_paragraph("Chief Vigilance Officer", style="List Bullet")
        doc.add_paragraph("Director (Personnel)", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_cvc_returns_docx(self, filename="CVC_Quarterly_Returns.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("CVC Quarterly Returns", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Central Vigilance Commission - Quarterly Compliance Report")
        doc.add_heading("Quarter: Q4 FY 2024-25 (January-March 2025)", level=1)
        doc.add_heading("1. Complaints Received", level=2)
        table = doc.add_table(rows=6, cols=3, style="Light Grid Accent 1")
        headers = ["Category", "Number", "Status"]
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
            for paragraph in table.rows[0].cells[i].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        data = [
            ["Corruption complaints", "5", "Under investigation"],
            ["Misconduct complaints", "8", "5 disposed, 3 pending"],
            ["Grievance (non-vigilance)", "12", "All disposed"],
            ["Public interest disclosure", "2", "Under investigation"],
            ["Anonymous complaints", "3", "Under preliminary inquiry"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("2. Vigilance Cases", level=2)
        doc.add_paragraph("Pending cases from previous quarter: 3", style="List Bullet")
        doc.add_paragraph("New cases during the quarter: 5", style="List Bullet")
        doc.add_paragraph("Cases disposed during the quarter: 4", style="List Bullet")
        doc.add_paragraph("Cases pending at quarter end: 4", style="List Bullet")
        doc.add_heading("3. Surprise Checks", level=2)
        doc.add_paragraph("Total surprise checks conducted: 15", style="List Bullet")
        doc.add_paragraph("Checks in cash handling: 8", style="List Bullet")
        doc.add_paragraph("Checks in stores/inventory: 5", style="List Bullet")
        doc.save(filepath)
        return filepath

    def generate_apr_return_docx(self, filename="Annual_Property_Return.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Annual Property Return (APR)", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Filed under CCS (Conduct) Rules / Company Policy")
        doc.add_heading("Officer Details", level=1)
        table = doc.add_table(rows=7, cols=2, style="Light Grid Accent 1")
        details = [
            ["Name", "[Officer Name]"], ["Designation", "[Designation]"],
            ["Employee ID", "[XXXXXX]"], ["Department", "[Department]"],
            ["Date of Joining", "[DD/MM/YYYY]"], ["Date of Filing APR", "31/12/2024"],
            ["Financial Year", "2024-25"],
        ]
        for r_idx, (label, value) in enumerate(details):
            table.rows[r_idx].cells[0].text = label
            table.rows[r_idx].cells[1].text = value
            for paragraph in table.rows[r_idx].cells[0].paragraphs:
                for run in paragraph.runs:
                    run.bold = True
        doc.add_heading("Immovable Property Details", level=1)
        table2 = doc.add_table(rows=4, cols=5, style="Light Grid Accent 1")
        headers = ["S.No", "Description", "Location", "Date of Acquisition", "Value (Rs)"]
        for i, h in enumerate(headers):
            table2.rows[0].cells[i].text = h
        data = [
            ["1", "Residential Flat", "Mumbai", "15/06/2015", "1,50,00,000"],
            ["2", "Agricultural Land", "Pune", "01/04/2018", "50,00,000"],
            ["3", "Commercial Shop", "Delhi", "20/12/2020", "80,00,000"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table2.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Movable Property Details", level=1)
        table3 = doc.add_table(rows=4, cols=4, style="Light Grid Accent 1")
        headers = ["S.No", "Description", "Date of Acquisition", "Value (Rs)"]
        for i, h in enumerate(headers):
            table3.rows[0].cells[i].text = h
        data = [
            ["1", "Car (SUV)", "01/03/2022", "25,00,000"],
            ["2", "Gold Jewelry", "25/11/2019", "8,00,000"],
            ["3", "Fixed Deposits", "Various", "45,00,000"],
        ]
        for r_idx, row_data in enumerate(data):
            for c_idx, val in enumerate(row_data):
                table3.rows[r_idx + 1].cells[c_idx].text = val
        doc.add_heading("Declaration", level=2)
        doc.add_paragraph("I hereby declare that the above information is true and complete to the best of my knowledge. I have not acquired any property out of proportion to my known sources of income.")
        doc.save(filepath)
        return filepath

    def generate_independent_directors_code_docx(self, filename="Code_Independent_Directors_Schedule_IV.docx"):
        filepath = self._get_output_path(filename)
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        title = doc.add_heading("Code for Independent Directors", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Schedule IV of the Companies Act, 2013 and Regulation 25 of SEBI LODR Regulations, 2015")
        doc.add_heading("1. Role and Responsibilities", level=2)
        doc.add_paragraph("Bring independent judgment on issues of strategy, performance, resources", style="List Bullet")
        doc.add_paragraph("Satisfy themselves on the integrity of financial information", style="List Bullet")
        doc.add_paragraph("Satisfy themselves that adequate systems of internal control exist", style="List Bullet")
        doc.add_paragraph("Safeguard the interests of minority shareholders", style="List Bullet")
        doc.add_heading("2. Responsibilities to the Company", level=2)
        doc.add_paragraph("Undertake appropriate induction and regularly update skills", style="List Bullet")
        doc.add_paragraph("Seek appropriate professional advice", style="List Bullet")
        doc.add_paragraph("Attend at least one meeting per quarter of independent directors", style="List Bullet")
        doc.add_heading("3. Manner of Appointment", level=2)
        doc.add_paragraph("Appointment by shareholders at general meeting", style="List Bullet")
        doc.add_paragraph("Term: 5 years, renewable up to 2 terms", style="List Bullet")
        doc.add_paragraph("Maximum age: 70 years", style="List Bullet")
        doc.add_heading("4. Separate Meetings", level=2)
        doc.add_paragraph("Independent directors shall meet at least once in a year without presence of non-independent directors and management.")
        doc.save(filepath)
        return filepath

    def generate_vigilance_review_pptx(self, filename="Vigilance_Review_Presentation.pptx"):
        filepath = self._get_output_path(filename)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(2.5), PptxInches(12), PptxInches(1.5))
        title.text = "Vigilance Review - FY 2024-25"
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = 1
            for run in paragraph.runs:
                run.font.size = PptxInches(0.7)
                run.bold = True
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Vigilance Dashboard"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Complaints Received: 25\n- Corruption: 5 (20%)\n- Misconduct: 8 (32%)\n- Grievance: 12 (48%)\n\nCases Disposed: 9\nCases Pending: 16\n\nSurprise Checks: 15\n- Cash handling: 8\n- Stores/inventory: 5"
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(0.8))
        title_box.text = "Anti-Bribery & Compliance"
        for paragraph in title_box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = PptxInches(0.5)
                run.bold = True
        content = slide.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        tf = content.text_frame
        tf.text = "Anti-Bribery Compliance:\n- ISO 37001 certified\n- Gift register: Quarterly review\n- Integrity pact: 100%\n\nAPR Filing:\n- 100% submitted\n- Verification: 100%\n\nCVC Returns:\n- Filed on time\n- No adverse observations"
        prs.save(filepath)
        return filepath
