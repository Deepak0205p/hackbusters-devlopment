"""
MRPL & ONGC Enterprise Legal & Agreements Deliverables Synthesis Engine
Implements authentic Maharatna PSU legal document generators:
- Word (.docx): Land Lease, Dealership, Technology License, Shareholders, MOU/JV,
                Crude Purchase, PPA, Pipeline ROW, Indemnity, Arbitration
- PowerPoint (.pptx): 16:9 Legal & Agreements Review Presentation
"""

import os
import time
from typing import Optional
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))


def ensure_output_directories():
    os.makedirs(os.path.join(OUTPUT_DIR, "docx"), exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, "pptx"), exist_ok=True)


class ExtendedGenerators_Legal:
    """
    Legal & Agreements Document Generator Engine for MRPL & ONGC.
    """

    def __init__(self):
        ensure_output_directories()

    # =========================================================================
    # 1. LAND LEASE AGREEMENT
    # =========================================================================

    def generate_land_lease_agreement_docx(
        self,
        filename: str = "MRPL_Land_Lease_Agreement.docx",
        lessor: str = "Mangalore Special Economic Zone Ltd.",
        lessee: str = "Mangalore Refinery and Petrochemicals Limited",
        area: str = "12.5 Acres (Survey No. 45/2, Boloor Village, Mangalore Taluk)",
        term: str = "30 Years (Commencement: 01-Apr-2026, Expiry: 31-Mar-2056)",
        annual_rent: str = "₹18,50,000/- (Rupees Eighteen Lakhs Fifty Thousand Only) per annum"
    ) -> str:
        """Generates Land Lease Agreement for MRPL/ONGC facility land acquisition."""
        doc = Document()
        title = doc.add_heading("LAND LEASE AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("This Deed of Lease is made and executed on this " + time.strftime("%d") + " day of " + time.strftime("%B, %Y"))
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Lessor:", lessor),
            ("Lessee:", lessee)
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Demised Premises & Area", level=1)
        doc.add_paragraph(f"The Lessor hereby grants and demises unto the Lessee the land measuring {area}, "
                          "together with all improvements, structures, and appurtenances thereunto belonging, "
                          "for the purpose of establishment and operation of refinery and petrochemical facilities.")

        doc.add_heading("2. Lease Term", level=1)
        doc.add_paragraph(f"The term of this Lease shall be {term}, unless sooner terminated in accordance with "
                          "the provisions herein contained.")

        doc.add_heading("3. Rental Consideration", level=1)
        doc.add_paragraph(f"3.1 The Lessee shall pay to the Lessor an annual rent of {annual_rent}, "
                          "payable in advance in four equal quarterly instalments on or before the 1st day of April, "
                          "July, October, and January of each year.")
        doc.add_paragraph("3.2 The rent shall be subject to an escalation of 5% (five percent) after every "
                          "five years of the Lease Term, or as mutually agreed upon in writing.")

        doc.add_heading("4. Permitted Use", level=1)
        doc.add_paragraph("The Lessee shall use the demised land exclusively for the construction, installation, "
                          "and operation of refinery units, petrochemical plants, tank farms, utilities, "
                          "and ancillary facilities connected with the refinery operations of ONGC-MRPL.")

        doc.add_heading("5. Obligations of the Lessee", level=1)
        doc.add_paragraph("5.1 The Lessee shall, at its own cost, maintain all environmental clearances "
                          "and statutory permissions required for operations.")
        doc.add_paragraph("5.2 The Lessee shall keep the demised premises in good and tenantable condition "
                          "and shall indemnify the Lessor against all losses arising from negligence or misuse.")
        doc.add_paragraph("5.3 The Lessee shall not assign, sub-let, or transfer any interest under this "
                          "Agreement without the prior written consent of the Lessor.")

        doc.add_heading("6. Termination", level=1)
        doc.add_paragraph("Either party may terminate this Agreement by giving not less than twelve (12) months "
                          "prior written notice to the other party. Upon termination, the Lessee shall vacate "
                          "the demised premises and restore the land to its original condition, reasonable wear "
                          "and tear excepted.")

        doc.add_heading("7. Governing Law & Dispute Resolution", level=1)
        doc.add_paragraph("This Agreement shall be governed by and construed in accordance with the laws of India. "
                          "Any dispute arising hereunder shall be referred to a sole arbitrator appointed by "
                          "mutual consent, and the arbitration shall be conducted in accordance with the "
                          "Arbitration and Conciliation Act, 1996.")

        doc.add_heading("8. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of the Lessor:", "For and on behalf of the Lessee:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 2. DEALERSHIP / DISTRIBUTORSHIP AGREEMENT
    # =========================================================================

    def generate_dealership_agreement_docx(
        self,
        filename: str = "MRPL_Dealership_Agreement.docx",
        dealer: str = "Bharat Petro-Marketing Solutions Pvt. Ltd.",
        territory: str = "State of Karnataka (Districts: Dakshina Kannada, Udupi, Chikkamagaluru, Hassan)",
        products: str = "MRPL Motor Spirit (MS), High Speed Diesel (HSD), Aviation Turbine Fuel (ATF), "
                        "Liquefied Petroleum Gas (LPG), Bitumen, and Petrochemical products"
    ) -> str:
        """Generates Dealership/Distributorship Agreement for MRPL fuel and petrochemical products."""
        doc = Document()
        title = doc.add_heading("DEALERSHIP / DISTRIBUTORSHIP AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("This Agreement is entered into on " + time.strftime("%d") + " day of " + time.strftime("%B, %Y"))
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Principal / Supplier:", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "A subsidiary of Oil and Natural Gas Corporation Limited (ONGC)"),
            ("Dealer / Distributor:", dealer)
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Appointment & Territory", level=1)
        doc.add_paragraph(f"1.1 The Principal hereby appoints the Dealer as an authorised dealer for the "
                          f"sale and distribution of the Products within the territory of {territory}.")
        doc.add_paragraph("1.2 This appointment is non-exclusive unless otherwise specified in writing.")

        doc.add_heading("2. Products & Specifications", level=1)
        doc.add_paragraph(f"The products covered under this Agreement are: {products}. "
                          "All products shall conform to the relevant Indian Standards (IS) and "
                          "Bureau of Indian Standards (BIS) specifications as applicable.")

        doc.add_heading("3. Ordering & Supply", level=1)
        doc.add_paragraph("3.1 The Dealer shall place purchase orders through the MRPL Order Management "
                          "System (OMS) or as mutually agreed.")
        doc.add_paragraph("3.2 The Principal shall endeavour to dispatch confirmed orders within "
                          "72 working hours from the nearest MRPL terminal or depot.")

        doc.add_heading("4. Pricing & Payment", level=1)
        doc.add_paragraph("4.1 The selling price of each product shall be determined by the Principal "
                          "from time to time, inclusive of applicable excise duty, GST, and other levies.")
        doc.add_paragraph("4.2 Payment shall be made by the Dealer within 15 (fifteen) days of the "
                          "date of invoice by way of electronic fund transfer (NEFT/RTGS).")

        doc.add_heading("5. Performance Obligations", level=1)
        doc.add_paragraph("5.1 The Dealer shall achieve minimum monthly off-take volumes as per "
                          "Schedule-A attached hereto.")
        doc.add_paragraph("5.2 The Dealer shall maintain adequate storage, dispensing, and "
                          "quality control facilities at its dealership outlets.")

        doc.add_heading("6. Term & Renewal", level=1)
        doc.add_paragraph("This Agreement shall be valid for a period of 5 (five) years from the "
                          "date of execution, renewable for successive periods of 3 (three) years "
                          "each, subject to satisfactory performance and mutual consent.")

        doc.add_heading("7. Termination", level=1)
        doc.add_paragraph("Either party may terminate this Agreement by giving not less than 90 "
                          "days' written notice. The Principal reserves the right to terminate "
                          "immediately in case of material breach, including default in payment, "
                          "quality deviation, or regulatory non-compliance.")

        doc.add_heading("8. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of India. Disputes shall "
                          "be subject to the exclusive jurisdiction of the courts at Mangalore, Karnataka.")

        doc.add_heading("9. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of MRPL:", f"For and on behalf of\n{dealer}:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 3. TECHNOLOGY LICENSE AGREEMENT
    # =========================================================================

    def generate_technology_license_docx(
        self,
        filename: str = "MRPL_Technology_License_Agreement.docx",
        licensor: str = "UOP LLC (A Honeywell Company), Des Plaines, Illinois, USA",
        technology: str = "UOP Penex Process for Isomerization of Light Naphtha (C5/C6 Fraction) "
                         "and UOP Merox Process for Mercaptan Oxidation of LPG and Naphtha",
        license_fee: str = "USD 4,500,000/- (US Dollars Four Million Five Hundred Thousand) "
                           "as lump-sum technology license fee plus 2.5% of net product sales revenue "
                           "as running royalty"
    ) -> str:
        """Generates Technology License Agreement for proprietary refinery processes."""
        doc = Document()
        title = doc.add_heading("TECHNOLOGY LICENSE AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED AS OF " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Licensor:", licensor),
            ("Licensee:", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "P.O. Mangalore Refinery, Mangalore – 575030, Karnataka, India")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Grant of License", level=1)
        doc.add_paragraph(f"1.1 The Licensor hereby grants to the Licensee the non-exclusive, non-transferable "
                          f"right and license to use the following proprietary technology: {technology}.")
        doc.add_paragraph("1.2 This License includes the right to use the Licensor's process design, "
                          "engineering specifications, catalyst formulations, operating manuals, and "
                          "technical know-how required for the design, construction, and operation of "
                          "the licensed unit(s).")

        doc.add_heading("2. Technology Fee Structure", level=1)
        doc.add_paragraph(f"2.1 License Fee: {license_fee}.")
        doc.add_paragraph("2.2 The Lump-sum Technology License Fee shall be payable in three (3) "
                          "instalments: (a) 30% upon execution of this Agreement; "
                          "(b) 40% upon Licensor's delivery of Front-End Engineering Design (FEED) "
                          "package; (c) 30% upon successful Performance Guarantee Test completion.")
        doc.add_paragraph("2.3 Running Royalty: 2.5% of net ex-factory sales revenue of products "
                          "manufactured using the licensed technology, payable quarterly in arrears.")

        doc.add_heading("3. Technical Services", level=1)
        doc.add_paragraph("3.1 The Licensor shall provide technical advisory services during the "
                          "detailed engineering, construction, pre-commissioning, and commissioning phases.")
        doc.add_paragraph("3.2 The Licensor shall guarantee the following performance metrics: "
                          "Isomerate RON ≥ 82 (Research Octane Number); Merox sweetened LPG mercaptan "
                          "sulphur ≤ 0.25 wt%.")

        doc.add_heading("4. Confidentiality", level=1)
        doc.add_paragraph("The Licensee shall maintain the confidentiality of all technical data, "
                          "process parameters, and proprietary information disclosed under this Agreement "
                          "for a period of 15 years from the date of expiration or termination.")

        doc.add_heading("5. Indemnification", level=1)
        doc.add_paragraph("The Licensor shall indemnify the Licensee against any third-party patent "
                          "infringement claims arising from the use of the licensed technology.")

        doc.add_heading("6. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of the State of Illinois, "
                          "United States of America, without regard to its conflict of laws principles. "
                          "Disputes shall be resolved by binding arbitration under ICC rules.")

        doc.add_heading("7. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of the Licensor:", "For and on behalf of the Licensee:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 4. SHAREHOLDERS AGREEMENT
    # =========================================================================

    def generate_shareholders_agreement_docx(
        self,
        filename: str = "MRPL_Shareholders_Agreement.docx",
        shareholding_pattern: str = ("Oil and Natural Gas Corporation Limited (ONGC): 51.00%\n"
                                     "MRPL: 12.32%\n"
                                     "Hindustan Petroleum Corporation Limited (HPCL): 16.68%\n"
                                     "Oil India Limited (OIL): 10.00%\n"
                                     "Public / Institutional Holding: 10.00%"),
        board_composition: str = ("Total Board Strength: 12 Directors\n"
                                  "ONGC Nominees: 6 Directors (including Chairman)\n"
                                  "HPCL Nominees: 2 Directors\n"
                                  "OIL Nominee: 1 Director\n"
                                  "Independent Directors: 3 Directors (as per SEBI LODR)")
    ) -> str:
        """Generates Shareholders Agreement with shareholding pattern and board composition."""
        doc = Document()
        title = doc.add_heading("SHAREHOLDERS' AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        doc.add_paragraph("1. Oil and Natural Gas Corporation Limited (ONGC)")
        doc.add_paragraph("2. Mangalore Refinery and Petrochemicals Limited (MRPL)")
        doc.add_paragraph("3. Hindustan Petroleum Corporation Limited (HPCL)")
        doc.add_paragraph("4. Oil India Limited (OIL)")
        doc.add_paragraph("(Collectively referred to as the \"Shareholders\")")

        doc.add_heading("1. Shareholding Pattern", level=1)
        doc.add_paragraph("The issued and paid-up share capital of MRPL shall be held in the "
                          "following proportions:")
        doc.add_paragraph(shareholding_pattern)

        doc.add_heading("2. Board of Directors Composition", level=1)
        doc.add_paragraph("The Board of Directors of MRPL shall comprise the following:")
        doc.add_paragraph(board_composition)

        doc.add_heading("3. Reserved Matters", level=1)
        doc.add_paragraph("The following matters shall require the affirmative vote of shareholders "
                          "holding not less than 75% of the issued share capital:")
        doc.add_paragraph("a) Amendment of the Memorandum or Articles of Association of MRPL.")
        doc.add_paragraph("b) Any change in the capital structure, including issuance of new equity shares.")
        doc.add_paragraph("c) Approval of annual capital expenditure exceeding ₹500 Crores.")
        doc.add_paragraph("d) Declaration of dividend exceeding the statutory minimum.")
        doc.add_paragraph("e) Appointment or removal of the Chief Financial Officer or Company Secretary.")

        doc.add_heading("4. Transfer Restrictions", level=1)
        doc.add_paragraph("4.1 No Shareholder shall transfer its shares without first offering them "
                          "to the other Shareholders on a pro-rata basis at fair market value determined "
                          "by an independent valuer.")
        doc.add_paragraph("4.2 Tag-Along Right: If any Shareholder proposes to transfer its shares "
                          "to a third party, the other Shareholders shall have the right to participate "
                          "in the sale on the same terms and conditions.")

        doc.add_heading("5. Dividend Policy", level=1)
        doc.add_paragraph("The Shareholders agree that a minimum dividend payout ratio of 30% "
                          "of the net profit shall be maintained, subject to statutory requirements "
                          "and reserves as required under the Companies Act, 2013.")

        doc.add_heading("6. Governing Law & Dispute Resolution", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of India. Disputes shall "
                          "be referred to a three-member arbitration panel with one arbitrator appointed "
                          "by each party and the third by mutual agreement, conducted under ICC rules.")

        doc.add_heading("7. Signatures", level=1)
        sig_table = doc.add_table(rows=4, cols=3)
        sigs = [
            ("ONGC:", "HPCL:", "OIL:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:"),
            ("MRPL:", "", ""),
            ("_________________________\nAuthorised Signatory\nDate:", "", "")
        ]
        for r_idx, row in enumerate(sigs):
            for c_idx, val in enumerate(row):
                sig_table.rows[r_idx].cells[c_idx].text = val

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 5. MOU / JOINT VENTURE AGREEMENT
    # =========================================================================

    def generate_mou_jv_agreement_docx(
        self,
        filename: str = "MRPL_JV_Agreement_MOU.docx",
        partners: str = ("Party A: Oil and Natural Gas Corporation Limited (ONGC)\n"
                         "Party B: Petronet LNG Limited (PLL)\n"
                         "Party C: Bharat Petroleum Corporation Limited (BPCL)"),
        spv: str = "ONGC-MRPL LNG Terminal SPV Private Limited",
        capital_structure: str = ("Equity Contribution:\n"
                                  "ONGC: 50% (₹2,500 Crores)\n"
                                  "PLL: 30% (₹1,500 Crores)\n"
                                  "BPCL: 20% (₹1,000 Crores)\n"
                                  "Total Project Cost: ₹5,000 Crores\n"
                                  "Debt-Equity Ratio: 70:30")
    ) -> str:
        """Generates MOU/Joint Venture Agreement for a Special Purpose Vehicle (SPV)."""
        doc = Document()
        title = doc.add_heading("MEMORANDUM OF UNDERSTANDING\n"
                                "JOINT VENTURE AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("ENTERED INTO ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        doc.add_paragraph(partners)
        doc.add_paragraph(f"\n(Hereinafter collectively referred to as the \"JV Partners\" or \"Promoters\")")

        doc.add_heading("1. Formation of Special Purpose Vehicle", level=1)
        doc.add_paragraph(f"1.1 The JV Partners hereby agree to incorporate a Special Purpose Vehicle "
                          f"(SPV) under the Companies Act, 2013, with the name: {spv}.")
        doc.add_paragraph("1.2 The SPV shall be registered under the Registrar of Companies, Karnataka, "
                          "and shall have its registered office in Mangalore, Karnataka.")

        doc.add_heading("2. Capital Structure", level=1)
        doc.add_paragraph(capital_structure)

        doc.add_heading("3. Object of the Joint Venture", level=1)
        doc.add_paragraph("The SPV shall be established for the following purposes:")
        doc.add_paragraph("a) Design, construction, commissioning, and operation of an LNG "
                          "Re-gasification Terminal at Mangalore SEZ.")
        doc.add_paragraph("b) Procurement of LNG on long-term and spot basis for supply to "
                          "the JV Partners and third-party customers.")
        doc.add_paragraph("c) Development of LNG bunkering infrastructure at Mangalore Port.")

        doc.add_heading("4. Governance & Management", level=1)
        doc.add_paragraph("4.1 The SPV shall be managed by a Board of Directors comprising nominees "
                          "of each JV Partner in proportion to their shareholding.")
        doc.add_paragraph("4.2 The Managing Director shall be nominated by ONGC. The CFO shall be "
                          "nominated by PLL. The Technical Director shall be nominated by BPCL.")

        doc.add_heading("5. Business Plan & Timeline", level=1)
        doc.add_paragraph("5.1 The SPV shall finalise its Detailed Project Report (DPR) within "
                          "6 months of incorporation.")
        doc.add_paragraph("5.2 Financial closure shall be achieved within 12 months of incorporation.")
        doc.add_paragraph("5.3 Commercial Operations Date (COD) shall be within 42 months of "
                          "financial closure.")

        doc.add_heading("6. Buy-Sell Arrangement", level=1)
        doc.add_paragraph("Each JV Partner shall have the right to off-take its proportionate share "
                          "of LNG throughput from the Terminal, subject to nomination of quantities "
                          "not less than 90 days prior to delivery.")

        doc.add_heading("7. Exit Provisions", level=1)
        doc.add_paragraph("7.1 No JV Partner may transfer its shares without the prior written "
                          "consent of the other JV Partners.")
        doc.add_paragraph("7.2 Tag-along and drag-along rights shall apply as per the terms "
                          "of the Shareholders' Agreement to be executed concurrently.")

        doc.add_heading("8. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of India. Disputes "
                          "shall be resolved by arbitration under ICC rules with the seat "
                          "of arbitration in New Delhi.")

        doc.add_heading("9. Signatures", level=1)
        sig_table = doc.add_table(rows=3, cols=3)
        sigs = [
            ("ONGC:", "PLL:", "BPCL:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:"),
        ]
        for r_idx, row in enumerate(sigs):
            for c_idx, val in enumerate(row):
                sig_table.rows[r_idx].cells[c_idx].text = val

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 6. CRUDE OIL PURCHASE & TRANSPORTATION AGREEMENT
    # =========================================================================

    def generate_crude_purchase_agreement_docx(
        self,
        filename: str = "MRPL_Crude_Purchase_Agreement.docx",
        supplier: str = "Saudi Aramco Trading Company (SATCO), Dhahran, Kingdom of Saudi Arabia",
        grade: str = "Arab Light Crude Oil (API Gravity: 33.5°, Sulphur: 1.78 wt%)",
        pricing: str = "FOB Ras Tanura at Platts OMD Crude Oil Marker Price (AOG) minus "
                       "a contractual discount of USD 0.85/bbl, with freight and insurance "
                       "on CFR Mangalore basis as per BAF and CAF adjustments"
    ) -> str:
        """Generates Crude Oil Purchase & Transportation Agreement."""
        doc = Document()
        title = doc.add_heading("CRUDE OIL PURCHASE & TRANSPORTATION AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Seller:", supplier),
            ("Buyer:", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "A subsidiary of ONGC, Mangalore SEZ, Karnataka, India")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Crude Oil Grade & Specifications", level=1)
        doc.add_paragraph(f"1.1 The Seller shall supply and the Buyer shall purchase the following crude oil:\n"
                          f"Grade: {grade}")
        doc.add_paragraph("1.2 The crude oil shall conform to the Seller's published specifications "
                          "and shall be delivered in a merchantable condition, free from water, "
                          "sediment, and other contaminants.")

        doc.add_heading("2. Volume & Delivery", level=1)
        doc.add_paragraph("2.1 Contract Quantity: 120,000 barrels per day (bpd) ± 10% at Seller's option.")
        doc.add_paragraph("2.2 Delivery: CFR (Cost and Freight) Mangalore Port, India.")
        doc.add_paragraph("2.3 The Seller shall deliver the crude oil in VLCC (Very Large Crude Carrier) "
                          "shipments of approximately 2,000,000 barrels per cargo.")
        doc.add_paragraph("2.4 Nominations shall be made not less than 60 days prior to the "
                          "first day of the loading window.")

        doc.add_heading("3. Pricing Mechanism", level=1)
        doc.add_paragraph(f"{pricing}.")
        doc.add_paragraph("3.2 Invoice Price: Based on Platts OMD assessment averaged over the "
                          "loading period, plus/minus the contractual differential and "
                          "applicable freight and insurance charges.")

        doc.add_heading("4. Transportation & Marine Insurance", level=1)
        doc.add_paragraph("4.1 The Seller shall arrange and bear the cost of marine transportation "
                          "from Ras Tanura to Mangalore Port.")
        doc.add_paragraph("4.2 Marine cargo insurance shall be effected by the Seller on 'Institute "
                          "Cargo Clauses (A)' basis for 110% of CIF value.")

        doc.add_heading("5. Quality Assurance", level=1)
        doc.add_paragraph("5.1 Independent inspection shall be conducted by SGS or Bureau Veritas "
                          "at the port of loading and discharge.")
        doc.add_paragraph("5.2 The average of the loading and discharge measurements shall be "
                          "final and binding for quantity and quality purposes.")

        doc.add_heading("6. Payment", level=1)
        doc.add_paragraph("Payment shall be made by the Buyer within 30 days of receipt of "
                          "the Seller's invoice, by telegraphic transfer in US Dollars "
                          "to the Seller's designated bank account.")

        doc.add_heading("7. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of England and Wales. "
                          "Disputes shall be resolved by arbitration under LCIA rules.")

        doc.add_heading("8. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of the Seller:", "For and on behalf of the Buyer:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 7. POWER PURCHASE AGREEMENT (PPA)
    # =========================================================================

    def generate_ppa_agreement_docx(
        self,
        filename: str = "MRPL_Power_Purchase_Agreement.docx",
        seller: str = "MRPL Cogeneration Power Plant (2x150 MW CFBC Units)",
        capacity: str = "Installed Capacity: 300 MW (2 x 150 MW Circulating Fluidised Bed Combustion boilers). "
                        "Contracted Capacity for internal consumption: 220 MW (73.3% plant load factor).",
        tariff: str = "Levelised Tariff: ₹5.85/kWh (Composite tariff comprising capacity charge of "
                      "₹2.15/kWh and energy charge of ₹3.70/kWh). Tariff subject to annual escalation "
                      "of 3.5% linked to CPI index as per CERC norms.",
        term: str = "25 years from the Commercial Operations Date (COD) of the power plant, "
                    "with provision for extension by mutual agreement"
    ) -> str:
        """Generates Power Purchase Agreement (PPA) for captive/co-generated power."""
        doc = Document()
        title = doc.add_heading("POWER PURCHASE AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Seller (Generator):", seller),
            ("Buyer (Off-taker):", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "For captive consumption at its refinery complex, Mangalore SEZ, Karnataka")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Capacity & Availability", level=1)
        doc.add_paragraph(f"{capacity}")
        doc.add_paragraph("1.2 The Seller shall maintain the plant at a minimum Plant Availability "
                          "Factor of 95% on a yearly basis, subject to planned and unplanned shutdowns.")

        doc.add_heading("2. Tariff & Payment", level=1)
        doc.add_paragraph(f"{tariff}")
        doc.add_paragraph("2.2 Monthly billing: The Seller shall raise monthly invoices based on "
                          "metered energy offtake, and the Buyer shall make payment within "
                          "15 days of receipt of invoice.")
        doc.add_paragraph("2.3 Late payment: Interest at 2% above SBI MCLR shall apply on "
                          "overdue amounts.")

        doc.add_heading("3. Term of Agreement", level=1)
        doc.add_paragraph(f"This Agreement shall be effective for a period of {term}.")

        doc.add_heading("4. Metering & Testing", level=1)
        doc.add_paragraph("4.1 Export and import metering shall be done through CT/PT operated "
                          "electronic energy meters at the 110 kV switchyard bus-coupler.")
        doc.add_paragraph("4.2 Both parties shall have the right to audit the metering system "
                          "annually, and the cost of such audit shall be borne by the requesting party.")

        doc.add_heading("5. Force Majeure", level=1)
        doc.add_paragraph("Neither party shall be liable for failure to perform its obligations "
                          "under this Agreement to the extent that such failure is caused by "
                          "Force Majeure events, including natural disasters, war, civil unrest, "
                          "government orders, or grid failure attributable to the State Transmission Utility.")

        doc.add_heading("6. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the Electricity Act, 2003, "
                          "and the Indian Contract Act, 1872. Disputes shall be resolved by "
                          "arbitration under the Arbitration and Conciliation Act, 1996.")

        doc.add_heading("7. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of the Seller:", "For and on behalf of the Buyer:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 8. PIPELINE RIGHT-OF-WAY (ROW) AGREEMENT
    # =========================================================================

    def generate_pipeline_row_agreement_docx(
        self,
        filename: str = "MRPL_Pipeline_ROW_Agreement.docx",
        pipeline: str = "MRPL–ONGC Crude Oil Cross-Country Pipeline\n"
                        "Pipeline No: MRPL-CP-01 | Diameter: 24\" (610 mm) | MOP: 85 bar(g) | "
                        "Material: API 5L X65 PSL-2",
        length: str = "Total Pipeline Length: 48.5 km (Km 0.00 at MRPL Refinery Gate to "
                      "Km 48.50 at ONGC Mangalore Beach Terminal)",
        compensation: str = "One-time compensation: ₹15,000/- per running metre of pipeline ROW width "
                            "(ROW Width: 30 metres = 15m on each side of pipeline centreline). "
                            "Annual recurring compensation: ₹2,500/- per running metre, "
                            "escalated at 5% every 3 years."
    ) -> str:
        """Generates Pipeline Right-of-Way (ROW) Agreement for cross-country pipelines."""
        doc = Document()
        title = doc.add_heading("PIPELINE RIGHT-OF-WAY (ROW) AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Pipeline Owner:", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "A subsidiary of ONGC, Mangalore SEZ, Karnataka"),
            ("Land Owner / Grantor:", "Karnataka Industrial Areas Development Board (KIADB)\n"
             "and individual landowners as listed in Schedule-I")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Pipeline & Right-of-Way Details", level=1)
        doc.add_paragraph(f"Pipeline: {pipeline}")
        doc.add_paragraph(f"Right-of-Way: {length}")
        doc.add_paragraph("ROW Width: 30 metres (15 metres on either side of pipeline centreline)")

        doc.add_heading("2. Grant of Right-of-Way", level=1)
        doc.add_paragraph("2.1 The Grantor hereby grants to the Pipeline Owner the perpetual, "
                          "irrevocable right-of-way across the land parcels described in Schedule-I "
                          "for the laying, construction, operation, maintenance, repair, replacement, "
                          "and abandonment of the pipeline.")
        doc.add_paragraph("2.2 The Pipeline Owner shall have the right to access the ROW at "
                          "all times for inspection, maintenance, and emergency repair purposes.")

        doc.add_heading("3. Compensation", level=1)
        doc.add_paragraph(f"{compensation}")

        doc.add_heading("4. Obligations of the Pipeline Owner", level=1)
        doc.add_paragraph("4.1 The Pipeline Owner shall restore the land to its original condition "
                          "immediately after completion of construction activities.")
        doc.add_paragraph("4.2 The Pipeline Owner shall maintain the pipeline at all times in "
                          "a safe and operable condition, in compliance with PNGRB regulations "
                          "and Petroleum and Natural Gas Regulatory Board (Pipeline) Regulations, 2006.")
        doc.add_paragraph("4.3 The Pipeline Owner shall carry comprehensive third-party liability "
                          "insurance coverage of not less than ₹500 Crores per incident.")

        doc.add_heading("5. Restriction on Grantor", level=1)
        doc.add_paragraph("5.1 The Grantor shall not construct any permanent structure within "
                          "the ROW without the prior written consent of the Pipeline Owner.")
        doc.add_paragraph("5.2 No heavy construction equipment or earthmoving machinery "
                          "shall be operated within the ROW without giving 48 hours advance notice "
                          "to the Pipeline Owner.")

        doc.add_heading("6. Governing Law", level=1)
        doc.add_paragraph("This Agreement shall be governed by the laws of India and the "
                          "Petroleum and Natural Gas Regulatory Board Act, 2006. Disputes "
                          "shall be resolved by arbitration under ICC rules.")

        doc.add_heading("7. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("For and on behalf of MRPL:", "For and on behalf of KIADB / Grantor:"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 9. INDEMNITY BOND
    # =========================================================================

    def generate_indemnity_bond_docx(
        self,
        filename: str = "MRPL_Indemnity_Bond.docx",
        indemnitor: str = "Mangalore Refinery and Petrochemicals Limited (MRPL), through its "
                          "Authorised Signatory",
        indemnity_terms: str = "Irrevocable and unconditional indemnity against all losses, damages, "
                               "claims, demands, actions, suits, costs, charges, and expenses of whatever "
                               "nature arising out of or in connection with the performance of the "
                               "Contract Work Order No. MRPL/WO/MAINT/2026/0847 dated 01-Apr-2026",
        insurance: str = "The Indemnitor shall maintain comprehensive insurance coverage as follows:\n"
                        "(a) Contractor's All Risk (CAR) Policy: ₹25 Crores\n"
                        "(b) Workmen's Compensation Policy: ₹50 Lakhs per worker\n"
                        "(c) Third-Party Liability Policy: ₹100 Crores\n"
                        "(d) Professional Indemnity Policy: ₹15 Crores"
    ) -> str:
        """Generates Indemnity Bond for contractor/subcontractor obligations."""
        doc = Document()
        title = doc.add_heading("INDEMNITY BOND", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTED ON " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        doc.add_heading("BETWEEN", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Indemnitor:", indemnitor),
            ("Indemnitee:", "Mangalore Refinery and Petrochemicals Limited (MRPL)\n"
             "A subsidiary of ONGC, Mangalore SEZ, Karnataka")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Scope of Indemnity", level=1)
        doc.add_paragraph(indemnity_terms)

        doc.add_heading("2. Undertaking of the Indemnitor", level=1)
        doc.add_paragraph("2.1 The Indemnitor hereby agrees to indemnify, defend, and hold "
                          "harmless the Indemnitee, its directors, officers, employees, and agents "
                          "from and against any and all losses, liabilities, damages, costs, and "
                          "expenses (including reasonable attorney's fees) arising out of or in "
                          "connection with the Contract Work.")
        doc.add_paragraph("2.2 This indemnity obligation shall survive the expiry or termination "
                          "of the Contract Work Order for a period of 7 (seven) years.")

        doc.add_heading("3. Insurance Requirements", level=1)
        doc.add_paragraph(insurance)
        doc.add_paragraph("3.2 Copies of all insurance policies shall be submitted to the "
                          "Indemnitee within 30 days of execution of this Bond.")
        doc.add_paragraph("3.3 The Indemnitor shall maintain insurance coverage throughout the "
                          "period of the Contract Work and for 3 years thereafter.")

        doc.add_heading("4. Limitation of Liability", level=1)
        doc.add_paragraph("4.1 The total aggregate liability of the Indemnitor under this Bond "
                          "shall not exceed ₹100 Crores (Rupees One Hundred Crores Only).")
        doc.add_paragraph("4.2 In no event shall the Indemnitor be liable for indirect, "
                          "consequential, incidental, or punitive damages.")

        doc.add_heading("5. Governing Law", level=1)
        doc.add_paragraph("This Bond shall be governed by the laws of India, specifically "
                          "the Indian Contract Act, 1872, and shall be subject to the "
                          "jurisdiction of the courts at Mangalore, Karnataka.")

        doc.add_heading("6. Signatures", level=1)
        sig_table = doc.add_table(rows=2, cols=2)
        sigs = [
            ("Indemnitor:", "Indemnitee (Witness):"),
            ("_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for idx, (k, v) in enumerate(sigs):
            sig_table.rows[0].cells[idx].text = k
            sig_table.rows[0].cells[idx].paragraphs[0].runs[0].font.bold = True
            sig_table.rows[1].cells[idx].text = v

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 10. ARBITRATION AWARD / DISPUTE RESOLUTION RECORD
    # =========================================================================

    def generate_arbitration_record_docx(
        self,
        filename: str = "MRPL_Arbitration_Award_Record.docx",
        reference: str = "Arbitration Reference No. MRPL/ARB/2026/007",
        claims: str = "Claim No. 1: ₹8,45,00,000/- (Rupees Eight Crores Forty-Five Lakhs Only) "
                      "on account of variation in crude oil quality parameters resulting in additional "
                      "refining costs and product yield losses.\n\n"
                      "Claim No. 2: ₹2,15,00,000/- (Rupees Two Crores Fifteen Lakhs Only) "
                      "on account of delayed shipment causing demurrage and port charges.\n\n"
                      "Counter-Claim: ₹3,80,00,000/- (Rupees Three Crores Eighty Lakhs Only) "
                      "on account of late payment interest and quality non-conformance penalty.",
        award: str = "The Arbitral Tribunal has unanimously awarded as follows:\n\n"
                     "On Claim No. 1: The Respondent shall pay the Claimant a sum of ₹6,25,00,000/- "
                     "(Rupees Six Crores Twenty-Five Lakhs Only) in full and final settlement of "
                     "this claim, after considering the quality adjustment differential.\n\n"
                     "On Claim No. 2: The Respondent shall pay the Claimant a sum of ₹1,80,00,000/- "
                     "(Rupees One Crore Eighty Lakhs Only) in full and final settlement of this claim, "
                     "considering the mitigated demurrage charges.\n\n"
                     "On Counter-Claim: The Counter-Claim is partially allowed to the extent of "
                     "₹1,25,00,000/- (Rupees One Crore Twenty-Five Lakhs Only), being the balance "
                     "late payment interest computed at SBI MCLR + 2% per annum.\n\n"
                     "Net Award Amount: ₹6,80,00,000/- (Rupees Six Crores Eighty Lakhs Only) "
                     "payable by the Respondent to the Claimant within 60 days of this Award."
    ) -> str:
        """Generates Arbitration Award / Dispute Resolution Record."""
        doc = Document()
        title = doc.add_heading("ARBITRATION AWARD\n"
                                "DISPUTE RESOLUTION RECORD", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("AWARD DATED " + time.strftime("%d").upper() + " " +
                                time.strftime("%B, %Y").upper())
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)

        meta_table = doc.add_table(rows=3, cols=2)
        meta_data = [
            ("Reference No:", reference),
            ("Seat of Arbitration:", "Mangalore, Karnataka, India"),
            ("Applicable Rules:", "Arbitration and Conciliation Act, 1996 (as amended 2015/2019)")
        ]
        for idx, (k, v) in enumerate(meta_data):
            c = meta_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("1. Parties to the Arbitration", level=1)
        party_table = doc.add_table(rows=2, cols=2)
        party_data = [
            ("Claimant:", "Mangalore Refinery and Petrochemicals Limited (MRPL)"),
            ("Respondent:", "Gulf Petro Trading LLC, Dubai, UAE")
        ]
        for idx, (k, v) in enumerate(party_data):
            c = party_table.rows[idx].cells
            c[0].text = k
            c[0].paragraphs[0].runs[0].font.bold = True
            c[1].text = v

        doc.add_heading("2. Claims & Counter-Claims", level=1)
        doc.add_paragraph(claims)

        doc.add_heading("3. Arbitral Tribunal", level=1)
        doc.add_paragraph("The Arbitral Tribunal was constituted with the following members:")
        doc.add_paragraph("Sole Arbitrator: Justice (Retd.) K. Ramaswamy, former Chief Justice "
                          "of the Karnataka High Court.")
        doc.add_paragraph("The Tribunal was duly constituted and both parties were given "
                          "adequate opportunity of hearing and submission of written arguments.")

        doc.add_heading("4. Proceedings Summary", level=1)
        doc.add_paragraph("4.1 The Claim was filed on 15-Jan-2026, and the Respondent submitted "
                          "its written statement on 15-Feb-2026.")
        doc.add_paragraph("4.2 Three rounds of hearings were conducted on 10-Mar-2026, "
                          "05-May-2026, and 20-Jul-2026.")
        doc.add_paragraph("4.3 Both parties submitted documentary evidence including "
                          "commercial invoices, quality certificates, BOL, and correspondence.")

        doc.add_heading("5. Award", level=1)
        doc.add_paragraph(award)

        doc.add_heading("6. Cost of Arbitration", level=1)
        doc.add_paragraph("The cost of arbitration, including arbitrator's fees and "
                          "administrative charges, shall be borne equally by both parties. "
                          "The Respondent shall additionally pay ₹15,00,000/- as cost of "
                          "arbitration payable to the Claimant within 60 days.")

        doc.add_heading("7. Enforcement", level=1)
        doc.add_paragraph("This Award is final and binding on both parties. It may be "
                          "enforced as a decree of a civil court in accordance with "
                          "Section 36 of the Arbitration and Conciliation Act, 1996.")

        doc.add_heading("8. Arbitrator's Signature", level=1)
        sig_table = doc.add_table(rows=2, cols=3)
        sigs = [
            ("Sole Arbitrator:", "Claimant:", "Respondent:"),
            ("_________________________\nJustice (Retd.) K. Ramaswamy\nDate:", "_________________________\nAuthorised Signatory\nDate:", "_________________________\nAuthorised Signatory\nDate:")
        ]
        for r_idx, row in enumerate(sigs):
            for c_idx, val in enumerate(row):
                sig_table.rows[r_idx].cells[c_idx].text = val
                if r_idx == 0:
                    sig_table.rows[r_idx].cells[c_idx].paragraphs[0].runs[0].font.bold = True

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    # =========================================================================
    # 11. LEGAL & AGREEMENTS REVIEW PRESENTATION (PPTX)
    # =========================================================================

    def generate_legal_review_pptx(
        self,
        filename: str = "MRPL_Legal_Agreements_Review.pptx"
    ) -> str:
        """Generates Legal & Agreements Review Presentation (16:9 widescreen)."""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        # Slide 1 - Title Slide
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL / ONGC LEGAL & AGREEMENTS REVIEW"
        s1.placeholders[1].text = (
            "Quarterly Legal Portfolio Review & Agreement Status Dashboard\n"
            "Legal & Regulatory Affairs Division — " + time.strftime("%B %Y")
        )

        # Slide 2 - Active Agreements Summary
        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Active Agreements Portfolio — Executive Summary"
        tf2 = s2.placeholders[1].text_frame
        tf2.text = "• Total Active Agreements: 127 agreements under management across MRPL & ONGC subsidiaries."
        tf2.add_paragraph().text = "• Land & Lease Agreements: 38 active leases covering 2,450 acres with ₹48.2 Cr annual outflow."
        tf2.add_paragraph().text = "• Technology License Agreements: 14 active licenses with international technology providers (UOP, Linde, CB&I)."
        tf2.add_paragraph().text = "• Crude Purchase & Transportation Agreements: 9 long-term crude supply contracts covering 85% of crude requirements."
        tf2.add_paragraph().text = "• Joint Venture & MOU Agreements: 6 active JV arrangements with ₹12,500 Cr cumulative capital commitment."
        tf2.add_paragraph().text = "• Power Purchase Agreements: 4 PPAs for captive and grid-connected renewable/cogeneration power."
        tf2.add_paragraph().text = "• Pipeline ROW Agreements: 12 ROW agreements for 485 km of cross-country pipeline network."

        # Slide 3 - Dispute Resolution & Litigation Status
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Dispute Resolution & Active Litigation Status"
        tf3 = s3.placeholders[1].text_frame
        tf3.text = "• Active Arbitration Cases: 5 pending arbitrations with total claim value of ₹142.5 Crores."
        tf3.add_paragraph().text = "• International Arbitration: 2 cases under ICC rules (Dubai & Singapore seats)."
        tf3.add_paragraph().text = "• Domestic Arbitration: 3 cases under Arbitration and Conciliation Act, 1996."
        tf3.add_paragraph().text = "• Court Litigation: 12 civil suits pending across High Courts of Karnataka, Mumbai, and Delhi."
        tf3.add_paragraph().text = "• Indemnity Claims: 7 indemnity bond invocations in process with total exposure of ₹28.3 Crores."
        tf3.add_paragraph().text = "• Awards Rendered in FY26: 3 awards in favour of MRPL with ₹6.8 Crores net recovery."

        # Slide 4 - Regulatory & Statutory Compliance
        s4 = prs.slides.add_slide(prs.slide_layouts[1])
        s4.shapes.title.text = "Regulatory & Statutory Compliance Overview"
        tf4 = s4.placeholders[1].text_frame
        tf4.text = "• PNGRB Pipeline Regulations: Full compliance with PNGRB (Pipeline) Regulations, 2006 and amendments."
        tf4.add_paragraph().text = "• Environmental Clearances: All 14 consent orders and environmental clearances renewed and valid."
        tf4.add_paragraph().text = "• SEZ Compliance: MRPL SEZ unit compliance report submitted to SEZ Authority, Karnataka."
        tf4.add_paragraph().text = "• CVC Integrity Pact: 100% coverage for all tenders exceeding ₹1 Crore."
        tf4.add_paragraph().text = "• FEMA / RBI Compliance: All cross-border technology license fees and crude purchase payments routed through AD Category-I banks."

        # Slide 5 - Key Agreements Under Review
        s5 = prs.slides.add_slide(prs.slide_layouts[1])
        s5.shapes.title.text = "Key Agreements Under Active Review"
        tf5 = s5.placeholders[1].text_frame
        tf5.text = "• Crude Purchase Agreement with Saudi Aramco: Under renewal negotiation — 5-year extension with revised pricing formula."
        tf5.add_paragraph().text = "• UOP Technology License: Running royalty audit for FY25-26 completed; pending final settlement of ₹1.8 Crores."
        tf5.add_paragraph().text = "• LNG Joint Venture SPV: Business Plan update and DPR revision for Mangalore LNG Terminal project."
        tf5.add_paragraph().text = "• Pipeline ROW Renewal: 6 ROW agreements due for renewal in FY26-27, covering 128 km of pipeline network."
        tf5.add_paragraph().text = "• Dealership Agreements: Performance review of 14 dealers with 3 agreements scheduled for termination/non-renewal."

        # Slide 6 - Action Items & Recommendations
        s6 = prs.slides.add_slide(prs.slide_layouts[1])
        s6.shapes.title.text = "Critical Action Items & Recommendations"
        tf6 = s6.placeholders[1].text_frame
        tf6.text = "1. Finalise Crude Purchase Agreement renewal with Saudi Aramco by Q2 FY27."
        tf6.add_paragraph().text = "2. Complete arbitration award enforcement for MRPL/ARB/2026/007 — ₹6.8 Crores recovery."
        tf6.add_paragraph().text = "3. Update 6 Pipeline ROW agreements before expiry to avoid operational disruption."
        tf6.add_paragraph().text = "4. Commission independent audit of UOP technology license running royalties."
        tf6.add_paragraph().text = "5. Review and strengthen indemnity bond insurance coverage limits for FY27 contracts."
        tf6.add_paragraph().text = "6. Submit consolidated legal compliance report to ONGC Board by 30-Sep-2026."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path


# =========================================================================
# CONVENIENCE FUNCTION FOR DIRECT EXECUTION
# =========================================================================

if __name__ == "__main__":
    generator = ExtendedGenerators_Legal()
    print("Generating MRPL/ONGC Legal & Agreements deliverables...")
    print(generator.generate_land_lease_agreement_docx())
    print(generator.generate_dealership_agreement_docx())
    print(generator.generate_technology_license_docx())
    print(generator.generate_shareholders_agreement_docx())
    print(generator.generate_mou_jv_agreement_docx())
    print(generator.generate_crude_purchase_agreement_docx())
    print(generator.generate_ppa_agreement_docx())
    print(generator.generate_pipeline_row_agreement_docx())
    print(generator.generate_indemnity_bond_docx())
    print(generator.generate_arbitration_record_docx())
    print(generator.generate_legal_review_pptx())
    print("All Legal & Agreements deliverables generated successfully.")
