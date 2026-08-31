import os
import time
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))


class ExtendedGenerators_QualityInfoSec:

    # ======================================================================
    # QUALITY (ISO) METHODS
    # ======================================================================

    def generate_iso9001_qms_manual_docx(self, filename="ISO9001_QMS_Manual.docx"):
        doc = Document()
        style = doc.styles["Title"]
        style.font.size = Pt(24)
        style.font.color.rgb = RGBColor(0, 51, 102)

        title = doc.add_heading("ISO 9001:2015 Quality Management System Manual", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph(f"Document Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Quality Policy", level=1)
        doc.add_paragraph(
            "We are committed to delivering refined petroleum products and petrochemicals "
            "that consistently meet or exceed customer and applicable statutory/regulatory "
            "requirements. We pursue continual improvement of our Quality Management System "
            "through risk-based thinking, leadership commitment, and engagement of all personnel."
        )
        doc.add_paragraph(
            "Our quality policy is communicated, understood, and implemented across all "
            "functional areas. Measurable quality objectives are established at relevant "
            "functions, levels, and processes as documented in the Quality Objectives Matrix."
        )

        doc.add_heading("2. Scope of QMS", level=1)
        doc.add_paragraph(
            "This manual defines the requirements for the Quality Management System applicable to:"
        )
        scope_items = [
            "Crude oil refining (atmospheric and vacuum distillation, FCC, hydrocracking)",
            "Petrochemical product manufacturing (polypropylene, LPG, bitumen, sulfur)",
            "Blending and distribution of finished petroleum products",
            "Terminal and pipeline operations for product dispatch",
            "Maintenance, turnaround, and capital project management",
            "Laboratory testing and quality assurance of feedstocks and products",
        ]
        for item in scope_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_paragraph(
            "Exclusions: Design and development of refining process technology (outsourced to "
            "OEMs/licensors) are excluded per Clause 8.3 justification."
        )

        doc.add_heading("3. Process Approach", level=1)
        doc.add_paragraph(
            "Our QMS is built on the process approach as mandated by ISO 9001:2015 Clause 4.4. "
            "The interaction of processes is mapped through a Turtle Diagram framework."
        )

        doc.add_heading("3.1 Core Process Map", level=2)
        process_map = [
            ("OP-01", "Crude Intake & Blending", "Feedstock evaluation, desalting, blending recipes"),
            ("OP-02", "Primary & Secondary Processing", "Distillation, FCC, hydrocracking, reforming"),
            ("OP-03", "Product Blending & Dispatch", "Finished product blending, LPG bottling, pipeline loading"),
            ("OP-04", "Laboratory & QA/QC", "Incoming/outgoing inspection, certification, NABL-accredited lab"),
            ("OP-05", "Maintenance & Turnaround", "Preventive/predictive maintenance, shutdown planning"),
            ("OP-06", "Capital Projects", "EPC execution, commissioning, performance guarantees"),
            ("OP-07", "HSE Operations", "Environmental compliance, safety, PESO regulations"),
            ("OP-08", "Procurement & Supply Chain", "Approved vendor management, material traceability"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Process ID"
        hdr[1].text = "Process Name"
        hdr[2].text = "Key Activities"
        for pid, pname, pact in process_map:
            row = table.add_row().cells
            row[0].text = pid
            row[1].text = pname
            row[2].text = pact

        doc.add_heading("3.2 PDCA Cycle Application", level=2)
        doc.add_paragraph(
            "Each process follows the Plan-Do-Check-Act cycle with defined process owners, "
            "key performance indicators (KPIs), risk assessments, and scheduled management reviews. "
            "Internal audits verify conformity on a 12-month cycle per the Annual Audit Plan."
        )

        doc.add_heading("4. Context of the Organization", level=1)
        doc.add_paragraph(
            "The organization has determined external and internal issues relevant to its purpose "
            "and strategic direction (Clause 4.1), needs and expectations of interested parties "
            "(Clause 4.2), and the scope of the QMS considering these factors."
        )

        doc.add_heading("5. Leadership & Commitment", level=1)
        doc.add_paragraph(
            "Top management demonstrates leadership and commitment per Clause 5.1 by: "
            "ensuring integration of QMS requirements, communicating the importance of effective "
            "QMS and conforming to requirements, promoting continual improvement, and supporting "
            "relevant roles to demonstrate leadership."
        )

        doc.add_heading("6. Document Control", level=1)
        doc.add_paragraph(
            "All documented information is controlled per Clause 7.5. Electronic document management "
            "system (eDMS) ensures version control, approval workflows, distribution, retention, "
            "and disposition. Master document register is maintained with revision history."
        )

        doc.add_heading("7. Risk & Opportunity Management", level=1)
        doc.add_paragraph(
            "Risk-based thinking is integrated per Clause 6.1. A corporate-level risk register "
            "is maintained with risk owners, controls, and residual risk ratings. Opportunities "
            "for improvement are captured through CAPA, audit findings, and management review."
        )

        doc.add_heading("8. Continual Improvement", level=1)
        doc.add_paragraph(
            "Objectives are set at functional level with SMART criteria. Performance is monitored "
            "through dashboards, KPI trending, customer satisfaction surveys, and process audits. "
            "Improvement actions are tracked through the CAPA system with root cause analysis."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_iso14001_ems_manual_docx(self, filename="ISO14001_EMS_Manual.docx"):
        doc = Document()
        title = doc.add_heading("ISO 14001:2015 Environmental Management System Manual", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph(f"Document Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Environmental Policy", level=1)
        doc.add_paragraph(
            "MRPL is committed to protecting the environment by preventing pollution, complying "
            "with all applicable legal and other requirements, and continually improving our "
            "Environmental Management System. We integrate environmental considerations into "
            "business planning, operations, and product lifecycle thinking."
        )
        doc.add_paragraph(
            "Our environmental policy provides a framework for setting and reviewing environmental "
            "objectives and targets, and is available to the public as required."
        )

        doc.add_heading("2. Scope of EMS", level=1)
        scope_items = [
            "All refining operations and petrochemical manufacturing processes",
            "Fugitive and point-source air emissions (SOx, NOx, PM, VOCs)",
            "Wastewater treatment and discharge (zero liquid discharge facilities)",
            "Solid and hazardous waste management (hydrocarbon sludge, spent catalysts)",
            "Land and soil contamination prevention",
            "Water resource consumption and conservation",
            "Energy consumption and greenhouse gas (GHG) emissions tracking",
            "Pipeline and terminal environmental operations",
        ]
        for item in scope_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("3. Environmental Aspects & Impacts", level=1)
        doc.add_paragraph(
            "A comprehensive environmental aspects and impacts register is maintained per "
            "Clause 6.1.2. The methodology follows a lifecycle perspective considering "
            "direct and indirect aspects."
        )

        aspects = [
            ("Air Emissions", "FCC regenerator flares, furnace stacks, loading operations", "High", "Scrubbers, flaring minimization, leak detection (LDAR)"),
            ("Wastewater", "Process water, cooling tower blowdown, stormwater", "High", "ZLD, API separator, biological treatment"),
            ("Solid Waste", "Spent catalyst, sludge, OWS skimmings", "Medium", "Authorized recycler, secure landfill, waste minimization"),
            ("Noise", "Compressors, cooling fans, turnaround activities", "Low", "Acoustic enclosures, scheduling, PPE"),
            ("Soil Contamination", "Underground storage tanks, pipe leaks", "Medium", "Cathodic protection, double-walled tanks, monitoring wells"),
            ("GHG Emissions", "CO2 from combustion, process emissions, F-gases", "High", "Energy efficiency, flare gas recovery, carbon credits"),
        ]
        table = doc.add_table(rows=1, cols=4, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Aspect"
        hdr[1].text = "Source"
        hdr[2].text = "Significance"
        hdr[3].text = "Controls"
        for asp, src, sig, ctrl in aspects:
            row = table.add_row().cells
            row[0].text = asp
            row[1].text = src
            row[2].text = sig
            row[3].text = ctrl

        doc.add_heading("4. Legal & Compliance Requirements", level=1)
        doc.add_paragraph(
            "A legal register is maintained covering: Environment Protection Act 1986, "
            "Air (Prevention and Control of Pollution) Act 1981, Water (Prevention and Control "
            "of Pollution) Act 1974, Hazardous Waste Rules 2016, PESO regulations, CPCB/SPCB "
            "consent conditions, and client-specific environmental requirements."
        )

        doc.add_heading("5. Objectives & Targets", level=1)
        obj_items = [
            "Reduce specific energy consumption by 2.5% year-on-year",
            "Achieve zero environmental incidents (Level 1 and above)",
            "Reduce freshwater intake by 5% through ZLD optimization",
            "Reduce carbon intensity per barrel of crude processed by 3%",
            "100% compliance with consent conditions - zero show cause notices",
            "Increase waste recycling/recovery rate to 95%",
        ]
        for item in obj_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("6. Operational Control", level=1)
        doc.add_paragraph(
            "Operational control is maintained through: standardized operating procedures (SOPs) "
            "for significant aspects, emergency preparedness and response plans, maintenance "
            "calibration schedules, and management of change (MOC) procedures."
        )

        doc.add_heading("7. Monitoring & Measurement", level=1)
        doc.add_paragraph(
            "Continuous emissions monitoring systems (CEMS) are installed on major sources. "
            "Stack monitoring, ambient air quality monitoring, effluent analysis, and groundwater "
            "monitoring are conducted per regulatory schedules. Results are reported to CPCB/SPCB."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_iso45001_ohs_manual_docx(self, filename="ISO45001_OHS_Manual.docx"):
        doc = Document()
        title = doc.add_heading("ISO 45001:2018 Occupational Health & Safety Management System Manual", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph(f"Document Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. OH&S Policy", level=1)
        doc.add_paragraph(
            "MRPL is committed to providing a safe and healthy workplace for all employees, "
            "contractors, and visitors. We proactively identify hazards, assess risks, and "
            "implement controls to prevent work-related injury and ill health. Worker consultation "
            "and participation is integral to our OH&S management system."
        )

        doc.add_heading("2. Scope of OH&S MS", level=1)
        scope_items = [
            "All refinery operations including processing units and utilities",
            "Petrochemical plant operations and product handling",
            "Maintenance activities including turnaround/shutdown",
            "Capital projects and construction activities",
            "Pipeline operations and terminal activities",
            "Contractor and subcontractor safety management",
            "Off-site emergency response and community interface",
        ]
        for item in scope_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("3. Hazard Identification & Risk Assessment", level=1)
        doc.add_paragraph(
            "A systematic hazard identification process is implemented per Clause 6.1.1 and 6.1.2 "
            "using HAZID, HAZOP, JSA, FMEA, and Bow-Tie methodologies."
        )

        hazards = [
            ("Fire & Explosion", "Hydrocarbon leaks, hot work, static electricity", "Catastrophic", "HAZOP, LEL monitoring, fire protection systems, PTW"),
            ("Toxic Gas Exposure", "H2S, SOx, benzene, ammonia", "Major", "GMDs, wind socks, escape routes, gas monitoring, PPE"),
            ("Pressure Systems", "Pressure vessels, heat exchangers, reactors", "Major", "API 510/570 inspection, relief valves, MI programs"),
            ("Heights & Falling Objects", "Tower work, scaffolding, crane operations", "Moderate", "Fall protection, scaffolding certification, lifting plans"),
            ("Confined Spaces", "Tanks, vessels, underground pits", "Major", "Gas testing, ventilation, confined space entry permit"),
            ("Electrical Hazards", "HT/LT systems, transformers, switchgear", "Major", "LOTO procedures, Arc Flash protection, electrical safety rules"),
            ("Vehicle Movement", "Trucks, forklifts, heavy equipment", "Moderate", "Traffic management plan, exclusion zones, reversing alarms"),
            ("Ergonomic Hazards", "Manual handling, repetitive tasks, vibration", "Low", "Workstation design, job rotation, health surveillance"),
        ]
        table = doc.add_table(rows=1, cols=4, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Hazard"
        hdr[1].text = "Source"
        hdr[2].text = "Risk Level"
        hdr[3].text = "Controls"
        for haz, src, risk, ctrl in hazards:
            row = table.add_row().cells
            row[0].text = haz
            row[1].text = src
            row[2].text = risk
            row[3].text = ctrl

        doc.add_heading("4. Risk Control Hierarchy", level=1)
        doc.add_paragraph(
            "Controls follow the hierarchy of controls: Elimination > Substitution > Engineering "
            "Controls > Administrative Controls > PPE. ALARP principle is applied to demonstrate "
            "risks are reduced to As Low As Reasonably Practicable."
        )

        doc.add_heading("5. Emergency Preparedness", level=1)
        doc.add_paragraph(
            "Emergency response plans cover: fire, explosion, toxic release, medical emergency, "
            "natural disaster, and environmental spill scenarios. Drills are conducted quarterly. "
            "On-site fire station operates 24/7 with trained firefighting personnel."
        )

        doc.add_heading("6. Safety Performance Indicators", level=1)
        kpis = [
            "Total Recordable Incident Rate (TRIR) target: < 0.1",
            "Lost Time Injury Frequency Rate (LTIFR) target: < 0.05",
            "Near-miss reporting rate: > 5 per employee per year",
            "Safety observation cards: > 100 per month per unit",
            "First aid cases: < 2 per 200,000 man-hours",
            "Zero major process safety incidents per API RP 754 Tier 1",
        ]
        for item in kpis:
            doc.add_paragraph(item, style="List Bullet")

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_iso50001_enms_docx(self, filename="ISO50001_ENMS_Manual.docx"):
        doc = Document()
        title = doc.add_heading("ISO 50001:2018 Energy Management System Manual", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph(f"Document Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Energy Policy", level=1)
        doc.add_paragraph(
            "MRPL is committed to energy efficiency and conservation in all operations. We will "
            "use energy efficiently to improve performance, reduce costs, and minimize "
            "environmental impact. This policy supports the procurement and use of energy-"
            "efficient products and services, and design activities that consider energy performance "
            "improvement."
        )

        doc.add_heading("2. Scope of EnMS", level=1)
        doc.add_paragraph(
            "The Energy Management System covers all energy sources and significant energy uses "
            "within the refinery complex."
        )
        energy_sources = [
            "Natural gas (primary fuel for furnaces and boilers)",
            "Furnace oil / LSHS (secondary fuel)",
            "Electricity (grid power - HT/LT distribution)",
            "Steam (high pressure, medium pressure, low pressure networks)",
            "Compressed air (instrument and utility air)",
            "Diesel (backup generators, mobile equipment)",
        ]
        for item in energy_sources:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("3. Significant Energy Uses (SEUs)", level=1)
        seu_items = [
            ("Crude distillation furnaces", "42% of total thermal energy", "Stack O2 optimization, coil outlet temperature control"),
            ("FCC regenerator & reactor", "18% of total thermal energy", "Catalyst circulation optimization, air blower efficiency"),
            ("Hydrogen plant (SMR)", "12% of total thermal energy", "Steam methane reforming efficiency, PSA optimization"),
            ("Cooling water system", "25% of total electrical energy", "Cooling tower fan VFDs, condenser cleaning"),
            ("Compressed air system", "8% of total electrical energy", "Leak reduction, pressure band optimization, VSD compressors"),
            ("Distillation column pumps", "6% of total electrical energy", "Pump efficiency, VFD installation, impeller trimming"),
            ("Lighting & HVAC", "5% of total electrical energy", "LED conversion, smart building controls, BMS"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "SEU"
        hdr[1].text = "Energy Share"
        hdr[2].text = "Improvement Measures"
        for seu, share, measures in seu_items:
            row = table.add_row().cells
            row[0].text = seu
            row[1].text = share
            row[2].text = measures

        doc.add_heading("4. Energy Performance Indicators (EnPIs)", level=1)
        enpi_items = [
            "Specific energy consumption (GJ/tonne of crude throughput)",
            "Energy intensity index (EII) - Solomon benchmarking",
            "Specific electricity consumption (kWh/tonne)",
            "Specific steam consumption (kg/tonne at different pressure levels)",
            "Furnace efficiency (%) - actual vs. design",
            "Cogeneration plant availability and efficiency (%)",
            "Energy cost per unit of production (INR/tonne)",
            "Renewable energy share (%) in total energy mix",
        ]
        for item in enpi_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("5. Energy Review Process", level=1)
        doc.add_paragraph(
            "The energy review is conducted annually per Clause 6.4 and includes: analysis of "
            "energy consumption data, identification of SEUs, evaluation of energy performance, "
            "identification of energy improvement opportunities, and review of operational "
            "control requirements."
        )

        doc.add_heading("6. Energy Objectives & Targets", level=1)
        obj_items = [
            "Reduce specific energy consumption by 3% annually",
            "Improve furnace thermal efficiency by 1% point",
            "Reduce electrical specific consumption by 2% annually",
            "Achieve 5 MW additional solar power installation",
            "Complete waste heat recovery projects on 3 major furnaces",
            "Convert 80% of motors to IE4 efficiency class",
        ]
        for item in obj_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("7. Competence & Awareness", level=1)
        doc.add_paragraph(
            "Personnel operating SEUs are trained on energy efficient practices. Energy champions "
            "are designated in each unit. Awareness programs cover energy policy, EnPIs, energy "
            "targets, and contribution to improvement."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_ims_policy_docx(self, filename="IMS_Integrated_Policy.docx"):
        doc = Document()
        title = doc.add_heading("Integrated Management System Policy", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph("Covering: ISO 9001:2015 | ISO 14001:2015 | ISO 45001:2018 | ISO 50001:2018")
        doc.add_paragraph(f"Effective Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Policy Statement", level=1)
        doc.add_paragraph(
            "MRPL/ONGC is committed to an integrated approach to quality, environmental protection, "
            "occupational health and safety, and energy management. This integrated policy applies "
            "to all employees, contractors, and stakeholders involved in our refining and "
            "petrochemical operations."
        )
        doc.add_paragraph(
            "We pledge to:\n"
            "1. Meet or exceed customer requirements and statutory/regulatory obligations\n"
            "2. Prevent pollution and minimize environmental impact throughout our operations\n"
            "3. Provide a safe and healthy workplace free from occupational hazards\n"
            "4. Use energy efficiently and seek opportunities for energy performance improvement\n"
            "5. Protect the health, safety, and welfare of our workers and the community\n"
            "6. Continually improve the effectiveness of our integrated management system"
        )

        doc.add_heading("2. Scope of Integration", level=1)
        doc.add_paragraph(
            "This policy integrates the requirements of four international standards into a "
            "single cohesive framework:"
        )
        standards = [
            ("ISO 9001:2015", "Quality Management System", "Product quality, customer satisfaction, process consistency"),
            ("ISO 14001:2015", "Environmental Management System", "Environmental protection, pollution prevention, legal compliance"),
            ("ISO 45001:2018", "OH&S Management System", "Worker safety, hazard control, emergency preparedness"),
            ("ISO 50001:2018", "Energy Management System", "Energy efficiency, performance monitoring, improvement"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Standard"
        hdr[1].text = "System"
        hdr[2].text = "Focus Areas"
        for std, sys_name, focus in standards:
            row = table.add_row().cells
            row[0].text = std
            row[1].text = sys_name
            row[2].text = focus

        doc.add_heading("3. Leadership & Commitment", level=1)
        doc.add_paragraph(
            "Top management demonstrates leadership and commitment to the IMS by:\n"
            "- Ensuring the policy and objectives are established and compatible with strategic direction\n"
            "- Integrating IMS requirements into business processes\n"
            "- Promoting awareness of applicable requirements\n"
            "- Ensuring resources are available\n"
            "- Communicating the importance of effective IMS and conforming to requirements\n"
            "- Supporting management reviews and continual improvement"
        )

        doc.add_heading("4. Risk-Based Thinking", level=1)
        doc.add_paragraph(
            "Risk-based thinking is embedded across all four standards. A unified risk register "
            "consolidates quality, environmental, safety, and energy risks. Risk treatment plans "
            "are integrated into operational planning and reviewed during management review meetings."
        )

        doc.add_heading("5. Interested Parties", level=1)
        parties = [
            ("Customers", "Quality, delivery, product specifications", "ISO 9001"),
            ("Regulators (CPCB, SPCB, PESO)", "Environmental, safety compliance", "ISO 14001/45001"),
            ("Employees", "Safe workplace, training, consultation", "ISO 45001"),
            ("Shareholders", "Operational efficiency, profitability", "All standards"),
            ("Community", "Environmental performance, safety", "ISO 14001/45001"),
            ("Energy regulators (BEE)", "Energy performance, conservation", "ISO 50001"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Interested Party"
        hdr[1].text = "Needs & Expectations"
        hdr[2].text = "Applicable Standard"
        for party, needs, std in parties:
            row = table.add_row().cells
            row[0].text = party
            row[1].text = needs
            row[2].text = std

        doc.add_heading("6. Documented Information", level=1)
        doc.add_paragraph(
            "The IMS is documented through a hierarchical structure:\n"
            "Level 1: Integrated Policy Manual (this document)\n"
            "Level 2: Integrated Procedures and Management Programs\n"
            "Level 3: Work Instructions, SOPs, and Forms\n"
            "Level 4: Records, Reports, and Logs\n\n"
            "Document control is managed through the centralized eDMS with electronic workflows "
            "for review, approval, and distribution."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_internal_iso_audit_report_docx(self, filename="Internal_Audit_Report.docx"):
        doc = Document()
        title = doc.add_heading("Internal Audit Report", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Integrated Management System Audit")
        doc.add_paragraph(f"Audit Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Audit Overview", level=1)
        audit_info = [
            ("Audit ID", "IA-2026-001"),
            ("Audit Type", "Internal Audit (Combined QMS + EMS + OHSMS + EnMS)"),
            ("Audit Period", "January 2026 - December 2026"),
            ("Lead Auditor", "[Lead Auditor Name], Certified ISO 9001/14001/45001 Auditor"),
            ("Audit Team", "[Auditor 1], [Auditor 2], [Technical Expert]"),
            ("Auditee Units", "Crude Distillation, FCC, Utilities, Laboratory, Maintenance"),
            ("Standards Audited", "ISO 9001:2015, ISO 14001:2015, ISO 45001:2018, ISO 50001:2018"),
            ("Audit Methodology", "Process-based audit per ISO 19011:2018 guidelines"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Item"
        hdr[1].text = "Details"
        for item, detail in audit_info:
            row = table.add_row().cells
            row[0].text = item
            row[1].text = detail

        doc.add_heading("2. Executive Summary", level=1)
        doc.add_paragraph(
            "The internal audit of MRPL/ONGC's Integrated Management System was conducted from "
            "the planned audit dates. The audit covered all four management system standards "
            "across the selected operational units. Overall, the IMS is functioning effectively "
            "with evidence of strong leadership commitment and worker engagement."
        )
        doc.add_paragraph(
            "Total findings: 12 (4 Major NCs, 5 Minor NCs, 3 Observations). "
            "No critical findings were identified. The organization has demonstrated "
            "a culture of continual improvement with 89% of previous audit actions closed "
            "within target dates."
        )

        doc.add_heading("3. Audit Findings Summary", level=1)
        findings = [
            ("F-01", "Major NC", "ISO 9001:8.5.1", "Control of Production", "FCC catalyst change procedure not followed - missing MOC approval", "Corrective Action Required"),
            ("F-02", "Major NC", "ISO 14001:9.1.1", "Monitoring & Measurement", "CEMS data validation not performed on scheduled dates for Stack 3", "Corrective Action Required"),
            ("F-03", "Major NC", "ISO 45001:8.1.2", "Procurement", "Contractor safety induction records incomplete for 15% of workforce", "Corrective Action Required"),
            ("F-04", "Major NC", "ISO 50001:6.1.1", "Energy Review", "Energy review missing updated SEU rankings after recent modifications", "Corrective Action Required"),
            ("F-05", "Minor NC", "ISO 9001:7.1.5", "Monitoring Resources", "Calibration due dates exceeded for 3 laboratory instruments", "Corrective Action Required"),
            ("F-06", "Minor NC", "ISO 14001:8.1", "Operational Control", "Hazardous waste manifest not updated for new waste streams", "Corrective Action Required"),
            ("F-07", "Minor NC", "ISO 45001:9.1", "Performance Evaluation", "Safety observation cards not reviewed for 2 months in Utilities", "Corrective Action Required"),
            ("F-08", "Minor NC", "ISO 50001:4.3", "Scope", "EnMS scope document not updated to include new solar installation", "Corrective Action Required"),
            ("F-09", "Minor NC", "ISO 9001:10.2", "Nonconformity", "CAPA closure evidence missing root cause verification", "Corrective Action Required"),
            ("F-10", "Observation", "ISO 45001:6.1.4", "Action Planning", "Risk assessment for turnaround activities not yet initiated (6 months out)", "Observation - Address before shutdown"),
            ("F-11", "Observation", "ISO 14001:7.3", "Awareness", "Environmental awareness quiz scores below target for contractor personnel", "Observation - Training improvement"),
            ("F-12", "Observation", "ISO 50001:7.2", "Competence", "Energy champion training plan pending for new unit operators", "Observation - Complete by Q2"),
        ]
        table = doc.add_table(rows=1, cols=6, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Finding #"
        hdr[1].text = "Severity"
        hdr[2].text = "ISO Clause"
        hdr[3].text = "Process"
        hdr[4].text = "Finding Description"
        hdr[5].text = "Action Required"
        for fno, sev, clause, process, desc, action in findings:
            row = table.add_row().cells
            row[0].text = fno
            row[1].text = sev
            row[2].text = clause
            row[3].text = process
            row[4].text = desc
            row[5].text = action

        doc.add_heading("4. Positive Findings", level=1)
        positives = [
            "Strong worker consultation mechanism with active safety committee",
            "Effective management of change process for all major modifications",
            "CEMS data integrated with regulatory reporting portal",
            "Energy performance tracking dashboard with real-time monitoring",
            "Proactive near-miss reporting culture with >95% closure rate",
            "Regular management review with cross-functional participation",
        ]
        for item in positives:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("5. Recommendations", level=1)
        doc.add_paragraph(
            "1. Prioritize closure of major non-conformances within 30 days\n"
            "2. Conduct follow-up verification audit within 60 days\n"
            "3. Enhance contractor management system with digital induction tracking\n"
            "4. Update CEMS calibration schedule and implement backup validation procedures\n"
            "5. Conduct energy review workshop to update SEU rankings\n"
            "6. Share audit findings across all units for awareness and prevention"
        )

        doc.add_heading("6. Audit Sign-off", level=1)
        doc.add_paragraph("Lead Auditor: ____________________  Date: ________________")
        doc.add_paragraph("Auditee Representative: ____________________  Date: ________________")
        doc.add_paragraph("HOD Quality: ____________________  Date: ________________")

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_capa_report_docx(self, filename="CAPA_Report.docx"):
        doc = Document()
        title = doc.add_heading("Corrective and Preventive Action (CAPA) Report", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Quality Assurance Department")
        doc.add_paragraph(f"Report Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. CAPA Information", level=1)
        capa_info = [
            ("CAPA Number", "CAPA-2026-001"),
            ("Date Initiated", time.strftime("%d %B %Y")),
            ("Initiated By", "[Quality Engineer Name]"),
            ("Source", "Internal Audit Finding F-01 (ISO 9001 Non-conformance)"),
            ("Priority", "High"),
            ("Target Closure Date", "30 days from initiation"),
            ("Responsible Department", "Operations - FCC Unit"),
            ("CAPA Owner", "[FCC Unit Manager]"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Item"
        hdr[1].text = "Details"
        for item, detail in capa_info:
            row = table.add_row().cells
            row[0].text = item
            row[1].text = detail

        doc.add_heading("2. Problem Statement / Non-Conformance Description", level=1)
        doc.add_paragraph(
            "During the internal audit of FCC unit operations, it was observed that the catalyst "
            "change procedure was not followed as per SOP-QA-FCC-003. Specifically:\n\n"
            "- The Management of Change (MOC) approval was not obtained before initiating catalyst "
            "change activity on [Date]\n"
            "- The pre-change safety review checklist was incomplete\n"
            "- Risk assessment for the new catalyst batch was not reviewed\n"
            "- Quality impact assessment on product octane number was not documented\n\n"
            "This constitutes a major non-conformance against ISO 9001:2015 Clause 8.5.1 "
            "(Control of production and service provision)."
        )

        doc.add_heading("3. Root Cause Analysis", level=1)
        doc.add_paragraph("Methodology: 5-Why Analysis + Fishbone (Ishikawa) Diagram")
        doc.add_paragraph("")

        doc.add_heading("5-Why Analysis:", level=2)
        why_items = [
            ("Why 1", "Why was MOC not obtained?", "Operator was not aware MOC was required for catalyst change"),
            ("Why 2", "Why was operator not aware?", "Training records show last MOC training was >18 months ago"),
            ("Why 3", "Why was training not refreshed?", "Annual training calendar did not include MOC refresher for FCC operators"),
            ("Why 4", "Why was training calendar incomplete?", "Training coordinator role was vacant for 6 months"),
            ("Why 5", "Why was role not filled?", "HR recruitment process delay and competing priorities"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Level"
        hdr[1].text = "Question"
        hdr[2].text = "Answer"
        for level, q, a in why_items:
            row = table.add_row().cells
            row[0].text = level
            row[1].text = q
            row[2].text = a

        doc.add_paragraph("")
        doc.add_paragraph(
            "Root Cause: Inadequate training management system with no automated reminders "
            "for refresher training and no backup resource for critical training coordinator role."
        )

        doc.add_heading("4. Corrective Actions (Immediate)", level=1)
        ca_items = [
            ("CA-01", "Conduct MOC awareness refresher training for all FCC operators", "FCC Unit Manager", "7 days", "Training records, attendance sheet"),
            ("CA-02", "Review and complete all pending MOC documents for recent catalyst changes", "Quality Engineer", "5 days", "MOC register, approval records"),
            ("CA-03", "Perform risk assessment for current catalyst batch and document quality impact", "HSE + Quality teams", "10 days", "Risk assessment report, quality analysis"),
            ("CA-04", "Verify all operator competencies for critical procedures", "Training Coordinator", "15 days", "Competency matrix, assessment records"),
        ]
        table = doc.add_table(rows=1, cols=5, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "CA #"
        hdr[1].text = "Action"
        hdr[2].text = "Responsible"
        hdr[3].text = "Target Date"
        hdr[4].text = "Evidence Required"
        for ca_no, action, resp, target, evidence in ca_items:
            row = table.add_row().cells
            row[0].text = ca_no
            row[1].text = action
            row[2].text = resp
            row[3].text = target
            row[4].text = evidence

        doc.add_heading("5. Preventive Actions (Systemic)", level=1)
        pa_items = [
            ("PA-01", "Implement automated training management module with expiry alerts in HRMS", "IT + HR", "60 days", "System implementation report"),
            ("PA-02", "Establish backup training coordinator role with cross-training plan", "HR", "30 days", "JD, cross-training records"),
            ("PA-03", "Update SOP-QA-FCC-003 with mandatory pre-activity MOC verification step", "Quality", "15 days", "Revised SOP, change log"),
            ("PA-04", "Implement digital MOC approval workflow with system-level blocking", "IT + Quality", "90 days", "Workflow configuration, UAT report"),
            ("PA-05", "Extend training compliance dashboard to all critical operational procedures", "HR + Operations", "45 days", "Dashboard access, monitoring reports"),
        ]
        table = doc.add_table(rows=1, cols=5, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "PA #"
        hdr[1].text = "Action"
        hdr[2].text = "Responsible"
        hdr[3].text = "Target Date"
        hdr[4].text = "Evidence Required"
        for pa_no, action, resp, target, evidence in pa_items:
            row = table.add_row().cells
            row[0].text = pa_no
            row[1].text = action
            row[2].text = resp
            row[3].text = target
            row[4].text = evidence

        doc.add_heading("6. Effectiveness Verification", level=1)
        doc.add_paragraph(
            "Effectiveness will be verified after 90 days of CAPA closure through:\n"
            "- Follow-up audit of FCC unit MOC compliance\n"
            "- Training compliance rate monitoring (target: 100%)\n"
            "- Review of training management system reports\n"
            "- Assessment of any recurrence of similar non-conformances"
        )

        doc.add_heading("7. Sign-off", level=1)
        doc.add_paragraph("CAPA Owner: ____________________  Date: ________________")
        doc.add_paragraph("Quality Manager: ____________________  Date: ________________")
        doc.add_paragraph("Plant Head: ____________________  Date: ________________")

        os.makedirs(OUTPUT_DIR, exist=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_ncr_report_docx(self, filename="NCR_Report.docx"):
        doc = Document()
        title = doc.add_heading("Non-Conformance Report (NCR)", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Quality Assurance Department")
        doc.add_paragraph(f"NCR Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. NCR Identification", level=1)
        ncr_info = [
            ("NCR Number", "NCR-2026-001"),
            ("NCR Type", "Product Non-Conformance"),
            ("Severity", "Major"),
            ("Detected By", "Quality Control Laboratory"),
            ("Detection Method", "Incoming raw material inspection"),
            ("Date Detected", time.strftime("%d %B %Y")),
            ("Affected Product", "Polypropylene Grade - IP-150"),
            ("Batch/Lot Number", "PP-2026-0156"),
            ("Quantity Affected", "45 Metric Tonnes"),
            ("Customer/End Use", "Domestic market - injection molding application"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Field"
        hdr[1].text = "Details"
        for field, detail in ncr_info:
            row = table.add_row().cells
            row[0].text = field
            row[1].text = detail

        doc.add_heading("2. Description of Non-Conformance", level=1)
        doc.add_paragraph(
            "During routine quality control testing of Polypropylene Grade IP-150 (Batch PP-2026-0156), "
            "the following non-conformances were identified:\n\n"
            "1. Melt Flow Index (MFI): 14.2 g/10min (Specification: 12.0 ± 1.0 g/10min) - EXCEEDS UPPER LIMIT\n"
            "2. Tensile Strength: 28.5 MPa (Specification: ≥ 30.0 MPa) - BELOW MINIMUM\n"
            "3. Ash Content: 0.08% (Specification: ≤ 0.05%) - EXCEEDS LIMIT\n\n"
            "The non-conformance was detected at the final product stage before dispatch. The batch "
            "has been quarantined and placed on hold pending disposition."
        )

        doc.add_heading("3. Immediate Containment Action", level=1)
        doc.add_paragraph(
            "1. Batch PP-2026-0156 quarantined in designated hold area\n"
            "2. Hold tag applied and inventory system updated\n"
            "3. Dispatch to customer suspended pending investigation\n"
            "4. Previous 3 batches (PP-2026-0153, 0154, 0155) placed on precautionary hold\n"
            "5. Customer notification sent for potential delivery delay"
        )

        doc.add_heading("4. Investigation Findings", level=1)
        doc.add_paragraph(
            "Investigation conducted using 5-Why methodology and process data analysis:\n\n"
            "Primary Cause: Catalyst batch variation - supplier delivered catalyst with higher "
            "activity than specified, leading to increased polymerization rate.\n\n"
            "Contributing Factors:\n"
            "- Incoming catalyst QC check did not include activity verification (only composition)\n"
            "- Supplier COA data showed catalyst activity at upper specification limit\n"
            "- Process parameters were not adjusted for catalyst activity variation"
        )

        doc.add_heading("5. Disposition Decision", level=1)
        disposition = [
            ("Disposition Option", "Decision", "Justification"),
            ("Use As Is", "REJECTED", "MFI and tensile strength do not meet customer specifications"),
            ("Rework/Reprocess", "NOT FEASIBLE", "Polymer chain degradation cannot be reversed"),
            ("Return to Supplier", "SELECTED", "Catalyst quality did not meet agreed specifications"),
            ("Scrap/Destroy", "ALTERNATIVE", "If supplier return not accepted, batch to be disposed as industrial waste"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Disposition Option"
        hdr[1].text = "Decision"
        hdr[2].text = "Justification"
        for option, decision, justification in disposition[1:]:
            row = table.add_row().cells
            row[0].text = option
            row[1].text = decision
            row[2].text = justification

        doc.add_heading("6. Corrective Actions", level=1)
        ca_items = [
            ("Issue supplier quality non-conformance notification", "Quality Manager", "3 days"),
            ("Conduct supplier audit for catalyst manufacturing process", "Supplier Quality Engineer", "20 days"),
            ("Update incoming catalyst QC specification to include activity test", "Quality Laboratory", "15 days"),
            ("Implement statistical process control for catalyst activity", "Process Engineering", "30 days"),
            ("Review and update catalyst procurement specifications", "Procurement + Quality", "25 days"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Corrective Action"
        hdr[1].text = "Responsible"
        hdr[2].text = "Target Date"
        for action, resp, target in ca_items:
            row = table.add_row().cells
            row[0].text = action
            row[1].text = resp
            row[2].text = target

        doc.add_heading("7. Financial Impact", level=1)
        doc.add_paragraph(
            "Estimated financial impact:\n"
            "- Batch value (45 MT @ Rs. 95/kg): Rs. 42,75,000\n"
            "- Supplier claim amount: Rs. 42,75,000\n"
            "- Investigation and testing costs: Rs. 1,25,000\n"
            "- Customer delay penalty (if applicable): Rs. 2,50,000 (TBD)\n"
            "- Total potential impact: Rs. 46,50,000"
        )

        doc.add_heading("8. Sign-off", level=1)
        doc.add_paragraph("Quality Manager: ____________________  Date: ________________")
        doc.add_paragraph("Plant Head: ____________________  Date: ________________")
        doc.add_paragraph("Commercial Head: ____________________  Date: ________________")

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_mgmt_review_minutes_docx(self, filename="Management_Review_Minutes.docx"):
        doc = Document()
        title = doc.add_heading("Management Review Meeting Minutes", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Integrated Management System Review")
        doc.add_paragraph(f"Meeting Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Meeting Details", level=1)
        meeting_info = [
            ("Meeting ID", "MR-2026-Q1-001"),
            ("Date & Time", f"{time.strftime('%d %B %Y')}, 10:00 AM - 12:30 PM"),
            ("Location", "Conference Room A, Admin Building"),
            ("Chairperson", "[Plant Head / Managing Director]"),
            ("Secretary", "[Quality Manager]"),
            ("Standards Covered", "ISO 9001, ISO 14001, ISO 45001, ISO 50001"),
            ("Meeting Type", "Quarterly Management Review"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Item"
        hdr[1].text = "Details"
        for item, detail in meeting_info:
            row = table.add_row().cells
            row[0].text = item
            row[1].text = detail

        doc.add_heading("2. Attendees", level=1)
        attendees = [
            "[Plant Head] - Chairperson",
            "[Quality Manager] - Secretary",
            "[Operations Head] - FCC/CDU/Petchem",
            "[HSE Head] - Health, Safety & Environment",
            "[Maintenance Head] - Reliability & Projects",
            "[HR Head] - Human Resources",
            "[IT Head] - Information Technology",
            "[Finance Head] - Finance & Commercial",
            "[Lab Head] - Quality Control Laboratory",
            "[Energy Manager] - Energy Management",
        ]
        for att in attendees:
            doc.add_paragraph(att, style="List Bullet")

        doc.add_heading("3. Review Inputs (as per ISO Clauses)", level=1)

        doc.add_heading("3.1 Status of Actions from Previous Review", level=2)
        doc.add_paragraph(
            "Previous review (Q4-2025) actions: 12 actions assigned, 10 closed (83%), "
            "2 in progress. The two open actions relate to contractor safety management system "
            "upgrade and energy monitoring dashboard implementation. Both on track for Q2 closure."
        )

        doc.add_heading("3.2 Changes in External/Internal Issues", level=2)
        doc.add_paragraph(
            "- External: New CPCB emission norms effective April 2026 (SOx limit tightened to 50 mg/Nm3)\n"
            "- External: DPDP Act 2023 implementation rules issued - data classification requirements\n"
            "- Internal: FCC unit catalyst change completed - new ZSM-5 catalyst installed\n"
            "- Internal: Solar plant Phase-2 (5 MW) commissioned - grid synchronization pending\n"
            "- Internal: New HRMS module deployed - training management integration ongoing"
        )

        doc.add_heading("3.3 Customer Satisfaction & Feedback", level=2)
        doc.add_paragraph(
            "Customer satisfaction survey results: Overall score 4.3/5.0 (Target: 4.0)\n"
            "- Product quality: 4.5/5.0\n"
            "- On-time delivery: 4.1/5.0\n"
            "- Technical support: 4.2/5.0\n"
            "- Complaint resolution: 3.8/5.0 (Needs improvement)\n\n"
            "Customer complaints received: 8 (4 product quality, 3 delivery, 1 documentation)\n"
            "Complaints closed: 7 (87.5%). 1 under investigation."
        )

        doc.add_heading("3.4 QMS Performance & Effectiveness", level=2)
        doc.add_paragraph(
            "- Internal audit findings: 12 (4 Major, 5 Minor, 3 Observations)\n"
            "- CAPA closure rate: 92% within target\n"
            "- Process KPI achievement: 88% of targets met\n"
            "- Document non-conformance rate: 2.3% (Target: <3%)\n"
            "- Supplier quality performance: 96% acceptance rate\n"
            "- Product rejection rate: 0.8% (Target: <1%)"
        )

        doc.add_heading("3.5 Environmental Performance", level=2)
        doc.add_paragraph(
            "- Emissions compliance: 100% within consent limits\n"
            "- Zero liquid discharge: 99.7% compliance\n"
            "- GHG emissions: Reduced by 2.8% year-on-year\n"
            "- Environmental incidents: 0 Level 1 or above\n"
            "- Waste recycling rate: 93.5% (Target: 95%)\n"
            "- Water consumption: Reduced by 4.2% through ZLD optimization"
        )

        doc.add_heading("3.6 OH&S Performance", level=2)
        doc.add_paragraph(
            "- TRIR: 0.08 (Target: <0.1) - On track\n"
            "- LTIFR: 0.03 (Target: <0.05) - On track\n"
            "- Near-miss reports: 1,247 (127 per 100 workers)\n"
            "- Safety observations: 2,891 (positive trend)\n"
            "- First aid cases: 3\n"
            "- Process safety events (API 754 Tier 1): 0\n"
            "- Process safety events (API 754 Tier 2): 1 (minor)"
        )

        doc.add_heading("3.7 Energy Performance", level=2)
        doc.add_paragraph(
            "- Specific energy consumption: 4.82 GJ/tonne (Target: 4.90) - Exceeds target\n"
            "- EII (Solomon benchmark): 96.8 (Industry average: 100)\n"
            "- Electricity consumption: Reduced by 3.1% year-on-year\n"
            "- Solar power generation: 4.2 GWh (cumulative Phase 1+2)\n"
            "- Energy improvement projects completed: 8 of 12 planned\n"
            "- Annual energy savings: Rs. 8.5 crores"
        )

        doc.add_heading("3.8 Supplier Performance", level=2)
        doc.add_paragraph(
            "- Approved vendors: 245 active\n"
            "- New vendors approved: 18\n"
            "- Vendor audit findings: 7 (2 major, 5 minor)\n"
            "- Supply chain disruption events: 3 (all resolved within 48 hours)\n"
            "- Vendor quality agreements updated: 100% of critical vendors"
        )

        doc.add_heading("4. Action Items", level=1)
        actions = [
            ("AI-01", "Complete CEMS calibration schedule for all stacks", "HSE Head", "30 days", "High"),
            ("AI-02", "Implement contractor digital safety induction system", "IT + HSE", "45 days", "High"),
            ("AI-03", "Update energy review with revised SEU rankings", "Energy Manager", "30 days", "Medium"),
            ("AI-04", "Launch customer complaint root cause analysis improvement program", "Quality Manager", "60 days", "Medium"),
            ("AI-05", "Complete waste recycling improvement to achieve 95% target", "HSE Head", "90 days", "Medium"),
            ("AI-06", "Finalize DPDP Act compliance roadmap", "IT Head + Legal", "45 days", "High"),
            ("AI-07", "Conduct turnaround risk assessment for upcoming shutdown", "Operations Head", "30 days", "High"),
            ("AI-08", "Complete energy champion training for new unit operators", "Energy Manager", "15 days", "Low"),
        ]
        table = doc.add_table(rows=1, cols=5, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Action #"
        hdr[1].text = "Action Description"
        hdr[2].text = "Responsible"
        hdr[3].text = "Due Date"
        hdr[4].text = "Priority"
        for act_no, desc, resp, due, priority in actions:
            row = table.add_row().cells
            row[0].text = act_no
            row[1].text = desc
            row[2].text = resp
            row[3].text = due
            row[4].text = priority

        doc.add_heading("5. Decisions", level=1)
        doc.add_paragraph(
            "1. Approve IMS Policy revision to include updated legal requirements\n"
            "2. Allocate budget of Rs. 1.5 crores for energy improvement projects\n"
            "3. Mandate digital safety induction for all contractors from Q2\n"
            "4. Approve expansion of internal audit scope to cover DPDP Act compliance\n"
            "5. Next management review scheduled for Q2-2026 (July)"
        )

        doc.add_heading("6. Sign-off", level=1)
        doc.add_paragraph("Chairperson: ____________________  Date: ________________")
        doc.add_paragraph("Secretary: ____________________  Date: ________________")

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_quality_review_pptx(self, filename="Quality_Review_Presentation.pptx"):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        def add_title_slide(title_text, subtitle_text):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = subtitle_text
            return slide

        def add_content_slide(title_text, bullets):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
            return slide

        def add_table_slide(title_text, headers, rows_data):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title_text
            rows = len(rows_data) + 1
            cols = len(headers)
            left = PptxInches(0.5)
            top = PptxInches(1.8)
            width = PptxInches(12.333)
            height = PptxInches(4.5)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            for i, h in enumerate(headers):
                table.cell(0, i).text = h
            for r, row_data in enumerate(rows_data):
                for c, cell_text in enumerate(row_data):
                    table.cell(r + 1, c).text = cell_text
            return slide

        add_title_slide(
            "Quality Management Review",
            f"MRPL/ONGC - Quarterly Review\n{time.strftime('%B %Y')}"
        )

        add_content_slide(
            "Agenda",
            [
                "1. Quality Policy & Objectives Status",
                "2. Customer Satisfaction Review",
                "3. Product Quality Performance",
                "4. Internal Audit Findings",
                "5. CAPA Status & Closure",
                "6. Supplier Quality Performance",
                "7. Process KPI Dashboard",
                "8. Risk Register Update",
                "9. Continual Improvement Initiatives",
                "10. Action Items & Next Steps",
            ],
        )

        add_content_slide(
            "Quality Policy Reminder",
            [
                "Committed to delivering refined products meeting/exceeding customer requirements",
                "Risk-based thinking integrated across all processes",
                "Measurable objectives at functional level with SMART criteria",
                "Annual policy review with top management commitment",
            ],
        )

        add_table_slide(
            "Quality Objectives - Q1 Status",
            ["Objective", "Target", "Actual", "Status", "Trend"],
            [
                ["Product rejection rate", "<1.0%", "0.8%", "MET", "Stable"],
                ["Customer complaints", "<10/quarter", "8", "MET", "Improving"],
                ["CAPA closure rate", ">90%", "92%", "MET", "Improving"],
                ["Internal audit findings", "<15", "12", "MET", "Stable"],
                ["Supplier acceptance rate", ">95%", "96%", "MET", "Stable"],
                ["Document NCR rate", "<3%", "2.3%", "MET", "Improving"],
            ],
        )

        add_table_slide(
            "Customer Satisfaction Dashboard",
            ["Parameter", "Score (out of 5)", "Target", "Gap"],
            [
                ["Overall satisfaction", "4.3", "4.0", "+0.3"],
                ["Product quality", "4.5", "4.0", "+0.5"],
                ["On-time delivery", "4.1", "4.0", "+0.1"],
                ["Technical support", "4.2", "4.0", "+0.2"],
                ["Complaint resolution", "3.8", "4.0", "-0.2"],
            ],
        )

        add_table_slide(
            "Internal Audit Summary",
            ["Finding Type", "Count", "% of Total", "Closure Status"],
            [
                ["Major Non-Conformance", "4", "33%", "3 closed, 1 in progress"],
                ["Minor Non-Conformance", "5", "42%", "4 closed, 1 in progress"],
                ["Observations", "3", "25%", "2 addressed, 1 pending"],
                ["Total", "12", "100%", "83% closed"],
            ],
        )

        add_content_slide(
            "CAPA Highlights",
            [
                "CAPA-2025-Q4-012: MOC training gap - CLOSED (FCC unit)",
                "CAPA-2025-Q4-013: CEMS calibration - CLOSED (HSE)",
                "CAPA-2025-Q4-014: Contractor induction - IN PROGRESS (IT+HSE)",
                "CAPA-2025-Q4-015: Energy review update - IN PROGRESS",
                "Overall CAPA closure rate: 92% within target",
                "Zero repeat non-conformances in Q1",
            ],
        )

        add_table_slide(
            "Process KPI Performance",
            ["Process", "KPI", "Target", "Actual", "Status"],
            [
                ["Crude Distillation", "Yield efficiency", "97.5%", "97.8%", "Exceeded"],
                ["FCC Unit", "Conversion rate", "72%", "73.1%", "Exceeded"],
                ["Product Blending", "First-pass success", "98%", "97.2%", "Below"],
                ["Laboratory", "Test turnaround time", "4 hours", "3.5 hours", "Exceeded"],
                ["Maintenance", "Equipment availability", "98%", "98.5%", "Exceeded"],
            ],
        )

        add_content_slide(
            "Risk Register Update",
            [
                "New risk: Catalyst supply chain disruption (Medium) - Mitigation: Dual sourcing",
                "Escalated risk: DPDP Act compliance (High) - Mitigation: Compliance roadmap",
                "Reduced risk: Process safety events (Low) - Monitoring continues",
                "Retained risk: Contractor safety management (Medium) - Digital system in progress",
                "Risk review conducted: 15 risks reviewed, 3 updated, 1 new",
            ],
        )

        add_content_slide(
            "Continual Improvement Initiatives",
            [
                "Lean Six Sigma: 3 projects completed, 2 in progress",
                "Kaizen events: 12 conducted in Q1",
                "Suggestion scheme: 47 suggestions, 28 implemented",
                "Digital transformation: eDMS upgrade, MOC workflow digitization",
                "Training: 100% completion of critical competency assessments",
            ],
        )

        add_table_slide(
            "Action Items",
            ["Action #", "Description", "Owner", "Due Date"],
            [
                ["AI-01", "CEMS calibration schedule update", "HSE Head", "30 days"],
                ["AI-02", "Contractor digital induction system", "IT + HSE", "45 days"],
                ["AI-03", "Customer complaint improvement program", "Quality Mgr", "60 days"],
                ["AI-04", "DPDP Act compliance roadmap", "IT + Legal", "45 days"],
            ],
        )

        add_title_slide(
            "Thank You",
            "Questions & Discussion\nNext Review: Q2-2026 (July)"
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        return filepath

    # ======================================================================
    # INFORMATION SECURITY METHODS
    # ======================================================================

    def generate_infosec_policy_docx(self, filename="InfoSec_Policy.docx"):
        doc = Document()
        title = doc.add_heading("Information Security Policy", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph("Based on ISO 27001:2022 Requirements")
        doc.add_paragraph(f"Effective Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Policy Statement", level=1)
        doc.add_paragraph(
            "MRPL/ONGC is committed to protecting the confidentiality, integrity, and availability "
            "of all information assets, including operational technology (OT), SCADA systems, "
            "business data, and customer information. This policy establishes the framework for "
            "information security management aligned with ISO 27001:2022."
        )

        doc.add_heading("2. Scope", level=1)
        doc.add_paragraph(
            "This policy applies to all information assets, systems, networks, and data processed "
            "or stored within MRPL/ONGC operations, including:"
        )
        scope_items = [
            "IT infrastructure (servers, networks, endpoints, cloud services)",
            "OT infrastructure (SCADA, DCS, PLC, RTU systems)",
            "Business applications (ERP, LIMS, MES, SAP)",
            "Communication systems (email, messaging, video conferencing)",
            "Physical and environmental security of data centers",
            "Mobile devices and remote access systems",
            "Third-party and contractor access to information systems",
            "Data processed by or on behalf of MRPL/ONGC",
        ]
        for item in scope_items:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("3. Information Security Controls (ISO 27001:2022 Annex A)", level=1)
        doc.add_paragraph(
            "The following control categories are implemented as per ISO 27001:2022:"
        )

        controls = [
            ("A.5", "Organizational Controls", "Policies, roles, segregation of duties, threat intelligence, cloud security, data filtering, secure coding"),
            ("A.6", "People Controls", "Screening, terms of employment, awareness, training, disciplinary, termination, remote work, event reporting"),
            ("A.7", "Physical Controls", "Perimeters, physical entry, offices, physical monitoring, equipment security, secure disposal, clear desk/screen"),
            ("A.8", "Technological Controls", "Endpoint, privileged access, source code, data leakage, configuration, vulnerability management, malware protection, backup, logging, testing, cryptography, SDLC, secure auth, network security, application security, secure transfer, information deletion, masking, data leakage prevention, monitoring, web filtering, secure coding"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Control ID"
        hdr[1].text = "Category"
        hdr[2].text = "Key Controls"
        for cid, cat, ctrl in controls:
            row = table.add_row().cells
            row[0].text = cid
            row[1].text = cat
            row[2].text = ctrl

        doc.add_heading("4. IT/OT Convergence Security", level=1)
        doc.add_paragraph(
            "Given the refinery's integrated IT/OT environment, specific controls are implemented "
            "for operational technology systems:\n"
            "- Network segmentation between IT and OT zones (Purdue model)\n"
            "- Industrial DMZ (IDMZ) for data exchange\n"
            "- Whitelisting on SCADA/DCS servers\n"
            "- USB device control and media sanitization\n"
            "- Patch management with OT-specific testing protocols\n"
            "- Remote access via jump servers with MFA\n"
            "- OT-specific incident response procedures\n"
            "- Regular OT vulnerability assessments"
        )

        doc.add_heading("5. Roles & Responsibilities", level=1)
        doc.add_paragraph(
            "- Chief Information Security Officer (CISO): Overall ISMS accountability\n"
            "- IT Security Manager: Day-to-day security operations\n"
            "- OT Security Lead: OT/SCADA security management\n"
            "- Data Protection Officer (DPO): DPDP Act compliance\n"
            "- Information Asset Owners: Risk assessment for assigned assets\n"
            "- All Employees: Compliance with security policies and procedures"
        )

        doc.add_heading("6. Compliance & Enforcement", level=1)
        doc.add_paragraph(
            "Non-compliance with this policy may result in disciplinary action, including termination "
            "of employment or contract. Willful violations may be subject to legal action under "
            "applicable laws including IT Act 2000, DPDP Act 2023, and Indian Penal Code."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_data_classification_docx(self, filename="Data_Classification_Policy.docx"):
        doc = Document()
        title = doc.add_heading("Data Classification Policy", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Information Security Department")
        doc.add_paragraph(f"Effective Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Purpose", level=1)
        doc.add_paragraph(
            "This policy defines the framework for classifying information assets based on their "
            "sensitivity, criticality, and legal/regulatory requirements. It ensures appropriate "
            "protection measures are applied proportionate to the classification level."
        )

        doc.add_heading("2. Classification Levels", level=1)
        levels = [
            ("Level 4 - RESTRICTED", "Highly Sensitive", "CRITICAL",
             "Information whose unauthorized disclosure could cause catastrophic damage to operations, "
             "safety, or legal standing. Includes: SCADA configurations, safety system parameters, "
             "trade secrets, M&A plans, employee PII (Aadhaar, PAN), financial statements (pre-release), "
             "cybersecurity incident details, encryption keys, credentials.",
             "Encryption at rest and in transit, MFA access, DLP monitoring, dedicated storage, "
             "annual audit, need-to-know basis, no removable media, watermarking."),
            ("Level 3 - CONFIDENTIAL", "Sensitive", "HIGH",
             "Information whose unauthorized disclosure could cause significant operational, "
             "financial, or reputational damage. Includes: Process parameters, production data, "
             "quality test results (pre-release), vendor contracts, audit reports, risk assessments, "
             "employee records, customer data, API integrations.",
             "Encryption at rest and in transit, role-based access, DLP monitoring, access logging, "
             "secure sharing protocols, annual review, approved devices only."),
            ("Level 2 - INTERNAL", "General Internal", "MEDIUM",
             "Information intended for internal use only. Unauthorized disclosure may cause "
             "minor operational inconvenience. Includes: Internal policies, meeting minutes, "
             "training materials, organizational charts, internal communications, SOPs, "
             "maintenance schedules, non-sensitive production reports.",
             "Access control (employee-only), standard network security, no public sharing "
             "without approval, acceptable use policy compliance."),
            ("Level 1 - PUBLIC", "Unrestricted", "LOW",
             "Information approved for public disclosure. Includes: Published financial reports, "
             "marketing materials, press releases, product specifications (public), "
             "career pages, CSR reports, environmental reports (public versions).",
             "No special protection required, content review before publication, "
             "brand guideline compliance."),
        ]
        table = doc.add_table(rows=1, cols=5, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Classification"
        hdr[1].text = "Sensitivity"
        hdr[2].text = "Risk"
        hdr[3].text = "Examples"
        hdr[4].text = "Protection Requirements"
        for cls, sens, risk, examples, protection in levels:
            row = table.add_row().cells
            row[0].text = cls
            row[1].text = sens
            row[2].text = risk
            row[3].text = examples
            row[4].text = protection

        doc.add_heading("3. Classification Process", level=1)
        doc.add_paragraph(
            "1. Information Asset Owner assigns initial classification at creation/collection\n"
            "2. Classification label applied (digital watermark, file property, label)\n"
            "3. Protection controls implemented per classification level\n"
            "4. Classification reviewed annually or upon change in sensitivity\n"
            "5. Downgrading requires approval from Information Asset Owner + CISO\n"
            "6. Disposal procedures followed per classification level"
        )

        doc.add_heading("4. Handling Requirements by Level", level=1)
        handling = [
            ("RESTRICTED", "AES-256", "MFA + RBAC", "Full DLP", "Encrypted backup", "Crypto erase + physical destruction"),
            ("CONFIDENTIAL", "AES-256", "RBAC + MFA recommended", "DLP monitoring", "Encrypted backup", "Secure deletion"),
            ("INTERNAL", "Not required", "Employee authentication", "Standard AV", "Regular backup", "Standard deletion"),
            ("PUBLIC", "Not required", "None", "None", "None", "Standard deletion"),
        ]
        table = doc.add_table(rows=1, cols=6, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Level"
        hdr[1].text = "Encryption"
        hdr[2].text = "Access Control"
        hdr[3].text = "DLP/Monitoring"
        hdr[4].text = "Backup"
        hdr[5].text = "Disposal"
        for level, enc, access, dlp, backup, disposal in handling:
            row = table.add_row().cells
            row[0].text = level
            row[1].text = enc
            row[2].text = access
            row[3].text = dlp
            row[4].text = backup
            row[5].text = disposal

        doc.add_heading("5. Marking & Labeling", level=1)
        doc.add_paragraph(
            "All documents and data must be marked with their classification level:\n"
            "- Digital documents: Header/footer watermark, file metadata property\n"
            "- Emails: Subject line prefix [RESTRICTED], [CONFIDENTIAL], [INTERNAL]\n"
            "- Physical documents: Printed label on top-right corner\n"
            "- Database records: Classification field in metadata\n"
            "- Network shares: Folder naming convention with classification prefix"
        )

        doc.add_heading("6. Breach Reporting", level=1)
        doc.add_paragraph(
            "Any suspected or confirmed data classification breach must be reported immediately "
            "to the IT Security team via the incident reporting hotline or email. "
            "Breaches involving RESTRICTED or CONFIDENTIAL data trigger the incident response "
            "plan and may require notification to regulators under DPDP Act 2023."
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_privacy_policy_docx(self, filename="Privacy_Policy.docx"):
        doc = Document()
        title = doc.add_heading("Privacy Policy", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Refinery & Petrochemical Division")
        doc.add_paragraph("In Compliance with Digital Personal Data Protection Act, 2023 (DPDP Act)")
        doc.add_paragraph(f"Effective Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Introduction", level=1)
        doc.add_paragraph(
            "MRPL/ONGC (\"Company\", \"we\", \"our\") is committed to protecting the privacy "
            "and personal data of individuals. This Privacy Policy outlines how we collect, "
            "process, store, and protect personal data in accordance with the Digital Personal "
            "Data Protection Act, 2023 and its rules."
        )

        doc.add_heading("2. Definitions (as per DPDP Act 2023)", level=1)
        defs = [
            ("Personal Data", "Any data about an individual who is identifiable by or in relation to such data"),
            ("Data Principal", "The individual to whom the personal data relates"),
            ("Data Fiduciary", "MRPL/ONGC - determines the purpose and means of processing personal data"),
            ("Data Processor", "Any person who processes personal data on behalf of the Data Fiduciary"),
            ("Consent", "Any freely given, specific, informed, and unconditional indication of agreement"),
            ("Processing", "Any operation performed on personal data including collection, storage, retrieval, use, disclosure"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Term"
        hdr[1].text = "Definition"
        for term, defn in defs:
            row = table.add_row().cells
            row[0].text = term
            row[1].text = defn

        doc.add_heading("3. Categories of Personal Data Collected", level=1)
        categories = [
            ("Employee Data", "Name, address, contact details, Aadhaar, PAN, bank details, emergency contacts, employment history, qualifications, performance records"),
            ("Contractor Data", "Name, company details, Aadhaar, PAN, safety training records, medical fitness certificates"),
            ("Vendor Data", "Contact person details, bank account information, GST details"),
            ("Customer Data", "Contact details, order information, delivery addresses, feedback"),
            ("Visitor Data", "Name, vehicle number, photo ID, entry/exit timestamps"),
            ("CCTV Footage", "Video recordings of premises for security purposes"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Category"
        hdr[1].text = "Data Elements"
        for cat, elements in categories:
            row = table.add_row().cells
            row[0].text = cat
            row[1].text = elements

        doc.add_heading("4. Purpose of Processing", level=1)
        purposes = [
            "Employment management and HR administration",
            "Contractor and vendor management",
            "Safety management and regulatory compliance",
            "Security of premises and assets",
            "Regulatory reporting to government authorities",
            "Customer relationship management",
            "Legal and contractual obligations",
            "Business continuity planning",
        ]
        for item in purposes:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("5. Consent Requirements", level=1)
        doc.add_paragraph(
            "Personal data shall be processed only with the consent of the Data Principal or "
            "for legitimate uses as specified under the DPDP Act. Consent must be:\n"
            "- Freely given with a clear affirmative action\n"
            "- Specific to the purpose\n"
            "- Informed of the data being collected and purpose\n"
            "- Unconditional (not bundled with other terms)\n"
            "- Withdrawable at any time by the Data Principal"
        )

        doc.add_heading("6. Data Principal Rights", level=1)
        rights = [
            "Right to Access: Obtain information about processing of personal data",
            "Right to Correction: Correct inaccurate or incomplete personal data",
            "Right to Erasure: Request deletion of personal data when no longer needed",
            "Right to Grievance Redressal: Raise complaints regarding data processing",
            "Right to Nominate: Nominate a person to exercise rights in case of death/incapacity",
        ]
        for item in rights:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("7. Data Security Measures", level=1)
        doc.add_paragraph(
            "Reasonable security safeguards are implemented including:\n"
            "- Encryption of personal data at rest and in transit\n"
            "- Access controls with role-based permissions\n"
            "- Regular security audits and vulnerability assessments\n"
            "- Employee training on data protection obligations\n"
            "- Incident response procedures for data breaches\n"
            "- Data retention policies with defined periods\n"
            "- Secure disposal of personal data"
        )

        doc.add_heading("8. Data Retention", level=1)
        doc.add_paragraph(
            "Personal data is retained only for as long as necessary for the purpose of processing "
            "or as required by law. Retention periods are defined in the Data Retention Schedule. "
            "Upon expiry of retention period, data is securely deleted or anonymized."
        )

        doc.add_heading("9. Cross-Border Transfer", level=1)
        doc.add_paragraph(
            "Personal data may be transferred outside India only to countries or territories "
            "as notified by the Central Government, and subject to appropriate safeguards "
            "as required under the DPDP Act."
        )

        doc.add_heading("10. Grievance Officer", level=1)
        doc.add_paragraph(
            "For any privacy-related queries or grievances, contact:\n\n"
            "Data Protection Officer\n"
            "MRPL/ONGC\n"
            "Email: dpo@mrplongc.com\n"
            "Phone: [DPO Contact Number]\n"
            "Grievance redressal within 30 days of receipt"
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_it_asset_register_docx(self, filename="IT_Asset_Register.docx"):
        doc = Document()
        title = doc.add_heading("IT & OT Asset Register", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - IT & OT Security Department")
        doc.add_paragraph(f"Last Updated: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Purpose", level=1)
        doc.add_paragraph(
            "This register maintains a comprehensive inventory of all IT and OT assets including "
            "SCADA, DCS, and industrial control systems. It supports asset management, security "
            "monitoring, vulnerability management, and regulatory compliance."
        )

        doc.add_heading("2. Asset Classification", level=1)
        asset_classes = [
            ("IT Infrastructure", "Servers, storage, network devices, endpoints"),
            ("OT - SCADA", "Supervisory Control and Data Acquisition systems"),
            ("OT - DCS", "Distributed Control Systems"),
            ("OT - PLC/RTU", "Programmable Logic Controllers, Remote Terminal Units"),
            ("OT - Instruments", "Field instruments, transmitters, analyzers"),
            ("Communication", "Switches, routers, firewalls, industrial DMZ"),
            ("Applications", "Business and industrial applications"),
            ("Cloud Services", "IaaS, PaaS, SaaS subscriptions"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Asset Class"
        hdr[1].text = "Description"
        for cls, desc in asset_classes:
            row = table.add_row().cells
            row[0].text = cls
            row[1].text = desc

        doc.add_heading("3. OT Asset Inventory - SCADA Systems", level=1)
        scada_assets = [
            ("SCADA-SRV-01", "Primary SCADA Server", "Honeywell Experion PKS", "v4.10", "Unit 1 - CDU", "Critical", "24x7", "Connected to IDMZ"),
            ("SCADA-SRV-02", "Redundant SCADA Server", "Honeywell Experion PKS", "v4.10", "Unit 1 - CDU", "Critical", "24x7", "Connected to IDMZ"),
            ("SCADA-WS-01", "SCADA Operator Workstation", "Dell OptiPlex 7090", "Windows 10 LTSC", "Control Room", "High", "Shift-based", "Hardened workstation"),
            ("SCADA-WS-02", "SCADA Operator Workstation", "Dell OptiPlex 7090", "Windows 10 LTSC", "Control Room", "High", "Shift-based", "Hardened workstation"),
            ("SCADA-WS-03", "SCADA Engineering Workstation", "Dell Precision 5820", "Windows 10 LTSC", "Eng. Office", "High", "Business hours", "Restricted access"),
            ("SCADA-HIST-01", "Historian Server", "OSIsoft PI System", "PI Server 2018", "Server Room", "Critical", "24x7", "Data archiving"),
        ]
        table = doc.add_table(rows=1, cols=8, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Asset ID"
        hdr[1].text = "Name"
        hdr[2].text = "Platform"
        hdr[3].text = "Version/OS"
        hdr[4].text = "Location"
        hdr[5].text = "Criticality"
        hdr[6].text = "Availability"
        hdr[7].text = "Notes"
        for asset in scada_assets:
            row = table.add_row().cells
            for i, val in enumerate(asset):
                row[i].text = val

        doc.add_heading("4. OT Asset Inventory - DCS Systems", level=1)
        dcs_assets = [
            ("DCS-CONT-01", "FCC DCS Controller", "Honeywell C300", "R520.1", "FCC Unit", "Critical", "24x7", "Triple modular redundancy"),
            ("DCS-CONT-02", "CDU DCS Controller", "Honeywell C300", "R520.1", "CDU Unit", "Critical", "24x7", "Triple modular redundancy"),
            ("DCS-CONT-03", "Utilities DCS Controller", "Honeywell C300", "R520.1", "Utilities", "Critical", "24x7", "Dual redundancy"),
            ("DCS-IO-01", "FCC DCS I/O Module", "Honeywell PMIO", "N/A", "FCC Unit", "High", "24x7", "128 channels"),
            ("DCS-IO-02", "CDU DCS I/O Module", "Honeywell PMIO", "N/A", "CDU Unit", "High", "24x7", "128 channels"),
        ]
        table = doc.add_table(rows=1, cols=8, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Asset ID"
        hdr[1].text = "Name"
        hdr[2].text = "Platform"
        hdr[3].text = "Version"
        hdr[4].text = "Location"
        hdr[5].text = "Criticality"
        hdr[6].text = "Availability"
        hdr[7].text = "Notes"
        for asset in dcs_assets:
            row = table.add_row().cells
            for i, val in enumerate(asset):
                row[i].text = val

        doc.add_heading("5. OT Asset Inventory - PLC/RTU", level=1)
        plc_assets = [
            ("PLC-01", "Catalyst Handling PLC", "Siemens S7-1500", "V2.8", "FCC Unit", "High", "24x7", "Safety interlocks"),
            ("PLC-02", "Loading Terminal PLC", "Allen-Bradley ControlLogix", "v33", "Terminal", "High", "24x7", "Loading automation"),
            ("RTU-01", "Pipeline RTU", "ABB RTU560", "FW 1.3", "Pipeline Station", "High", "24x7", "Telemetry"),
            ("RTU-02", "Tank Farm RTU", "ABB RTU560", "FW 1.3", "Tank Farm", "High", "24x7", "Level monitoring"),
        ]
        table = doc.add_table(rows=1, cols=8, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Asset ID"
        hdr[1].text = "Name"
        hdr[2].text = "Platform"
        hdr[3].text = "Version"
        hdr[4].text = "Location"
        hdr[5].text = "Criticality"
        hdr[6].text = "Availability"
        hdr[7].text = "Notes"
        for asset in plc_assets:
            row = table.add_row().cells
            for i, val in enumerate(asset):
                row[i].text = val

        doc.add_heading("6. IT Infrastructure Assets", level=1)
        it_assets = [
            ("IT-SRV-01", "ERP Server (SAP)", "SAP S/4HANA", "2023 FPS02", "Data Center", "Critical", "24x7", "Cluster node 1"),
            ("IT-SRV-02", "ERP Server (SAP)", "SAP S/4HANA", "2023 FPS02", "Data Center", "Critical", "24x7", "Cluster node 2"),
            ("IT-SRV-03", "LIMS Server", "LabWare LIMS v7", "v7.4", "Data Center", "High", "24x7", "Laboratory management"),
            ("IT-SRV-04", "Email Server", "Microsoft Exchange", "2019 CU12", "Data Center", "High", "24x7", "On-premise"),
            ("IT-FW-01", "Perimeter Firewall", "Palo Alto PA-5260", "PAN-OS 11.1", "DMZ", "Critical", "24x7", "Internet gateway"),
            ("IT-FW-02", "IT/OT Firewall", "Palo Alto PA-3260", "PAN-OS 11.1", "IDMZ", "Critical", "24x7", "IT/OT boundary"),
            ("IT-SW-01", "Core Switch", "Cisco Nexus 9500", "NX-OS 10.3", "Data Center", "Critical", "24x7", "Core network"),
        ]
        table = doc.add_table(rows=1, cols=8, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Asset ID"
        hdr[1].text = "Name"
        hdr[2].text = "Platform"
        hdr[3].text = "Version"
        hdr[4].text = "Location"
        hdr[5].text = "Criticality"
        hdr[6].text = "Availability"
        hdr[7].text = "Notes"
        for asset in it_assets:
            row = table.add_row().cells
            for i, val in enumerate(asset):
                row[i].text = val

        doc.add_heading("7. Asset Lifecycle Management", level=1)
        doc.add_paragraph(
            "- New asset procurement: Security review and hardening before deployment\n"
            "- Commissioning: Baseline configuration documented and backed up\n"
            "- Operations: Regular patching, monitoring, vulnerability scanning\n"
            "- Decommissioning: Data sanitization, certificate revocation, register update\n"
            "- End-of-life: Secure disposal with certificate of destruction"
        )

        doc.add_heading("8. Vulnerability Management", level=1)
        doc.add_paragraph(
            "- IT assets: Monthly vulnerability scanning, quarterly penetration testing\n"
            "- OT assets: Semi-annual vulnerability assessment (non-disruptive)\n"
            "- Patch management: IT patches monthly, OT patches per change management\n"
            "- Critical vulnerabilities: Emergency patching within 72 hours (IT) or next maintenance window (OT)"
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_access_control_log_docx(self, filename="Access_Control_Log.docx"):
        doc = Document()
        title = doc.add_heading("Access Control Log & Quarterly Review Report", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Information Security Department")
        doc.add_paragraph(f"Review Period: Q1 2026 (January - March)")
        doc.add_paragraph(f"Report Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Executive Summary", level=1)
        doc.add_paragraph(
            "This quarterly access control review covers all user accounts, privileged access, "
            "physical access, and system access across IT and OT environments. The review ensures "
            "compliance with the principle of least privilege and identifies orphaned, dormant, "
            "or excessive access rights."
        )

        doc.add_heading("2. Access Control Summary", level=1)
        access_summary = [
            ("Total Active User Accounts", "2,847"),
            ("IT User Accounts", "1,923"),
            ("OT User Accounts", "924"),
            ("Privileged Accounts", "156"),
            ("Service Accounts", "89"),
            ("Dormant Accounts (>90 days)", "47"),
            ("Orphaned Accounts (terminated)", "12"),
            ("Accounts with MFA", "2,156 (75.7%)"),
            ("Physical Access Cards Active", "3,214"),
        ]
        table = doc.add_table(rows=1, cols=2, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Metric"
        hdr[1].text = "Count/Status"
        for metric, count in access_summary:
            row = table.add_row().cells
            row[0].text = metric
            row[1].text = count

        doc.add_heading("3. Privileged Access Review", level=1)
        doc.add_paragraph(
            "Privileged accounts are reviewed quarterly with enhanced scrutiny:"
        )
        priv_accounts = [
            ("SVC-SAP-ADMIN", "SAP Administrator", "Service Account", "IT", "Approved", "MFA + PAM"),
            ("SVC-SCADA-ENG", "SCADA Engineering", "Service Account", "OT", "Approved", "MFA + PAM + MFA"),
            ("USR-DBA-001", "Database Administrator", "User Account", "IT", "Approved", "MFA + PAM"),
            ("USR-CISO-001", "CISO Admin", "User Account", "IT", "Approved", "MFA + PAM + Privilege Review"),
            ("USR-OT-ENG-03", "OT Engineer", "User Account", "OT", "Under Review", "Excessive privileges identified"),
            ("SVC-BACKUP", "Backup Service", "Service Account", "IT", "Approved", "Password rotation automated"),
            ("USR-CONTRACTOR-12", "Contractor DBA", "User Account", "IT", "REVOKED", "Contract expired - access removed"),
        ]
        table = doc.add_table(rows=1, cols=6, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Account ID"
        hdr[1].text = "Description"
        hdr[2].text = "Type"
        hdr[3].text = "Environment"
        hdr[4].text = "Status"
        hdr[5].text = "Controls"
        for account in priv_accounts:
            row = table.add_row().cells
            for i, val in enumerate(account):
                row[i].text = val

        doc.add_heading("4. Dormant & Orphaned Account Remediation", level=1)
        doc.add_paragraph(
            "Findings from this quarter's review:"
        )
        findings = [
            "47 dormant accounts identified (>90 days no login) - 38 disabled, 9 verified active (seasonal workers)",
            "12 orphaned accounts from terminated employees - ALL disabled within 24 hours of HR notification",
            "3 contractor accounts with expired contracts - ALL revoked",
            "1 service account with expired password - rotated and monitoring enhanced",
            "2 accounts with shared credentials identified - separated, users trained",
        ]
        for item in findings:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("5. Physical Access Control Review", level=1)
        physical = [
            ("Data Center", "124", "52", "Biometric + card", "Monthly", "Completed"),
            ("Control Room", "89", "34", "Card access", "Monthly", "Completed"),
            ("Server Room", "67", "28", "Biometric + card", "Monthly", "Completed"),
            ("SCADA Engineering", "23", "12", "Biometric + card", "Monthly", "Completed"),
            ("Admin Building", "456", "156", "Card access", "Quarterly", "Completed"),
            ("Plant Area", "2,134", "845", "Card access", "Quarterly", "Completed"),
            ("Tank Farm", "156", "67", "Card + PIN", "Monthly", "Completed"),
        ]
        table = doc.add_table(rows=1, cols=7, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Zone"
        hdr[1].text = "Total Cards"
        hdr[2].text = "Active Users"
        hdr[3].text = "Access Method"
        hdr[4].text = "Review Frequency"
        hdr[5].text = "Review Status"
        hdr[6].text = "Exceptions"
        for zone, cards, users, method, freq, status in physical:
            row = table.add_row().cells
            row[0].text = zone
            row[1].text = cards
            row[2].text = users
            row[3].text = method
            row[4].text = freq
            row[5].text = status
            row[6].text = "None"

        doc.add_heading("6. MFA Compliance Status", level=1)
        doc.add_paragraph(
            "Multi-Factor Authentication (MFA) implementation status:"
        )
        mfa_status = [
            ("IT Administrators", "100%", "Target met"),
            ("OT Administrators", "100%", "Target met"),
            ("All IT Users", "82%", "Target: 90% by Q2"),
            ("All OT Users", "61%", "Target: 80% by Q3"),
            ("VPN Access", "100%", "Target met"),
            ("Cloud Services", "100%", "Target met"),
            ("Email Access", "94%", "Target met"),
            ("ERP Access", "78%", "Target: 95% by Q2"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Category"
        hdr[1].text = "MFA Coverage"
        hdr[2].text = "Status"
        for cat, coverage, status in mfa_status:
            row = table.add_row().cells
            row[0].text = cat
            row[1].text = coverage
            row[2].text = status

        doc.add_heading("7. Recommendations", level=1)
        recs = [
            "Expedite MFA rollout for OT users - currently at 61%, target 80% by Q3",
            "Implement automated account lifecycle management integrated with HRMS",
            "Review and reduce privileged account count (currently 156, target <120)",
            "Implement session recording for all privileged access sessions",
            "Enhance monitoring for service accounts with behavior analytics",
            "Conduct access certification campaign for all IT/OT user accounts",
            "Implement just-in-time (JIT) privileged access for administrative tasks",
        ]
        for item in recs:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("8. Sign-off", level=1)
        doc.add_paragraph("IT Security Manager: ____________________  Date: ________________")
        doc.add_paragraph("CISO: ____________________  Date: ________________")
        doc.add_paragraph("IT Head: ____________________  Date: ________________")

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_incident_response_plan_docx(self, filename="Incident_Response_Plan.docx"):
        doc = Document()
        title = doc.add_heading("Information Security Incident Response Plan", 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("MRPL/ONGC - Information Security Department")
        doc.add_paragraph(f"Version: 2.0 | Effective Date: {time.strftime('%d %B %Y')}")
        doc.add_paragraph("")

        doc.add_heading("1. Purpose & Scope", level=1)
        doc.add_paragraph(
            "This plan establishes the framework for detecting, responding to, containing, "
            "eradicating, and recovering from information security incidents. It covers IT systems, "
            "OT/SCADA systems, and all connected infrastructure within MRPL/ONGC operations."
        )

        doc.add_heading("2. Incident Classification (SEV Levels)", level=1)
        sev_levels = [
            ("SEV-1 (Critical)", "Critical operational impact",
             "SCADA/DCS compromise, refinery shutdown, safety system manipulation, ransomware on production systems, data breach of >10,000 records, regulatory data breach notification required",
             "Immediate", "CSIRT + CISO + Plant Head + Legal", "4 hours (response), 24 hours (containment)"),
            ("SEV-2 (High)", "Significant operational impact",
             "IT system compromise, unauthorized access to OT network, malware outbreak affecting multiple systems, denial of service on critical applications, data breach of 1,000-10,000 records",
             "1 hour", "CSIRT + CISO + IT Head", "8 hours (response), 48 hours (containment)"),
            ("SEV-3 (Medium)", "Limited operational impact",
             "Single system compromise, phishing campaign with credential theft, unauthorized access attempt (blocked), suspicious network activity, data breach of <1,000 records",
             "4 hours", "CSIRT + IT Security", "24 hours (response), 72 hours (containment)"),
            ("SEV-4 (Low)", "No/minimal operational impact",
             "Policy violation, minor malware infection (contained), vulnerability scan detection, social engineering attempt (no impact), spam campaign",
             "Next business day", "IT Security Team", "72 hours (response), 1 week (resolution)"),
        ]
        table = doc.add_table(rows=1, cols=6, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Severity"
        hdr[1].text = "Impact"
        hdr[2].text = "Examples"
        hdr[3].text = "Response Time"
        hdr[4].text = "Escalation"
        hdr[5].text = "SLA"
        for sev, impact, examples, resp_time, escalation, sla in sev_levels:
            row = table.add_row().cells
            row[0].text = sev
            row[1].text = impact
            row[2].text = examples
            row[3].text = resp_time
            row[4].text = escalation
            row[5].text = sla

        doc.add_heading("3. Computer Security Incident Response Team (CSIRT)", level=1)
        csirt = [
            ("CSIRT Lead", "CISO / IT Security Manager", "Overall incident coordination, decision authority"),
            ("Technical Lead", "Senior Security Analyst", "Technical investigation, forensics, containment"),
            ("OT Security Lead", "OT Security Engineer", "SCADA/DCS incident handling, OT containment"),
            ("Network Lead", "Network Security Engineer", "Network forensics, traffic analysis, isolation"),
            ("Systems Lead", "System Administrator", "System recovery, backup restoration"),
            ("Legal Advisor", "Legal Department", "Regulatory compliance, legal implications"),
            ("Communications Lead", "Corporate Communications", "Internal/external communications"),
            ("Business Liaison", "Operations Representative", "Business impact assessment, operational coordination"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Role"
        hdr[1].text = "Responsible Person"
        hdr[2].text = "Responsibilities"
        for role, person, resp in csirt:
            row = table.add_row().cells
            row[0].text = role
            row[1].text = person
            row[2].text = resp

        doc.add_heading("4. Incident Response Phases", level=1)

        doc.add_heading("Phase 1: Detection & Identification", level=2)
        doc.add_paragraph(
            "- SIEM alerts from security monitoring tools\n"
            "- SCADA/DCS anomaly detection systems\n"
            "- User reports (helpdesk, security hotline)\n"
            "- Threat intelligence feeds\n"
            "- Vulnerability scan findings\n"
            "- Third-party notifications\n"
            "- Log analysis and correlation"
        )

        doc.add_heading("Phase 2: Containment", level=2)
        doc.add_paragraph(
            "- Short-term: Isolate affected systems, block malicious IPs/accounts\n"
            "- OT-specific: Engage process control team, assess safety implications\n"
            "- Evidence preservation (disk images, memory dumps, logs)\n"
            "- Network segmentation enforcement\n"
            "- Credential reset for compromised accounts"
        )

        doc.add_heading("Phase 3: Eradication", level=2)
        doc.add_paragraph(
            "- Remove malware and persistence mechanisms\n"
            "- Patch exploited vulnerabilities\n"
            "- Rebuild compromised systems from known-good backups\n"
            "- Update security controls to prevent recurrence\n"
            "- Verify OT system integrity against known baselines"
        )

        doc.add_heading("Phase 4: Recovery", level=2)
        doc.add_paragraph(
            "- Restore systems from validated backups\n"
            "- Gradual return to production with monitoring\n"
            "- Verify system functionality and integrity\n"
            "- Monitor for signs of residual compromise\n"
            "- Document recovery procedures and timelines"
        )

        doc.add_heading("Phase 5: Post-Incident Activities", level=2)
        doc.add_paragraph(
            "- Post-incident review meeting within 5 business days\n"
            "- Lessons learned documentation\n"
            "- Incident report preparation\n"
            "- Security control improvements\n"
            "- Training updates based on findings\n"
            "- Trend analysis and metrics reporting"
        )

        doc.add_heading("5. Communication Plan", level=1)
        comms = [
            ("Internal - CSIRT", "Encrypted channel / secure messaging", "Immediate"),
            ("Internal - Management", "Phone + encrypted email", "Within 1 hour (SEV-1/2)"),
            ("Internal - All employees", "Internal communication portal", "As needed"),
            ("External - Regulators", "Official communication", "As required by law (DPDP Act)"),
            ("External - Customers", "Customer notification portal", "If customer data affected"),
            ("External - Law enforcement", "Official channels", "If criminal activity suspected"),
            ("Media", "Corporate communications only", "If public disclosure required"),
        ]
        table = doc.add_table(rows=1, cols=3, style="Light Grid Accent 1")
        hdr = table.rows[0].cells
        hdr[0].text = "Audience"
        hdr[1].text = "Channel"
        hdr[2].text = "Timing"
        for audience, channel, timing in comms:
            row = table.add_row().cells
            row[0].text = audience
            row[1].text = channel
            row[2].text = timing

        doc.add_heading("6. OT-Specific Incident Handling", level=1)
        doc.add_paragraph(
            "OT/SCADA incidents require additional considerations:\n"
            "1. SAFETY FIRST: Assess impact on process safety before any containment action\n"
            "2. Do NOT disconnect OT systems without process control team approval\n"
            "3. Engage OT security lead and process control engineer immediately\n"
            "4. Use OT-specific forensics tools (no standard IT tools on OT networks)\n"
            "5. Coordinate with plant operations for any system isolation\n"
            "6. Document all actions for regulatory reporting (CPCB, PESO, MeitY)\n"
            "7. Verify safety instrumented system (SIS) integrity"
        )

        doc.add_heading("7. Incident Response Toolkit", level=1)
        toolkit = [
            "Forensic imaging tools (FTK Imager, dd)",
            "Memory analysis tools (Volatility)",
            "Network analysis (Wireshark, tcpdump)",
            "Malware analysis sandbox",
            "Secure communication channels",
            "Evidence bags and chain-of-custody forms",
            "OT-specific forensics toolkit",
            "Backup and recovery procedures",
            "Contact lists and escalation matrices",
        ]
        for item in toolkit:
            doc.add_paragraph(item, style="List Bullet")

        doc.add_heading("8. Testing & Drills", level=1)
        doc.add_paragraph(
            "The incident response plan is tested through:\n"
            "- Tabletop exercises: Quarterly (simulated scenarios)\n"
            "- Technical drills: Semi-annual (technical response validation)\n"
            "- Full-scale simulation: Annual (enterprise-wide)\n"
            "- OT-specific drill: Annual (SCADA incident scenario)\n"
            "- Results documented and plan updated based on findings"
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        doc.save(filepath)
        return filepath

    def generate_infosec_review_pptx(self, filename="InfoSec_Review_Presentation.pptx"):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        def add_title_slide(title_text, subtitle_text):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = subtitle_text
            return slide

        def add_content_slide(title_text, bullets):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title_text
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
            return slide

        def add_table_slide(title_text, headers, rows_data):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title_text
            rows = len(rows_data) + 1
            cols = len(headers)
            left = PptxInches(0.5)
            top = PptxInches(1.8)
            width = PptxInches(12.333)
            height = PptxInches(4.5)
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            for i, h in enumerate(headers):
                table.cell(0, i).text = h
            for r, row_data in enumerate(rows_data):
                for c, cell_text in enumerate(row_data):
                    table.cell(r + 1, c).text = cell_text
            return slide

        add_title_slide(
            "Information Security Review",
            f"MRPL/ONGC - Quarterly Security Review\n{time.strftime('%B %Y')}"
        )

        add_content_slide(
            "Agenda",
            [
                "1. Security Posture Overview",
                "2. Incident Summary & Trends",
                "3. Vulnerability Management Status",
                "4. Access Control Review",
                "5. OT/SCADA Security Update",
                "6. DPDP Act Compliance Progress",
                "7. Security Awareness & Training",
                "8. Third-Party Risk Management",
                "9. Key Metrics & KPIs",
                "10. Action Items & Roadmap",
            ],
        )

        add_content_slide(
            "Security Posture Dashboard",
            [
                "Overall Security Score: 78/100 (Target: 85)",
                "Critical vulnerabilities open: 3 (down from 7 last quarter)",
                "Mean time to detect (MTTD): 4.2 hours (Target: <4 hours)",
                "Mean time to respond (MTTR): 12.5 hours (Target: <12 hours)",
                "Security incidents: 23 total (1 SEV-2, 5 SEV-3, 17 SEV-4)",
                "Phishing click rate: 4.2% (Target: <3%)",
                "MFA coverage: 75.7% (Target: 90% by Q4)",
            ],
        )

        add_table_slide(
            "Incident Summary - Q1 2026",
            ["Category", "Count", "SEV-1/2", "SEV-3", "SEV-4", "Trend"],
            [
                ["Malware/Ransomware", "5", "0", "2", "3", "Decreasing"],
                ["Phishing/Social Engineering", "8", "0", "2", "6", "Stable"],
                ["Unauthorized Access", "3", "1", "2", "0", "Decreasing"],
                ["Data Leakage", "2", "0", "1", "1", "New"],
                ["Policy Violation", "4", "0", "0", "4", "Increasing"],
                ["DDoS/Network Attack", "1", "0", "0", "1", "Decreasing"],
                ["Total", "23", "1", "7", "15", "Stable"],
            ],
        )

        add_table_slide(
            "Vulnerability Management",
            ["Asset Type", "Critical", "High", "Medium", "Low", "Patch Rate"],
            [
                ["IT Servers", "2", "15", "45", "120", "94%"],
                ["IT Endpoints", "0", "8", "35", "85", "88%"],
                ["OT - SCADA", "1", "3", "8", "15", "72%"],
                ["OT - DCS", "0", "2", "5", "12", "68%"],
                ["Network Devices", "0", "4", "12", "28", "91%"],
                ["Cloud Services", "0", "1", "6", "18", "96%"],
            ],
        )

        add_content_slide(
            "OT/SCADA Security Update",
            [
                "IT/OT network segmentation completed - Purdue model implemented",
                "Industrial DMZ (IDMZ) operational with dual firewalls",
                "SCADA server hardening: 85% complete (target 100% by Q2)",
                "USB device control deployed on 90% of OT workstations",
                "OT vulnerability assessment completed - 3 critical findings remediated",
                "Remote access: Jump server + MFA enforced for all OT remote access",
                "OT security monitoring: SIEM integration for SCADA/DCS logs",
                "Safety Instrumented System (SIS) integrity verified - no findings",
            ],
        )

        add_table_slide(
            "Access Control Metrics",
            ["Metric", "Current", "Target", "Status"],
            [
                ["Total user accounts", "2,847", "N/A", "Baseline"],
                ["Dormant accounts", "47", "<20", "Behind"],
                ["Orphaned accounts", "12", "0", "Behind"],
                ["Privileged accounts", "156", "<120", "Behind"],
                ["MFA coverage", "75.7%", "90%", "Behind"],
                ["Access review completion", "100%", "100%", "Met"],
                ["Contractor access review", "100%", "100%", "Met"],
            ],
        )

        add_content_slide(
            "DPDP Act Compliance Progress",
            [
                "Data Protection Officer (DPO) appointed and registered",
                "Data classification policy implemented across all departments",
                "Privacy impact assessments completed for 12 critical systems",
                "Consent management framework deployed for employee data",
                "Data retention schedules defined and implementation started",
                "Cross-border data transfer agreements reviewed (8 pending)",
                "Breach notification procedures documented and tested",
                "Employee privacy training: 78% completion (target: 100% by Q2)",
            ],
        )

        add_content_slide(
            "Security Awareness & Training",
            [
                "Mandatory security awareness training: 92% completion",
                "Phishing simulation campaigns: 3 conducted, click rate trending down",
                "Role-based training: IT security (100%), OT security (85%), developers (78%)",
                "New employee security induction: 100% coverage",
                "Contractor security orientation: 100% coverage",
                "Security champion program: 24 champions across departments",
                "Top attack vector: Phishing (35% of incidents)",
            ],
        )

        add_table_slide(
            "Key Security KPIs",
            ["KPI", "Q1 2026", "Q4 2025", "Target", "Trend"],
            [
                ["Security incidents", "23", "28", "<25", "Improving"],
                ["Critical vulnerabilities", "3", "7", "<5", "Improving"],
                ["MTTD (hours)", "4.2", "5.1", "<4", "Improving"],
                ["MTTR (hours)", "12.5", "15.3", "<12", "Improving"],
                ["Phishing click rate", "4.2%", "5.8%", "<3%", "Improving"],
                ["MFA coverage", "75.7%", "62%", "90%", "Improving"],
                ["Patch compliance", "89%", "82%", "95%", "Improving"],
                ["Training completion", "92%", "85%", "100%", "Improving"],
            ],
        )

        add_table_slide(
            "Action Items",
            ["Action #", "Description", "Owner", "Due Date", "Priority"],
            [
                ["AI-01", "Complete MFA rollout for OT users", "OT Security Lead", "60 days", "High"],
                ["AI-02", "Remediate 3 critical SCADA vulnerabilities", "CSIRT Lead", "30 days", "Critical"],
                ["AI-03", "Reduce dormant accounts to <20", "IT Security", "30 days", "High"],
                ["AI-04", "Complete DPDP Act employee training", "DPO + HR", "45 days", "High"],
                ["AI-05", "Implement automated account lifecycle", "IT + HR", "90 days", "Medium"],
                ["AI-06", "Deploy OT security monitoring (Phase 2)", "OT Security", "60 days", "Medium"],
            ],
        )

        add_title_slide(
            "Thank You",
            "Questions & Discussion\nNext Review: Q2-2026 (July)\nSecurity Hotline: [Extension]\nEmail: security@mrplongc.com"
        )

        os.makedirs(OUTPUT_DIR, exist_ok=True)
        filepath = os.path.join(OUTPUT_DIR, filename)
        prs.save(filepath)
        return filepath
