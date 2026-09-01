"""
MRPL & ONGC Enterprise Deliverables Synthesis Engine
Implements authentic Maharatna PSU document generators across 9 operational departments:
- Word (.docx): DoP Note Sheets, OISD-STD-105 Permits, CVC Integrity Pact, FIIR, BRSR
- Excel (.xlsx): Dynamic formulas for API 570 Remaining Life, API 610 Hydraulic Power, e-MB RA Billing, GeM CST L1, GHG Scope 1/2
- PowerPoint (.pptx): 16:9 Executive Briefings, Apex Safety Reviews, TAR Decks
"""

import os
import time
from typing import Dict, Any, List, Optional
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor

try:
    from apps.admin_backend.generators.deliverables_ep_project import ExtendedGenerators_EpProject
    from apps.admin_backend.generators.deliverables_quality_infosec import ExtendedGenerators_QualityInfoSec
    from apps.admin_backend.generators.deliverables_legal import ExtendedGenerators_Legal
    from apps.admin_backend.generators.deliverables_hse_ops import ExtendedGenerators_HseOps
    from apps.admin_backend.generators.deliverables_mech_materials import ExtendedGenerators_MechMaterials
    from apps.admin_backend.generators.deliverables_finance_esg import ExtendedGenerators_FinanceEsg
except ImportError:
    try:
        from .deliverables_ep_project import ExtendedGenerators_EpProject
        from .deliverables_quality_infosec import ExtendedGenerators_QualityInfoSec
        from .deliverables_legal import ExtendedGenerators_Legal
        from .deliverables_hse_ops import ExtendedGenerators_HseOps
        from .deliverables_mech_materials import ExtendedGenerators_MechMaterials
        from .deliverables_finance_esg import ExtendedGenerators_FinanceEsg
    except ImportError:
        from generators.deliverables_ep_project import ExtendedGenerators_EpProject
        from generators.deliverables_quality_infosec import ExtendedGenerators_QualityInfoSec
        from generators.deliverables_legal import ExtendedGenerators_Legal
        from generators.deliverables_hse_ops import ExtendedGenerators_HseOps
        from generators.deliverables_mech_materials import ExtendedGenerators_MechMaterials
        from generators.deliverables_finance_esg import ExtendedGenerators_FinanceEsg

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))

def ensure_output_directories():
    os.makedirs(os.path.join(OUTPUT_DIR, "docx"), exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, "xlsx"), exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, "pptx"), exist_ok=True)
    os.makedirs(os.path.join(OUTPUT_DIR, "scripts"), exist_ok=True)

class DeliverableGenerator(
    ExtendedGenerators_EpProject,
    ExtendedGenerators_QualityInfoSec,
    ExtendedGenerators_Legal,
    ExtendedGenerators_HseOps,
    ExtendedGenerators_MechMaterials,
    ExtendedGenerators_FinanceEsg,
):
    """
    Complete 47-Deliverable Enterprise Generator Engine for MRPL & ONGC.
    """
    def __init__(self):
        ensure_output_directories()

    def generate_approval_note_docx(
        self,
        filename: str = "approval_note.docx",
        findings: Optional[List[Dict[str, Any]]] = None,
        sop_violations: Optional[List[str]] = None
    ) -> str:
        """Generates formal executive approval note. Requires caller-provided findings; no hardcoded demo data."""
        if not findings:
            raise ValueError("findings is required — caller must provide inspection findings; no hardcoded demo fallback")
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("EXECUTIVE APPROVAL NOTE / MOP&NG STATUTORY SUBMISSION")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)
        meta_table = doc.add_table(rows=4, cols=2)
        meta_data = [
            ("Unit / Location:", "As per inspection report"),
            ("Inspection Date:", time.strftime("%Y-%m-%d")),
            ("Refinery SOP Reference:", "As per applicable SOP"),
            ("Classification:", "As per inspection findings")
        ]
        for idx, (label, val) in enumerate(meta_data):
            r = meta_table.rows[idx].cells
            r[0].text = label
            r[0].paragraphs[0].runs[0].font.bold = True
            r[1].text = val

        doc.add_heading("1. Executive Summary", level=1)
        doc.add_paragraph(
            "This approval note is generated from the provided inspection findings. "
            "All parameters below are taken directly from caller input — no hardcoded demo values."
        )

        doc.add_heading("2. Verified Inspection Findings", level=1)
        findings_table = doc.add_table(rows=1, cols=4)
        hdr_cells = findings_table.rows[0].cells
        hdr_cells[0].text = "Parameter / Instrument"
        hdr_cells[1].text = "Measured Value"
        hdr_cells[2].text = "Safe Operating Limit"
        hdr_cells[3].text = "Compliance Status"
        for cell in hdr_cells:
            cell.paragraphs[0].runs[0].font.bold = True

        for item in findings:
            row_cells = findings_table.add_row().cells
            # Support both tuple and dict findings
            if isinstance(item, dict):
                row_cells[0].text = str(item.get("parameter", item.get("key", "")))
                row_cells[1].text = str(item.get("measured", item.get("value", "")))
                row_cells[2].text = str(item.get("limit", item.get("sop_clause", "")))
                row_cells[3].text = str(item.get("status", ""))
            else:
                row_cells[0].text = str(item[0]) if len(item) > 0 else ""
                row_cells[1].text = str(item[1]) if len(item) > 1 else ""
                row_cells[2].text = str(item[2]) if len(item) > 2 else ""
                row_cells[3].text = str(item[3]) if len(item) > 3 else ""

        doc.add_heading("3. Standard Operating Procedure (SOP) Compliance Verdict", level=1)
        p_sop = doc.add_paragraph()
        r_sop = p_sop.add_run(
            "MANDATORY CLAUSE TRIGGERED:\n"
            "Applicable SOP dictates that measured values exceeding safe operating limits mandate immediate "
            "corrective action per the defined compliance protocol."
        )
        r_sop.font.bold = True
        r_sop.font.color.rgb = RGBColor(180, 40, 40)

        doc.add_heading("4. Recommended Authority Action", level=1)
        doc.add_paragraph("1. Authority action as per applicable SOP and inspection findings.")
        doc.add_paragraph("2. Schedule corrective action per reviewer recommendation.")

        doc.add_heading("5. Approval Authority Sign-offs", level=1)
        sig_table = doc.add_table(rows=2, cols=3)
        sigs = [
            ("Initiated By:", "Senior Operations Engineer"),
            ("Reviewed By:", "Chief General Manager (Technical Services)"),
            ("Approved By:", "Executive Director (Operations)")
        ]
        for idx, (role, title_str) in enumerate(sigs):
            c = sig_table.rows[0].cells[idx]
            c.text = f"{role}\n\n___________________\n{title_str}"

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_hydraulic_register_xlsx(self, filename: str = "hydraulic_calculation_register.xlsx", calc_params: Optional[Dict[str, Any]] = None) -> str:
        """Generates Excel workbook with hydraulic calculations. Requires calc_params; no hardcoded demo values."""
        if not calc_params:
            raise ValueError("calc_params is required — caller must provide hydraulic parameters; no hardcoded demo fallback")
        wb = Workbook()
        ws = wb.active
        ws.title = "Calculation"
        title_font = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        regular_font = Font(name="Calibri", size=11)
        bold_font = Font(name="Calibri", size=11, bold=True)

        header_fill = PatternFill(start_color="1A365D", end_color="1A365D", fill_type="solid")
        section_fill = PatternFill(start_color="2B6CB0", end_color="2B6CB0", fill_type="solid")
        zebra_fill = PatternFill(start_color="F7FAFC", end_color="F7FAFC", fill_type="solid")
        pass_fill = PatternFill(start_color="C6F6D5", end_color="C6F6D5", fill_type="solid")

        thin_border = Border(
            left=Side(style='thin', color='E2E8F0'),
            right=Side(style='thin', color='E2E8F0'),
            top=Side(style='thin', color='E2E8F0'),
            bottom=Side(style='thin', color='E2E8F0')
        )

        ws.merge_cells("A1:E1")
        ws["A1"] = "HYDRAULIC CALCULATION REGISTER"
        ws["A1"].font = title_font
        ws["A1"].fill = header_fill
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")

        # All values come from caller-provided calc_params — no hardcoded equipment demo
        calc_rows = [
            ("Equipment Tag", calc_params.get("equipment_tag", ""), calc_params.get("equipment_desc", ""), "", ""),
            ("Design Standard", calc_params.get("standard", ""), calc_params.get("standard_desc", ""), "", ""),
            ("Fluid Handled", calc_params.get("fluid", ""), calc_params.get("fluid_desc", ""), "", ""),
            ("", "", "", "", ""),
            ("Parameter", "Symbol", "Input Value", "Unit", "Verification / Design Envelope"),
            ("Volumetric Flow Rate", "Q", calc_params.get("Q", ""), calc_params.get("Q_unit", ""), calc_params.get("Q_note", "")),
            ("Differential Head", "H", calc_params.get("H", ""), calc_params.get("H_unit", ""), calc_params.get("H_note", "")),
            ("Fluid Density", "ρ", calc_params.get("rho", ""), calc_params.get("rho_unit", ""), calc_params.get("rho_note", "")),
            ("Gravitational Acceleration", "g", calc_params.get("g", 9.81), "m/s²", "Standard gravity"),
            ("Flow Rate (SI Units)", "Q_si", calc_params.get("Q_si", ""), "m³/s", calc_params.get("Q_si_formula", "")),
            ("Hydraulic Power Generated", "Ph", calc_params.get("Ph", ""), "kW", calc_params.get("Ph_formula", "")),
            ("Electrical Motor Input Power", "Pin", calc_params.get("Pin", ""), "kW", calc_params.get("Pin_note", "")),
            ("Operating Hydraulic Efficiency", "η", calc_params.get("eta", ""), "%", calc_params.get("eta_formula", "")),
            ("Compliance Verdict", "VERDICT", calc_params.get("verdict", ""), "-", calc_params.get("verdict_note", ""))
        ]

        for r_idx, row in enumerate(calc_rows, start=2):
            for c_idx, val in enumerate(row, start=1):
                cell = ws.cell(row=r_idx, column=c_idx, value=val)
                cell.font = regular_font
                cell.border = thin_border
                if r_idx == 6:
                    cell.font = header_font
                    cell.fill = section_fill
                elif r_idx == 15:
                    cell.font = bold_font
                    cell.fill = pass_fill
                elif r_idx % 2 == 0 and r_idx > 6:
                    cell.fill = zebra_fill

        for col_idx in range(1, 6):
            col_letter = get_column_letter(col_idx)
            max_len = max(len(str(ws.cell(row=r, column=col_idx).value or '')) for r in range(2, 16))
            ws.column_dimensions[col_letter].width = max(max_len + 4, 16)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_asset_register_xlsx(self, filename: str = "asset_register.xlsx", asset_tags: Optional[List[tuple]] = None) -> str:
        """Generates ISA 5.1 instrumentation asset register Excel workbook. Requires asset_tags; no hardcoded demo tags."""
        wb = Workbook()
        ws = wb.active
        ws.title = "P&ID Asset Register"
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        regular_font = Font(name="Calibri", size=10)
        header_fill = PatternFill(start_color="1A365D", end_color="1A365D", fill_type="solid")
        zebra_fill = PatternFill(start_color="F7FAFC", end_color="F7FAFC", fill_type="solid")

        headers = ["ISA Tag", "Equipment Name", "Subsystem", "Design Rating", "Calibration / Inspection Interval", "Verification Status"]
        for col_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx, value=h)
            cell.font = header_font
            cell.fill = header_fill

        tags_data = asset_tags if asset_tags else []
        if not tags_data:
            raise ValueError("asset_tags is required — caller must provide asset register tags; no hardcoded demo tags")

        for r_idx, row in enumerate(tags_data, start=2):
            for c_idx, val in enumerate(row, start=1):
                cell = ws.cell(row=r_idx, column=c_idx, value=val)
                cell.font = regular_font
                if r_idx % 2 == 0:
                    cell.fill = zebra_fill

        for col_idx in range(1, len(headers) + 1):
            col_letter = get_column_letter(col_idx)
            max_len = max(len(str(ws.cell(row=r, column=col_idx).value or '')) for r in range(1, len(tags_data) + 2))
            ws.column_dimensions[col_letter].width = max(max_len + 4, 16)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_turnaround_briefing_pptx(self, filename: str = "turnaround_briefing.pptx", briefing_data: Optional[Dict[str, Any]] = None) -> str:
        """Generates 16:9 widescreen executive briefing presentation. Requires briefing_data; no hardcoded demo content."""
        if not briefing_data:
            raise ValueError("briefing_data is required — caller must provide briefing content; no hardcoded demo fallback")
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        title_slide_layout = prs.slide_layouts[0]
        slide1 = prs.slides.add_slide(title_slide_layout)
        slide1.shapes.title.text = briefing_data.get("title", "OPERATIONS BRIEFING")
        slide1.placeholders[1].text = briefing_data.get("subtitle", "Generated by Sovereign Agentic AI Workbench")

        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = briefing_data.get("anomaly_title", "Anomaly Summary")
        tf2 = slide2.placeholders[1].text_frame
        findings = briefing_data.get("findings", [])
        if findings:
            tf2.text = f"• {findings[0]}"
            for f in findings[1:]:
                tf2.add_paragraph().text = f"• {f}"
        else:
            tf2.text = "• No findings provided"

        slide3 = prs.slides.add_slide(bullet_layout)
        slide3.shapes.title.text = "7-Day Decoking Turnaround Milestone Schedule"
        tf3 = slide3.placeholders[1].text_frame
        tf3.text = "• Day 1-2: Firing de-rating and crude throughput rebalancing."
        tf3.add_paragraph().text = "• Day 3-4: Unit isolation, steam-air decoking, and pigging of radiant passes."
        tf3.add_paragraph().text = "• Day 5-6: Hydrostatic re-testing, ultrasonic gauge verification, and refractory inspection."
        tf3.add_paragraph().text = "• Day 7: Safe unit start-up and ramp-up under CGM sign-off protocol."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 1. HSE, FIRE & DISASTER MANAGEMENT (ERDMP)
    # =========================================================================

    def generate_hot_work_permit_docx(self, filename: str = "HW-B-OISD105.docx") -> str:
        """OISD-STD-105 Form B Hot Work Permit with Gas Testing Matrix"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("HOT WORK & NAKED FLAME PERMIT — FORM B (OISD-STD-105)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)
        meta_table = doc.add_table(rows=4, cols=2)
        meta = [
            ("Permit Serial No:", f"HW-MRPL-{time.strftime('%Y%m%d')}-042"),
            ("Plant / Location:", "[Unit Name] ([Unit ID]) - [Location]"),
            ("Validity Window:", f"{time.strftime('%Y-%m-%d')} 08:00 to 18:00 hrs"),
            ("Designated Firewatch:", "Mr. Ramesh Kumar (ID: FW-8841)")
        ]
        for idx, (k, v) in enumerate(meta):
            r = meta_table.rows[idx].cells
            r[0].text = k
            r[0].paragraphs[0].runs[0].font.bold = True
            r[1].text = v

        doc.add_heading("1. Multi-Gas Atmospheric Test Matrix (Pre-Entry & Periodic)", level=1)
        gt_table = doc.add_table(rows=5, cols=4)
        headers = ["Atmospheric Gas Parameter", "Permissible Safe Limit", "Measured Test Value", "Status"]
        for c_idx, h in enumerate(headers):
            gt_table.rows[0].cells[c_idx].text = h
            gt_table.rows[0].cells[c_idx].paragraphs[0].runs[0].font.bold = True
        
        gas_data = [
            ("Combustible Hydrocarbons (%LEL)", "0.0% (Strictly Nil)", "0.0%", "PASS - SAFE"),
            ("Oxygen Content (O₂ %)", "19.5% - 23.5%", "20.8%", "PASS - SAFE"),
            ("Hydrogen Sulphide (H₂S ppm)", "< 10.0 ppm (Ceiling)", "0.0 ppm", "PASS - SAFE"),
            ("Carbon Monoxide (CO ppm)", "< 35.0 ppm", "2.1 ppm", "PASS - SAFE")
        ]
        for r_idx, row in enumerate(gas_data, start=1):
            for c_idx, val in enumerate(row):
                gt_table.rows[r_idx].cells[c_idx].text = val

        doc.add_heading("2. Mandatory Safety & Isolation Checkpoints", level=1)
        doc.add_paragraph("[X] Equipment positively blinded / spaded with blind register entry.")
        doc.add_paragraph("[X] Sewers, catch basins, and oily water drains within 15m covered with fire blankets.")
        doc.add_paragraph("[X] Continuous pressurized fire hose and two 10kg DCP extinguishers stationed at site.")
        doc.add_paragraph("[X] LOTO Certificate issued for electrical drive isolation.")

        doc.add_heading("3. Issuing & Receiving Authority Sign-off Ladder", level=1)
        sig_table = doc.add_table(rows=2, cols=3)
        sigs = [
            ("Gas Analyst / QA-QC:", "Instrument QC Specialist"),
            ("Issuing Authority:", "Shift In-Charge (Operations)"),
            ("Accepting Authority:", "Executing Engineer (Maintenance)")
        ]
        for idx, (role, title_str) in enumerate(sigs):
            c = sig_table.rows[0].cells[idx]
            c.text = f"{role}\n\n___________________\n{title_str}"

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_confined_space_permit_docx(self, filename: str = "CSE-C-OISD105.docx") -> str:
        """OISD-STD-105 Form C Confined Space Entry Permit"""
        doc = Document()
        title = doc.add_heading("ONGC - MRPL ENERGY COMPLEX", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("CONFINED SPACE ENTRY & BOX-UP PERMIT — FORM C (OISD-STD-105)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("─" * 55)
        doc.add_heading("1. Confined Vessel Particulars", level=1)
        doc.add_paragraph("Vessel Tag: C-101 (Crude Fractionator Column) | Height: 48.5m | Volume: 320 m³")
        doc.add_paragraph("Last Contained Fluid: Heavy Crude Gas Oil / Naphtha Fractions")
        doc.add_paragraph("Standby Attendant (Outside Manway): Mr. Suresh Naik (Cert # SA-2026)")

        doc.add_heading("2. Continuous 2-Hour Atmospheric Surveillance Log", level=1)
        log_table = doc.add_table(rows=4, cols=5)
        cols = ["Time (Hrs)", "O₂ % (19.5-23.5)", "LEL % (0.0%)", "H₂S ppm (<10)", "Analyst Initials"]
        for i, h in enumerate(cols):
            log_table.rows[0].cells[i].text = h
            log_table.rows[0].cells[i].paragraphs[0].runs[0].font.bold = True
        
        entries = [
            ("08:00", "20.9%", "0.0%", "0.0 ppm", "QC-Analyst 1"),
            ("10:00", "20.8%", "0.0%", "0.0 ppm", "QC-Analyst 1"),
            ("12:00", "20.9%", "0.0%", "0.0 ppm", "QC-Analyst 2")
        ]
        for r_idx, row in enumerate(entries, 1):
            for c_idx, val in enumerate(row):
                log_table.rows[r_idx].cells[c_idx].text = val

        doc.add_heading("3. Standby & Rescue Verification", level=1)
        doc.add_paragraph("Full-body safety harness with life retrieval line anchored at manway entrance.")
        doc.add_paragraph("SCBA (Self-Contained Breathing Apparatus) reserve cylinders (2x 300 bar) on standby.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_erdmp_master_docx(self, filename: str = "ERDMP-MRPL-Master.docx") -> str:
        """PNGRB Emergency Response and Disaster Management Plan Master Volume"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("ON-SITE EMERGENCY RESPONSE & DISASTER MANAGEMENT PLAN (PNGRB ERDMP 2010/2026)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Tiered Emergency Classification", level=1)
        doc.add_paragraph("• Level I: Localized incident controllable within the unit boundaries (Shift In-Charge Incident Commander).")
        doc.add_paragraph("• Level II: Plant-wide incident requiring internal mutual aid and site main control room activation.")
        doc.add_paragraph("• Level III: Catastrophic off-site disaster requiring District Disaster Management Authority (DDMA) & NDRF escalation.")

        doc.add_heading("2. Emergency Communication Cascade & Siren Codes", level=1)
        doc.add_paragraph("• Emergency Siren: Wailing sound for 2 minutes (interrupted tone).")
        doc.add_paragraph("• All Clear Siren: Continuous straight pitch tone for 2 minutes.")
        doc.add_paragraph("• Designated Assembly Points: Primary Assembly Point #1 (Main Admin Gate), Secondary (North Flare Perimeter).")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_hse_kpi_dashboard_xlsx(self, filename: str = "HSE-KPI-Dash-FY26.xlsx") -> str:
        """Safety KPI Dashboard with Dynamic LTIFR and TRIR Formulas"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Monthly_LTIFR_TRIR"

        header_fill = PatternFill(start_color="1A365D", end_color="1A365D", fill_type="solid")
        header_font = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
        regular_font = Font(name="Calibri", size=10)

        ws.merge_cells("A1:K1")
        ws["A1"] = "MRPL REFINERY - MONTHLY SAFETY KPI & LTIFR/TRIR DASHBOARD (FY 2025-26)"
        ws["A1"].font = Font(name="Calibri", size=13, bold=True, color="FFFFFF")
        ws["A1"].fill = header_fill
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")

        headers = ["Month", "Own Man-Hours", "Contract Man-Hours", "Total Hours", "LTI Count", "RWC Count", "MTC Count", "Recordable Incidents", "LTIFR (per 1M hrs)", "TRIR (per 200k hrs)", "Status"]
        for c_idx, h in enumerate(headers, 1):
            cell = ws.cell(row=2, column=c_idx, value=h)
            cell.font = header_font
            cell.fill = header_fill

        months_data = [
            ("Apr-25", 220000, 480000, 0, 0, 1),
            ("May-25", 225000, 510000, 0, 1, 0),
            ("Jun-25", 218000, 495000, 0, 0, 0),
            ("Jul-25", 230000, 530000, 0, 0, 0),
            ("Aug-25", 222000, 510000, 0, 0, 2),
            ("Sep-25", 228000, 540000, 0, 0, 0),
        ]

        for idx, (m, own, con, lti, rwc, mtc) in enumerate(months_data, start=3):
            ws.cell(row=idx, column=1, value=m).font = regular_font
            ws.cell(row=idx, column=2, value=own).font = regular_font
            ws.cell(row=idx, column=3, value=con).font = regular_font
            ws.cell(row=idx, column=4, value=f"=B{idx}+C{idx}").font = regular_font
            ws.cell(row=idx, column=5, value=lti).font = regular_font
            ws.cell(row=idx, column=6, value=rwc).font = regular_font
            ws.cell(row=idx, column=7, value=mtc).font = regular_font
            ws.cell(row=idx, column=8, value=f"=E{idx}+F{idx}+G{idx}").font = regular_font
            ws.cell(row=idx, column=9, value=f"=(E{idx}*1000000)/D{idx}").font = regular_font
            ws.cell(row=idx, column=10, value=f"=(H{idx}*200000)/D{idx}").font = regular_font
            ws.cell(row=idx, column=11, value=f'=IF(I{idx}<=0.10,"TARGET MET","ATTENTION")').font = regular_font

        for col_idx in range(1, 12):
            col_letter = get_column_letter(col_idx)
            ws.column_dimensions[col_letter].width = 18

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_gas_monitoring_register_xlsx(self, filename: str = "Gas_Monitoring_Register.xlsx") -> str:
        """Daily Gas Monitoring & Detector Calibration Register"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Daily_Gas_Test_Log"
        ws["A1"] = "MRPL REFINERY - DAILY MULTI-GAS SURVEILLANCE & SENSOR CALIBRATION REGISTER"
        ws["A1"].font = Font(name="Calibri", size=12, bold=True)

        headers = ["Date", "Permit No", "Unit / Area", "Equipment Tag", "O2 %", "LEL %", "H2S ppm", "CO ppm", "Detector S/N", "Calibration Due", "Verdict"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        sample = [
            (time.strftime("%Y-%m-%d"), "PERMIT-001", "Unit Area", "Equipment-001", 20.9, 0.0, 0.0, 1.2, "DET-001", "2026-12-31", "=IF(AND(E3>=19.5,E3<=23.5,F3=0,G3<10),'PASS','STOP')"),
            (time.strftime("%Y-%m-%d"), "PERMIT-002", "Unit Area", "Equipment-002", 20.8, 0.0, 0.0, 0.0, "DET-002", "2026-12-31", "=IF(AND(E4>=19.5,E4<=23.5,F4=0,G4<10),'PASS','STOP')")
        ]
        for r_idx, row in enumerate(sample, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_hazop_worksheet_xlsx(self, filename: str = "HAZOP_Worksheet_IEC61882.xlsx") -> str:
        """13-Column IEC 61882 HAZOP Deviation & CAPA RAG Matrix"""
        wb = Workbook()
        ws = wb.active
        ws.title = "HAZOP_Node_01"
        ws["A1"] = "HAZOP STUDY WORKSHEET — IEC 61882 (CDU CRUDE CHARGE TRAIN)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Node ID", "Description", "Parameter", "Guide Word", "Deviation", "Causes", "Consequences", "Sev (1-5)", "Lik (1-5)", "Risk (S*L)", "Existing Safeguards", "Recommendation", "Status"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        row_data = [
            ("N-01", "Charge Line to Equipment", "Flow", "LESS", "Low Flow", "Pump Trip / Strainer Clog", "Equipment overheating", 4, 3, "=H3*I3", "Low Flow Alarm", "Install automated low flow interlock trip", "OPEN"),
            ("N-01", "Charge Line to Equipment", "Pressure", "MORE", "High Pressure Surge", "Control valve fail closed", "Upstream flange leakage", 3, 2, "=H4*I4", "PSV relief to flare", "Quarterly PSV pop test calibration", "CLOSED")
        ]
        for r_idx, row in enumerate(row_data, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_apex_safety_review_pptx(self, filename: str = "Apex_Safety_Committee_Review.pptx") -> str:
        """Monthly Apex Safety Committee Review Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL APEX SAFETY COMMITTEE REVIEW"
        s1.placeholders[1].text = f"Monthly Executive Safety Scorecard & OISD Compliance\nChaired by Director (Technical) — {time.strftime('%B %Y')}"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Executive Safety Scorecard — FY 2025-26"
        tf2 = s2.placeholders[1].text_frame
        tf2.text = "• Cumulative Safe Man-Hours: 4.82 Million Hours achieved without Lost Time Injury (LTI)."
        tf2.add_paragraph().text = "• LTIFR Metric: 0.00 (Target: ≤0.10 per 1M man-hours)."
        tf2.add_paragraph().text = "• Near-Miss & Hazard Hunt: 142 reports logged; 94.3% closed within 14-day SLA."
        tf2.add_paragraph().text = "• High Potential (HIPO) Findings: 2 cases reviewed with completed CAPA."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 2. REFINERY OPERATIONS & PROCESS ENGINEERING
    # =========================================================================

    def generate_unit_commissioning_plan_docx(self, filename: str = "COMM-PLAN-Unit-RevX.docx") -> str:
        """Unit Commissioning Plan & Start-up Operating Manual"""
        doc = Document()
        title = doc.add_heading("MRPL REFINERY OPERATIONS DIVISION", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("UNIT PRE-COMMISSIONING & START-UP MASTER PLAN")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Operating Envelope & Feed Slate Envelope", level=1)
        doc.add_paragraph("Unit: [Unit Name] | Design Capacity: [Capacity]")
        doc.add_paragraph("Crude Slate: [Crude Type] / [Blend] (API: [API]°, Sulphur: [Sulphur] wt%).")

        doc.add_heading("2. Pre-Startup Safety Review (PSSR) Punch-list Clearance", level=1)
        doc.add_paragraph("• Category A Punch Items: 0 Open (Mandatory for hydrocarbon introduction).")
        doc.add_paragraph("• Nitrogen Purging & Blind De-isolation: System O₂ < 0.5% verified by gas chromatography.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_dop_process_notesheet_docx(self, filename: str = "NoteSheet-DoP-Proc-01.docx") -> str:
        """Process Change & Operating Envelope Modification Note Sheet (DoP Schedule-II)"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("NOTE SHEET — PROCESS ENVELOPE MODIFICATION (DoP SCHEME-II)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_paragraph("File Ref: MRPL/OPS/DoP-II/2026/089 | Date: " + time.strftime("%Y-%m-%d"))
        doc.add_heading("1. Proposal & Operational Background", level=1)
        doc.add_paragraph("Proposal for temporary de-rating of [Unit] throughput and re-routing of heavy naphtha cuts to Aromatic Complex.")
        doc.add_paragraph("Financial Implication: ₹18.5 Lakhs (Well within CGM DoP Schedule-II limit of ₹50 Lakhs).")

        doc.add_heading("2. DoP Competent Authority Approval Sign-off", level=1)
        doc.add_paragraph("Submitted for formal sanction under Schedule-II Sub-clause 4.3.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_production_yield_recon_xlsx(self, filename: str = "Prod_Yield_Recon_FY26.xlsx") -> str:
        """Daily Refinery Mass Balance & Blend Assay Yield Reconciliation"""
        wb = Workbook()
        ws = wb.active
        ws.title = "CDU_Yield"
        ws["A1"] = "MRPL REFINERY — DAILY MASS BALANCE & CRUDE YIELD RECONCILIATION"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Fraction / Stream", "Design Yield %", "Actual Yield %", "Actual MT/Day", "Cum MTD (MT)", "Variance Δ%", "API Gravity", "Sulphur wt%", "Status"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        streams = [
            ("LPG & Off-Gas", 2.8, 3.1, 380, 11400, "=(C3-B3)/B3*100", 0.55, 0.01, "ON-SPEC"),
            ("Light Naphtha", 7.5, 7.8, 955, 28650, "=(C4-B4)/B4*100", 72.0, 0.04, "ON-SPEC"),
            ("Motor Spirit (Gasoline)", 18.2, 17.9, 2195, 65850, "=(C5-B5)/B5*100", 58.5, 0.001, "ON-SPEC"),
            ("Aviation Turbine Fuel (ATF)", 12.0, 12.4, 1520, 45600, "=(C6-B6)/B6*100", 41.5, 0.08, "ON-SPEC"),
            ("High Speed Diesel (HSD)", 42.5, 41.8, 5120, 153600, "=(C7-B7)/B7*100", 35.0, 0.001, "ON-SPEC"),
            ("Reduced Crude / Bitumen", 17.0, 17.0, 2085, 62550, "=(C8-B8)/B8*100", 12.0, 3.50, "ON-SPEC")
        ]
        for r_idx, row in enumerate(streams, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_energy_mbi_xlsx(self, filename: str = "Energy_MBI_MBtu_per_t.xlsx") -> str:
        """Solomon Modified Energy Intensity Index (MBN) Workbook"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Energy_MBN"
        ws["A1"] = "MRPL SPECIFIC ENERGY CONSUMPTION & SOLOMON MBN DASHBOARD"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Energy Source", "Daily Consumption", "Unit", "NCV / Factor", "Heat Energy (MMBtu)", "Specific Energy (MBN)", "Benchmark"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        data = [
            ("Refinery Fuel Gas", 420.0, "MT/day", 46.5, "=B3*D3*0.9478", 54.2, 58.0),
            ("Natural Gas (Imported)", 180000.0, "Sm³/day", 38.2, "=B4*D4*0.0009478", 22.1, 25.0),
            ("Grid Electricity", 340.0, "MWh/day", 3.412, "=B5*D5", 8.4, 9.5)
        ]
        for r_idx, row in enumerate(data, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_tar_kickoff_pptx(self, filename: str = "Refinery_M&I_TAR_Kickoff.pptx") -> str:
        """Refinery Turnaround M&I Master Kickoff Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL M&I REFINERY TURNAROUND KICKOFF"
        s1.placeholders[1].text = "Master Shutdown Schedule & Critical Path Execution Plan"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Turnaround Scope & High-Risk Milestone Gates"
        tf2 = s2.placeholders[1].text_frame
        tf2.text = "• Shutdown Duration: [Duration] Days (Critical Path: [Critical Path])."
        tf2.add_paragraph().text = "• Total Equipment In-Scope: 142 Vessels, 88 Heat Exchangers, 320 PSVs, 18 Control Valves."
        tf2.add_paragraph().text = "• Safety Goal: Zero Lost Time Injuries (Goal Zero)."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 3. MECHANICAL, ELECTRICAL & ASSET RELIABILITY (M&I)
    # =========================================================================

    def generate_fiir_vessel_inspection_docx(self, filename: str = "FIIR-SE-01-API510.docx") -> str:
        """Static Pressure Vessel Final Internal Inspection Report (API 510)"""
        doc = Document()
        title = doc.add_heading("MRPL ASSET INTEGRITY & INSPECTION DIVISION", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("FINAL INTERNAL INSPECTION REPORT — API 510 (FIIR-SE-01)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Equipment & Metallurgy Particulars", level=1)
        doc.add_paragraph("Vessel Tag: [Tag] ([Service]) | Material: [Material]")
        doc.add_paragraph("Design Temp / Press: 220°C @ 4.5 bar-g | Joint Efficiency E = 1.0 (Full Radiography)")

        doc.add_heading("2. Ultrasonic Thickness & Corrosion Rate Evaluation", level=1)
        doc.add_paragraph("• Nominal Thickness: 12.0 mm | Minimum Required (t_min): 8.4 mm")
        doc.add_paragraph("• Actual Measured Thickness (t_act): 10.2 mm | Long-Term Corrosion Rate: 0.12 mm/year")
        doc.add_paragraph("• Remaining Life: (10.2 - 8.4) / 0.12 = 15.0 Years | Next Inspection Due: 7.5 Years (API 510 Half-Life)")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_pump_runin_cert_docx(self, filename: str = "ROT-RUN-API610-Cert.docx") -> str:
        """API 610 Centrifugal Pump Run-in & Acceptance Certificate"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("API 610 ROTATING EQUIPMENT RUN-IN & ACCEPTANCE CERTIFICATE")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. 4-Hour Continuous Mechanical Run Test Record", level=1)
        doc.add_paragraph("Pump Tag: [Equipment Tag] | Motor Rating: [Rating] | Speed: [Speed]")
        doc.add_paragraph("• Overall Vibration (DE/NDE Bearings): 1.85 mm/s RMS (ISO 10816 Zone A Compliant, <2.8 mm/s limit).")
        doc.add_paragraph("• Bearing Temperature Rise: DE 58°C, NDE 54°C (Limit: <75°C).")
        doc.add_paragraph("• Mechanical Seal Plan 53A: Zero leakage observed during full pressure hold.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_api570_ut_cr_register_xlsx(self, filename: str = "API570_UT_CR_Register.xlsx") -> str:
        """API 570 Piping UT Survey, Short/Long-Term CR & Remaining Life Workbook"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Corrosion_Rate_Calc"

        header_fill = PatternFill(start_color="1A365D", end_color="1A365D", fill_type="solid")
        header_font = Font(name="Calibri", size=10, bold=True, color="FFFFFF")

        ws.merge_cells("A1:N1")
        ws["A1"] = "MRPL REFINERY — API 570 PIPING CIRCUIT ULTRASONIC CORROSION REGISTER"
        ws["A1"].font = Font(size=12, bold=True, color="FFFFFF")
        ws["A1"].fill = header_fill
        ws["A1"].alignment = Alignment(horizontal="center", vertical="center")

        headers = ["Line ID", "Service Fluid", "Material", "CML No", "Nominal (mm)", "t_min (mm)", "t_prev (mm)", "t_act (mm)", "ΔYears", "Short CR (mm/y)", "Long CR (mm/y)", "Gov CR", "Remaining Life (y)", "Next Insp Due"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = header_font
            ws.cell(row=2, column=c).fill = header_fill

        cmls = [
            ("L-001", "Service Fluid", "Carbon Steel", "CML-01", 9.52, 4.20, 8.10, 7.65, 2.0, "=(G3-H3)/I3", "=(E3-H3)/8.0", "=MAX(J3,K3)", "=(H3-F3)/L3", "=MIN(M3/2, 5.0)"),
            ("L-002", "Service Fluid", "CS + Clad", "CML-02", 8.18, 3.80, 7.40, 7.22, 2.0, "=(G4-H4)/I4", "=(E4-H4)/8.0", "=MAX(J4,K4)", "=(H4-F4)/L4", "=MIN(N4/2, 5.0)"),
            ("L-003", "Service Fluid", "SS 316L", "CML-03", 12.70, 5.50, 11.80, 11.50, 2.0, "=(G5-H5)/I5", "=(E5-H5)/8.0", "=MAX(J5,K5)", "=(H5-F5)/L5", "=MIN(N5/2, 3.0)")
        ]
        for r_idx, row in enumerate(cmls, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_api581_rbi_matrix_xlsx(self, filename: str = "API581_RBI_Matrix_Plant.xlsx") -> str:
        """API 581 Risk-Based Inspection 5x5 Heatmap Matrix"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Risk_Matrix_Heatmap"
        ws["A1"] = "MRPL ASSET INTEGRITY — API 581 RISK-BASED INSPECTION (RBI) MASTER REGISTER"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Equipment Tag", "Damage Mechanism", "PoF (1-5)", "CoF (1-5)", "Risk Score", "Risk Category", "Recommended Interval"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        data = [
            ("Equipment-001", "High Temp Oxidation / Creep", 4, 4, "=C3*D3", "HIGH RISK (RED)", "1.0 Year (Mandatory Online Pyrometry)"),
            ("Equipment-002", "High Pressure H2 Attack (Nelson)", 2, 5, "=C4*D4", "MEDIUM-HIGH", "2.0 Years (Advanced TOFD/PAUT)"),
            ("Equipment-003", "Cavitation / Erosion", 2, 2, "=C5*D5", "LOW RISK", "5.0 Years (Standard PM)")
        ]
        for r_idx, row in enumerate(data, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_asset_integrity_pptx(self, filename: str = "Refinery_M&I_Asset_Integrity.pptx") -> str:
        """Asset Integrity & Static Equipment Inspection Review Deck"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL ASSET RELIABILITY & INTEGRITY BRIEFING"
        s1.placeholders[1].text = "API 510/570/653 Mechanical Integrity Survey & Half-Life Inspection Planner"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Critical Static Equipment Findings"
        tf = s2.placeholders[1].text_frame
        tf.text = "• 100% of Class 1 Piping Circuits surveyed via automated Ultrasonic Testing."
        tf.add_paragraph().text = "• Zero assets below minimum required structural wall thickness (t_min)."
        tf.add_paragraph().text = "• Next statutory turnaround window locked per API 570 Table 1 intervals."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 4. MATERIALS MANAGEMENT, CONTRACTS & GeM PROCUREMENT
    # =========================================================================

    def generate_purchase_indent_docx(self, filename: str = "Indent-MM-01.docx") -> str:
        """Purchase Indent & Tender Authorization Note Sheet (ONGC IMMM)"""
        doc = Document()
        title = doc.add_heading("ONGC — MATERIALS MANAGEMENT DIVISION", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("PURCHASE INDENT & TENDER MEMORANDUM (ONGC IMMM MANUAL)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Indent & Requirement Justification", level=1)
        doc.add_paragraph("Indent Ref: ONGC/MM/IND/2026/0411 | Estimated Value: ₹1.45 Crores")
        doc.add_paragraph("Item Description: API 610 Centrifugal Pump Spares & High-Pressure Mechanical Seal Cartridges.")
        doc.add_paragraph("Procurement Mode: Open Tender via Government e-Marketplace (GeM) with Make in India Preference.")

        doc.add_heading("2. Mandatory CVC & Statutory Tender Undertakings", level=1)
        doc.add_paragraph("[X] Integrity Pact applicable (Value > ₹1.00 Crore).")
        doc.add_paragraph("[X] Rule 144(xi) GFR Land Border Sharing declaration required.")
        doc.add_paragraph("[X] Minimum 50% Local Content for Class-I Local Supplier status.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_vendor_registration_form_docx(self, filename: str = "VRF-EFT-Ver9-MRPL.docx") -> str:
        """MRPL Vendor Registration Format & EFT Details Form"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("VENDOR REGISTRATION FORMAT & EFT DETAILS (VRF-EFT-Ver9)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Vendor Commercial & Statutory Particulars", level=1)
        doc.add_paragraph("Company Name: Bharat Heavy Spares & Engineering Pvt. Ltd.")
        doc.add_paragraph("GSTIN: 29AAACB1234F1Z5 | PAN: AAACB1234F | MSME Udyam: UDYAM-KR-03-0012345")
        doc.add_paragraph("Bank IFSC Code: SBIN0001234 | Bank Account: 334455667788 (State Bank of India)")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_integrity_pact_docx(self, filename: str = "Integrity-Pact-CVC.docx") -> str:
        """CVC Mandated Integrity Pact Agreement"""
        doc = Document()
        title = doc.add_heading("ONGC / MRPL ENTERPRISE TENDER AGREEMENT", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("INTEGRITY PACT — CVC MANDATED MODEL (OFFICE ORDER 41/12/07)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Commitments of the Principal & the Bidder", level=1)
        doc.add_paragraph("1. The Principal commits to take all measures necessary to prevent corruption and ensure transparency.")
        doc.add_paragraph("2. The Bidder commits not to offer or give any bribe, gift, or advantage in exchange for tender award.")
        doc.add_paragraph("3. Independent External Monitors (IEMs) nominated by CVC shall oversee compliance and address grievances.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_gem_cst_matrix_xlsx(self, filename: str = "GeM_CST_Matrix.xlsx") -> str:
        """Comparative Statement of Tenders (CST) & Total Landed Cost (TLC) Matrix"""
        wb = Workbook()
        ws = wb.active
        ws.title = "TLC_Arrival_CST"
        ws["A1"] = "GeM TENDER EVALUATION — COMPARATIVE STATEMENT & TOTAL LANDED COST (TLC) MATRIX"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Bidder Name", "MSE Status", "MII Content %", "Base Quoted (₹)", "Freight & P&F (₹)", "GST %", "GST Amount (₹)", "Total Landed Cost (₹)", "Evaluated Rank"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        bids = [
            ("L&T Hydrocarbon Engineering", "Non-MSE", 78.0, 12000000, 250000, 18.0, "=(D3+E3)*F3/100", "=D3+E3+G3", '=RANK(H3,$H$3:$H$5,1)'),
            ("BHEL Heavy Engineering", "Class-I MII", 85.0, 11500000, 300000, 18.0, "=(D4+E4)*F4/100", "=D4+E4+G4", '=RANK(H4,$H$3:$H$5,1)'),
            ("Triveni Turbine Works", "MSE-Medium", 65.0, 12800000, 200000, 18.0, "=(D5+E5)*F5/100", "=D5+E5+G5", '=RANK(H5,$H$3:$H$5,1)')
        ]
        for r_idx, row in enumerate(bids, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_vendor_performance_rating_xlsx(self, filename: str = "Vendor_Performance_Rating.xlsx") -> str:
        """Vendor Performance Evaluation Model (40/40/10/10 Weighting)"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Rating_Formula_Engine"
        ws["A1"] = "ONGC / MRPL VENDOR PERFORMANCE RATING ENGINE"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Vendor Name", "PO Value (₹)", "Delivery Score (40%)", "Quality Score (40%)", "HSE Score (10%)", "Resp Score (10%)", "Composite Score", "Status"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        vendors = [
            ("Tata Projects Ltd", 45000000, 38.5, 39.0, 9.5, 9.0, "=C3+D3+E3+F3", '=IF(G3>=80,"APPROVED VENDOR","UNDER WATCH")'),
            ("Thermax Babcock", 28000000, 36.0, 37.5, 9.0, 8.5, "=C4+D4+E4+F4", '=IF(G4>=80,"APPROVED VENDOR","UNDER WATCH")')
        ]
        for r_idx, row in enumerate(vendors, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_scm_procurement_review_pptx(self, filename: str = "SCM_Quarterly_Procurement_Review.pptx") -> str:
        """Quarterly SCM & GeM Procurement Review Deck"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL / ONGC SCM PROCUREMENT REVIEW"
        s1.placeholders[1].text = "Quarterly GeM Spend & Vendor Performance Audit"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Procurement Key Performance Indicators"
        tf = s2.placeholders[1].text_frame
        tf.text = "• GeM Procurement Share: 91.4% of total eligible procurement executed via GeM portal."
        tf.add_paragraph().text = "• MSME Procurement Compliance: 34.2% (Mandatory Target: 25.0%)."
        tf.add_paragraph().text = "• Integrity Pact Coverage: 100% of tenders >₹1 Crore covered under CVC IEM oversight."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 5. FINANCE, ACCOUNTS & e-MEASUREMENT BOOK (e-MB)
    # =========================================================================

    def generate_ra_bill_cert_memo_docx(self, filename: str = "RA-Bill-Cert-Memo.docx") -> str:
        """Running Account (RA) Bill Claim Compilation & Certification Memo"""
        doc = Document()
        title = doc.add_heading("ONGC / MRPL FINANCE & ACCOUNTS DIVISION", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("RUNNING ACCOUNT (RA) BILL CERTIFICATION MEMORANDUM")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Bill Abstract & Contract Details", level=1)
        doc.add_paragraph("Contract Work Order No: MRPL/WO/CIVIL/2025/1102 | RA Bill No: 04 (Running)")
        doc.add_paragraph("Contractor: M/s Coastal Infrastructure Builders Pvt. Ltd.")
        doc.add_paragraph("Gross Value of Work Done: ₹84,50,000.00 | Net Payable this Bill: ₹72,67,000.00")

        doc.add_heading("2. Mandatory Statutory & Contractual Deductions", level=1)
        doc.add_paragraph("• Retention Money (5.0% capped): ₹4,22,500.00 (Held in Defect Liability Escrow).")
        doc.add_paragraph("• Income Tax TDS @ Sec 194C (2.0%): ₹1,69,000.00.")
        doc.add_paragraph("• GST TDS @ Sec 51 (2.0%): ₹1,69,000.00.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_mb_open_close_format_docx(self, filename: str = "MB_Open_Close_Format.docx") -> str:
        """Measurement Book Opening / Closing Certificate (CPWD Form 7)"""
        doc = Document()
        title = doc.add_heading("CENTRAL PUBLIC WORKS / ONGC WORKS SYSTEM", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("MEASUREMENT BOOK (e-MB) OPENING & CLOSING CERTIFICATE")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Joint Physical Measurement Verification", level=1)
        doc.add_paragraph("This is to certify that the measurements recorded from Page No. 01 to Page No. 120 have been jointly taken in the presence of the Contractor's authorized engineer and the Departmental Engineer-in-Charge (EIC).")
        doc.add_paragraph("No unauthorized overwriting or erasure exists in this Measurement Book.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_emb_ra_bill_ledger_xlsx(self, filename: str = "eMB_RA_Bill_Ledger_FY26.xlsx") -> str:
        """e-Measurement Book Dimension Log & RA Billing Ledger with Deductions"""
        wb = Workbook()
        ws = wb.active
        ws.title = "RA_Bill_Abstract"
        ws["A1"] = "ONGC / MRPL — e-MEASUREMENT BOOK (e-MB) & RA BILLING LEDGER"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["BoQ Item", "Item Description", "Unit", "DSR Rate (₹)", "Length (m)", "Breadth (m)", "Depth (m)", "Qty Executed", "Gross Amount (₹)", "TDS 194C (2%)", "GST TDS (2%)", "Retention (5%)", "Net Payable (₹)"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        items = [
            ("BoQ-01", "RCC M-30 Equipment Foundation", "m³", 8500.0, 12.0, 6.5, 1.8, "=E3*F3*G3", "=D3*H3", "=I3*0.02", "=I3*0.02", "=I3*0.05", "=I3-J3-K3-L3"),
            ("BoQ-02", "Structural Steel Fabrication", "MT", 72000.0, 1.0, 1.0, 24.5, "=G4", "=D4*H4", "=I4*0.02", "=I4*0.02", "=I4*0.05", "=I4-J4-K4-L4")
        ]
        for r_idx, row in enumerate(items, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_contract_cashflow_xlsx(self, filename: str = "ContractCashFlow_Milestone.xlsx") -> str:
        """Contract Earned Value & Milestone Linked Cash Flow Tracker"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Milestone_Tracker"
        ws["A1"] = "MRPL CONTRACT CASH FLOW & EARNED VALUE ANALYSIS (EVA)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Contract Ref", "Contractor", "Total Value (₹)", "Physical Progress %", "Financial Progress %", "Planned Value BCWS", "Earned Value BCWP", "Actual Cost ACWP", "Schedule Var SV", "Cost Var CV"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        sample = [
            ("MRPL-C-88", "Punj Lloyd Ltd", 150000000, 75.0, 70.0, 120000000, "=C3*D3/100", 110000000, "=G3-F3", "=G3-H3")
        ]
        for r_idx, row in enumerate(sample, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_finance_cashflow_pptx(self, filename: str = "Finance_Contract_CashFlow_Review.pptx") -> str:
        """Quarterly Contract Billing & Cash Flow Review Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL / ONGC CONTRACT CASH FLOW & e-MB AUDIT"
        s1.placeholders[1].text = "Quarterly Vendor Payments & Retention Ledger Review"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "e-MB Digital Billing Highlights"
        tf = s2.placeholders[1].text_frame
        tf.text = "• 100% of contractor invoices processed through authenticated e-Measurement Book."
        tf.add_paragraph().text = "• Average Bill Cycle Time: 6.8 Days from measurement entry to NEFT release."
        tf.add_paragraph().text = "• Total Retention Escrow: ₹18.4 Crores secured against Defect Liability Periods."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 6. ESG, ENVIRONMENTAL CELL & BRSR SUSTAINABILITY
    # =========================================================================

    def generate_brsr_core_annex_docx(self, filename: str = "BRSR-Core-Annex-Principlewise.docx") -> str:
        """SEBI BRSR Core Principles 1-9 Narrative Annexure"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("BUSINESS RESPONSIBILITY & SUSTAINABILITY REPORT (SEBI BRSR CORE FY25)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Principle 6: Environmental Performance & Net-Zero 2038 Target", level=1)
        doc.add_paragraph("• Total GHG Emissions (Scope 1 + Scope 2): 1.42 Million tCO₂e.")
        doc.add_paragraph("• Specific GHG Intensity: 0.098 tCO₂e per tonne of crude processed.")
        doc.add_paragraph("• Decarbonization Roadmap: ₹2,400 Crore capital commitment toward Green Hydrogen and 20MW Solar Farms.")

        doc.add_heading("2. Principle 3: Employee Well-being & Zero Harassment", level=1)
        doc.add_paragraph("100% of permanent and contract workforce covered under comprehensive health and accident insurance.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_ec_compliance_docx(self, filename: str = "EC-Compliance-Half-Yearly.docx") -> str:
        """Environmental Clearance (EC) 6-Monthly Compliance Report"""
        doc = Document()
        title = doc.add_heading("MRPL ENVIRONMENTAL MANAGEMENT CELL", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("SIX-MONTHLY ENVIRONMENTAL CLEARANCE COMPLIANCE REPORT (MoEF&CC)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Continuous Ambient Air & Stack Emissions Status", level=1)
        doc.add_paragraph("• Stack SO₂ Emissions: 142 mg/Nm³ (Strictly below CPCB limit of 250 mg/Nm³).")
        doc.add_paragraph("• Ambient PM₂.₅: 38.4 µg/m³ (CPCB Limit: 60 µg/m³).")
        doc.add_paragraph("• Treated Effluent BOD: 12 mg/L (Limit: 30 mg/L) — 100% recycled back into cooling water makeup.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_ghg_inventory_xlsx(self, filename: str = "GHG_Inventory_Scope1_Scope2.xlsx") -> str:
        """GHG Protocol Scope 1 & Scope 2 Emissions Inventory Workbook"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Scope1_Stationary"
        ws["A1"] = "MRPL REFINERY — GREENHOUSE GAS (GHG) INVENTORY (SCOPE 1 & SCOPE 2)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Emission Source", "Fuel Type", "Activity Quantity", "Unit", "Net Calorific Value", "Emission Factor (kg CO2/MMBtu)", "Scope 1 tCO2e", "Data Quality"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        emissions = [
            ("Fired Heater Equipment-001", "Refinery Fuel Gas", 18500.0, "MT/yr", 46.5, 53.06, "=C3*D3*E3*F3/1000000", "Continuous Metered"),
            ("Captive Power Plant Boilers", "Natural Gas", 45000.0, "kSm³/yr", 38.2, 53.06, "=C4*D4*E4*F4/1000000", "Continuous Metered"),
            ("Emergency DG Sets", "HSD (Diesel)", 120.0, "kL/yr", 36.8, 74.02, "=C5*D5*E5*F5/1000000", "Batch Logged")
        ]
        for r_idx, row in enumerate(emissions, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_water_waste_brsr_xlsx(self, filename: str = "Energy_Water_Waste_BRSR_P6.xlsx") -> str:
        """Water Neutrality & Waste Management Principle 6 Dashboard"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Water_Drawal_Discharge"
        ws["A1"] = "MRPL WATER BALANCE & EFFLUENT RECYCLING REGISTER (BRSR PRINCIPLE 6)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Water Source", "Fresh Intake (ML)", "Treated Effluent (ML)", "Recycled Makeup (ML)", "Specific Water Intensity (m3/t)", "Recycle %"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        data = [
            ("Netravati River Pipeline", 8500.0, 0.0, 0.0, 0.58, "-"),
            ("Municipal Sewage STP Treated", 4200.0, 3950.0, 3950.0, 0.28, "=D4/B4*100"),
            ("Desalination Plant", 2100.0, 1950.0, 1950.0, 0.14, "=D5/B5*100")
        ]
        for r_idx, row in enumerate(data, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_board_brsr_pptx(self, filename: str = "Board_BRSR_Sustainability_Review.pptx") -> str:
        """Annual Board BRSR Sustainability & Net-Zero Strategy Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL BOARD SUSTAINABILITY REVIEW"
        s1.placeholders[1].text = "Business Responsibility and Sustainability Report (SEBI BRSR Core FY25)"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Decarbonization & Water Neutrality Milestones"
        tf = s2.placeholders[1].text_frame
        tf.text = "• Scope 1 & 2 Emissions reduced by 4.2% Year-on-Year."
        tf.add_paragraph().text = "• Net-Zero 2038 Target: Signed onto Oil & Gas Decarbonization Charter (OGDC)."
        tf.add_paragraph().text = "• Water Neutrality: 72.8% of refinery makeup water drawn from treated municipal sewage."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 7. STATUTORY CAG & INTERNAL AUDIT COMPLIANCE
    # =========================================================================

    def generate_internal_audit_charter_docx(self, filename: str = "IA-Charter-Annual.docx") -> str:
        """Internal Audit Charter (IIA Global Standards / CVC Guidelines)"""
        doc = Document()
        title = doc.add_heading("ONGC / MRPL AUDIT COMMITTEE OF BOARD", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("INTERNAL AUDIT CHARTER (IIA GLOBAL STANDARDS & CVC)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Purpose, Authority & Dual Reporting", level=1)
        doc.add_paragraph("The Chief Internal Auditor reports functionally to the Audit Committee of the Board (ACB) and administratively to the Managing Director.")
        doc.add_paragraph("Internal Audit has unrestricted access to all company systems, personnel, and physical refinery assets.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_stat_audit_reply_docx(self, filename: str = "StatAuditPara-Reply.docx") -> str:
        """CAG Supplementary Audit Observation Response Note Sheet"""
        doc = Document()
        title = doc.add_heading("MANGALORE REFINERY AND PETROCHEMICALS LIMITED", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("STATUTORY CAG AUDIT OBSERVATION RESPONSE NOTE")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. CAG Audit Observation (Para 14.2)", level=1)
        doc.add_paragraph("Observation: Installation and calibration intervals of corrosion monitoring facilities (ER probes & UT sensors).")
        doc.add_paragraph("Company Response: 100% of designated corrosion probes calibrated and integrated into real-time DCS telemetry per OISD-156.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_cag_para_tracker_xlsx(self, filename: str = "AuditParaTracker_CAG_Internal.xlsx") -> str:
        """CAG SAR & PAC Audit Para Ageing & Financial Risk Tracker"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Paras_Master"
        ws["A1"] = "MRPL CAG & STATUTORY AUDIT OBSERVATION TRACKER"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Para ID", "Audit Body", "Subject", "Financial Risk (₹ Cr)", "Action Owning HoD", "Target Date", "Closure Status", "Ageing Days"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        paras = [
            ("CAG-2025-01", "CAG SAR", "Furnace Decoking Cycle Optimization", 1.85, "CGM Operations", "2026-06-30", "IN PROGRESS", 45),
            ("INT-2025-04", "Internal Audit", "e-MB Contractor Billing Reconciliation", 0.42, "Controller F&A", "2026-03-31", "CLOSED-VERIFIED", 0)
        ]
        for r_idx, row in enumerate(paras, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_statutory_compliance_calendar_xlsx(self, filename: str = "StatutoryCompliance_Calendar.xlsx") -> str:
        """Multi-Act Statutory Compliance Master Calendar"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Compliance_Master"
        ws["A1"] = "ONGC / MRPL STATUTORY COMPLIANCE CALENDAR"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Compliance Item", "Governing Statute", "Filing Frequency", "Last Filed Date", "Next Due Date", "Responsible Officer", "Status RAG"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        items = [
            ("Annual Safety Return", "Factories Act 1948", "Annual", "2025-01-15", "2026-01-15", "DGM HSE", "GREEN (COMPLIANT)"),
            ("Water Cess Return", "Water Act 1974", "Monthly", "2026-01-31", "2026-02-28", "Manager Env", "GREEN (COMPLIANT)"),
            ("Contractor PF Return", "EPF Act 1952", "Monthly", "2026-01-15", "2026-02-15", "HR CONTRACT", "GREEN (COMPLIANT)")
        ]
        for r_idx, row in enumerate(items, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_acb_board_pptx(self, filename: str = "Audit_Committee_Board_Update.pptx") -> str:
        """Half-Yearly Audit Committee of Board (ACB) Update Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL AUDIT COMMITTEE OF THE BOARD"
        s1.placeholders[1].text = "Statutory CAG & Internal Audit Progress Update"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Audit Universe & Para Closure Status"
        tf = s2.placeholders[1].text_frame
        tf.text = "• Risk-Based Internal Audit Plan: 88% of planned audits completed on schedule."
        tf.add_paragraph().text = "• CAG Audit Paras: 0 Significant Non-Conformances pending beyond 90 days."
        tf.add_paragraph().text = "• Internal Financial Controls (IFC): Certified effective with zero material weaknesses."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 8. HUMAN RESOURCES, LABOUR WELFARE & CONTRACTOR REGULATIONS
    # =========================================================================

    def generate_contract_worker_induction_docx(self, filename: str = "CW-Induction-Authorisation.docx") -> str:
        """Contract Worker Induction & Safety Pass Authorization Letter"""
        doc = Document()
        title = doc.add_heading("MRPL CONTRACTOR SAFETY CELL", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("CONTRACT WORKER SAFETY INDUCTION & RFID PASS CLEARANCE")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Worker & Medical Fitness Particulars", level=1)
        doc.add_paragraph("Worker Name: Anil Kumar | UAN / Aadhaar: 100988776655 | Blood Group: O+ve")
        doc.add_paragraph("Contractor Agency: M/s Coastal Plant Services Pvt. Ltd. | Trade: High-Pressure Welder (6G)")
        doc.add_paragraph("• Pre-Employment Medical Examination: Fit for Hydrocarbon Plant & Height Work (>1.8m).")
        doc.add_paragraph("• 8-Hour Safety Induction & JSA Walkthrough: Completed with 100% test score.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_clra_form_vi_docx(self, filename: str = "CLRA-Form-VI-Application.docx") -> str:
        """Contract Labour License Form VI/VII Application Packet (CLRA 1970)"""
        doc = Document()
        title = doc.add_heading("CENTRAL LABOUR COMMISSIONERATE", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("FORM VI — APPLICATION FOR CONTRACT LABOUR LICENSE (CLRA 1970)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Principal Employer & Establishment Details", level=1)
        doc.add_paragraph("Principal Employer: Mangalore Refinery and Petrochemicals Limited (MRPL)")
        doc.add_paragraph("Registration Certificate No: CLRA/PE/MRPL/2018/009 | Maximum Permitted Workers: 1,500")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_clra_wage_compliance_xlsx(self, filename: str = "CLRA_Wage_Compliance_Register.xlsx") -> str:
        """Contractor Form XII/XIII Deployment, Minimum Wage & EPF/ESI Register"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Monthly_Wage_Payment_Tracker"
        ws["A1"] = "MRPL CONTRACT LABOUR WAGE, EPF & ESI COMPLIANCE REGISTER (FORM XIII)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Contractor Name", "Work Order Ref", "Workers Deployed", "Wage Rate Paid (₹/day)", "Central Min Wage (₹/day)", "Wage Deficiency Flag", "Total Wages Paid (₹)", "EPF (12% ₹)", "ESI (3.25% ₹)", "Status"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        contractors = [
            ("Apex Industrial Services", "MRPL-M-01", 120, 945.0, 890.0, '=IF(D3<E3,"DEFICIENT","COMPLIANT")', "=C3*D3*26", "=G3*0.12", "=G3*0.0325", "VERIFIED"),
            ("Coastal Mech Maintenance", "MRPL-M-02", 85, 945.0, 890.0, '=IF(D4<E4,"DEFICIENT","COMPLIANT")', "=C4*D4*26", "=G4*0.12", "=G4*0.0325", "VERIFIED")
        ]
        for r_idx, row in enumerate(contractors, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_hr_ir_review_pptx(self, filename: str = "HR_IR_Annual_Review.pptx") -> str:
        """HR Committee Workforce Demographics & Contractor Welfare Review Deck"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL HUMAN RESOURCES & INDUSTRIAL RELATIONS"
        s1.placeholders[1].text = "Annual Workforce Health, CLRA Welfare & Training Review"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Labour Welfare & Minimum Wage Adherence"
        tf = s2.placeholders[1].text_frame
        tf.text = "• 100% compliance with Central Minimum Wages notified rates."
        tf.add_paragraph().text = "• EPF / ESI Direct Bank Remittance verified for all 2,400+ contractor personnel."
        tf.add_paragraph().text = "• Industrial Relations: Zero man-days lost due to industrial dispute in FY25."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # 9. VIGILANCE, ETHICS, WHISTLE BLOWER & ANTI-BRIBERY
    # =========================================================================

    def generate_wb_complaint_reg_docx(self, filename: str = "WB-Complaint-Reg-01.docx") -> str:
        """Whistle Blower Complaint Registration Note (CVC PIDPI Resolution)"""
        doc = Document()
        title = doc.add_heading("ONGC / MRPL VIGILANCE DIVISION", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub = doc.add_paragraph("CONFIDENTIAL WHISTLE BLOWER COMPLAINT RECORD (PIDPI / CVC)")
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub.runs[0].font.bold = True

        doc.add_heading("1. Confidential Receipt Metadata", level=1)
        doc.add_paragraph("Registration Code: VIG/WB/2026/014 (Masked under PIDPI Identity Protection)")
        doc.add_paragraph("Subject: Allegation regarding tender qualification in materials procurement.")
        doc.add_paragraph("Action Taken: Forwarded to Independent External Monitor (IEM) and Preliminary Enquiry (PE) team.")

        out_path = os.path.join(OUTPUT_DIR, "docx", filename)
        doc.save(out_path)
        return out_path

    def generate_vigilance_case_register_xlsx(self, filename: str = "Vigilance_Case_Register_FY.xlsx") -> str:
        """CVC Preliminary Enquiry & Disciplinary Case Management Register"""
        wb = Workbook()
        ws = wb.active
        ws.title = "Cases_Master"
        ws["A1"] = "ONGC / MRPL VIGILANCE CASE MANAGEMENT REGISTER (CVC GUIDELINES)"
        ws["A1"].font = Font(size=12, bold=True)

        headers = ["Case ID", "Domain", "Receipt Date", "Stage", "Value at Stake (₹ Lakhs)", "IEM Involved", "Status RAG"]
        for c, h in enumerate(headers, 1):
            ws.cell(row=2, column=c, value=h).font = Font(bold=True)

        cases = [
            ("VIG-2025-01", "Procurement / GeM", "2025-11-10", "Detailed Investigation (DI)", 45.0, "YES", "ACTIVE INQUIRY"),
            ("VIG-2025-02", "Works / Civil", "2025-12-04", "Preliminary Enquiry (PE)", 12.5, "NO", "STAGE-1 ADVICE")
        ]
        for r_idx, row in enumerate(cases, 3):
            for c_idx, val in enumerate(row, 1):
                ws.cell(row=r_idx, column=c_idx, value=val)

        out_path = os.path.join(OUTPUT_DIR, "xlsx", filename)
        wb.save(out_path)
        return out_path

    def generate_vigilance_awareness_pptx(self, filename: str = "Vigilance_Awareness_Week_Kickoff.pptx") -> str:
        """Annual Vigilance Awareness Week Kickoff Presentation"""
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)

        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s1.shapes.title.text = "MRPL / ONGC VIGILANCE AWARENESS WEEK"
        s1.placeholders[1].text = "Preventive Vigilance & Transparent Public Procurement Governance"

        s2 = prs.slides.add_slide(prs.slide_layouts[1])
        s2.shapes.title.text = "Preventive Vigilance Initiatives"
        tf = s2.placeholders[1].text_frame
        tf.text = "• 100% E-Tendering & GeM Reverse Auction deployment."
        tf.add_paragraph().text = "• Integrity Pact operationalized for all major tenders >₹1 Crore."
        tf.add_paragraph().text = "• Online Vendor Grievance Redressal portal active with 15-day SLA."

        out_path = os.path.join(OUTPUT_DIR, "pptx", filename)
        prs.save(out_path)
        return out_path

    # =========================================================================
    # BATCH GENERATION & DISPATCH HELPER
    # =========================================================================

    # generate_all_enterprise_deliverables removed — no demo file generation
        generated.append(self.generate_scm_procurement_review_pptx())
        # Finance
        generated.append(self.generate_ra_bill_cert_memo_docx())
        generated.append(self.generate_mb_open_close_format_docx())
        generated.append(self.generate_emb_ra_bill_ledger_xlsx())
        generated.append(self.generate_contract_cashflow_xlsx())
        generated.append(self.generate_finance_cashflow_pptx())
        # ESG
        generated.append(self.generate_brsr_core_annex_docx())
        generated.append(self.generate_ec_compliance_docx())
        generated.append(self.generate_ghg_inventory_xlsx())
        generated.append(self.generate_water_waste_brsr_xlsx())
        generated.append(self.generate_board_brsr_pptx())
        # Audit
        generated.append(self.generate_internal_audit_charter_docx())
        generated.append(self.generate_stat_audit_reply_docx())
        generated.append(self.generate_cag_para_tracker_xlsx())
        generated.append(self.generate_statutory_compliance_calendar_xlsx())
        generated.append(self.generate_acb_board_pptx())
        # HR
        generated.append(self.generate_contract_worker_induction_docx())
        generated.append(self.generate_clra_form_vi_docx())
        generated.append(self.generate_clra_wage_compliance_xlsx())
        generated.append(self.generate_hr_ir_review_pptx())
        # Vigilance
        generated.append(self.generate_wb_complaint_reg_docx())
        generated.append(self.generate_vigilance_case_register_xlsx())
        generated.append(self.generate_vigilance_awareness_pptx())

        # === EXTENDED ~203 GENERATORS ===
        # E&P (ONGC)
        generated.append(self.generate_well_control_docx())
        generated.append(self.generate_cementing_report_docx())
        generated.append(self.generate_drilling_daily_report_docx())
        generated.append(self.generate_wireline_log_report_docx())
        generated.append(self.generate_workover_report_docx())
        generated.append(self.generate_geophysical_ops_docx())
        generated.append(self.generate_simops_report_docx())
        generated.append(self.generate_offshore_pipeline_inspection_docx())
        generated.append(self.generate_oilfield_explosive_safety_docx())
        generated.append(self.generate_pipeline_corrosion_monitoring_docx())
        generated.append(self.generate_production_safety_report_docx())
        generated.append(self.generate_ep_safety_pptx())
        # Project & Construction
        generated.append(self.generate_detailed_engineering_pkg_docx())
        generated.append(self.generate_as_built_drawings_docx())
        generated.append(self.generate_mechanical_completion_cert_docx())
        generated.append(self.generate_pre_commissioning_checklist_docx())
        generated.append(self.generate_performance_test_run_docx())
        generated.append(self.generate_punch_list_xlsx())
        generated.append(self.generate_construction_safety_audit_docx())
        generated.append(self.generate_project_completion_report_docx())
        generated.append(self.generate_tef_report_docx())
        generated.append(self.generate_sow_docx())
        generated.append(self.generate_rfp_docx())
        generated.append(self.generate_rfq_docx())
        generated.append(self.generate_bom_docx())
        generated.append(self.generate_material_submittal_docx())
        generated.append(self.generate_qaqc_plan_docx())
        generated.append(self.generate_project_review_pptx())
        # Quality (ISO)
        generated.append(self.generate_iso9001_qms_manual_docx())
        generated.append(self.generate_iso14001_ems_manual_docx())
        generated.append(self.generate_iso45001_ohs_manual_docx())
        generated.append(self.generate_iso50001_enms_docx())
        generated.append(self.generate_ims_policy_docx())
        generated.append(self.generate_internal_iso_audit_report_docx())
        generated.append(self.generate_capa_report_docx())
        generated.append(self.generate_ncr_report_docx())
        generated.append(self.generate_mgmt_review_minutes_docx())
        generated.append(self.generate_quality_review_pptx())
        # Info Security
        generated.append(self.generate_infosec_policy_docx())
        generated.append(self.generate_data_classification_docx())
        generated.append(self.generate_privacy_policy_docx())
        generated.append(self.generate_it_asset_register_docx())
        generated.append(self.generate_access_control_log_docx())
        generated.append(self.generate_incident_response_plan_docx())
        generated.append(self.generate_infosec_review_pptx())
        # Land & Legal
        generated.append(self.generate_land_lease_agreement_docx())
        generated.append(self.generate_dealership_agreement_docx())
        generated.append(self.generate_technology_license_docx())
        generated.append(self.generate_shareholders_agreement_docx())
        generated.append(self.generate_mou_jv_agreement_docx())
        generated.append(self.generate_crude_purchase_agreement_docx())
        generated.append(self.generate_ppa_agreement_docx())
        generated.append(self.generate_pipeline_row_agreement_docx())
        generated.append(self.generate_indemnity_bond_docx())
        generated.append(self.generate_arbitration_record_docx())
        generated.append(self.generate_legal_review_pptx())
        # Extended HSE
        generated.append(self.generate_cold_work_permit_docx())
        generated.append(self.generate_electrical_isolation_permit_docx())
        generated.append(self.generate_excavation_permit_docx())
        generated.append(self.generate_radiation_safety_permit_docx())
        generated.append(self.generate_lifting_permit_docx())
        generated.append(self.generate_work_at_height_permit_docx())
        generated.append(self.generate_vehicle_entry_permit_docx())
        generated.append(self.generate_loto_certificate_docx())
        generated.append(self.generate_safety_induction_cert_docx())
        generated.append(self.generate_gas_testing_certificate_docx())
        generated.append(self.generate_safety_manual_docx())
        generated.append(self.generate_drill_report_docx())
        generated.append(self.generate_jsa_document_docx())
        generated.append(self.generate_sds_sheet_docx())
        generated.append(self.generate_incident_investigation_docx())
        generated.append(self.generate_near_miss_report_docx())
        generated.append(self.generate_ppe_compliance_register_xlsx())
        generated.append(self.generate_contractor_worker_safety_docx())
        generated.append(self.generate_ptw_register_xlsx())
        generated.append(self.generate_hse_composite_permit_pptx())
        # Extended Refinery Ops
        generated.append(self.generate_operating_manual_docx())
        generated.append(self.generate_pid_document_docx())
        generated.append(self.generate_cause_effect_diagram_docx())
        generated.append(self.generate_hmb_sheet_docx())
        generated.append(self.generate_moc_document_docx())
        generated.append(self.generate_pssr_checklist_docx())
        generated.append(self.generate_psmdocument_docx())
        generated.append(self.generate_sop_document_docx())
        generated.append(self.generate_turnaround_plan_docx())
        generated.append(self.generate_design_basis_docx())
        generated.append(self.generate_operations_review_pptx())
        # Extended Mechanical
        generated.append(self.generate_corrosion_control_docx())
        generated.append(self.generate_vibration_report_docx())
        generated.append(self.generate_pm_schedule_docx())
        generated.append(self.generate_equipment_history_docx())
        generated.append(self.generate_hx_inspection_report_docx())
        generated.append(self.generate_tank_inspection_report_docx())
        generated.append(self.generate_boiler_inspection_report_docx())
        generated.append(self.generate_fired_heater_inspection_docx())
        generated.append(self.generate_psv_test_certificate_docx())
        generated.append(self.generate_electrical_equipment_inspection_docx())
        generated.append(self.generate_insulation_inspection_docx())
        generated.append(self.generate_painting_inspection_docx())
        generated.append(self.generate_ndt_report_docx())
        generated.append(self.generate_idle_equipment_preservation_docx())
        generated.append(self.generate_maintenance_review_pptx())
        # Extended Materials
        generated.append(self.generate_tender_evaluation_docx())
        generated.append(self.generate_purchase_order_docx())
        generated.append(self.generate_contract_agreement_docx())
        generated.append(self.generate_bank_guarantee_docx())
        generated.append(self.generate_delivery_challan_docx())
        generated.append(self.generate_gem_procurement_review_docx())
        generated.append(self.generate_pqc_document_docx())
        generated.append(self.generate_materials_review_pptx())
        # Extended Finance
        generated.append(self.generate_annual_financial_stmt_docx())
        generated.append(self.generate_statutory_audit_report_docx())
        generated.append(self.generate_internal_audit_report_docx())
        generated.append(self.generate_tax_audit_report_docx())
        generated.append(self.generate_cost_audit_report_docx())
        generated.append(self.generate_income_tax_return_docx())
        generated.append(self.generate_board_report_docx())
        generated.append(self.generate_annual_return_mgt7_docx())
        generated.append(self.generate_register_contracts_docx())
        generated.append(self.generate_register_charges_docx())
        generated.append(self.generate_dividend_distribution_docx())
        generated.append(self.generate_finance_review_pptx())
        # Extended ESG
        generated.append(self.generate_crz_clearance_docx())
        generated.append(self.generate_hazwaste_authorization_docx())
        generated.append(self.generate_biomedical_waste_auth_docx())
        generated.append(self.generate_ewaste_return_docx())
        generated.append(self.generate_energy_audit_report_docx())
        generated.append(self.generate_biodiversity_register_docx())
        generated.append(self.generate_eia_report_docx())
        generated.append(self.generate_esg_review_pptx())
        # Extended Audit
        generated.append(self.generate_governance_report_docx())
        generated.append(self.generate_rpt_policy_docx())
        generated.append(self.generate_vigil_mechanism_docx())
        generated.append(self.generate_insider_trading_code_docx())
        generated.append(self.generate_csr_policy_docx())
        generated.append(self.generate_secretarial_compliance_docx())
        generated.append(self.generate_material_subsidiary_policy_docx())
        generated.append(self.generate_risk_management_docx())
        generated.append(self.generate_ifc_report_docx())
        generated.append(self.generate_statutory_returns_docx())
        generated.append(self.generate_audit_review_pptx())
        # Extended HR
        generated.append(self.generate_hr_policy_docx())
        generated.append(self.generate_transfer_policy_docx())
        generated.append(self.generate_promotion_policy_docx())
        generated.append(self.generate_medical_benefit_policy_docx())
        generated.append(self.generate_standing_orders_docx())
        generated.append(self.generate_human_rights_policy_docx())
        generated.append(self.generate_diversity_report_docx())
        generated.append(self.generate_training_register_docx())
        generated.append(self.generate_hr_review_pptx())
        # Extended Vigilance
        generated.append(self.generate_anti_bribery_policy_docx())
        generated.append(self.generate_gift_policy_docx())
        generated.append(self.generate_vigilance_clearance_docx())
        generated.append(self.generate_cvc_returns_docx())
        generated.append(self.generate_apr_return_docx())
        generated.append(self.generate_independent_directors_code_docx())
        generated.append(self.generate_vigilance_review_pptx())
        return generated

# Global Singleton
deliverable_generator = DeliverableGenerator()
