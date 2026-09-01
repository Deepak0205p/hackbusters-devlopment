import os
from datetime import datetime
from docx import Document
from docx.shared import Pt, RGBColor
from docx.enum.table import WD_TABLE_ALIGNMENT
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))


class ExtendedGenerators_EpProject:

    def __init__(self):
        os.makedirs(os.path.join(OUTPUT_DIR, "docx"), exist_ok=True)
        os.makedirs(os.path.join(OUTPUT_DIR, "xlsx"), exist_ok=True)
        os.makedirs(os.path.join(OUTPUT_DIR, "pptx"), exist_ok=True)

    @staticmethod
    def _h(doc, text, level=1):
        h = doc.add_heading(text, level=level)
        for r in h.runs:
            r.font.color.rgb = RGBColor(0, 51, 102)
        return h

    @staticmethod
    def _t(doc, rows, cols, data=None):
        table = doc.add_table(rows=rows, cols=cols, style="Light Grid Accent 1")
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        if data:
            for i, rd in enumerate(data):
                for j, ct in enumerate(rd):
                    table.rows[i].cells[j].text = str(ct)
        return table

    @staticmethod
    def _m(doc, meta):
        for label, value in meta.items():
            p = doc.add_paragraph()
            rl = p.add_run(f"{label}: ")
            rl.bold = True
            rl.font.size = Pt(10)
            rv = p.add_run(str(value))
            rv.font.size = Pt(10)

    @staticmethod
    def _p(doc, text):
        p = doc.add_paragraph(text)
        p.style.font.size = Pt(10)
        return p

    @staticmethod
    def _n(doc, items):
        for item in items:
            doc.add_paragraph(item, style="List Number")

    def generate_well_control_docx(self, filename=None):
        fn = filename or f"OISD174_Well_Control_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-174 Well Control Operations Document", 0)
        self._m(d, {"Document No": "MRPL/E&P/WC/2026-0041", "Well": "MH-4A-12", "Platform": "Mumbai High North",
                     "MASP": "520 bar", "Formation Pressure": "385 bar @ 2850 m TVDSS", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Purpose & Scope", 1)
        self._p(d, "This document establishes well control procedures for Well MH-4A-12 per OISD-STD-174. Covers BOP configuration, mud weight windows, kick detection, kill procedures, and contingency planning.")
        self._h(d, "2. BOP Stack Configuration", 1)
        self._t(d, 8, 3, [["Component", "Rating", "Status"],
                           ["18 3/4in Annular Preventer", "10,000 psi", "Certified"],
                           ["18 3/4in Double Ram (Pipe)", "10,000 psi", "Certified"],
                           ["18 3/4in Double Ram (Blind/Shear)", "10,000 psi", "Certified"],
                           ["Kill Manifold", "10,000 psi", "Certified"],
                           ["Choke Manifold", "10,000 psi", "Certified"],
                           ["BOP Control System", "Dual Accumulator - 3000 psi", "Functional"],
                           ["Diverter System", "21in - 500 psi", "Certified"]])
        self._h(d, "3. Mud Weight Window", 1)
        self._t(d, 7, 4, [["Interval (m)", "Pore Pressure (SG)", "Fracture Gradient (SG)", "Mud Weight (SG)"],
                           ["0 - 350", "Seawater", "1.25", "1.03"],
                           ["350 - 1200", "0.95", "1.40", "1.10 - 1.15"],
                           ["1200 - 2400", "1.02", "1.55", "1.15 - 1.25"],
                           ["2400 - 2700", "1.15", "1.70", "1.25 - 1.30"],
                           ["2700 - 2900", "1.22", "1.80", "1.30 - 1.32"],
                           ["2900 - 3200", "1.28", "1.85", "1.32 - 1.35"]])
        self._h(d, "4. Kill Procedures", 1)
        self._n(d, ["Stop drilling, pick up off bottom", "Shut-in using hard shut-in method",
                     "Record SIDPP, SICP, and pit gain", "Circulate at original mud weight to remove kick",
                     "Maintain constant BHP equal to formation pressure", "Kill weight = Original MW + (SIDPP / (0.0981 x TVD))"])
        self._h(d, "5. Sign-Off", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Drilling Superintendent", "Shri R.K. Sharma", "________________"],
                           ["Well Control Lead", "Shri A.V. Deshmukh", "________________"], ["MRPL Drilling Manager", "Shri P.K. Mehra", "________________"]])
        d.save(fp)
        return fp

    def generate_cementing_report_docx(self, filename=None):
        fn = filename or f"OISD175_Cementing_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-175 Cementing Operations Report", 0)
        self._m(d, {"Document No": "MRPL/CEM/RPT/2026-0118", "Well": "MH-4A-12", "Casing": "9 5/8in @ 2900 m",
                     "Cement Type": "Class G + Silica Flour", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Slurry Design", 1)
        self._t(d, 8, 2, [["Property", "Specification"], ["Base Cement", "API Class G + 35% Silica Flour"],
                           ["Density", "1.92 - 1.95 SG"], ["Thickening Time", "> 6 hrs @ 128C/385 bar"],
                           ["Free Water", "< 1.0 mL"], ["Fluid Loss", "< 50 mL/30 min"],
                           ["24-hr Strength", "> 25 MPa"], ["Volume", "58.5 m3 (247 sacks)"]])
        self._h(d, "2. Pumping Schedule", 1)
        self._t(d, 8, 5, [["Stage", "Fluid", "Vol (m3)", "Rate (bpm)", "Time (min)"],
                           ["Preflush", "Seawater", "12.0", "10", "14"],
                           ["Spacer Heavy", "3.16 SG Spacer", "8.5", "6", "16"],
                           ["Lead Slurry", "1.92 SG Cement", "35.0", "8", "51"],
                           ["Tail Slurry", "1.95 SG Cement", "23.5", "8", "34"],
                           ["Displacement", "Seawater", "42.0", "6", "80"],
                           ["Plug Down", "Seawater", "8.0", "4", "23"],
                           ["Total", "-", "138.0", "-", "245"]])
        self._h(d, "3. CBL Results", 1)
        self._t(d, 5, 3, [["Interval", "CBL (mV)", "Assessment"], ["2400-2650", "5-12", "Excellent Bond"],
                           ["2650-2750", "10-18", "Good Bond"], ["2750-2850", "6-10", "Excellent Bond"],
                           ["Overall", "-", "ACCEPTABLE"]])
        d.save(fp)
        return fp

    def generate_drilling_daily_report_docx(self, filename=None):
        fn = filename or f"Drilling_Daily_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Drilling Daily Progress Report", 0)
        self._m(d, {"Well": "MH-4A-12", "Rig": "Sagar Ratna MRJU-04", "Hole Section": "17 1/2in (1200-2400 m)",
                     "Report No": "DR-MH4A12-2026-0038", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Drilling Parameters", 1)
        self._t(d, 8, 2, [["Parameter", "Value"], ["Depth Start", "2145 m MD"], ["Depth End", "2210 m MD"],
                           ["Interval Drilled", "65 m"], ["ROP (avg)", "8.1 m/hr"],
                           ["WOB", "12-16 klbs"], ["RPM", "120-140"], ["Flow Rate", "48 L/s"]])
        self._h(d, "2. Mud Properties", 1)
        self._t(d, 7, 3, [["Property", "Morning", "Evening"], ["MW (SG)", "1.18", "1.20"],
                           ["Funnel Visc (s/qt)", "48", "50"], ["PV (cP)", "22", "24"],
                           ["YP (lb/100ft2)", "14", "16"], ["pH", "9.5", "9.4"], ["Sand (%)", "0.4", "0.5"]])
        self._h(d, "3. Operations Summary", 1)
        self._n(d, ["06:00 - Commenced drilling at 2145 m", "08:30 - Connection at 2165 m, gas 2.1% C1",
                     "11:30 - Reamed hole to clean shale pack-off", "15:30 - Connection at 2200 m",
                     "17:00 - Reached 2210 m MD", "18:00 - End of shift"])
        d.save(fp)
        return fp

    def generate_wireline_log_report_docx(self, filename=None):
        fn = filename or f"OISD183_Wireline_Log_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-183 Wireline Log Report", 0)
        self._m(d, {"Well": "MH-4A-12", "Hole Section": "8 1/2in (2400-3200 m)", "Logging Co": "Schlumberger",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Tool Suite", 1)
        self._t(d, 7, 4, [["Tool", "Service", "Interval", "Serial"],
                           ["Platform Express", "Triple Combo", "2400-3200 m", "PE-0451"],
                           ["DSL", "Resistivity", "2400-3200 m", "DL-0328"],
                           ["CMR", "NMR Porosity", "2500-3100 m", "CMR-0192"],
                           ["Litho-Density", "Density/PEF", "2400-3200 m", "LDT-0415"],
                           ["CNL", "Neutron", "2400-3200 m", "CNL-0387"],
                           ["MDT", "Pressure/Sampling", "2550,2820 m", "MDT-0201"]])
        self._h(d, "2. QC Assessment", 1)
        self._t(d, 5, 3, [["Log", "Check", "Result"], ["GR/Caliper", "Gauge hole 8.4-8.7in", "PASS"],
                           ["Density", "Neutron crossover consistent", "PASS"],
                           ["Resistivity", "LLD>LLS>MSFL invasion", "PASS"], ["NMR", "Bimodal T2 in reservoir", "PASS"]])
        self._h(d, "3. Reservoir Intervals", 1)
        self._t(d, 4, 5, [["Interval (m)", "Net Pay (m)", "Porosity (%)", "Sw (%)", "Remarks"],
                           ["2545-2580", "28", "22.5", "18", "Oil-bearing L-I Sand"],
                           ["2710-2750", "32", "24.1", "15", "Best reservoir"],
                           ["2880-2900", "8", "15.2", "80", "Water-bearing"]])
        d.save(fp)
        return fp

    def generate_workover_report_docx(self, filename=None):
        fn = filename or f"Workover_Stimulation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Workover / Well Stimulation Report", 0)
        self._m(d, {"Well": "MH-12A-05", "Job Type": "Matrix Acidizing", "Co": "Halliburton",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Job Objective", 1)
        self._p(d, "Restore well productivity by removing near-wellbore formation damage through matrix acidizing of L-I Sand (2710-2750 m). Pre-job PI: 0.8 bopd/psi, Skin: +12.5. Target: restore to 80% of original.")
        self._h(d, "2. Stimulation Design", 1)
        self._t(d, 7, 3, [["Stage", "Fluid", "Volume"], ["Pre-Flush", "Diesel + Mutual Solvent", "8.0 m3"],
                           ["Preflush", "15% HCl + Inhibitor", "6.0 m3"],
                           ["Main Acid 1", "28% HCl + HF (12:3)", "18.0 m3"],
                           ["Main Acid 2", "28% HCl + HF (9:1)", "12.0 m3"],
                           ["Overflush", "5% NH4Cl Brine", "8.0 m3"],
                           ["Displacement", "5% NH4Cl Brine", "15.0 m3"]])
        self._h(d, "3. Post-Job Results", 1)
        self._t(d, 5, 2, [["Parameter", "Result"], ["Production Rate", "385 bopd (from 120)"],
                           ["Water Cut", "9.4% (from 22.5%)"], ["Skin Factor", "-1.2 (from +12.5)"],
                           ["PI", "3.2 bopd/psi (from 0.8)"]])
        self._p(d, "Conclusion: Job successful. 86% recovery of original PI achieved.")
        d.save(fp)
        return fp

    def generate_geophysical_ops_docx(self, filename=None):
        fn = filename or f"OISD181_Geophysical_Ops_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-181 Geophysical Operations Report", 0)
        self._m(d, {"Survey": "Mumbai High 3D Seismic Phase III", "Vessel": "Ramform Vanguard (PGS)",
                     "Area": "1,850 km2", "Period": "Jan-Mar 2026"})
        self._h(d, "1. Acquisition Parameters", 1)
        self._t(d, 8, 2, [["Parameter", "Specification"], ["Survey Type", "Full Azimuth 3D Marine"],
                           ["Source", "Airgun array - 8 guns, 2550 in3"], ["Streamer Length", "8,000 m (10 cables)"],
                           ["Streamer Depth", "12 m"], ["Shot Interval", "25 m"],
                           ["Fold", "60-80"], ["Record Length", "8 seconds / 2 ms"]])
        self._h(d, "2. QC Summary", 1)
        self._t(d, 5, 2, [["QC Metric", "Result"], ["Navigation Accuracy", "+/- 2.5 m"],
                           ["Streamer feathering", "Max 8 deg"], ["Source timing", "All guns within 2 ms"],
                           ["Bandwidth", "5-80 Hz - ACHIEVED"]])
        d.save(fp)
        return fp

    def generate_simops_report_docx(self, filename=None):
        fn = filename or f"OISD186_SIMOPS_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-186 Simultaneous Operations Report", 0)
        self._m(d, {"Asset": "Mumbai High North Complex", "Period": "Q1 2026", "Coordinator": "Mr. Deepak Iyer"})
        self._h(d, "1. Concurrent Activities", 1)
        self._t(d, 5, 4, [["Activity", "Location", "Party", "Risk"],
                           ["Drilling MH-4A-12", "WHP-A", "MRPL Drilling", "HIGH"],
                           ["Production Wells", "PP-B", "MRPL Production", "MEDIUM"],
                           ["Compressor Overhaul", "CPP-1", "MRPL Maintenance", "MEDIUM"],
                           ["Diving Ops", "Subsea PP-B", "Subsea India", "MEDIUM"]])
        self._h(d, "2. Controls", 1)
        self._n(d, ["Weekly SIMOPS meetings (Monday 07:00)", "Simultaneous Operations Permit for each activity",
                     "Dedicated radio channel Ch. 7", "Exclusion zones with barriers",
                     "Real-time gas monitoring with auto-shutdown at 10% LEL",
                     "NO-GO: wind >25 kts, visibility <500m, sea state >3"])
        self._h(d, "3. Incidents", 1)
        self._p(d, "No SIMOPS-related incidents in Q1 2026. One near-miss: supply vessel within 200m of PP-B without notification. Corrective action taken.")
        d.save(fp)
        return fp

    def generate_offshore_pipeline_inspection_docx(self, filename=None):
        fn = filename or f"OISD139_Pipeline_Inspection_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-139 Offshore Pipeline Inspection Report", 0)
        self._m(d, {"Pipeline": "Mumbai High - Uran Trunk Line", "Length": "128 km", "Diameter": "24in",
                     "ILI Tool": "ROSEN ROCDP", "Date": "28-Feb-2026"})
        self._h(d, "1. ILI Findings", 1)
        self._t(d, 6, 4, [["Defect", "Count", "Max Depth", "Action"],
                           ["Metal Loss (General)", "142", "18%", "Monitor"],
                           ["Metal Loss (Local)", "38", "32%", "Dig 12, repair 2"],
                           ["Dents", "15", "N/A", "Dig 5 critical"],
                           ["Girth Weld Defects", "8", "12%", "Monitor"],
                           ["Coating Disbondment", "24", "N/A", "CP survey"]])
        self._h(d, "2. CP Status", 1)
        self._t(d, 5, 3, [["Section", "CP (V)", "Status"],
                           ["KM 0-30", "-0.82 to -0.95", "ADEQUATE"],
                           ["KM 30-80", "-0.72 to -0.88", "MARGINAL"],
                           ["KM 80-120", "-0.68 to -0.82", "INADEQUATE"],
                           ["KM 120-128", "-0.85 to -0.95", "ADEQUATE"]])
        self._h(d, "3. Recommendations", 1)
        self._n(d, ["Repair 2 defects >30% wt", "CP rehabilitation KM 30-120 (Rs 4.2 Cr)", "Full ILI in 3 years"])
        d.save(fp)
        return fp

    def generate_oilfield_explosive_safety_docx(self, filename=None):
        fn = filename or f"OISD191_Explosive_Safety_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-191 Oil Field Explosive Safety Document", 0)
        self._m(d, {"Asset": "Mumbai High North WHP-A & PP-B", "Standard": "IEC 60079-10 / OISD-STD-191", "Rev": "Rev 3"})
        self._h(d, "1. Hazardous Area Classification", 1)
        self._t(d, 8, 4, [["Location", "Zone", "Gas Group", "Temp Class"],
                           ["Wellhead (3m)", "Zone 0", "IIC", "T3 (200C)"],
                           ["Separator area", "Zone 1", "IIC", "T3"],
                           ["Produced water", "Zone 1", "IIB", "T3"],
                           ["Tank (internal)", "Zone 0", "IIC", "T3"],
                           ["Compressor room", "Zone 1", "IIC", "T2"],
                           ["Fuel gas skid", "Zone 2", "IIC", "T3"],
                           ["Open deck", "Safe Area", "-", "-"]])
        self._h(d, "2. Equipment Requirements", 1)
        self._t(d, 4, 3, [["Zone", "Category", "Examples"],
                           ["Zone 0", "ia/ma - EPL Ga", "IS transmitters, submersible pumps"],
                           ["Zone 1", "ib/mb - EPL Gb", "Ex d motors, Ex e junction boxes"],
                           ["Zone 2", "n - EPL Gc", "Ex n non-sparking motors"]])
        d.save(fp)
        return fp

    def generate_pipeline_corrosion_monitoring_docx(self, filename=None):
        fn = filename or f"OISD188_Corrosion_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "OISD-STD-188 Pipeline Corrosion Monitoring Report", 0)
        self._m(d, {"Pipeline": "Mumbai High - Uran 24in", "Period": "Q4 2025-Q1 2026"})
        self._h(d, "1. Corrosion Rates", 1)
        self._t(d, 6, 4, [["Location", "Rate (mm/yr)", "Coupon (mg)", "Status"],
                           ["KM 15", "0.025", "18.2", "ACCEPTABLE"],
                           ["KM 45", "0.085", "61.8", "MONITORING"],
                           ["KM 78", "0.120", "87.4", "ACTION REQUIRED"],
                           ["KM 105", "0.065", "47.3", "ACCEPTABLE"],
                           ["KM 125", "0.030", "21.9", "ACCEPTABLE"]])
        self._h(d, "2. ILI Summary", 1)
        self._t(d, 5, 3, [["Defect", "Count", "Max Depth"],
                           ["General Metal Loss", "142", "18%"],
                           ["Pitting", "38", "32%"],
                           ["Erosion at Bend", "8", "15%"],
                           ["Weld Corrosion", "12", "10%"]])
        d.save(fp)
        return fp

    def generate_production_safety_report_docx(self, filename=None):
        fn = filename or f"Production_Safety_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Production Installation Safety Report", 0)
        self._m(d, {"Platform": "Mumbai High North PP-B", "Period": "Q1 2026", "HSE Mgr": "Mr. Vikram Joshi"})
        self._h(d, "1. Safety Systems", 1)
        self._t(d, 7, 3, [["System", "Last Inspection", "Status"],
                           ["Fire & Gas", "15-Mar-2026", "Operational"],
                           ["ESD", "15-Mar-2026", "SIL 2 certified"],
                           ["Fire Water", "10-Mar-2026", "Operational"],
                           ["Lifeboats", "15-Feb-2026", "All 4 serviced"],
                           ["Helideck", "01-Mar-2026", "Foam serviced"],
                           ["SCBA", "01-Mar-2026", "48 units ready"]])
        self._h(d, "2. Incident Stats", 1)
        self._t(d, 5, 2, [["Metric", "Q1 2026"], ["LTI", "0"], ["Recordable Incidents", "2"],
                           ["Near-Miss Reports", "14"], ["Safe Man-Hours", "145,000"]])
        d.save(fp)
        return fp

    def generate_ep_safety_pptx(self, filename=None):
        fn = filename or f"EP_Safety_Review_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        fp = os.path.join(OUTPUT_DIR, "pptx", fn)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = "E&P Offshore Safety Review"
        sl.placeholders[1].text = f"Q1 2026 - Mumbai High North\nMRPL HSE Division\n{datetime.now().strftime('%d %B %Y')}"
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = "Agenda"
        tf = sl2.placeholders[1].text_frame
        for item in ["Platform Safety Overview", "Incident Statistics", "Safety Systems Status", "Emergency Drills", "Action Items"]:
            p = tf.add_paragraph(); p.text = item; p.level = 0
        sl3 = prs.slides.add_slide(prs.slide_layouts[5])
        sl3.shapes.title.text = "Safety Statistics"
        ts = sl3.shapes.add_table(6, 3, PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(4)).table
        for i, rd in enumerate([["Metric", "Q1 2026", "Target"], ["LTI", "0", "0"], ["Recordable", "2", "<3"],
                                 ["Near-Miss", "14", ">10"], ["Safe Man-Hrs", "145,000", "150,000"], ["Env Incidents", "0", "0"]]):
            for j, ct in enumerate(rd): ts.cell(i, j).text = ct
        sl4 = prs.slides.add_slide(prs.slide_layouts[1])
        sl4.shapes.title.text = "Action Items"
        tf2 = sl4.placeholders[1].text_frame
        for item in ["Install 2 gas detectors (Apr-2026)", "Update ERP for well control (May-2026)",
                      "SCBA refresher training (Apr-2026)", "HAZOP review Q3-2026", "Target: Zero LTI 2026"]:
            p = tf2.add_paragraph(); p.text = item; p.level = 0
        prs.save(fp)
        return fp

    def generate_detailed_engineering_pkg_docx(self, filename=None):
        fn = filename or f"Engineering_Package_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Detailed Engineering Package", 0)
        self._m(d, {"Project": "MRPL Crude FPSO - Flowline Upgrade", "Contractor": "L&T Hydrocarbon", "Status": "IFC"})
        self._h(d, "1. Design Basis", 1)
        self._t(d, 7, 2, [["Parameter", "Specification"], ["Design Code", "ASME B31.4 / B31.8"],
                           ["Design Pressure", "150 bar (flowline)"], ["Pipeline Material", "API 5L X65 PSL2"],
                           ["Flowline Material", "Duplex SS UNS S31803"], ["Design Life", "25 years"],
                           ["CP", "Aluminium anodes - 15yr life"]])
        self._h(d, "2. Deliverables Checklist", 1)
        self._t(d, 8, 3, [["#", "Deliverable", "Status"], ["1", "P&ID (IFC)", "Issued"],
                           ["2", "PFD", "Issued"], ["3", "Equipment Layout", "Issued"],
                           ["4", "Pipe Stress Report", "Issued"], ["5", "Pipeline Route Survey", "Issued"],
                           ["6", "Electrical SLD", "Issued"], ["7", "Instrument Index", "Issued"]])
        d.save(fp)
        return fp

    def generate_as_built_drawings_docx(self, filename=None):
        fn = filename or f"As_Built_Register_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "As-Built Drawings Register", 0)
        self._m(d, {"Project": "MRPL Crude FPSO", "Register No": "MRPL/ABD/REG/2026-001"})
        self._h(d, "1. Drawing Register", 1)
        self._t(d, 8, 6, [["Drawing No", "Title", "Discipline", "Orig Rev", "As-Built", "Status"],
                           ["CFG-PID-001", "P&ID - FPSO", "Process", "C", "D", "Final"],
                           ["CFG-PLD-001", "Pipeline Route", "Pipeline", "B", "C", "Final"],
                           ["CFG-ELD-001", "Electrical Layout", "Electrical", "A", "B", "Final"],
                           ["CFG-INS-001", "Instrument Index", "Instrument", "C", "D", "Final"],
                           ["CFG-STR-001", "Structural Steel", "Structural", "B", "C", "Final"],
                           ["CFG-CIV-001", "Foundation Plan", "Civil", "A", "A", "Final"],
                           ["CFG-HSE-001", "Safety Layout", "HSE", "B", "C", "Final"]])
        d.save(fp)
        return fp

    def generate_mechanical_completion_cert_docx(self, filename=None):
        fn = filename or f"MC_Certificate_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Mechanical Completion Certificate", 0)
        self._m(d, {"Certificate No": "MRPL/MCC/2026/CFG-04/003", "System": "10in Crude Flowline", "Contractor": "L&T Hydrocarbon"})
        self._h(d, "1. Scope Verification", 1)
        self._t(d, 6, 3, [["#", "Work Package", "Status"], ["1", "Pipeline Install (12.8 km)", "COMPLETE"],
                           ["2", "Hydro Test (Section 1)", "COMPLETE"], ["3", "Hydro Test (Section 2)", "COMPLETE"],
                           ["4", "Anode Installation", "COMPLETE"], ["5", "Coating & Insulation", "COMPLETE"]])
        self._h(d, "2. Punch List", 1)
        self._t(d, 3, 3, [["Category", "Count", "Status"], ["Cat A (Safety)", "0", "ALL CLOSED"],
                           ["Cat B (Non-safety)", "3", "2 CLOSED, 1 OPEN"]])
        d.save(fp)
        return fp

    def generate_pre_commissioning_checklist_docx(self, filename=None):
        fn = filename or f"PreComm_Checklist_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Pre-Commissioning Checklist", 0)
        self._m(d, {"System": "10in Crude Flowline", "Area": "FPSO Main Deck"})
        self._h(d, "1. Hydrostatic Test", 1)
        self._t(d, 5, 3, [["#", "Item", "Result"], ["1.1", "Test pressure (150 bar)", "PASS"],
                           ["1.2", "Hold time (4 hr)", "PASS - 4.25 hr"], ["1.3", "Pressure drop <0.5%", "PASS - 0.12%"],
                           ["1.4", "Visual - no leaks", "PASS"]])
        self._h(d, "2. Loop Check", 1)
        self._t(d, 5, 3, [["#", "Item", "Result"], ["2.1", "Instrument loops (42)", "PASS"],
                           ["2.2", "Transmitter cal", "PASS"], ["2.3", "Control valve stroke", "PASS"],
                           ["2.4", "ESD function", "PASS"]])
        d.save(fp)
        return fp

    def generate_performance_test_run_docx(self, filename=None):
        fn = filename or f"Perf_Test_Run_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Performance Test Run Report", 0)
        self._m(d, {"System": "10in Crude Transfer Pumps", "Duration": "72 hours", "Contractor": "L&T Hydrocarbon"})
        self._h(d, "1. Guarantees vs Actual", 1)
        self._t(d, 8, 4, [["Parameter", "Guarantee", "Actual", "Status"],
                           ["Flow Rate", ">=2,500 m3/hr", "2,620 m3/hr", "PASS"],
                           ["Discharge Pressure", ">=120 bar", "125 bar", "PASS"],
                           ["Pump Efficiency", ">=82%", "84.3%", "PASS"],
                           ["Motor Power", "<=1,200 kW", "1,155 kW", "PASS"],
                           ["Vibration", "<=4.5 mm/s", "2.8 mm/s", "PASS"],
                           ["Bearing Temp", "<=85C", "68C", "PASS"],
                           ["Seal Leakage", "Zero", "Zero", "PASS"]])
        self._p(d, "Conclusion: All guarantees met. Equipment recommended for handover.")
        d.save(fp)
        return fp

    def generate_punch_list_xlsx(self, filename=None):
        fn = filename or f"Punch_List_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        fp = os.path.join(OUTPUT_DIR, "xlsx", fn)
        wb = Workbook()
        ws = wb.active; ws.title = "Category A"
        hf = Font(bold=True, color="FFFFFF", size=11)
        hrf = PatternFill(start_color="CC0000", end_color="CC0000", fill_type="solid")
        hbf = PatternFill(start_color="CC7700", end_color="CC7700", fill_type="solid")
        gf = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        of = PatternFill(start_color="FF6666", end_color="FF6666", fill_type="solid")
        headers = ["Punch ID", "Drawing Ref", "Location", "Description", "Discipline", "Assigned To", "Due Date", "Status", "Remarks"]
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=c, value=h); cell.font = hf; cell.fill = hrf; cell.alignment = Alignment(horizontal="center", wrap_text=True)
        a_data = [
            ["PL-A-001", "CFG-PPE-001", "Equipment-001", "Missing PSV on header", "Piping", "Eng. Prasad", "05-Apr-2026", "OPEN", "Vendor delay"],
            ["PL-A-002", "CFG-ELD-001", "S-02", "Incomplete earthing at MCC", "Electrical", "Eng. Mehta", "01-Apr-2026", "OPEN", "In progress"],
            ["PL-A-003", "CFG-INS-001", "Control Room", "ESD trip not verified Pump", "Instrument", "Eng. Desai", "10-Apr-2026", "CLOSED", "Verified"],
            ["PL-A-004", "CFG-STR-001", "PR-01", "Missing fireproofing", "Structural", "Eng. Singh", "15-Mar-2026", "CLOSED", "Completed"],
        ]
        for ri, rd in enumerate(a_data, 2):
            for ci, v in enumerate(rd, 1):
                cell = ws.cell(row=ri, column=ci, value=v); cell.alignment = Alignment(wrap_text=True)
                if ci == 8: cell.fill = of if v == "OPEN" else gf; cell.font = Font(bold=(v == "OPEN"))
        ws2 = wb.create_sheet("Category B")
        for c, h in enumerate(headers, 1):
            cell = ws2.cell(row=1, column=c, value=h); cell.font = hf; cell.fill = hbf; cell.alignment = Alignment(horizontal="center", wrap_text=True)
        b_data = [
            ["PL-B-001", "CFG-CIV-001", "Equipment-001", "Surface cracks on foundation", "Civil", "Eng. Prasad", "15-Apr-2026", "OPEN", "Monitoring"],
            ["PL-B-002", "CFG-PPE-002", "Inj Line", "Pipe labels missing 3 locations", "Piping", "Eng. Mehta", "05-Apr-2026", "OPEN", "Ordered"],
            ["PL-B-003", "CFG-ELD-001", "L-03", "Missing cable gland 2 cables", "Electrical", "Eng. Mehta", "01-Apr-2026", "OPEN", "Scheduled"],
            ["PL-B-004", "CFG-COR-001", "KM 5.2", "Coating damage", "Corrosion", "Eng. Singh", "10-Apr-2026", "OPEN", "Repair scheduled"],
            ["PL-B-005", "CFG-INS-001", "Analyzer", "Calibration label missing", "Instrument", "Eng. Desai", "20-Apr-2026", "CLOSED", "Completed"],
        ]
        for ri, rd in enumerate(b_data, 2):
            for ci, v in enumerate(rd, 1):
                cell = ws2.cell(row=ri, column=ci, value=v); cell.alignment = Alignment(wrap_text=True)
                if ci == 8: cell.fill = PatternFill(start_color="FFCC66", end_color="FFCC66", fill_type="solid") if v == "OPEN" else gf; cell.font = Font(bold=(v == "OPEN"))
        for wsx in [ws, ws2]:
            for col in wsx.columns:
                ml = max(len(str(c.value or "")) for c in col)
                wsx.column_dimensions[col[0].column_letter].width = min(ml + 2, 40)
        wb.save(fp)
        return fp

    def generate_construction_safety_audit_docx(self, filename=None):
        fn = filename or f"Safety_Audit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Construction Safety Audit Report", 0)
        self._m(d, {"Audit No": "MRPL/CSA/2026/003", "Project": "MRPL Crude FPSO", "Auditor": "Mr. Vikram Joshi + Mr. Lars Petersen"})
        self._h(d, "1. Findings Summary", 1)
        self._t(d, 5, 3, [["Category", "Count", "Severity"],
                           ["Major NC", "2", "Safety-critical"],
                           ["Minor NC", "8", "Action in 14 days"],
                           ["Observations", "12", "Recommendations"],
                           ["Positive Findings", "18", "-"]])
        self._h(d, "2. Major NCs", 1)
        self._n(d, ["MNC-01: Scaffolding tag system not enforced - 3 untaged scaffolds. Action: Taken out of service.",
                     "MNC-02: Lifting without banksman - 2T pipe spool lift. Action: Stopped, crew re-briefed."])
        self._p(d, "Overall Rating: SATISFACTORY WITH RESERVATIONS")
        d.save(fp)
        return fp

    def generate_project_completion_report_docx(self, filename=None):
        fn = filename or f"Completion_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Project Completion Report", 0)
        self._m(d, {"Project": "MRPL Crude FPSO - Flowline Upgrade", "Contractor": "L&T Hydrocarbon", "Value": "Rs 192.3 Cr"})
        self._h(d, "1. Cost Summary", 1)
        self._t(d, 5, 3, [["Category", "Budget (Rs Cr)", "Actual (Rs Cr)"],
                           ["EPC Contract", "185.00", "192.30"],
                           ["Owner Costs", "12.50", "11.80"],
                           ["Contingency", "18.50", "8.20"],
                           ["Total", "220.00", "216.45"]])
        self._h(d, "2. Schedule", 1)
        self._t(d, 5, 3, [["Milestone", "Planned", "Actual"],
                           ["Contract Award", "01-Jan-2025", "15-Jan-2025"],
                           ["FEED", "15-Mar-2025", "20-Mar-2025"],
                           ["Engineering", "30-Jun-2025", "10-Jul-2025"],
                           ["MC", "15-Mar-2026", "25-Mar-2026"]])
        self._p(d, "Result: 1.6% under budget, 10 days late. All commissioning complete.")
        d.save(fp)
        return fp

    def generate_tef_report_docx(self, filename=None):
        fn = filename or f"TEF_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Techno-Economic Feasibility Report", 0)
        self._m(d, {"Project": "MRPL Crude FPSO Upgrade", "Prepared": "MRPL Projects Dept"})
        self._h(d, "1. Financial Indicators", 1)
        self._t(d, 6, 2, [["Indicator", "Value"], ["NPV", "Rs 342 Crore @ 12%"],
                           ["IRR", "22.8%"], ["Payback", "4.2 years"],
                           ["ROCE", "28.5%"], ["Break-even Oil Price", "$38/bbl"]])
        self._h(d, "2. CAPEX", 1)
        self._t(d, 5, 2, [["Category", "Amount (Rs Cr)"],
                           ["EPC Contract", "192.3"], ["Owner Costs", "11.8"],
                           ["Spares & Testing", "11.7"], ["Total CAPEX", "243.15"]])
        self._h(d, "3. Sensitivity", 1)
        self._t(d, 4, 3, [["Scenario", "NPV", "IRR"],
                           ["Base Case", "342", "22.8%"],
                           ["Oil Price -20%", "185", "15.2%"],
                           ["CAPEX +15%", "295", "19.6%"]])
        self._p(d, "RECOMMENDATION: PROCEED WITH PROJECT")
        d.save(fp)
        return fp

    def generate_sow_docx(self, filename=None):
        fn = filename or f"SOW_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Scope of Work Document", 0)
        self._m(d, {"SOW No": "MRPL/SOW/2026/CFG-04/001", "Client": "MRPL", "Contractor": "L&T Hydrocarbon"})
        self._h(d, "1. Scope Overview", 1)
        self._p(d, "EPC services for upgrade of crude oil transfer system on Mumbai High North FPSO, including 10in crude flowline (12.8 km) and associated modifications.")
        self._h(d, "2. In-Scope Work", 1)
        self._n(d, ["Detailed engineering per MRPL specs", "Procurement of materials and equipment",
                     "Fabrication of pipe spools and skids", "Installation of 10in flowline (12.8 km)",
                     "Hook-up to FPSO process system", "Electrical and instrumentation modifications",
                     "Pre-commissioning and commissioning support", "As-built documentation"])
        self._h(d, "3. Deliverables", 1)
        self._t(d, 7, 2, [["#", "Deliverable"], ["1", "Engineering Package"],
                           ["2", "Material Requisitions"], ["3", "Construction Method Statements"],
                           ["4", "Test Packages"], ["5", "MC Certificate"],
                           ["6", "As-Built Register"]])
        d.save(fp)
        return fp

    def generate_rfp_docx(self, filename=None):
        fn = filename or f"RFP_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Request for Proposal (RFP)", 0)
        self._m(d, {"RFP No": "MRPL/RFP/2026/CFG-04/001", "Due Date": "30-Apr-2026"})
        self._h(d, "1. Introduction", 1)
        self._p(d, "MRPL invites proposals from qualified EPC contractors for Crude FPSO Flowline & Pipeline Upgrade at Mumbai High North.")
        self._h(d, "2. Prequalification", 1)
        self._n(d, ["Min 10 years offshore EPC experience", "3+ similar projects in last 5 years",
                     "Turnover >= Rs 500 Cr", "ISO 9001, 14001, 45001 certified", "Zero fatalities in 3 years"])
        self._h(d, "3. Evaluation Criteria", 1)
        self._t(d, 5, 2, [["Criterion", "Weight"], ["Technical Compliance", "40%"],
                           ["Past Performance", "20%"], ["Project Team", "15%"],
                           ["Commercial", "20%"]])
        d.save(fp)
        return fp

    def generate_rfq_docx(self, filename=None):
        fn = filename or f"RFQ_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Request for Quotation (RFQ)", 0)
        self._m(d, {"RFQ No": "MRPL/RFQ/2026/MAT-012", "Material": "API 5L X65 24in Line Pipe",
                     "Quantity": "12,800 m", "Due": "20-Apr-2026"})
        self._h(d, "1. Material Specification", 1)
        self._t(d, 7, 2, [["Parameter", "Specification"], ["Standard", "API 5L PSL2"],
                           ["Grade", "X65"], ["OD", "24in (609.6 mm)"],
                           ["WT", "0.500in (12.7 mm)"], ["Length", "50 m joints"],
                           ["Coating", "FBE 350 um min"]])
        self._h(d, "2. Terms", 1)
        self._n(d, ["FOB Mumbai Port", "Payment: 30% advance, 70% on delivery",
                     "LD: 0.5%/week, max 10%", "Warranty: min 5 years"])
        d.save(fp)
        return fp

    def generate_bom_docx(self, filename=None):
        fn = filename or f"BOM_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Bill of Material (BOM)", 0)
        self._m(d, {"BOM No": "MRPL/BOM/2026/CFG-04/BOM-005", "System": "Chemical Injection Skid"})
        self._h(d, "1. Equipment List", 1)
        self._t(d, 8, 6, [["#", "Description", "Qty", "Unit", "Spec", "Vendor"],
                           ["1", "Metering Pump", "2", "Nos", "API 675, 50 L/hr", "Milton Roy"],
                           ["2", "Storage Tank SS316", "1", "No", "5 m3, API 650", "Permian"],
                           ["3", "Check Valve SS316", "4", "Nos", "2in CL150", "Velan"],
                           ["4", "Globe Valve SS316", "6", "Nos", "2in CL150", "Velan"],
                           ["5", "PRV", "2", "Nos", "2in CL150", "Fisher"],
                           ["6", "Pressure Transmitter", "2", "Nos", "0-160 bar", "Rosemount"],
                           ["7", "Magnetic Flow Meter", "1", "No", "2in DN50", "E+H"]])
        d.save(fp)
        return fp

    def generate_material_submittal_docx(self, filename=None):
        fn = filename or f"Material_Submittal_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Material Submittal / Vendor Document Review", 0)
        self._m(d, {"Submittal No": "MRPL/MS/2026/MS-023", "Material": "API 5L X65 24in Pipe", "Vendor": "Jindal Saw Ltd"})
        self._h(d, "1. Documents Received", 1)
        self._t(d, 5, 3, [["#", "Document", "Review"],
                           ["1", "Mill Test Certificate", "APPROVED"],
                           ["2", "Material Test Report", "APPROVED"],
                           ["3", "Dimensional Report", "APPROVED"],
                           ["4", "FBE Coating Cert", "APPROVED WITH COMMENTS"]])
        self._h(d, "2. Comments", 1)
        self._n(d, ["Charpy impact at -46C to be provided", "Holiday test voltage verify at 5 kV",
                     "Heat number cross-reference required"])
        d.save(fp)
        return fp

    def generate_qaqc_plan_docx(self, filename=None):
        fn = filename or f"QAQC_Plan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "QA/QC Plan", 0)
        self._m(d, {"Plan No": "MRPL/QA/2026/QAP-007", "Scope": "10in Crude Flowline Welding & NDE"})
        self._h(d, "1. Inspection & Test Plan (ITP)", 1)
        self._t(d, 8, 5, [["#", "Activity", "Hold Pt", "Contractor", "DNV GL"],
                           ["1", "Material Receiving", "Witness", "Yes", "-"],
                           ["2", "WPS/PQR Qualification", "Hold", "Yes", "Yes"],
                           ["3", "Welder Qualification", "Hold", "Yes", "Yes"],
                           ["4", "Fit-up Inspection", "Witness", "Yes", "-"],
                           ["5", "Visual (Post-Weld)", "Hold", "Yes", "Yes"],
                           ["6", "UT - 100%", "Hold", "Yes", "Yes"],
                           ["7", "Hydrostatic Test", "Hold", "Yes", "Yes"]])
        self._h(d, "2. Welding Requirements", 1)
        self._t(d, 5, 2, [["Parameter", "Spec"],
                           ["WPS", "ASME IX / AWS D1.1"],
                           ["Process", "GTAW Root + SMAW/FCAW Fill"],
                           ["Preheat", "100C min"],
                           ["Interpass", "250C max"]])
        d.save(fp)
        return fp

    def generate_project_review_pptx(self, filename=None):
        fn = filename or f"Project_Review_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        fp = os.path.join(OUTPUT_DIR, "pptx", fn)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = "Project Review - MRPL Crude FPSO Upgrade"
        sl.placeholders[1].text = f"Q1 2026 Progress Review\nL&T Hydrocarbon Engineering\n{datetime.now().strftime('%d %B %Y')}"
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = "Agenda"
        tf = sl2.placeholders[1].text_frame
        for item in ["Project Status Overview", "Cost & Schedule Update", "Construction Progress", "Commissioning Status", "HSE Performance", "Key Risks & Mitigations"]:
            p = tf.add_paragraph(); p.text = item; p.level = 0
        sl3 = prs.slides.add_slide(prs.slide_layouts[5])
        sl3.shapes.title.text = "Cost & Schedule"
        ts = sl3.shapes.add_table(5, 3, PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(4)).table
        for i, rd in enumerate([["Milestone", "Planned", "Actual"], ["Contract Award", "Jan-2025", "Jan-2025"],
                                 ["FEED", "Mar-2025", "Mar-2025"], ["Engineering", "Jun-2025", "Jul-2025"],
                                 ["MC", "Mar-2026", "Mar-2026"]]):
            for j, ct in enumerate(rd): ts.cell(i, j).text = ct
        sl4 = prs.slides.add_slide(prs.slide_layouts[1])
        sl4.shapes.title.text = "Key Risks"
        tf2 = sl4.placeholders[1].text_frame
        for item in ["Monsoon delays (Jun-Sep) - mitigated with float", "Vendor delivery - 6-month lead time",
                      "Commissioning overlap - resource leveling applied", "Remaining punch list (7 items) - targeted closure"]:
            p = tf2.add_paragraph(); p.text = item; p.level = 0
        sl5 = prs.slides.add_slide(prs.slide_layouts[0])
        sl5.shapes.title.text = "Thank You"
        sl5.placeholders[1].text = "Questions & Discussion"
        prs.save(fp)
        return fp
