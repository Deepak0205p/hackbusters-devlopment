import os
import time
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))


class ExtendedGenerators_HseOps:

    def __init__(self):
        os.makedirs(os.path.join(OUTPUT_DIR, "docx"), exist_ok=True)
        os.makedirs(os.path.join(OUTPUT_DIR, "xlsx"), exist_ok=True)
        os.makedirs(os.path.join(OUTPUT_DIR, "pptx"), exist_ok=True)

    @staticmethod
    def _h(doc, text, level=1):
        h = doc.add_heading(text, level=level)
        for r in h.runs:
            r.font.color.rgb = RGBColor(0, 51, 102)
        return h

    @staticmethod
    def _t(doc, rows, cols, data=None):
        table = doc.add_table(rows=rows, cols=cols, style="Light Grid Accent 1")
        table.alignment = WD_TABLE_ALIGNMENT.CENTER
        if data:
            for i, rd in enumerate(data):
                for j, ct in enumerate(rd):
                    table.rows[i].cells[j].text = str(ct)
        return table

    @staticmethod
    def _m(doc, meta):
        for label, value in meta.items():
            p = doc.add_paragraph()
            rl = p.add_run(f"{label}: ")
            rl.bold = True
            rl.font.size = Pt(10)
            rv = p.add_run(str(value))
            rv.font.size = Pt(10)

    @staticmethod
    def _p(doc, text):
        p = doc.add_paragraph(text)
        p.style.font.size = Pt(10)
        return p

    @staticmethod
    def _n(doc, items):
        for item in items:
            doc.add_paragraph(item, style="List Number")

    # ── HSE Methods ───────────────────────────────────────────────────────

    def generate_cold_work_permit_docx(self, filename=None):
        fn = filename or f"Cold_Work_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Cold Work Permit (OISD-STD-105)", 0)
        self._m(d, {"Permit No": "MRPL/CWP/2026/00412", "Area": "CDU-III Pipe Rack Level 2",
                     "Activity": "Grinding & Cutting on Decommissioned 6in Line",
                     "Date": datetime.now().strftime("%d-%b-%Y"), "Valid From": "08:00", "Valid To": "17:00"})
        self._h(d, "1. Hazard Assessment", 1)
        self._t(d, 6, 3, [["Hazard", "Risk", "Control"], ["Sparks/Hot Particles", "Fire/Ignition", "Fire blanket, spark arrestor, fire watch"],
                           ["Slips/Falls", "Personnel Injury", "Barricades, warning signs, scaffold inspection"],
                           ["Noise >85 dBA", "Hearing Damage", "Ear protection mandatory"],
                           ["Dust/Fumes", "Respiratory", "Local exhaust ventilation, RPE"],
                           ["Sharp Edges", "Cuts/Lacerations", "Cut-resistant gloves"]])
        self._h(d, "2. Pre-Work Checklist", 1)
        self._n(d, ["Area gas-tested and within limits (<1% LEL)", "Adjacent work isolated and secured",
                     "Fire extinguisher (4 kg DCP) within 5 m", "Non-combustible shields in place",
                     "Fire watch personnel briefed and assigned", "PTW signed by Area Authority"])
        self._h(d, "3. Authorizations", 1)
        self._t(d, 5, 3, [["Role", "Name", "Signature"], ["Permit Issuer", "Shri K.R. Nair", "________________"],
                           ["Area Authority", "Shri S.V. Patil", "________________"],
                           ["Contractor Supervisor", "Shri R.M. Gupta", "________________"],
                           ["Safety Officer", "Shri A.D. Mistry", "________________"]])
        d.save(fp)
        return fp

    def generate_electrical_isolation_permit_docx(self, filename=None):
        fn = filename or f"Electrical_Isolation_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Electrical Isolation & Energization Permit", 0)
        self._m(d, {"Permit No": "MRPL/EIP/2026/00089", "Equipment": "Motor M-4401B (350 kW, 11 kV)",
                     "Location": "CDU-III Extractor Pump Area", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Equipment Details", 1)
        self._t(d, 5, 2, [["Parameter", "Detail"], ["Motor Tag", "M-4401B"],
                           ["Rating", "350 kW, 11 kV, 3-phase, 50 Hz"], ["MCC", "MCC-CDU3-07, Feeder F-22"],
                           ["Lock Position", "OFF / ISOLATED"]])
        self._h(d, "2. Isolation Procedure", 1)
        self._n(d, ["Notify Area Authority and control room", "De-energize feeder at MCC-CDU3-07 F-22",
                     "Apply lockout device (LOTO Tag No: LOTO-2026-0187)", "Verify dead with voltage detector",
                     "Apply earthing on both sides", "Confirm zero-energy state to Area Authority"])
        self._h(d, "3. Energization Procedure", 1)
        self._n(d, ["All personnel clear of equipment", "Work completion certificate signed", "Earthing removed",
                     "Area Authority authorizes energization", "Control room informed, breaker closed",
                     "Motor started and observed for 30 minutes"])
        self._h(d, "4. Sign-Off", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Electrical Engineer", "Shri P.V. Rao", "________________"],
                           ["Area Authority", "Shri S.V. Patil", "________________"],
                           ["Control Room Operator", "Shri T.K. Verma", "________________"]])
        d.save(fp)
        return fp

    def generate_excavation_permit_docx(self, filename=None):
        fn = filename or f"Excavation_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Excavation Permit", 0)
        self._m(d, {"Permit No": "MRPL/EXC/2026/00056", "Location": "CDU-III Utility Corridor, Grid Ref: N19-04-E12",
                     "Depth": "1.5 m", "Area": "20 m x 6 m", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Underground Services Survey", 1)
        self._t(d, 6, 4, [["Service", "Depth (m)", "Status", "Action"], ["11 kV Cable", "0.8", "LIVE", "Hand-dig only, maintain 0.5 m clearance"],
                           ["Fire Water 8in", "1.2", "LIVE", "Mark and protect"],
                           ["Drainage 12in", "1.8", "Abandoned", "Confirm with records"],
                           ["Instrument Tray", "0.6", "LIVE", "Support before excavation"],
                           ["Fiber Optic", "0.5", "LIVE", "Reroute before excavation"]])
        self._h(d, "2. Safety Requirements", 1)
        self._n(d, ["Barricade excavation area with signage", "Shoring/lagging for depth >1.2 m",
                     "No mechanical excavation within 0.5 m of live services", "Atmospheric monitoring for confined excavation",
                     "Means of egress within 7.5 m for depth >1.2 m", "Backfill material pre-staged"])
        self._h(d, "3. Authorizations", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Permit Issuer", "Shri K.R. Nair", "________________"],
                           ["Electrical Engineer", "Shri P.V. Rao", "________________"],
                           ["Contractor Supervisor", "Shri R.M. Gupta", "________________"]])
        d.save(fp)
        return fp

    def generate_radiation_safety_permit_docx(self, filename=None):
        fn = filename or f"Radiation_Safety_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Radiation Safety Permit (AERB License)", 0)
        self._m(d, {"Permit No": "MRPL/RSP/2026/00012", "Source": "Ir-192 (40 Ci) Gamma Radiography",
                     "Equipment": "GF-1000 X-Ray Cabinet / Radiography Set", "AERB License": "AERB/ROS/R/2024-2027/0156",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Source Details", 1)
        self._t(d, 5, 2, [["Parameter", "Detail"], ["Isotope", "Iridium-192"],
                           ["Activity", "40 Ci (1.48 TBq)"], ["Half Period", "73.83 days"],
                           ["IPEV Category", "Category II"]])
        self._h(d, "2. Radiation Protection Measures", 1)
        self._n(d, ["Controlled area established with 30 m exclusion zone",
                     "Radiation warning signs and barrier tape deployed",
                     "Personal dosimeters (TLD) issued to all personnel within 10 m",
                     "Survey meter calibrated within 6 months (Calib Due: Sep-2026)",
                     "Area dose rate <0.02 mSv/hr at boundary",
                     "Temporary worker exposure limit: 1 mSv/week"])
        self._h(d, "3. Personnel Monitoring", 1)
        self._t(d, 5, 4, [["Personnel", "TLD Badge", "Monthly Dose", "Annual Dose"], ["Radiographer A. Kumar", "TLD-4401", "0.18 mSv", "1.42 mSv"],
                           ["Asst. B. Singh", "TLD-4402", "0.12 mSv", "0.98 mSv"],
                           ["Safety Officer D. Patil", "TLD-4403", "0.06 mSv", "0.45 mSv"]])
        self._h(d, "4. Emergency Contact", 1)
        self._m(d, {"Radiation Safety Officer": "Shri V.K. Deshmukh", "Emergency No": "+91-98765-43210",
                     "AERB Emergency": "022-2550-1234"})
        d.save(fp)
        return fp

    def generate_lifting_permit_docx(self, filename=None):
        fn = filename or f"Lifting_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Lifting Permit - Critical Lift", 0)
        self._m(d, {"Permit No": "MRPL/LP/2026/00234", "Lift Type": "Critical (>20 Tonne)",
                     "Equipment": "Liebherr LTM 1250-8.1 (250T)", "Load": "Exchanger E-4201 (28.5 MT)",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Lift Details", 1)
        self._t(d, 7, 2, [["Parameter", "Detail"], ["Crane Capacity", "250 MT"],
                           ["Load Weight", "28.5 MT (with rigging)"], ["Radius", "18 m"],
                           ["Boom Length", "42 m"], ["Rigging", "4-leg chain sling, 36T SWL"],
                           ["Ground Bearing Pressure", "12.5 bar (checked)"]])
        self._h(d, "2. Risk Assessment", 1)
        self._t(d, 5, 3, [["Hazard", "Risk", "Control"], ["Crane Tip-Over", "Fatality", "Ground survey, outriggers on pads, wind <35 km/h"],
                           ["Dropped Load", "Fatality", "Certified rigging, pre-use inspection, 2-man signal"],
                           ["Overloaded Crane", "Fatality", "Load chart verification, zero-moment method"],
                           ["Electrical Contact", "Fatality", "Exclusion zone, height survey, radio comms"]])
        self._h(d, "3. Lift Sequence", 1)
        self._n(d, ["Pre-lift briefing (all personnel)", "Crane positioned, outriggers extended",
                     "Load trial lift to 10% SWL", "Load lifted to 0.5 m, hold 5 min, brake check",
                     "Load slewed to final position", "Load lowered, secured, demobilize"])
        self._h(d, "4. Authorizations", 1)
        self._t(d, 5, 3, [["Role", "Name", "Signature"], ["Lift Supervisor", "Shri G.S. Kulkarni", "________________"],
                           ["Crane Operator", "Shri M.N. Jadhav", "________________"],
                           ["Rigging Supervisor", "Shri A.R. Patil", "________________"],
                           ["HSE Officer", "Shri D.V. Sharma", "________________"]])
        d.save(fp)
        return fp

    def generate_work_at_height_permit_docx(self, filename=None):
        fn = filename or f"Work_at_Height_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Work at Height Permit", 0)
        self._m(d, {"Permit No": "MRPL/WAH/2026/00478", "Location": "CDU-III Column C-301, Level 4 (18.5 m)",
                     "Activity": "Insulation repair on 8in steam line", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Fall Protection Plan", 1)
        self._t(d, 5, 3, [["Protection Type", "Specification", "Status"], ["Scaffold", "Full access, green tag", "INSPECTED"],
                           ["Fall Arrest Harness", "Full body, 3-point", "INSPECTED"],
                           ["Anchor Point", "15 kN rated, overhead", "TESTED"],
                           ["Rescue Plan", "Horizontal retrieval", "IN PLACE"]])
        self._h(d, "2. Pre-Work Checklist", 1)
        self._n(d, ["Scaffold inspection tag GREEN and valid", "Harness and lanyard inspected (no defects)",
                     "Weather check: wind <40 km/h, no lightning, no rain",
                     "Tool lanyards attached for all hand tools",
                     "Barricade ground level exclusion zone",
                     "Rescue team on standby"])
        self._h(d, "3. Personnel", 1)
        self._t(d, 4, 3, [["Role", "Name", "Certification"], ["Worker", "Shri P.R. Deshmukh", "WAH Certified"],
                           ["Scaffold Inspector", "Shri S.K. Mane", "Scaffold Inspector Cert"],
                           ["Safety Watch", "Shri A.V. Kulkarni", "First Aid + Rescue"]])
        self._h(d, "4. Emergency Contact", 1)
        self._m(d, {"Emergency": "Ext. 5555", "Nearest Hospital": "MRPL Occupational Health Centre, 8 min"})
        d.save(fp)
        return fp

    def generate_vehicle_entry_permit_docx(self, filename=None):
        fn = filename or f"Vehicle_Entry_Permit_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Vehicle Entry Permit (Hazardous Area)", 0)
        self._m(d, {"Permit No": "MRPL/VEP/2026/01156", "Vehicle": "Tata Prima 4038 Tanker (MH-12-AB-3456)",
                     "Driver": "Shri H.R. Pawar (Lic: MH/2020/445567)", "Area": "CDU-III Flare Area, Zone 2",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Vehicle Details", 1)
        self._t(d, 5, 2, [["Parameter", "Detail"], ["Vehicle Type", "Fuel Tanker, Diesel, 16,000 L"],
                           ["Explosive Group", "IIA, T3"], ["Ex Rating", "Ex nA IIA T3 Gc"],
                           ["Spark Arrestor", "Fitted and functional"], ["Emergency Kit", "2x DCP, 1x blanket"]])
        self._h(d, "2. Driver Checks", 1)
        self._n(d, ["Valid driving license verified", "Safety induction completed (Induction Cert No: IND-2026-0891)",
                     "Vehicle inspection checklist passed", "Engine oil, coolant, tyre pressure OK",
                     "Fuel cap secured and sealed", "Mobile phone switched off"])
        self._h(d, "3. Route & Restrictions", 1)
        self._n(d, ["Entry via Gate 3 only (speed limit 15 km/h)", "No parking within 10 m of process equipment",
                     "Engine to be switched off at unloading point", "No hot work within 30 m during fueling",
                     "Return route via Gate 4 (one-way flow)"])
        self._h(d, "4. Authorizations", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Gate Security", "Shri T.K. Verma", "________________"],
                           ["Area Authority", "Shri S.V. Patil", "________________"],
                           ["HSE Officer", "Shri D.V. Sharma", "________________"]])
        d.save(fp)
        return fp

    def generate_loto_certificate_docx(self, filename=None):
        fn = filename or f"LOTO_Certificate_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Lock Out Tag Out (LOTO) Certificate", 0)
        self._m(d, {"Certificate No": "MRPL/LOTO/2026/00312", "Equipment": "Heat Exchanger E-4201 (38 bar / 320C)",
                     "Purpose": "Tube bundle replacement", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Isolation Points", 1)
        self._t(d, 7, 5, [["#", "Isolation Point", "Type", "Lock No", "Person"], ["1", "8in Inlet Valve V-4201A", "Gate Valve + LOTO", "L-001", "Shri R.M. Gupta"],
                           ["2", "8in Outlet Valve V-4201B", "Gate Valve + LOTO", "L-002", "Shri R.M. Gupta"],
                           ["3", "Steam Supply 4in V-4201C", "Gate Valve + LOTO", "L-003", "Shri R.M. Gupta"],
                           ["4", "Condensate Return 3in", "Gate Valve + LOTO", "L-004", "Shri R.M. Gupta"],
                           ["5", "MCC-CDU3-12 F-05 (Pump)", "Electrical Isolation", "L-005", "Shri P.V. Rao"],
                           ["6", "Instrument Air 1in", "Ball Valve + LOTO", "L-006", "Shri R.M. Gupta"]])
        self._h(d, "2. Verification", 1)
        self._n(d, ["Zero energy verified at all isolation points", "Process fluid drained and flushed",
                     "Electrical dead verified with voltage detector", "Stored energy (springs, thermal) relieved",
                     "LOTO tags applied and serial numbers recorded"])
        self._h(d, "3. Removal Authorization", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["LOTO Owner", "Shri R.M. Gupta", "________________"],
                           ["Area Authority", "Shri S.V. Patil", "________________"],
                           ["Control Room", "Shri T.K. Verma", "________________"]])
        self._p(d, "Note: All locks must be removed by the person who applied them. Group LOTO requires all tags removed before re-energization.")
        d.save(fp)
        return fp

    def generate_safety_induction_cert_docx(self, filename=None):
        fn = filename or f"Safety_Induction_Certificate_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Safety Induction Certificate", 0)
        self._m(d, {"Certificate No": "MRPL/IND/2026/00891", "Name": "Shri Anil K. Sharma",
                     "Company": "Tata Projects Ltd (Contractor)", "ID No": "TPL-4401",
                     "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Induction Details", 1)
        self._t(d, 7, 2, [["Module", "Duration", "Status"], ["General Safety Rules", "30 min", "PASS"],
                           ["Fire Safety & Emergency", "30 min", "PASS"],
                           ["Hazardous Area Awareness", "20 min", "PASS"],
                           ["Permit to Work System", "20 min", "PASS"],
                           ["PPE Requirements", "15 min", "PASS"],
                           ["LOTO Procedures", "15 min", "PASS"],
                           ["Environmental Awareness", "10 min", "PASS"]])
        self._h(d, "2. Assessment", 1)
        self._t(d, 4, 2, [["Question", "Result"], ["Q1: Emergency assembly point", "Correct"],
                           ["Q2: Fire extinguisher use", "Correct"], ["Q3: Permit requirements", "Correct"]])
        self._p(d, "Score: 100% - PASS (Minimum 80% required)")
        self._h(d, "3. Authorization", 1)
        self._m(d, {"Induction Conducted By": "Shri D.V. Sharma (HSE Officer)", "Valid Until": "31-Dec-2026",
                     "Badge Color": "GREEN (Authorized for Zone 1 & 2 areas)"})
        d.save(fp)
        return fp

    def generate_gas_testing_certificate_docx(self, filename=None):
        fn = filename or f"Gas_Testing_Certificate_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Gas Testing Certificate", 0)
        self._m(d, {"Certificate No": "MRPL/GTC/2026/00567", "Location": "CDU-III Exchanger E-4201 (Confined Space Entry)",
                     "Tester": "Shri V.K. Deshmukh (Gas Testing Certified)", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Test Results", 1)
        self._t(d, 7, 4, [["Gas", "Unit", "Result", "Limit"], ["Oxygen (O2)", "%v/v", "20.9", "19.5 - 23.5"],
                           ["Lower Explosive Limit (LEL)", "%v/v", "0.0", "< 1.0"],
                           ["Hydrogen Sulphide (H2S)", "ppm", "0.0", "< 10"],
                           ["Carbon Monoxide (CO)", "ppm", "2", "< 25"],
                           ["VOC (Total)", "ppm", "18", "< 50"],
                           ["Benzene", "ppm", "0.2", "< 0.5"]])
        self._h(d, "2. Instrument Details", 1)
        self._t(d, 5, 2, [["Instrument", "Details"], ["Multi-Gas Detector", "Drager X-am 8000, S/N: 8000-4401"],
                           ["Calibration Due", "15-Sep-2026"], ["Bump Test", "Passed 07:30 today"],
                           ["Sampling Method", "Direct probe, 3 locations"]])
        self._h(d, "3. Atmospheric Conditions", 1)
        self._t(d, 4, 2, [["Parameter", "Value"], ["Temperature", "38C"], ["Humidity", "72%"], ["Wind", "NW 8 km/h"]])
        self._p(d, "RESULT: ATMOSPHERE SAFE FOR ENTRY. Valid for 4 hours from time of test (expires 12:45 hrs).")
        d.save(fp)
        return fp

    def generate_safety_manual_docx(self, filename=None):
        fn = filename or f"Fire_Safety_Manual_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Fire Safety Manual", 0)
        self._m(d, {"Document No": "MRPL/HSE/FSM/2026-003", "Unit": "Unit Complex",
                     "Revision": "Rev 5", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Fire Detection Systems", 1)
        self._t(d, 6, 3, [["System", "Coverage", "Status"], ["Flame Detector (UV/IR)", "Process area, 48 nos", "Operational"],
                           ["Heat Detector", "Control room, MCC rooms", "Operational"],
                           ["Gas Detector (LEL)", "Enclosed areas, 62 nos", "Operational"],
                           ["Manual Call Point", "All escape routes, 36 nos", "Operational"],
                           ["Smoke Detector", "Office/control room", "Operational"]])
        self._h(d, "2. Fire Fighting Equipment", 1)
        self._t(d, 7, 3, [["Equipment", "Quantity", "Location"], ["DCP Extinguisher (6 kg)", "48", "Every 15 m in process area"],
                           ["CO2 Extinguisher (5 kg)", "24", "Electrical rooms"],
                           ["Foam Monitor (fixed)", "8", "Tank farm perimeter"],
                           ["Fire Water Monitor", "6", "Process platform elevations"],
                           ["Hydrant (pillar type)", "12", "Road access points"],
                           ["Deluge System", "3 zones", "Column & heater areas"]])
        self._h(d, "3. Emergency Response", 1)
        self._n(d, ["Raise alarm: activate nearest manual call point", "Call Emergency Control Room: Ext. 5555",
                     "Evacuate upwind to muster point (North Parking)",
                     "Fire team to assemble at Fire Station (3 min response)",
                     "Account for all personnel via roll call",
                     "Do NOT re-enter until ALL-CLEAR by Incident Commander"])
        self._h(d, "4. Drill Schedule", 1)
        self._t(d, 5, 3, [["Drill Type", "Frequency", "Last Conducted"], ["Fire Drill", "Monthly", "15-Mar-2026"],
                           ["Hazmat Drill", "Quarterly", "01-Feb-2026"],
                           ["Full Emergency", "Bi-annual", "15-Jan-2026"],
                           ["Night Drill", "Annual", "20-Dec-2025"]])
        d.save(fp)
        return fp

    def generate_drill_report_docx(self, filename=None):
        fn = filename or f"Drill_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Emergency Drill Report", 0)
        self._m(d, {"Drill No": "MRPL/MD/2026/0006", "Scenario": "Major Fire in CDU-III Heat Exchanger Area",
                     "Date": datetime.now().strftime("%d-%b-%Y"), "Duration": "42 minutes (07:00 - 07:42)",
                     "Participants": "124 personnel (all on shift)"})
        self._h(d, "1. Timeline", 1)
        self._t(d, 9, 3, [["Time", "Event", "Target"], ["07:00", "Alarm activated (manual call point)", "-"],
                           ["07:02", "Control room acknowledges, calls fire team", "< 1 min"],
                           ["07:05", "Fire team arrives at scene (6 members)", "< 5 min"],
                           ["07:08", "Muster point roll call started", "< 8 min"],
                           ["07:12", "Fire team deploys 2x foam monitors", "< 10 min"],
                           ["07:18", "All 124 personnel accounted for", "< 15 min"],
                           ["07:30", "Rescue team extracts simulated casualty", "< 20 min"],
                           ["07:42", "ALL-CLEAR signal, drill concluded", "42 min"]])
        self._h(d, "2. Performance Assessment", 1)
        self._t(d, 6, 3, [["Criterion", "Result", "Rating"], ["Alarm response time", "2 min", "GOOD"],
                           ["Fire team response", "5 min", "GOOD"],
                           ["Muster assembly", "18 min", "ACCEPTABLE"],
                           ["Casualty extraction", "30 min", "ACCEPTABLE"],
                           ["Communication clarity", "Good", "GOOD"]])
        self._h(d, "3. Observations & Actions", 1)
        self._n(d, ["OBS-01: 3 persons did not hear alarm (noise area) - ACTION: Install strobe lights",
                     "OBS-02: Muster point crowded at Gate B - ACTION: Define secondary muster point",
                     "OBS-03: Fire hose coupling delay 90 sec - ACTION: Replace old couplings"])
        self._p(d, "Overall Rating: SATISFACTORY. All critical objectives met.")
        d.save(fp)
        return fp

    def generate_jsa_document_docx(self, filename=None):
        fn = filename or f"JSA_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Job Safety Analysis (JSA)", 0)
        self._m(d, {"JSA No": "MRPL/JSA/2026/00892", "Job": "Replacement of Gasket on 10in Flange (E-4201)",
                     "Location": "CDU-III Level 3, Pipe Rack", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Job Steps & Hazards", 1)
        self._t(d, 9, 4, [["Step", "Job Description", "Potential Hazard", "Control Measure"],
                           ["1", "Obtain PTW and brief crew", "Unauthorized work", "Verify permit, attend toolbox talk"],
                           ["2", "Isolate and drain line (LOTO)", "Residual pressure/fluid", "Depressure, drain, verify zero energy"],
                           ["3", "Loosen flange bolts", "Falling wrench, struck-by", "Use safety harness, tool lanyards"],
                           ["4", "Remove old gasket", "Sharp edges, chemical", "Cut-resistant gloves, RPE"],
                           ["5", "Clean flange faces", "Solvent exposure, sparks", "Non-flammable solvent, no ignition sources"],
                           ["6", "Install new gasket", "Pinch points, heavy lifting", "Mechanical aids, 2-person lift"],
                           ["7", "Re-torque bolts", "Musculoskeletal", "Torque wrench, proper stance"],
                           ["8", "Remove LOTO, pressurize", "Leak, caught off-guard", "Stand clear, leak test at 1.5x"]])
        self._h(d, "2. PPE Requirements", 1)
        self._n(d, ["Safety helmet with chin strap", "Safety glasses (side shields)", "Hearing protection (if >85 dBA)",
                     "Cut-resistant gloves (Level 3)", "Steel toe safety boots", "Coveralls (FR for hydrocarbon area)",
                     "Fall arrest harness (if >1.8 m)"])
        self._h(d, "3. Emergency Arrangements", 1)
        self._m(d, {"Nearest Fire Point": "6 m from work area", "Assembly Point": "North Muster Point",
                     "Emergency Contact": "Ext. 5555", "First Aider": "Shri A.V. Kulkarni"})
        self._h(d, "4. Sign-Off", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Job Supervisor", "Shri R.M. Gupta", "________________"],
                           ["Safety Representative", "Shri A.V. Kulkarni", "________________"],
                           ["All Workers (List attached)", "-", "________________"]])
        d.save(fp)
        return fp

    def generate_sds_sheet_docx(self, filename=None):
        fn = filename or f"SDS_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Safety Data Sheet (SDS) - GHS Format", 0)
        self._m(d, {"Product Name": "Hydrochloric Acid (HCl) 30% w/w", "CAS No": "7647-01-0",
                     "Supplier": "MRPL Chemical Store", "Date": datetime.now().strftime("%d-%b-%Y"),
                     "GHS Pictograms": "GHS05 (Corrosion), GHS04 (Gas Cylinder)"})
        self._h(d, "1. Identification", 1)
        self._m(d, {"Product Identifier": "Hydrochloric Acid 30%", "Use": "Acidizing, pH adjustment, cleaning",
                     "Supplier Emergency": "+91-98765-43210"})
        self._h(d, "2. Hazard Identification", 1)
        self._n(d, ["H314: Causes severe skin burns and eye damage", "H290: May be corrosive to metals",
                     "Signal Word: DANGER", "Prevention: P260, P264, P280, P301+P330+P331",
                     "Response: P301+P330+P331, P303+P361+P353, P305+P351+P338"])
        self._h(d, "3. Composition / Information on Ingredients", 1)
        self._t(d, 4, 3, [["Component", "CAS No", "Concentration"], ["Hydrochloric Acid", "7647-01-0", "30% w/w"],
                           ["Water", "7732-18-5", "68-70%"], ["Iron (impurity)", "7439-89-6", "<0.01%"]])
        self._h(d, "4. First Aid Measures", 1)
        self._n(d, ["Inhalation: Remove to fresh air. Seek medical attention.",
                     "Skin Contact: Remove contaminated clothing. Rinse with water for 20 min.",
                     "Eye Contact: Rinse cautiously with water for 15 min. Seek ophthalmologist.",
                     "Ingestion: Rinse mouth. Do NOT induce vomiting. Seek immediate medical attention."])
        self._h(d, "5. Fire-Fighting Measures", 1)
        self._n(d, ["Non-flammable. Use CO2, dry chemical, or foam for surrounding fires.",
                     "Wear SCBA and full protective gear.",
                     "Neutralize spill with soda ash or lime."])
        self._h(d, "6. Storage & Handling", 1)
        self._n(d, ["Store in compatible (HDPE/FRP) containers, away from metals",
                     "Keep in ventilated area, below 35C", "Use only with adequate ventilation",
                     "Ground containers during transfer to prevent static buildup"])
        d.save(fp)
        return fp

    def generate_incident_investigation_docx(self, filename=None):
        fn = filename or f"Incident_Investigation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Incident/Accident Investigation Report (5-Why Analysis)", 0)
        self._m(d, {"Incident No": "MRPL/INC/2026/00023", "Type": "Recordable - Lost Time Injury",
                     "Date of Incident": "18-Mar-2026", "Location": "CDU-III Pump P-4401 Area",
                     "Injured Person": "Shri P.R. Deshmukh (TPL Contractor)"})
        self._h(d, "1. Incident Description", 1)
        self._p(d, "Worker sustained a crush injury to left hand while reconnecting a flexible hose to pump P-4401. The hose coupling slipped under pressure causing the worker's hand to be caught between the coupling and the pipe support. Worker was wearing gloves but they were not cut-resistant rated for pinch-point hazards.")
        self._h(d, "2. 5-Why Analysis", 1)
        self._t(d, 6, 2, [["Level", "Finding"], ["Why 1: Why did the injury occur?", "Worker's hand was caught between hose coupling and pipe support during reconnection."],
                           ["Why 2: Why was the hand in the crush zone?", "Worker was manually aligning the coupling while pressurizing the hose."],
                           ["Why 3: Why was the hose pressurized during alignment?", "Procedure did not specify de-pressurize-before-alignment step for this task."],
                           ["Why 4: Why was the procedure inadequate?", "Original SOP written for rigid pipe only; flexible hose modification added without SOP update."],
                           ["Why 5: Why was the SOP not updated?", "MOC process did not trigger SOP review for this minor modification (2022)."]])
        self._h(d, "3. Root Causes", 1)
        self._n(d, ["Incomplete MOC process for equipment modifications", "Inadequate SOP for flexible hose connection procedure",
                     "Task-specific risk assessment not performed for modified task"])
        self._h(d, "4. Corrective Actions", 1)
        self._t(d, 5, 4, [["#", "Action", "Owner", "Due Date"], ["1", "Update SOP for flexible hose connections", "Shri S.V. Patil", "05-Apr-2026"],
                           ["2", "Conduct task JSA for all modified equipment", "Shri R.M. Gupta", "15-Apr-2026"],
                           ["3", "Provide cut-resistant gloves (Level 5) for coupling work", "Shri D.V. Sharma", "25-Mar-2026"],
                           ["4", "MOC trigger review for all equipment modifications", "Shri K.R. Nair", "30-Apr-2026"]])
        self._h(d, "5. Sign-Off", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Investigation Lead", "Shri D.V. Sharma", "________________"],
                           ["HSE Manager", "Shri A.D. Mistry", "________________"],
                           ["Unit Head", "Shri S.V. Patil", "________________"]])
        d.save(fp)
        return fp

    def generate_near_miss_report_docx(self, filename=None):
        fn = filename or f"Near_Miss_Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Near-Miss Report", 0)
        self._m(d, {"Report No": "MRPL/NM/2026/00089", "Date": datetime.now().strftime("%d-%b-%Y"),
                     "Location": "CDU-III Column C-301, Level 4"})
        self._h(d, "1. Description", 1)
        self._p(d, "A 2 kg wrench fell from Level 4 scaffold (18.5 m) and landed within the barricaded exclusion zone. No personnel were in the drop zone at the time. If the barricade had not been in place, a person walking by could have been struck.")
        self._h(d, "2. Potential Consequences", 1)
        self._n(d, ["Fatality or serious head injury to personnel below",
                     "Damage to equipment/piping below",
                     "Potential for secondary incident (startling, evasive action)"])
        self._h(d, "3. Root Causes", 1)
        self._n(d, ["Tool was not secured with lanyard", "Worker was working at height without tool tethering",
                     "Tool bag was open and not secured to scaffold"])
        self._h(d, "4. Immediate Actions", 1)
        self._n(d, ["Area barricaded, tool recovered", "Worker reminded of tool lanyard requirement",
                     "Tool lanyards checked - 2 more unsecured tools found and corrected"])
        self._h(d, "5. Corrective Actions", 1)
        self._t(d, 4, 3, [["#", "Action", "Owner"], ["1", "Mandatory tool lanyards for all WAH >3m", "Shri A.V. Kulkarni"],
                           ["2", "Tool box audit before each WAH shift", "Shri R.M. Gupta"],
                           ["3", "Tool tethering training refresher", "Shri D.V. Sharma"]])
        self._p(d, "Severity: Potential Fatal / Actual: Near-Miss. Lucky escape - barricade prevented exposure.")
        d.save(fp)
        return fp

    def generate_ppe_compliance_register_xlsx(self, filename=None):
        fn = filename or f"PPE_Compliance_Register_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        fp = os.path.join(OUTPUT_DIR, "xlsx", fn)
        wb = Workbook()
        ws = wb.active; ws.title = "PPE Compliance"
        hf = Font(bold=True, color="FFFFFF", size=11)
        hrf = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
        gf = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        rf = PatternFill(start_color="FF6666", end_color="FF6666", fill_type="solid")
        yf = PatternFill(start_color="FFCC66", end_color="FFCC66", fill_type="solid")
        headers = ["Emp ID", "Name", "Company", "Helmet", "Safety Glasses", "Gloves", "Boots", "Hearing", "Harness", "Coverall", "Overall Status"]
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=c, value=h); cell.font = hf; cell.fill = hrf; cell.alignment = Alignment(horizontal="center", wrap_text=True)
        data = [
            ["MRPL-001", "Shri K.R. Nair", "MRPL", "OK", "OK", "OK", "OK", "OK", "N/A", "OK", "COMPLIANT"],
            ["MRPL-002", "Shri S.V. Patil", "MRPL", "OK", "OK", "OK", "OK", "OK", "N/A", "OK", "COMPLIANT"],
            ["TPL-4401", "Shri P.R. Deshmukh", "Tata Projects", "OK", "OK", "EXPIRED", "OK", "OK", "OK", "OK", "NON-COMPLIANT"],
            ["TPL-4402", "Shri A.V. Kulkarni", "Tata Projects", "OK", "OK", "OK", "OK", "OK", "OK", "OK", "COMPLIANT"],
            ["L&T-1101", "Shri R.M. Gupta", "L&T Hydrocarbon", "OK", "OK", "OK", "EXPIRED", "OK", "OK", "OK", "NON-COMPLIANT"],
            ["MRPL-005", "Shri D.V. Sharma", "MRPL", "OK", "OK", "OK", "OK", "OK", "N/A", "OK", "COMPLIANT"],
            ["TPL-4403", "Shri S.K. Mane", "Tata Projects", "OK", "EXPIRED", "OK", "OK", "OK", "OK", "OK", "NON-COMPLIANT"],
            ["MRPL-008", "Shri A.D. Mistry", "MRPL", "OK", "OK", "OK", "OK", "OK", "N/A", "OK", "COMPLIANT"],
        ]
        for ri, rd in enumerate(data, 2):
            for ci, v in enumerate(rd, 1):
                cell = ws.cell(row=ri, column=ci, value=v); cell.alignment = Alignment(horizontal="center", wrap_text=True)
                if ci == 11:
                    cell.fill = gf if v == "COMPLIANT" else rf; cell.font = Font(bold=True)
                elif ci > 2 and ci < 11:
                    if v == "OK": cell.fill = gf
                    elif v == "EXPIRED": cell.fill = rf; cell.font = Font(bold=True)
                    elif v == "N/A": cell.fill = PatternFill(start_color="D3D3D3", end_color="D3D3D3", fill_type="solid")
        for col in ws.columns:
            ml = max(len(str(c.value or "")) for c in col)
            ws.column_dimensions[col[0].column_letter].width = min(ml + 2, 20)
        wb.save(fp)
        return fp

    def generate_contractor_worker_safety_docx(self, filename=None):
        fn = filename or f"Contractor_Safety_Pass_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Contractor Worker Safety Pass", 0)
        self._m(d, {"Pass No": "MRPL/CSP/2026/00891", "Name": "Shri Anil K. Sharma",
                     "Company": "Tata Projects Ltd", "ID No": "TPL-4401", "Valid From": "01-Jan-2026",
                     "Valid To": "31-Dec-2026"})
        self._h(d, "1. Worker Details", 1)
        self._t(d, 6, 2, [["Field", "Detail"], ["Full Name", "Anil Kumar Sharma"],
                           ["Date of Birth", "15-Mar-1988"], ["Trade", "Pipe Fitter (Certified)"],
                           ["Experience", "8 years"], ["Medical Fitness", "FIT (Dr. Kulkarni, 01-Jan-2026)"]])
        self._h(d, "2. Safety Certifications", 1)
        self._t(d, 7, 3, [["Certification", "Valid Until", "Status"], ["Safety Induction (MRPL)", "31-Dec-2026", "VALID"],
                           ["Work at Height", "30-Jun-2026", "VALID"],
                           ["Confined Space Entry", "30-Jun-2026", "VALID"],
                           ["Fire Fighting (Basic)", "31-Dec-2026", "VALID"],
                           ["First Aid", "31-Dec-2026", "VALID"],
                           ["BBS Observer", "31-Dec-2026", "VALID"]])
        self._h(d, "3. Authorized Areas", 1)
        self._n(d, ["Zone 1: Process Area (with escort for first 3 visits)", "Zone 2: Utilities & Fire System Area",
                     "Zone 3: Scaffold Erection Areas", "NOT AUTHORIZED: Zone 0 (Wellhead / Reactor Area)"])
        self._h(d, "4. PPE Issued", 1)
        self._t(d, 5, 3, [["PPE Item", "Issue Date", "Due for Replacement"], ["Safety Helmet", "01-Jan-2026", "01-Jan-2027"],
                           ["Safety Glasses", "01-Jan-2026", "01-Jul-2026"],
                           ["Cut-Resistant Gloves (L3)", "15-Mar-2026", "15-Jun-2026"],
                           ["Steel Toe Boots", "01-Jan-2026", "01-Jul-2026"]])
        d.save(fp)
        return fp

    def generate_ptw_register_xlsx(self, filename=None):
        fn = filename or f"PTW_Register_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        fp = os.path.join(OUTPUT_DIR, "xlsx", fn)
        wb = Workbook()
        ws = wb.active; ws.title = "PTW Register"
        hf = Font(bold=True, color="FFFFFF", size=11)
        hrf = PatternFill(start_color="003366", end_color="003366", fill_type="solid")
        gf = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        rf = PatternFill(start_color="FF6666", end_color="FF6666", fill_type="solid")
        yf = PatternFill(start_color="FFCC66", end_color="FFCC66", fill_type="solid")
        bf = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        headers = ["PTW No", "Type", "Location", "Activity", "Permit Holder", "Start Time", "End Time", "Status", "Issuer"]
        for c, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=c, value=h); cell.font = hf; cell.fill = hrf; cell.alignment = Alignment(horizontal="center", wrap_text=True)
        data = [
            ["PTW-2026/0561", "Cold Work", "CDU-III Level 2", "Grinding on 6in line", "Shri R.M. Gupta", "08:00", "17:00", "ACTIVE", "Shri K.R. Nair"],
            ["PTW-2026/0562", "Hot Work", "CDU-III Pipe Rack", "Welding on bracket", "Shri P.R. Deshmukh", "08:30", "16:30", "ACTIVE", "Shri K.R. Nair"],
            ["PTW-2026/0563", "Confined Space", "V-4201 Tank", "Internal inspection", "Shri S.K. Mane", "09:00", "15:00", "ACTIVE", "Shri S.V. Patil"],
            ["PTW-2026/0564", "Excavation", "Utility Corridor", "Trench digging", "Shri A.V. Kulkarni", "07:00", "17:00", "COMPLETED", "Shri K.R. Nair"],
            ["PTW-2026/0565", "Electrical", "MCC Room", "Panel modification", "Shri T.K. Verma", "10:00", "14:00", "ACTIVE", "Shri P.V. Rao"],
            ["PTW-2026/0558", "Cold Work", "Unit Level 3", "Flange gasket change", "Shri R.M. Gupta", "08:00", "12:00", "CLOSED", "Shri K.R. Nair"],
            ["PTW-2026/0559", "Hot Work", "Unit Pipe Rack", "Support welding", "Shri A.V. Kulkarni", "08:00", "16:00", "CLOSED", "Shri K.R. Nair"],
        ]
        for ri, rd in enumerate(data, 2):
            for ci, v in enumerate(rd, 1):
                cell = ws.cell(row=ri, column=ci, value=v); cell.alignment = Alignment(horizontal="center", wrap_text=True)
                if ci == 8:
                    if v == "ACTIVE": cell.fill = gf
                    elif v == "COMPLETED": cell.fill = bf; cell.font = Font(color="FFFFFF", bold=True)
                    elif v == "CLOSED": cell.fill = PatternFill(start_color="A0A0A0", end_color="A0A0A0", fill_type="solid")
                    elif v == "EXPIRED": cell.fill = rf; cell.font = Font(bold=True)
        for col in ws.columns:
            ml = max(len(str(c.value or "")) for c in col)
            ws.column_dimensions[col[0].column_letter].width = min(ml + 2, 22)
        wb.save(fp)
        return fp

    def generate_hse_composite_permit_pptx(self, filename=None):
        fn = filename or f"HSE_Composite_Permit_Review_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        fp = os.path.join(OUTPUT_DIR, "pptx", fn)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = "HSE Composite Permit Review"
        sl.placeholders[1].text = f"Unit Complex\nDaily Permit Review Meeting\n{datetime.now().strftime('%d %B %Y')}"
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = "Agenda"
        tf = sl2.placeholders[1].text_frame
        for item in ["Active Permits Summary", "High-Risk Permits Review", "Expired/Overdue Permits", "Gas Testing Status", "HSE Observations", "Action Items"]:
            p = tf.add_paragraph(); p.text = item; p.level = 0
        sl3 = prs.slides.add_slide(prs.slide_layouts[5])
        sl3.shapes.title.text = "Active Permits by Type"
        ts = sl3.shapes.add_table(8, 3, PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(4.5)).table
        for i, rd in enumerate([["Permit Type", "Active", "Total Today"], ["Cold Work", "4", "7"], ["Hot Work", "2", "5"],
                                 ["Confined Space", "1", "2"], ["Excavation", "1", "3"], ["Electrical", "1", "2"],
                                 ["Working at Height", "3", "6"], ["Lifting", "2", "4"]]):
            for j, ct in enumerate(rd): ts.cell(i, j).text = ct
        sl4 = prs.slides.add_slide(prs.slide_layouts[1])
        sl4.shapes.title.text = "High-Risk Permits"
        tf2 = sl4.placeholders[1].text_frame
        for item in ["PTW-0563: Confined Space E-4201 - 24-hr continuous gas monitoring active",
                      "PTW-0565: Electrical isolation MCC Room - LOTO verified, standby maintained",
                      "PTW-0561: Cold Work near flare - Fire watch on standby, wind 12 km/h (limit 35)"]:
            p = tf2.add_paragraph(); p.text = item; p.level = 0
        sl5 = prs.slides.add_slide(prs.slide_layouts[1])
        sl5.shapes.title.text = "Action Items"
        tf3 = sl5.placeholders[1].text_frame
        for item in ["PTW-0562 expiry at 16:30 - confirm extension or closure",
                      "Gas test due for PTW-0563 at 13:00", "Toolbox talk for lifting ops scheduled 06:30 tomorrow",
                      "Night shift permit handover at 18:00"]:
            p = tf3.add_paragraph(); p.text = item; p.level = 0
        prs.save(fp)
        return fp

    # ── Refinery Operations Methods ────────────────────────────────────────

    def generate_operating_manual_docx(self, filename=None):
        fn = filename or f"Operating_Manual_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Operating Manual - CDU-III Crude Distillation Unit", 0)
        self._m(d, {"Document No": "MRPL/OPS/OM/2026/CDU-III-001", "Unit": "CDU-III",
                     "Capacity": "6 MTPA (120,000 BPD)", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Unit Description", 1)
        self._p(d, "CDU-III is a atmospheric crude distillation unit designed to process 120,000 BPD of mixed Indian crudes (Bombay High, Krishna-Godavari, Rajasthan). The unit includes crude preheat train, desalter, atmospheric column with side strippers, and naphtha/kerosene/diesel product systems.")
        self._h(d, "2. Process Parameters", 1)
        self._t(d, 8, 3, [["Parameter", "Design", "Normal Range"], ["Crude Feed Rate", "120,000 BPD", "100,000 - 120,000 BPD"],
                           ["Column Top Pressure", "1.2 kg/cm2g", "1.0 - 1.4 kg/cm2g"],
                           ["Stripper Temperature", "345C", "335 - 355C"],
                           ["Pumparound Rate", "280 m3/hr", "250 - 300 m3/hr"],
                           ["Desalter Temperature", "130C", "120 - 140C"],
                           ["Desalter Pressure", "8 kg/cm2g", "6 - 10 kg/cm2g"],
                           ["Overhead Temperature", "105C", "95 - 115C"]])
        self._h(d, "3. Operating Procedures", 1)
        self._n(d, ["Normal Start-up sequence (16-hour timeline)", "Normal Shutdown sequence (12-hour timeline)",
                     "Emergency Shutdown (ESD) procedures", "Emergency response procedures for leaks, fires, upsets",
                     "Product quality control (flash point, pour point, sulfur)",
                     "Alarm management and operator response"])
        self._h(d, "4. Safety Systems", 1)
        self._t(d, 5, 3, [["System", "Set Point", "Action"], ["Column High Pressure", "2.5 kg/cm2g", "ESD - Feed cut"],
                           ["Column High Temperature", "380C", "Alarm + Trip pumparound"],
                           ["Desalter High Level", "85%", "Alarm + Feed rate reduction"],
                           ["Overhead Receiver High Level", "80%", "Alarm + Pump start"]])
        d.save(fp)
        return fp

    def generate_pid_document_docx(self, filename=None):
        fn = filename or f"PID_Document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Process & Instrumentation Diagram (P&ID) Document", 0)
        self._m(d, {"Document No": "MRPL/PID/CDU-III/001-008", "Unit": "CDU-III Crude Distillation",
                     "Drawing Series": "CDU3-PID-001 through CDU3-PID-008", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Drawing Index", 1)
        self._t(d, 9, 3, [["Drawing No", "Description", "Status"], ["CDU3-PID-001", "Crude Feed & Preheat Train", "IFC Rev D"],
                           ["CDU3-PID-002", "Desalter System", "IFC Rev C"],
                           ["CDU3-PID-003", "Atmospheric Column (Main)", "IFC Rev D"],
                           ["CDU3-PID-004", "Side Strippers (Naphtha/Kero/Diesel)", "IFC Rev C"],
                           ["CDU3-PID-005", "Overhead System & Reflux", "IFC Rev B"],
                           ["CDU3-PID-006", "Pumparound Systems", "IFC Rev B"],
                           ["CDU3-PID-007", "Product Storage & Blending", "IFC Rev C"],
                           ["CDU3-PID-008", "Utility Connections (Steam, CW, N2)", "IFC Rev A"]])
        self._h(d, "2. Instrument Index Summary", 1)
        self._t(d, 6, 3, [["Instrument Type", "Count", "Range"], ["Temperature (TI/TIC)", "48", "-20C to 500C"],
                           ["Pressure (PI/PIC)", "32", "0 - 50 kg/cm2g"],
                           ["Flow (FI/FIC)", "26", "0 - 500 m3/hr"],
                           ["Level (LI/LIC)", "18", "0 - 100%"],
                           ["Analyzers (AI)", "8", "Various"]])
        self._h(d, "3. Valve Index Summary", 1)
        self._t(d, 5, 3, [["Valve Type", "Count", "Material"], ["Gate Valve", "128", "CS / SS316"],
                           ["Globe Valve", "64", "CS / SS316"], ["Check Valve", "36", "CS"],
                           ["Control Valve", "42", "CS / Alloy"]])
        self._h(d, "4. Safety Valve Settings", 1)
        self._t(d, 5, 3, [["Tag", "Location", "Set Pressure"], ["PSV-301A/B", "Column Overhead", "3.0 kg/cm2g"],
                           ["PSV-302A/B", "Pumparound Circuit", "12.0 kg/cm2g"],
                           ["PSV-303A/B", "Desalter", "15.0 kg/cm2g"]])
        d.save(fp)
        return fp

    def generate_cause_effect_diagram_docx(self, filename=None):
        fn = filename or f"Cause_Effect_Diagram_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Cause & Effect Diagram (SIS) - CDU-III", 0)
        self._m(d, {"Document No": "MRPL/SIS/CE/2026/CDU-III-001", "Unit": "CDU-III",
                     "SIL Rating": "SIL 2 / SIL 3", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. ESD Function Matrix", 1)
        self._t(d, 7, 5, [["Function", "Initiator", "SIL", "Action", "Final Element"],
                           ["ESD-001: Unit Shutdown", "High column pressure >2.5 kg/cm2g", "SIL 3", "Full feed cut, all pumps off", "XV-301, XV-302, XV-303"],
                           ["ESD-002: Column Trip", "High column temperature >380C", "SIL 2", "Feed cut, reflux max", "XV-301, XV-302"],
                           ["ESD-003: Desalter Trip", "High desalter level >90%", "SIL 2", "Desalter isolation", "XV-310, XV-311"],
                           ["ESD-004: Fire ESD", "F&G detection (2oo3)", "SIL 3", "Full unit isolation", "All XVs"],
                           ["ESD-005: seismic", "Seismic trigger >0.5g", "SIL 3", "Full unit isolation", "All XVs"],
                           ["ESD-006: Gas Detection", "LEL >25% (2oo3)", "SIL 2", "Isolate source, alert", "XV-320, XV-321"]])
        self._h(d, "2. Cause & Effect Logic", 1)
        self._n(d, ["All SIS functions per IEC 61511 / OISD-STD-178", "Voting logic: 2oo3 for critical functions, 1oo2 for SIL 2",
                     "Final elements: automated shutdown valves with 15-second stroke time",
                     "Annual proof testing required (documented in SIS register)",
                     "Bypass management: no bypass >24 hours without Plant Manager approval"])
        self._h(d, "3. SIF Summary", 1)
        self._t(d, 5, 4, [["SIF ID", "Description", "SIL", "PFDavg"], ["SIF-001", "High Pressure Trip", "SIL 3", "0.00098"],
                           ["SIF-002", "High Temperature Trip", "SIL 2", "0.0042"],
                           ["SIF-003", "Desalter Level Trip", "SIL 2", "0.0038"],
                           ["SIF-004", "Fire & Gas ESD", "SIL 3", "0.00085"]])
        d.save(fp)
        return fp

    def generate_hmb_sheet_docx(self, filename=None):
        fn = filename or f"HMB_Sheet_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Heat & Material Balance Sheet - CDU-III", 0)
        self._m(d, {"Document No": "MRPL/HMB/CDU-III/2026-001", "Basis": "100% Design Capacity (120,000 BPD)",
                     "Crude Blend": "60% Bombay High + 25% KG + 15% Rajasthan", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Material Balance Summary", 1)
        self._t(d, 9, 4, [["Stream", "Rate (MT/hr)", "Rate (BPD)", "Density (kg/m3)"], ["Crude Feed", "5,000", "120,000", "852"],
                           ["Light Naphtha (C5-80C)", "350", "12,600", "660"],
                           ["Heavy Naphtha (80-180C)", "520", "15,800", "740"],
                           ["Kerosene (180-250C)", "680", "17,200", "810"],
                           ["Light Diesel (250-350C)", "1,050", "24,000", "840"],
                           ["Atmospheric Residue (>350C)", "2,280", "48,000", "960"],
                           ["Light Ends (C1-C4)", "120", "N/A (gas)", "-"],
                           ["Losses (0.3%)", "15", "-", "-"]])
        self._h(d, "2. Heat Balance Summary", 1)
        self._t(d, 6, 3, [["Item", "Heat Duty (GCal/hr)", "Temperature (C)"], ["Crude Preheat (fired heater)", "185.0", "345 (outlet)"],
                           ["Column Overhead Condenser", "-52.0", "105 (overhead)"],
                           ["Pumparound 1 (Naphtha)", "-28.5", "150 (return)"],
                           ["Pumparound 2 (Kerosene)", "-35.2", "220 (return)"],
                           ["Stripper Steam (MDM)", "12.5", "350 (superheated)"]])
        self._h(d, "3. Energy Efficiency", 1)
        self._t(d, 5, 2, [["Indicator", "Value"], ["Specific Energy Consumption", "8.2 GCal/MT crude"],
                           ["Fired Heater Efficiency", "92%"], ["Heat Recovery Factor", "78%"],
                           ["CO2 Emission Intensity", "32 kg CO2/MT crude"]])
        d.save(fp)
        return fp

    def generate_moc_document_docx(self, filename=None):
        fn = filename or f"MOC_Document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Management of Change (OISD-STD-178)", 0)
        self._m(d, {"MOC No": "MRPL/MOC/2026/00045", "Unit": "CDU-III", "Title": "Replacement of Control Valve CV-3012 with Smart Positioner",
                     "Category": "Temporary", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Description of Change", 1)
        self._p(d, "Replace existing Fisher Positioner (pneumatic) with Emerson DeltaV smart positioner on Control Valve CV-3012 (Naphtha Draw). Includes addition of digital position feedback transmitter and integration with DCS for advanced diagnostics.")
        self._h(d, "2. Technical Justification", 1)
        self._n(d, ["Improved valve diagnostics (HART protocol)", "Predictive maintenance capability",
                     "Reduced air consumption (30% savings)", "Enhanced process control (faster response time)",
                     "Compatibility with existing DeltaV DCS"])
        self._h(d, "3. Risk Assessment (HAZOP Review)", 1)
        self._t(d, 5, 4, [["Hazard", "Consequence", "Safeguard", "Residual Risk"], ["Smart positioner failure", "Loss of valve control", "Manual bypass valve", "LOW"],
                           ["DCS communication loss", "Valve fails to last position", "Fail-last positioner", "LOW"],
                           ["Electrical interference", "Erratic valve movement", "Shielded cables", "LOW"],
                           ["Software configuration error", "Incorrect valve action", "Pre-commissioning test", "LOW"]])
        self._h(d, "4. Approvals", 1)
        self._t(d, 5, 3, [["Role", "Name", "Signature"], ["Unit Engineer", "Shri S.V. Patil", "________________"],
                           ["Instrument Lead", "Shri T.K. Verma", "________________"],
                           ["HSE Lead", "Shri D.V. Sharma", "________________"],
                           ["Plant Manager", "Shri A.D. Mistry", "________________"]])
        self._h(d, "5. Implementation Plan", 1)
        self._t(d, 5, 3, [["Phase", "Activity", "Target Date"], ["1", "Engineering review & procurement", "15-Apr-2026"],
                           ["2", "Installation (during planned shutdown)", "25-Apr-2026"],
                           ["3", "Loop check & commissioning", "28-Apr-2026"],
                           ["4", "Performance monitoring (30 days)", "28-May-2026"]])
        d.save(fp)
        return fp

    def generate_pssr_checklist_docx(self, filename=None):
        fn = filename or f"PSSR_Checklist_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Pre-Startup Safety Review (PSSR) Checklist", 0)
        self._m(d, {"Checklist No": "MRPL/PSSR/2026/00012", "Unit": "CDU-III (Post-Turnaround Restart)",
                     "Turnaround Period": "01-Mar-2026 to 15-Mar-2026", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Engineering Verification", 1)
        self._t(d, 7, 3, [["#", "Item", "Status"], ["1.1", "P&ID markups reconciled with as-built", "COMPLETE"],
                           ["1.2", "Equipment nameplates verified", "COMPLETE"],
                           ["1.3", "Instrument loop checks (142 loops)", "COMPLETE"],
                           ["1.4", "Control valve stroke tests", "COMPLETE"],
                           ["1.5", "Safety valve set pressure verified", "COMPLETE"],
                           ["1.6", "Electrical system checks (MCC, motors)", "COMPLETE"]])
        self._h(d, "2. Safety System Verification", 1)
        self._t(d, 7, 3, [["#", "Item", "Status"], ["2.1", "ESD function test (all SIFs)", "COMPLETE"],
                           ["2.2", "Fire & Gas system functional test", "COMPLETE"],
                           ["2.3", "Fire water system test", "COMPLETE"],
                           ["2.4", "Deluge system test", "COMPLETE"],
                           ["2.5", "BPCS alarm setpoints verified", "COMPLETE"],
                           ["2.6", "Operator training completed", "COMPLETE"]])
        self._h(d, "3. Operational Readiness", 1)
        self._n(d, ["SOP updated and approved", "Emergency response plan updated",
                     "Operator competency assessment complete", "Turnaround punch list closed (0 Cat A, 2 Cat B)",
                     "Spare parts inventory verified"])
        self._h(d, "4. Sign-Off", 1)
        self._t(d, 5, 3, [["Role", "Name", "Signature"], ["PSSR Chairman", "Shri A.D. Mistry", "________________"],
                           ["Operations Lead", "Shri S.V. Patil", "________________"],
                           ["Maintenance Lead", "Shri K.R. Nair", "________________"],
                           ["HSE Lead", "Shri D.V. Sharma", "________________"]])
        self._p(d, "RESULT: PSSR APPROVED. Unit ready for safe start-up.")
        d.save(fp)
        return fp

    def generate_psmdocument_docx(self, filename=None):
        fn = filename or f"PSM_Document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Process Safety Management Document - CDU-III", 0)
        self._m(d, {"Document No": "MRPL/PSM/2026/CDU-III-001", "Unit": "CDU-III Crude Distillation",
                     "Standard": "OISD-STD-178 / OISD-STD-153", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Process Safety Information", 1)
        self._n(d, ["Chemical inventory and SDS for all process fluids", "Process flow data: temperatures, pressures, flow rates",
                     "Equipment design data (material of construction, design code)",
                     "P&IDs (8 drawings, IFC Rev C/D)", "Cause & Effect diagrams (6 SIFs)"])
        self._h(d, "2. Process Hazard Analysis (PHA)", 1)
        self._t(d, 6, 3, [["#", "PHA Study", "Date", "Status"], ["1", "HAZOP (Original)", "Mar-2018", "COMPLETE"],
                           ["2", "HAZOP Revalidation", "Jan-2024", "COMPLETE"],
                           ["3", "LOPA (SIL Assessment)", "Jan-2024", "COMPLETE"],
                           ["4", "What-If Analysis (Minor)", "Mar-2025", "COMPLETE"],
                           ["5", "SIL Verification (Proof Test)", "Dec-2025", "COMPLETE"]])
        self._h(d, "3. Operating Procedures", 1)
        self._t(d, 5, 2, [["Procedure", "Last Review", "Next Due"], ["Start-up", "Mar-2026", "Mar-2027"],
                           ["Normal Operation", "Mar-2026", "Mar-2027"],
                           ["Shutdown", "Mar-2026", "Mar-2027"],
                           ["Emergency", "Mar-2026", "Mar-2027"]])
        self._h(d, "4. Mechanical Integrity", 1)
        self._t(d, 6, 3, [["Equipment", "Inspection Method", "Next Due"], ["Column C-301", "UT thickness / visual", "Mar-2027"],
                           ["Exchangers (12 nos)", "Tube bundle inspection", "Mar-2027"],
                           ["Piping (critical)", "UT / ACFM", "Mar-2028"],
                           ["PSVs (8 nos)", "Bench test", "Mar-2027"],
                           ["Storage Tanks", "API 653 inspection", "Dec-2026"]])
        self._h(d, "5. Incident Investigation", 1)
        self._t(d, 5, 3, [["Year", "Incidents", "Near-Miss", "Root Cause Analysis"], ["2025", "1 (LTI)", "12", "5-Why (all)"],
                           ["2024", "0", "18", "5-Why (all)"],
                           ["2023", "2 (LTI)", "15", "5-Why + Bowtie"],
                           ["2026 YTD", "1 (LTI)", "8", "5-Why (ongoing)"]])
        d.save(fp)
        return fp

    def generate_sop_document_docx(self, filename=None):
        fn = filename or f"SOP_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Standard Operating Procedure (SOP) - CDU-III Start-up", 0)
        self._m(d, {"SOP No": "MRPL/SOP/CDU-III/SU-001", "Unit": "CDU-III",
                     "Title": "Normal Start-up from Cold State", "Rev": "Rev 5", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Purpose & Scope", 1)
        self._p(d, "This SOP defines the step-by-step procedure for safe start-up of CDU-III from a cold (depressured, depressurized) state to normal operating conditions. Applicable to Operations, Control Room, and Field personnel.")
        self._h(d, "2. Prerequisites", 1)
        self._n(d, ["PSSR completed and signed", "All SIS functions tested and verified", "Operator competency verified",
                     "Emergency response plan current", "Utility systems available (steam, cooling water, nitrogen)"])
        self._h(d, "3. Start-up Sequence", 1)
        self._t(d, 9, 4, [["Phase", "Step", "Action", "Duration"], ["1", "1.1", "Establish utility systems (CW, steam, N2)", "2 hr"],
                           ["1", "1.2", "Nitrogen purge crude column", "3 hr"],
                           ["2", "2.1", "Introduce crude at 20% rate (24,000 BPD)", "2 hr"],
                           ["2", "2.2", "Light fired heater (low fire)", "1 hr"],
                           ["3", "3.1", "Increase to 50% rate, establish pumparounds", "3 hr"],
                           ["3", "3.2", "Establish reflux and product draw", "2 hr"],
                           ["4", "4.1", "Increase to 80% rate, optimize temperatures", "2 hr"],
                           ["4", "4.2", "Sample all products, adjust specs", "1 hr"],
                           ["5", "5.1", "Ramp to 100% rate, stabilize", "3 hr"],
                           ["5", "5.2", "Final tuning, document all parameters", "1 hr"]])
        self._h(d, "4. Key Operating Parameters", 1)
        self._t(d, 6, 3, [["Parameter", "Start-up Target", "Normal Target"], ["Feed Rate", "24,000 - 60,000 BPD", "100,000 - 120,000 BPD"],
                           ["Column Temp", "150C - 250C", "345C (bottom)"],
                           ["Column Pressure", "0.5 - 1.0 kg/cm2g", "1.2 kg/cm2g"],
                           ["Reflux Ratio", "3:1 - 5:1", "1.8:1"],
                           ["Naphtha Cut", "10% - 15%", "15.5%"]])
        self._h(d, "5. Emergency Procedures", 1)
        self._n(d, ["If any SIS trips: Follow ESD procedure (SOP-ESD-001)",
                     "If column flooding detected: Reduce feed rate immediately",
                     "If heater tube leak: Emergency shutdown, isolate fuel gas",
                     "Control Room to maintain constant communication with field"])
        self._h(d, "6. Approval", 1)
        self._t(d, 4, 3, [["Role", "Name", "Signature"], ["Prepared by", "Shri S.V. Patil", "________________"],
                           ["Reviewed by", "Shri A.D. Mistry", "________________"],
                           ["Approved by", "Shri Plant Manager", "________________"]])
        d.save(fp)
        return fp

    def generate_turnaround_plan_docx(self, filename=None):
        fn = filename or f"Turnaround_Plan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Turnaround Master Plan - CDU-III", 0)
        self._m(d, {"Plan No": "MRPL/TA/2026/CDU-III-001", "Unit": "CDU-III",
                     "Planned Shutdown": "01-Mar-2026 to 15-Mar-2026 (15 days)",
                     "Contractor": "L&T Hydrocarbon", "Budget": "Rs 28.5 Crore", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Scope of Work", 1)
        self._n(d, ["Column C-301 internal inspection & tray repair ( trays 1-25)",
                     "Heat exchanger bundle pull, cleaning & re-tube (6 units)",
                     "PSV overhaul and recertification (8 units)",
                     "Piping and valve inspection (142 critical joints)",
                     "Instrument calibration and loop check (142 loops)",
                     "Corrosion coupon replacement (12 locations)",
                     "Turnaround of 3 pumps (P-4401, P-4402, P-4403)"])
        self._h(d, "2. Schedule", 1)
        self._t(d, 8, 3, [["Phase", "Duration", "Key Activities"], ["Pre-Shutdown", "2 weeks prior", "Material staging, scaffolding, isolation"],
                           ["Shutdown", "Day 1-2", "Depressure, drain, purge, gas-free"],
                           ["Inspection", "Day 3-5", "NDE, thickness measurement, visual"],
                           ["Repair", "Day 5-11", "Welding, tray replacement, bundle work"],
                           ["Reassembly", "Day 12-13", "Valve reinstallation, flange up"],
                           ["Pre-Comm", "Day 14", "Loop checks, hydro test, function test"],
                           ["Restart", "Day 15", "PSSR, start-up, stabilize"]])
        self._h(d, "3. Resource Plan", 1)
        self._t(d, 6, 3, [["Trade", "Peak Manpower", "Total Man-Hours"], ["Fitters", "24", "2,880"],
                           ["Welders (Certified)", "16", "1,920"],
                           ["Scaffolders", "12", "1,440"],
                           ["Instrument Technicians", "8", "960"],
                           ["Supervisors", "6", "720"]])
        self._h(d, "4. HSE Plan", 1)
        self._n(d, ["Daily safety briefing at 06:30", "BBS observations target: 100/shift",
                     "Hot work permit for all welding", "Confined space entry with continuous gas monitoring",
                     "Emergency response: Ambulance on standby, fire team alert",
                     "Target: Zero LTI, Zero Environmental Incident"])
        self._h(d, "5. Cost Summary", 1)
        self._t(d, 6, 2, [["Category", "Amount (Rs Cr)"], ["Contractor (L&T)", "22.50"],
                           ["Material & Spares", "3.20"], ["Owner Costs", "1.80"],
                           ["Contingency (10%)", "2.50"], ["Total", "28.50"]])
        d.save(fp)
        return fp

    def generate_design_basis_docx(self, filename=None):
        fn = filename or f"Design_Basis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        fp = os.path.join(OUTPUT_DIR, "docx", fn)
        d = Document()
        self._h(d, "Design Basis Memorandum - CDU-III Revamp", 0)
        self._m(d, {"Document No": "MRPL/DBM/2026/CDU-III-001", "Project": "CDU-III Throughput Increase (6 to 7.5 MTPA)",
                     "Client": "MRPL", "Date": datetime.now().strftime("%d-%b-%Y")})
        self._h(d, "1. Project Objectives", 1)
        self._n(d, ["Increase CDU-III throughput from 6 MTPA (120,000 BPD) to 7.5 MTPA (150,000 BPD)",
                     "Maintain product slate and quality specifications",
                     "Minimal modifications to existing equipment (reuse where possible)",
                     "Compliance with OISD and Petroleum Safety Guidelines"])
        self._h(d, "2. Design Basis - Process", 1)
        self._t(d, 8, 3, [["Parameter", "Original Design", "Revamp Design"], ["Throughput", "120,000 BPD", "150,000 BPD"],
                           ["Crude Blend", "Bombay High dominant", "60% BHC + 25% KG + 15% RJ"],
                           ["Column Diameter", "8.4 m", "9.2 m (new)"],
                           ["Fired Heater Duty", "52 MW", "65 MW (replacement)"],
                           ["Pumparound Duty", "65 MW", "82 MW"],
                           ["Overhead Condenser", "32 MW", "42 MW"],
                           ["Desalter Capacity", "120,000 BPD", "150,000 BPD"]])
        self._h(d, "3. Design Basis - Mechanical", 1)
        self._t(d, 6, 3, [["Equipment", "Code", "Material"], ["Column", "ASME VIII Div 1", "CS SA-516 Gr 70"],
                           ["Exchangers", "TEMA R / ASME", "SS316L / CS"],
                           ["Piping", "ASME B31.3", "CS / SS316"],
                           ["Pumps", "API 610", "CS / Alloy"],
                           ["Vessels", "ASME VIII Div 1", "CS SA-516 Gr 70"]])
        self._h(d, "4. Design Basis - Utilities", 1)
        self._t(d, 5, 3, [["Utility", "Existing Capacity", "Revamp Requirement"], ["Cooling Water", "8,000 m3/hr", "10,500 m3/hr"],
                           ["Steam (6 kg/cm2g)", "45 TPH", "58 TPH"],
                           ["Nitrogen", "500 Nm3/hr", "500 Nm3/hr (no change)"],
                           ["Electrical", "12 MW", "15.5 MW"]])
        self._h(d, "5. Estimated Cost", 1)
        self._t(d, 5, 2, [["Category", "Amount (Rs Cr)"], ["EPC Contract", "320.00"],
                           ["Owner Costs", "18.50"], ["Contingency", "32.00"],
                           ["Total Project Cost", "370.50"]])
        d.save(fp)
        return fp

    def generate_operations_review_pptx(self, filename=None):
        fn = filename or f"Operations_Review_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        fp = os.path.join(OUTPUT_DIR, "pptx", fn)
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        sl = prs.slides.add_slide(prs.slide_layouts[0])
        sl.shapes.title.text = "Operations Review - Unit Operations"
        sl.placeholders[1].text = f"Q1 2026 Performance Review\nMRPL Refinery Operations\n{datetime.now().strftime('%d %B %Y')}"
        sl2 = prs.slides.add_slide(prs.slide_layouts[1])
        sl2.shapes.title.text = "Agenda"
        tf = sl2.placeholders[1].text_frame
        for item in ["Production Performance", "Energy & Utilities", "Reliability & Maintenance", "HSE & Process Safety", "Turnaround Status", "Key Action Items"]:
            p = tf.add_paragraph(); p.text = item; p.level = 0
        sl3 = prs.slides.add_slide(prs.slide_layouts[5])
        sl3.shapes.title.text = "Production Performance"
        ts = sl3.shapes.add_table(7, 3, PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(4.5)).table
        for i, rd in enumerate([["Parameter", "Q1 2026 Actual", "Target"], ["Throughput (BPD)", "118,500", "120,000"],
                                 ["Energy Intensity (GCal/MT)", "8.4", "<8.5"], ["Yield - Naphtha (%)", "15.8", "15.5"],
                                 ["Yield - Diesel (%)", "22.1", "22.0"], ["Yield - Residue (%)", "42.5", "43.0"],
                                 ["Uptime (%)", "97.2", ">95%"]]):
            for j, ct in enumerate(rd): ts.cell(i, j).text = ct
        sl4 = prs.slides.add_slide(prs.slide_layouts[5])
        sl4.shapes.title.text = "HSE Performance"
        ts2 = sl4.shapes.add_table(6, 3, PptxInches(1.5), PptxInches(2), PptxInches(10), PptxInches(4)).table
        for i, rd in enumerate([["Metric", "Q1 2026", "Target"], ["LTI", "1", "0"],
                                 ["Recordable Incidents", "2", "<3"], ["Near-Miss Reports", "8", ">5"],
                                 ["Environmental Incidents", "0", "0"], ["Safe Man-Hours", "142,000", "150,000"]]):
            for j, ct in enumerate(rd): ts2.cell(i, j).text = ct
        sl5 = prs.slides.add_slide(prs.slide_layouts[1])
        sl5.shapes.title.text = "Action Items"
        tf2 = sl5.placeholders[1].text_frame
        for item in ["CDU-III turnaround: 15-Mar start, on schedule", "Fired heater upgrade RFQ to be issued by Apr-10",
                      "Desalter chemistry optimization trial - Apr 2026", "PSV recertification: 8 units due Jun-2026",
                      "Target: 150,000 BPD by Q3 2026"]:
            p = tf2.add_paragraph(); p.text = item; p.level = 0
        prs.save(fp)
        return fp
