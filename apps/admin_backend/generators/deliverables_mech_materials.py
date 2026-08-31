from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches as PptxInches
import os
import time

OUTPUT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "data", "outputs"))
os.makedirs(OUTPUT_DIR, exist_ok=True)


class ExtendedGenerators_MechMaterials:

    def _add_meta_table(self, doc, data):
        meta_table = doc.add_table(rows=len(data), cols=2, style="Table Grid")
        for i, (k, v) in enumerate(data):
            meta_table.rows[i].cells[0].text = k
            meta_table.rows[i].cells[1].text = v
        return meta_table

    def _add_table(self, doc, headers, rows):
        table = doc.add_table(rows=len(rows) + 1, cols=len(headers), style="Table Grid")
        for i, h in enumerate(headers):
            table.rows[0].cells[i].text = h
        for r, row_data in enumerate(rows):
            for c, val in enumerate(row_data):
                table.rows[r + 1].cells[c].text = val
        return table

    def _add_sign_table(self, doc, roles):
        sign_table = doc.add_table(rows=len(roles) + 1, cols=3, style="Table Grid")
        for i, h in enumerate(["Role", "Name", "Signature / Date"]):
            sign_table.rows[0].cells[i].text = h
        for r, role in enumerate(roles):
            sign_table.rows[r + 1].cells[0].text = role
            sign_table.rows[r + 1].cells[1].text = ""
            sign_table.rows[r + 1].cells[2].text = ""
        return sign_table

    def _init_doc(self):
        doc = Document()
        style = doc.styles["Normal"]
        style.font.name = "Calibri"
        style.font.size = Pt(11)
        return doc

    # ────────────────────────────────────────────────
    # MECHANICAL & INSPECTION METHODS
    # ────────────────────────────────────────────────

    def generate_corrosion_control_docx(self, filename="Corrosion_Control_Document.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Corrosion Control Document (CCD)", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("Document No.", "CCD/MECH/2026-XX"), ("Plant / Area", "Refinery Unit / Sector"),
            ("Prepared By", "Corrosion Engineer"), ("Revision", "Rev. 00"),
        ])
        doc.add_heading("1. Purpose & Scope", level=1)
        doc.add_paragraph("This document defines the corrosion management strategy for static and rotating equipment in compliance with OISD-STD-149, OISD-STD-166, and NACE standards.")
        doc.add_heading("2. Corrosion Mechanisms Identified", level=1)
        self._add_table(doc, ["S.No.", "Corrosion Mechanism", "Affected Equipment", "Severity"], [
            ("1", "High Temperature Sulphidation", "Crude Column overhead, Reboiler", "High"),
            ("2", "Naphthenic Acid Corrosion", "TAN > 1.0 streams, Furnace tubes", "High"),
            ("3", "Ammonium Chloride Corrosion", "Crude Column top, Overhead system", "Medium"),
            ("4", "Wet H2S / HCN Corrosion", "LCFS, Amine system, Sour water", "High"),
            ("5", "Hydrogen Embrittlement", "High pressure reactors, Bolting", "Medium"),
        ])
        doc.add_heading("3. Material Selection Criteria", level=1)
        doc.add_paragraph("Material selection based on Corrosion Rate (mpy), Expected Life, Cost optimization, and compatibility with process fluids. Ref: NACE MR0175, ASTM.")
        doc.add_heading("4. Corrosion Monitoring Plan", level=1)
        self._add_table(doc, ["Monitoring Method", "Location", "Frequency"], [
            ("UT Thickness Survey", "All critical piping and vessels", "Annual"),
            ("Corrosion Coupons", "Cooling water, Crude overhead", "Quarterly"),
            ("Online Corrosion Probes", "High TAN streams, Amine circuits", "Continuous"),
        ])
        doc.add_heading("5. Inspection Intervals", level=1)
        doc.add_paragraph("Based on remaining life: >10yr=5yr interval; 5-10yr=3yr; <5yr=1yr.")
        doc.add_heading("6. Remedial Measures", level=1)
        doc.add_paragraph("Chemical inhibition, Internal coatings, Linings, Cathodic protection, Material upgrade, Operational changes, Cladding / weld overlay.")
        doc.add_heading("7. Approval", level=1)
        self._add_sign_table(doc, ["Corrosion Engineer", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_vibration_report_docx(self, filename="Vibration_Analysis_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Vibration Analysis Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Standard: ISO 10816-3 / ISO 20816-3 (Industrial Machines)")
        self._add_meta_table(doc, [
            ("Report No.", "VIB/MECH/2026-XX"), ("Equipment", "Centrifugal Pump / Compressor"),
            ("Tag No.", "P-XXXX / C-XXXX"), ("Date of Test", "DD/MM/YYYY"),
            ("Analyst", "Vibration Analyst (ISO 18436-2 Cat-II)"),
        ])
        doc.add_heading("1. Equipment Details", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Type", "Centrifugal Compressor"), ("Make / Model", "Manufacturer / Model"),
            ("Rated Speed", "XXXX RPM"), ("Bearing Type", "Tilting Pad Journal / Thrust"),
            ("Coupling Type", "Flexible Diaphragm"),
        ])
        doc.add_heading("2. Vibration Data Summary", level=1)
        self._add_table(doc, ["Location", "Overall (mm/s RMS)", "1X (mm/s)", "2X (mm/s)", "Class (ISO 10816)"], [
            ("DE Horizontal", "2.1", "1.5", "0.4", "A - Good"),
            ("DE Vertical", "2.3", "1.6", "0.3", "A - Good"),
            ("NDE Horizontal", "4.8", "3.2", "1.0", "B - Acceptable"),
            ("NDE Vertical", "5.1", "3.5", "1.1", "B - Acceptable"),
        ])
        doc.add_heading("3. ISO 10816-3 Classification", level=1)
        self._add_table(doc, ["Zone", "Velocity Range (mm/s RMS)", "Meaning"], [
            ("A", "0 - 2.3", "Good - Newly commissioned"),
            ("B", "2.3 - 4.5", "Acceptable - Unrestricted long-term use"),
            ("C", "4.5 - 7.1", "Restricted - Not suitable for long-term continuous"),
            ("D", "> 7.1", "Danger - Damage may occur"),
        ])
        doc.add_heading("4. Spectrum Analysis Findings", level=1)
        doc.add_paragraph("Dominant 1X peak at NDE bearing indicating residual unbalance. 2X within limits. No sub-synchronous or bearing defect frequencies detected.")
        doc.add_heading("5. Recommendations", level=1)
        self._add_table(doc, ["S.No.", "Observation", "Action"], [
            ("1", "NDE vibration trending upward", "Plan field balancing at next shutdown"),
            ("2", "Oil temperature marginally high", "Check oil cooler, clean if needed"),
            ("3", "No bearing defects", "Continue monitoring quarterly"),
        ])
        doc.add_heading("6. Signature", level=1)
        self._add_sign_table(doc, ["Vibration Analyst", "Maintenance Incharge"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_pm_schedule_docx(self, filename="Preventive_Maintenance_Schedule.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Preventive Maintenance Schedule", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("Document No.", "PM/PLAN/2026-XX"), ("Plant / Unit", "Refinery Unit"),
            ("Prepared By", "Maintenance Planning Engineer"), ("Effective Period", "Jan 2026 - Dec 2026"),
        ])
        doc.add_heading("1. Objective", level=1)
        doc.add_paragraph("To establish a systematic preventive maintenance schedule for all critical static, rotating, and electrical equipment.")
        doc.add_heading("2. PM Schedule - Rotating Equipment", level=1)
        self._add_table(doc, ["Tag No.", "Equipment", "PM Task", "Interval", "Due Date"], [
            ("P-1001", "Crude Feed Pump", "Bearing greasing, Alignment check", "3 Months", "01/04/2026"),
            ("C-2001", "Main Air Compressor", "Oil change, Filter replacement", "6 Months", "01/06/2026"),
            ("T-3001", "Cooling Tower Fan", "Belt replacement, Lubrication", "12 Months", "01/01/2026"),
            ("F-4001", "Induced Draft Fan", "Balancing, Bearing replacement", "12 Months", "01/07/2026"),
            ("M-5001", "Mixer Motor", "Insulation resistance test", "6 Months", "01/03/2026"),
        ])
        doc.add_heading("3. PM Schedule - Static Equipment", level=1)
        self._add_table(doc, ["Tag No.", "Equipment", "PM Task", "Interval", "Due Date"], [
            ("V-1001", "Crude Column", "Internal inspection, Thickness survey", "5 Years", "01/01/2030"),
            ("HX-2001", "Crude Preheat Exchanger", "Tube inspection, Hydro test", "3 Years", "01/01/2029"),
            ("TK-3001", "FRTK-01", "Shell settlement, Deck inspection", "5 Years", "01/01/2030"),
            ("B-4001", "Process Heater", "Tube inspection, Firebrick check", "2 Years", "01/01/2028"),
        ])
        doc.add_heading("4. Escalation Matrix", level=1)
        self._add_table(doc, ["Overdue Days", "Escalation To", "Action Required"], [
            ("7 Days", "Maintenance Supervisor", "Expedite and reschedule"),
            ("15 Days", "Maintenance Manager", "Written explanation + revised plan"),
            ("30 Days", "Plant Manager", "Report with root cause and CAPA"),
        ])
        doc.add_heading("5. Approval", level=1)
        self._add_sign_table(doc, ["Maintenance Planning Engineer", "Maintenance Manager", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_equipment_history_docx(self, filename="Equipment_History_Sheet.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Equipment History Sheet", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("Document No.", "EH/MECH/2026-XX"), ("Equipment", "Centrifugal Compressor"),
            ("Tag No.", "C-1001"), ("Location", "Compressor House, Unit 1"), ("Date Compiled", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Equipment Identification", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Manufacturer", "OEM Name"), ("Model / Type", "Centrifugal, Multi-stage"),
            ("Serial No.", "MFG-XXXX-XXXX"), ("Year of Installation", "20XX"),
            ("Capacity / Rating", "XXXX TPH / XXXX kW"), ("Operating Pressure", "XX kg/cm2"),
            ("Operating Temperature", "XX deg C"), ("MOC", "SS 316 / CS / Alloy"),
        ])
        doc.add_heading("2. Maintenance History", level=1)
        self._add_table(doc, ["Date", "Work Done", "Parts Replaced", "Cost (Rs)", "Done By"], [
            ("01/03/2024", "Bearing replacement (DE & NDE)", "SKF 6312-2RS", "45,000", "Maint. Team"),
            ("15/09/2024", "Coupling alignment & balancing", "Coupling element", "12,000", "External"),
            ("01/01/2025", "Major overhaul", "Gaskets, O-rings, Bearings", "2,80,000", "OEM"),
            ("15/06/2025", "Oil change & filter replacement", "Synthetic oil 68 cSt", "8,500", "Maint. Team"),
            ("01/01/2026", "Minor overhaul - seal replacement", "Mechanical seal set", "1,20,000", "OEM"),
        ])
        doc.add_heading("3. Breakdown History", level=1)
        self._add_table(doc, ["Date", "Failure Mode", "Root Cause", "Downtime", "Impact"], [
            ("10/04/2024", "Bearing failure", "Contaminated oil", "36 hrs", "Unit trip"),
            ("22/11/2024", "Seal leak", "Mechanical seal wear", "12 hrs", "Reduced throughput"),
            ("05/08/2025", "Coupling bolt shearing", "Fatigue / Misalignment", "24 hrs", "Unit trip"),
        ])
        doc.add_heading("4. Inspection Records", level=1)
        self._add_table(doc, ["Date", "Type", "Findings", "Action Taken"], [
            ("01/03/2025", "UT Thickness", "All readings within limits", "None"),
            ("15/09/2025", "Vibration survey", "NDE bearing marginally elevated", "Planned balancing"),
            ("01/01/2026", "Borescope inspection", "Impeller leading edge erosion", "Planned replacement"),
        ])
        doc.add_heading("5. Spare Parts Inventory", level=1)
        self._add_table(doc, ["Part", "Part No.", "Qty in Stock", "Min Required"], [
            ("Mechanical Seal Set", "MS-XXXX", "1", "1"),
            ("Bearing Assembly (DE)", "SKF 6312-2RS", "2", "1"),
            ("Coupling Element", "CU-XXXX", "1", "1"),
            ("Gasket Kit", "GK-XXXX", "3", "2"),
        ])
        doc.add_heading("6. Summary", level=1)
        doc.add_paragraph("Equipment running XXXXX hours since last major overhaul. Bearing replacement and seal changes are primary cost drivers.")
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_hx_inspection_report_docx(self, filename="HX_Inspection_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Heat Exchanger Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference Standards: API 570, API 510, TEMA Standards")
        self._add_meta_table(doc, [
            ("Report No.", "HX/INSP/2026-XX"), ("Equipment", "Shell & Tube Heat Exchanger"),
            ("Tag No.", "HX-1001"), ("Service", "Crude / Preheat"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Design Data", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("TEMA Type", "BEM / AES / BEU"), ("Shell ID", "XX mm"),
            ("Tube OD / Thickness", "XX mm / XX mm"), ("Tube Length", "XXXX mm"),
            ("No. of Tubes", "XXX"), ("Tube MOC", "SS 316 / CS / Cu-Ni"),
            ("Shell MOC", "CS SA-516 Gr. 70"), ("Design Pressure (Shell/Tube)", "XX / XX kg/cm2"),
            ("Design Temperature (Shell/Tube)", "XX / XX deg C"), ("Baffles", "Single segmental, XX% cut"),
        ])
        doc.add_heading("2. Inspection Findings", level=1)
        self._add_table(doc, ["Area", "Finding", "Severity", "Action"], [
            ("Shell Interior", "General corrosion, CR = 0.12 mpy", "Low", "Monitor"),
            ("Tube Bundle", "3 tubes plugged (pitting)", "Medium", "Replace bundle"),
            ("Tubesheet", "Fouling in tube-side passages", "Medium", "Chemical cleaning"),
            ("Nozzle Flanges", "Gasket seating surface pitting", "Low", "Re-machine"),
            ("Baffles", "Erosion on baffle edges", "Low", "No immediate action"),
        ])
        doc.add_heading("3. Tube Thickness Survey", level=1)
        self._add_table(doc, ["Tube Row", "Measured (mm)", "Nominal (mm)", "Remaining Life"], [
            ("Row 1 (Inlet)", "1.82", "2.11", "> 10 years"),
            ("Row 2", "1.95", "2.11", "> 10 years"),
            ("Row 3", "2.01", "2.11", "> 10 years"),
            ("Row 10 (Mid)", "2.08", "2.11", "> 10 years"),
        ])
        doc.add_heading("4. Hydro Test Results", level=1)
        doc.add_paragraph("Shell Side: XX kg/cm2 for 30 min - No leakage. Tube Side: XX kg/cm2 for 30 min - No leakage.")
        doc.add_heading("5. Recommendations", level=1)
        doc.add_paragraph("1. Chemical cleaning of tube-side. 2. Plan bundle replacement. 3. Monitor shell-side CR annually. 4. Replace all gaskets.")
        doc.add_heading("6. Approval", level=1)
        self._add_sign_table(doc, ["Inspector (API 510/570)", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_tank_inspection_report_docx(self, filename="Tank_Inspection_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Storage Tank Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: API 653, OISD-STD-129, IS 803")
        self._add_meta_table(doc, [
            ("Report No.", "TK/INSP/2026-XX"), ("Tank", "FRTK-01 (Floating Roof Tank)"),
            ("Capacity", "50,000 KL"), ("Product", "Crude Oil / MS / HSD"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Tank Identification", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Tank Type", "External Floating Roof (EFR)"), ("Diameter", "XXXX mm"),
            ("Height", "XXXX mm"), ("Shell Thickness", "XX mm (Bottom) / XX mm (Top)"),
            ("Year Built", "20XX"), ("MOC", "CS SA-516 Gr. 70"),
            ("Roof Type", "Pontoon / Deck"), ("Seal Type", "Primary shoe + Secondary resilient seal"),
        ])
        doc.add_heading("2. Shell Inspection", level=1)
        self._add_table(doc, ["Course", "Nominal (mm)", "Measured (mm)", "CR (mpy)", "Remaining Life"], [
            ("1st (Bottom)", "12.00", "11.52", "0.12", "> 15 years"),
            ("2nd", "12.00", "11.65", "0.09", "> 15 years"),
            ("3rd", "12.00", "11.78", "0.06", "> 15 years"),
            ("4th (Top)", "8.00", "7.85", "0.04", "> 15 years"),
        ])
        doc.add_heading("3. Bottom Plate Inspection", level=1)
        doc.add_paragraph("MFL scan conducted. 3 areas with metal loss > 20%. Minimum remaining thickness: 4.2 mm (nominal 6.0 mm).")
        doc.add_heading("4. Floating Roof Inspection", level=1)
        self._add_table(doc, ["Component", "Finding", "Action"], [
            ("Pontoon compartments", "No water ingress detected", "None"),
            ("Deck plates", "Minor surface corrosion", "Coat at next opportunity"),
            ("Primary seal", "Wear within limits", "Replace in 12 months"),
            ("Secondary seal", "Intact", "Monitor"),
            ("Drain system", "Functional, no blockage", "Clean annually"),
        ])
        doc.add_heading("5. Settlement Survey", level=1)
        doc.add_paragraph("Max differential settlement: 25 mm (allowable: 50 mm per API 653). Shell within verticality tolerance.")
        doc.add_heading("6. Recommendations", level=1)
        doc.add_paragraph("1. Bottom plate repair for 3 areas. 2. Replace primary seal. 3. Schedule internal cleaning. 4. Annual inspection per OISD-STD-129.")
        doc.add_heading("7. Approval", level=1)
        self._add_sign_table(doc, ["API 653 Inspector", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_boiler_inspection_report_docx(self, filename="Boiler_Inspection_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Boiler Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-131, IBR, IS 12238, API 573")
        self._add_meta_table(doc, [
            ("Report No.", "BOILER/INSP/2026-XX"), ("Boiler", "Process Heater / Steam Boiler"),
            ("Capacity", "XX TPH @ XX kg/cm2"), ("Fuel", "LNG / FO / AG / Combination"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Boiler Data", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Type", "Water Tube / Fire Tube"), ("Manufacturer", "OEM Name"),
            ("Year of Installation", "20XX"), ("Design Pressure", "XX kg/cm2"),
            ("Operating Pressure", "XX kg/cm2"), ("Steam Temperature", "XX deg C"),
            ("MOC (Drum)", "SA-516 Gr. 70"), ("MOC (Tubes)", "SA-213 T11 / T22"),
            ("Last Inspection Date", "DD/MM/YYYY"),
        ])
        doc.add_heading("2. External Inspection", level=1)
        self._add_table(doc, ["Area", "Observation", "Status"], [
            ("Drum shell & heads", "No visible distortion, no leaks", "OK"),
            ("Headers & downcomers", "Minor surface corrosion", "Monitor"),
            ("Burner zone tubes", "Scale buildup, thickness within limits", "Monitor"),
            ("Safety valves", "Tested, set pressure confirmed", "OK"),
        ])
        doc.add_heading("3. Internal Inspection", level=1)
        self._add_table(doc, ["Component", "Finding", "Action"], [
            ("Drum internals", "Mdl collectors clean, no damage", "None"),
            ("Tube bundle (radiant)", "Mild oxidation, no bulging", "Monitor CR"),
            ("Tube bundle (convective)", "Soot/blowdown fouling", "Chemical cleaning"),
            ("Water treatment", "Dosing rates within limits", "None"),
        ])
        doc.add_heading("4. Tube Thickness Survey", level=1)
        self._add_table(doc, ["Location", "Measured (mm)", "Nominal (mm)", "CR (mpy)"], [
            ("Radiant (bottom)", "4.85", "5.00", "0.04"),
            ("Radiant (top)", "4.72", "5.00", "0.07"),
            ("Convective", "3.65", "3.90", "0.06"),
            ("Return bend", "4.50", "5.00", "0.12"),
        ])
        doc.add_heading("5. Combustion & Efficiency", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Excess Air", "XX% (Target: 3-5%)"), ("Flue Gas Temperature", "XX deg C"),
            ("CO2", "XX%"), ("Stack Loss", "XX%"), ("Boiler Efficiency", "XX%"),
        ])
        doc.add_heading("6. Recommendations", level=1)
        doc.add_paragraph("1. Chemical cleaning of convective bank. 2. Monitor return bend CR. 3. Optimize excess air. 4. Eddy current testing. 5. Safety valve per IBR.")
        doc.add_heading("7. Approval", level=1)
        self._add_sign_table(doc, ["Boiler Inspector (IBR)", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_fired_heater_inspection_docx(self, filename="Fired_Heater_Inspection_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Fired Heater Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-133, API 573, API 556, API 530")
        self._add_meta_table(doc, [
            ("Report No.", "FH/INSP/2026-XX"), ("Heater", "Process Fired Heater"),
            ("Tag No.", "H-1001"), ("Capacity", "XX Gcal/hr"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Heater Data", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Type", "Cabin / Vertical / Horizontal"), ("Firebox Duty", "XX Gcal/hr"),
            ("Tube OD / Thickness", "XX mm / XX mm"), ("Tube MOC", "Incoloy 800HT / HP Mod."),
            ("Number of Passes", "X"), ("Design Film Temperature", "XX deg C"),
            ("TMT", "XX deg C"), ("Fuel", "LNG + FO (Dual fired)"),
            ("Burner Type", "Ultra-low NOx / Regenerative"), ("Radiant Coil Length", "XXXX mm"),
        ])
        doc.add_heading("2. External Inspection", level=1)
        self._add_table(doc, ["Area", "Observation", "Status"], [
            ("Tube supports", "No distortion, proper clearance", "OK"),
            ("Refractory (Casing)", "Minor hot spot on west wall (180 deg C)", "Monitor"),
            ("Burner tips", "Nozzle erosion on B3, B5", "Replace at shutdown"),
            ("Stack & damper", "Functional, no binding", "OK"),
            ("Convection section", "Slight fouling of fins", "Clean at shutdown"),
        ])
        doc.add_heading("3. Tube Inspection", level=1)
        self._add_table(doc, ["Location", "TMT (deg C)", "Measured (mm)", "Nominal (mm)", "CR (mpy)"], [
            ("Pass 1 (Bottom)", "850", "5.12", "6.00", "0.18"),
            ("Pass 2 (Mid)", "870", "4.95", "6.00", "0.22"),
            ("Pass 3 (Top)", "890", "4.78", "6.00", "0.26"),
            ("Return bend", "910", "4.60", "6.00", "0.30"),
        ])
        doc.add_heading("4. Refractory Inspection", level=1)
        doc.add_paragraph("Castable refractory: Minor cracking on west wall near B3. Hot spot 180 deg C (limit: 200 deg C).")
        doc.add_heading("5. Combustion Analysis", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("O2 in Flue Gas", "XX% (Target: 1.5-3%)"), ("CO", "< 50 ppm"),
            ("NOx", "XX ppm"), ("Flue Gas Temp", "XX deg C"), ("Thermal Efficiency", "XX%"),
        ])
        doc.add_heading("6. Recommendations", level=1)
        doc.add_paragraph("1. Repair refractory. 2. Replace burner tips. 3. Eddy current on return bends. 4. Clean convection fins. 5. Increase TMT monitoring.")
        doc.add_heading("7. Approval", level=1)
        self._add_sign_table(doc, ["Heater Inspector", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_psv_test_certificate_docx(self, filename="PSV_Test_Certificate.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Pressure Safety Valve (PSV) Test Certificate", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-132, API 576, ASME Section VIII")
        self._add_meta_table(doc, [
            ("Certificate No.", "PSV/TEST/2026-XX"), ("Valve Tag No.", "PSV-XXXX"),
            ("Equipment Protected", "Vessel / Heat Exchanger / Column"), ("Date of Test", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Valve Identification", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Manufacturer", "Crosby / Consolidated / Pentair"), ("Model / Size", "XXXX / X inch"),
            ("Valve Type", "Conventional / Balanced Bellows"), ("Inlet Connection", "XX inch Flanged"),
            ("Outlet Connection", "XX inch Flanged"), ("Set Pressure", "XX.X kg/cm2"),
            ("Blowdown", "XX% of set pressure"), ("Orifice Size", "XX mm"),
            ("Seat Material", "Stellite / Monel / SS 316"), ("Body Material", "CS / SS 316 / Alloy"),
        ])
        doc.add_heading("2. Test Results", level=1)
        self._add_table(doc, ["Parameter", "Set Value", "Actual Value", "Pass/Fail"], [
            ("Set Pressure", "XX.X kg/cm2", "XX.X kg/cm2", "PASS"),
            ("Blowdown", "XX%", "XX%", "PASS"),
            ("Leak Test (Seat)", "No visible leakage", "No leakage at 90% set", "PASS"),
            ("Re-seat Pressure", "XX.X kg/cm2", "XX.X kg/cm2", "PASS"),
            ("Lift (Visual)", "Smooth, no chatter", "Smooth, no chatter", "PASS"),
            ("Back Pressure Effect", "Within design limits", "No effect observed", "PASS"),
        ])
        doc.add_heading("3. Seat Leakage Test", level=1)
        doc.add_paragraph("Seat leakage test per API 527. Leakage rate: < 0.01 cc/min (PASS criteria: < 0.01 cc/min).")
        doc.add_heading("4. Visual Inspection", level=1)
        self._add_table(doc, ["Component", "Observation"], [
            ("Body & Bonnet", "No cracks, no corrosion, no deformation"),
            ("Spring", "No corrosion, correct color code"),
            ("Seat & Disc", "Minor lapping marks - acceptable"),
            ("Stem & Guide", "Smooth operation, no binding"),
        ])
        doc.add_heading("5. Disposition", level=1)
        doc.add_paragraph("Valve tested satisfactory and released for service. Next test: DD/MM/YYYY (12 months).")
        doc.add_heading("6. Certifier Details", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Tested By", "API 576 Certified Technician"),
            ("Certificate No.", "API 576-XXXX"),
            ("Workshop", "Authorized PSV Service Shop"),
        ])
        doc.add_heading("7. Signature", level=1)
        self._add_sign_table(doc, ["PSV Technician", "Maintenance Incharge"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_electrical_equipment_inspection_docx(self, filename="Electrical_Equipment_Inspection.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Electrical Equipment Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-137, IS 3427, IS 1681, IEEE Standards")
        self._add_meta_table(doc, [
            ("Report No.", "ELEC/INSP/2026-XX"), ("Equipment", "Transformer / Switchgear / Motor / MCC"),
            ("Tag No.", "TR-XXXX / SWG-XXXX"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Equipment Details", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Type", "Power Transformer / VCB / LT Motor"), ("Rating", "XX MVA / XX Amp / XX kW"),
            ("Voltage", "XX kV / XX V"), ("Manufacturer", "ABB / Siemens / BHEL"),
            ("Year of Manufacture", "20XX"), ("Location", "Switchyard / MCC Room"),
            ("Last Maintenance", "DD/MM/YYYY"), ("Protection Class", "IP55 / IP65"),
        ])
        doc.add_heading("2. Inspection Checklist", level=1)
        self._add_table(doc, ["S.No.", "Item", "Observation", "Status"], [
            ("1", "Visual condition", "No damage, no corrosion", "OK"),
            ("2", "Terminal connections", "All tight, no hotspots", "OK"),
            ("3", "IR test", "XX MOhm", "PASS"),
            ("4", "Polarization Index", "XX (>2.0 required)", "PASS"),
            ("5", "Winding resistance", "XX mOhm", "OK"),
            ("6", "Oil test (transformer)", "BDV: XX kV, Moisture: XX ppm", "OK"),
            ("7", "Bearing insulation", "XX MOhm", "OK"),
            ("8", "Earth connections", "All leads intact", "OK"),
            ("9", "Alarm/trip settings", "Confirmed per relay settings", "OK"),
            ("10", "Space heater", "Functional", "OK"),
        ])
        doc.add_heading("3. Insulation Resistance Test Summary", level=1)
        self._add_table(doc, ["Winding", "IR (MOhm) @ 500V", "PI", "Acceptance"], [
            ("R-Y", "500+", "2.8", "PASS"), ("Y-B", "500+", "2.7", "PASS"),
            ("B-R", "500+", "2.9", "PASS"), ("R-Y-B-E", "500+", "2.8", "PASS"),
        ])
        doc.add_heading("4. Thermographic Survey", level=1)
        doc.add_paragraph("IR thermography: No abnormal temperature rise. Max differential at joints: 8 deg C (limit: <15 deg C).")
        doc.add_heading("5. Recommendations", level=1)
        doc.add_paragraph("1. Schedule winding resistance test. 2. Oil reclamation. 3. Re-torque busbar joints. 4. Verify earth grid resistance.")
        doc.add_heading("6. Approval", level=1)
        self._add_sign_table(doc, ["Electrical Engineer", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_insulation_inspection_docx(self, filename="Insulation_Inspection.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Insulation Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-177, ASTM C585, ASTM C795, MSS-SP-58")
        self._add_meta_table(doc, [
            ("Report No.", "INS/INSP/2026-XX"), ("Area / Equipment", "Unit / Pipeline / Vessel"),
            ("Insulation Type", "Mineral Wool / Calcium Silicate / Aerogel"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Inspection Findings", level=1)
        self._add_table(doc, ["Location", "Insulation Type", "Condition", "Surface Temp (deg C)", "Action"], [
            ("Hot oil header (12 inch)", "Mineral Wool 80mm", "Good - no moisture", "42", "None"),
            ("Steam line (6 inch)", "Calcium Silicate 50mm", "Minor jacket damage", "38", "Repair jacket"),
            ("Crude column (bottom)", "Mineral Wool 100mm", "Wet insulation", "65", "Replace insulation"),
            ("Reboiler nozzles", "Aerogel blanket 40mm", "Good", "55", "None"),
            ("FG header", "Mineral Wool 50mm", "Damage on support", "68", "Repair cladding"),
            ("HX shell", "Mineral Wool 80mm", "Oil staining", "58", "Investigate leak"),
        ])
        doc.add_heading("2. Cladding Inspection", level=1)
        self._add_table(doc, ["Location", "Condition", "Action"], [
            ("All hot piping", "SS 304 cladding - good", "None"),
            ("Crude column", "Aluminum cladding - minor denting", "Straighten"),
            ("FG header", "SS 304 - paint peeling", "Repaint at shutdown"),
            ("Steam line", "SS 304 - joint gap > 3mm", "Re-seal joints"),
        ])
        doc.add_heading("3. Recommendations", level=1)
        doc.add_paragraph("1. Replace wet insulation. 2. Repair damaged cladding. 3. CUI inspection at next shutdown. 4. Improve jacketing QC. 5. Thermal imaging quarterly.")
        doc.add_heading("4. Approval", level=1)
        self._add_sign_table(doc, ["Insulation Engineer", "Maintenance Head"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_painting_inspection_docx(self, filename="Painting_Coating_Inspection.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Painting & Coating Inspection Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-177, SSPC-SP, ISO 12944, IS 1080")
        self._add_meta_table(doc, [
            ("Report No.", "PAINT/INSP/2026-XX"), ("Area / Equipment", "Structural Steel / Pipeline / Tank"),
            ("Coating System", "3-Coat Epoxy + Polyurethane / Zinc Rich"), ("Date of Inspection", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Surface Preparation", level=1)
        self._add_table(doc, ["Parameter", "Specification", "Actual"], [
            ("Surface Profile", "SSPC-SP10 / Sa 2.5", "Achieved Sa 2.5"),
            ("DFT", "50-75 microns", "62 microns"),
            ("Dust Level", "ISO 8502-3: Rating 2", "Rating 1 (Good)"),
        ])
        doc.add_heading("2. Coating Application Inspection", level=1)
        self._add_table(doc, ["Coat", "Product", "NFT (um)", "DFT Achieved (um)", "Status"], [
            ("Primer", "Zinc Rich Epoxy", "75", "82", "PASS"),
            ("Mid Coat", "Epoxy MIO", "125", "135", "PASS"),
            ("Finish Coat", "Aliphatic PU", "50", "55", "PASS"),
            ("Total DFT", "-", "250", "272", "PASS"),
        ])
        doc.add_heading("3. Adhesion Test", level=1)
        self._add_table(doc, ["Location", "Pull-off (MPa)", "Failure Mode"], [
            ("Beam B-12", "5.2 MPa (Min: 3.5)", "Cohesive failure in substrate"),
            ("Pipe P-08", "4.8 MPa (Min: 3.5)", "Adhesive + Cohesive"),
        ])
        doc.add_heading("4. Recommendations", level=1)
        doc.add_paragraph("1. Touch-up painting. 2. Monitor adhesion on pipe supports. 3. Ensure ambient conditions within spec. 4. Re-inspect after 12 months.")
        doc.add_heading("5. Approval", level=1)
        self._add_sign_table(doc, ["Painting Inspector (BGAS/PCN Level 2)", "Maintenance Head"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_ndt_report_docx(self, filename="NDT_Report.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Non-Destructive Testing (NDT) Report", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Methods: UT, MT, RT, PT - Per ASME Section V, ASME Section VIII")
        self._add_meta_table(doc, [
            ("Report No.", "NDT/2026-XX"), ("Equipment", "Pressure Vessel / Piping / Weldment"),
            ("Tag No.", "V-XXXX / P-XXXX"), ("Weld / Area", "Longitudinal / Circumferential / Nozzle"),
            ("Date of Testing", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. NDT Personnel", level=1)
        self._add_table(doc, ["Technician", "Method", "Cert Level", "Certificate No."], [
            ("Technician A", "UT", "ASNT Level II", "ASNT-II-XXXX"),
            ("Technician B", "MT", "ASNT Level II", "ASNT-II-XXXX"),
            ("Technician C", "RT", "ASNT Level II", "ASNT-II-XXXX"),
        ])
        doc.add_heading("2. Ultrasonic Testing (UT)", level=1)
        self._add_table(doc, ["Point", "Nominal (mm)", "Measured (mm)", "Min Allowable (mm)", "Result"], [
            ("UT-01", "12.00", "11.82", "9.50", "PASS"),
            ("UT-02", "12.00", "11.65", "9.50", "PASS"),
            ("UT-03", "12.00", "10.10", "9.50", "PASS"),
            ("UT-04", "12.00", "9.65", "9.50", "PASS"),
            ("UT-05", "12.00", "11.45", "9.50", "PASS"),
        ])
        doc.add_heading("3. Magnetic Particle Testing (MT)", level=1)
        self._add_table(doc, ["Weld", "Method", "Indication", "Severity"], [
            ("W-01 (Long.)", "Wet fluorescent", "No indications", "-"),
            ("W-02 (Circ.)", "Wet fluorescent", "2 linear indications", "Repair & re-test"),
            ("N-01", "Wet fluorescent", "No indications", "-"),
        ])
        doc.add_heading("4. Radiographic Testing (RT)", level=1)
        self._add_table(doc, ["Weld", "Technique", "Sensitivity", "Acceptance"], [
            ("W-01", "X-ray, single wall", "Fe 10, 2%", "ACCEPT"),
            ("W-02", "X-ray, single wall", "Fe 10, 2%", "Reject - Porosity"),
            ("N-01", "Gamma (Ir-192)", "Fe 8, 1.5%", "ACCEPT"),
        ])
        doc.add_heading("5. Summary", level=1)
        doc.add_paragraph("UT: PASS. MT: W-02 repaired, re-tested PASS. RT: W-02 repaired, re-tested PASS. PT: All PASS. Equipment released for service.")
        doc.add_heading("6. Approval", level=1)
        self._add_sign_table(doc, ["NDT Level III", "QC Engineer", "Maintenance Head"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_idle_equipment_preservation_docx(self, filename="Idle_Equipment_Preservation.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Idle Equipment Preservation Plan", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Reference: OISD-STD-171, API 670")
        self._add_meta_table(doc, [
            ("Document No.", "PRES/PLAN/2026-XX"), ("Equipment", "Centrifugal Pump / Compressor / Motor"),
            ("Tag No.", "P-XXXX / C-XXXX"), ("Idle Since", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Equipment Inventory", level=1)
        self._add_table(doc, ["Tag No.", "Equipment", "Idle Date", "Class", "Next Action"], [
            ("P-1001", "Crude Feed Pump", "01/01/2025", "Class A", "Monthly rotation"),
            ("C-2001", "Spare Compressor", "15/06/2024", "Class B", "Quarterly inspection"),
            ("M-3001", "Spare Motor", "01/03/2024", "Class C", "Annual overhaul"),
            ("HX-4001", "Spare Heat Exchanger", "01/01/2023", "Class C", "Annual inspection"),
        ])
        doc.add_heading("2. Preservation Methods", level=1)
        doc.add_paragraph("Rotating: Manual rotation weekly, maintain oil level, flush seals, cover openings.\nStatic: Drain and dry, rust preventive, seal nozzles, nitrogen blanket.\nElectric: Space heaters on, IR test quarterly, rotate shaft.")
        doc.add_heading("3. Inspection Schedule", level=1)
        self._add_table(doc, ["Class", "Interval", "Activities", "Responsible"], [
            ("Class A (< 6mo)", "Monthly", "Visual, rotation, oil level", "Mech. Technician"),
            ("Class B (6-24mo)", "Quarterly", "Full check, IR, rotation", "Maintenance Engineer"),
            ("Class C (> 24mo)", "Annual", "Overhaul, NDT, calibration", "Reliability Engineer"),
        ])
        doc.add_heading("4. Return-to-Service Checklist", level=1)
        self._add_table(doc, ["S.No.", "Item", "Status"], [
            ("1", "Remove blind flanges", ""),
            ("2", "Flush / nitrogen purge", ""),
            ("3", "Check alignment", ""),
            ("4", "IR test", ""),
            ("5", "Bearing lubrication", ""),
            ("6", "Operational test (1 hr)", ""),
        ])
        doc.add_heading("5. Approval", level=1)
        self._add_sign_table(doc, ["Reliability Engineer", "Maintenance Head", "Plant Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_maintenance_review_pptx(self, filename="MI_Review_Presentation.pptx"):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide_layout = prs.slide_layouts[5]
        slide1 = prs.slides.add_slide(slide_layout)
        title_shape = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(2))
        tf = title_shape.text_frame
        tf.text = "Mechanical & Inspection Review"
        tf.paragraphs[0].font.size = PptxInches(1.2)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = 1
        sub = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(4), PptxInches(12), PptxInches(1))
        sub.text_frame.text = "Monthly Review - Jan 2026\nMaintenance Department"
        sub.text_frame.paragraphs[0].font.size = PptxInches(0.6)
        sub.text_frame.paragraphs[0].alignment = 1
        slide2 = prs.slides.add_slide(slide_layout)
        tb = slide2.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
        tb.text_frame.text = "Equipment Reliability KPIs"
        tb.text_frame.paragraphs[0].font.size = PptxInches(0.8)
        tb.text_frame.paragraphs[0].font.bold = True
        kpi = slide2.shapes.add_table(7, 4, PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(5)).table
        for i, h in enumerate(['KPI', 'Target', 'Actual', 'Status']):
            kpi.cell(0, i).text = h
        for r, d in enumerate([
            ('OEE', '> 90%', '87.5%', 'Below Target'), ('PM Compliance', '> 95%', '92%', 'Below Target'),
            ('Breakdown Hours', '< 48 hrs/mo', '36 hrs', 'On Track'), ('MTBF', '> 4000 hrs', '4250 hrs', 'On Track'),
            ('MTTR', '< 4 hrs', '3.2 hrs', 'On Track'), ('Inspection Compliance', '100%', '98%', 'On Track'),
        ]):
            for c, v in enumerate(d):
                kpi.cell(r+1, c).text = v
        slide3 = prs.slides.add_slide(slide_layout)
        tb3 = slide3.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
        tb3.text_frame.text = "Action Items & Decisions"
        tb3.text_frame.paragraphs[0].font.size = PptxInches(0.8)
        tb3.text_frame.paragraphs[0].font.bold = True
        ab = slide3.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        af = ab.text_frame
        for i, a in enumerate([
            "1. Expedite spare compressor bearing procurement",
            "2. Complete laser alignment on all pumps by end of Feb",
            "3. Install vibration monitoring on critical compressors",
            "4. Schedule API 570 inspection (April shutdown)",
            "5. Review PM frequency for cooling tower gearbox",
            "6. Update equipment criticality ranking (due March)",
            "7. Approve NDT contractor budget",
        ]):
            if i == 0:
                af.text = a
            else:
                p = af.add_paragraph()
                p.text = a
                p.space_after = PptxInches(0.2)
                p.font.size = PptxInches(0.45)
        path = os.path.join(OUTPUT_DIR, filename)
        prs.save(path)
        return path

    # ────────────────────────────────────────────────
    # MATERIALS & PROCUREMENT METHODS
    # ────────────────────────────────────────────────

    def generate_tender_evaluation_docx(self, filename="Tender_Evaluation.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Tender Evaluation / Comparative Statement", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("Tender No.", "TENDER/2026-XXX"), ("Description", "Supply of Centrifugal Pump Spares"),
            ("NIT Ref.", "NIT/2026-XXX"), ("Opening Date", "DD/MM/YYYY"), ("Evaluation Date", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Bidders Participated", level=1)
        self._add_table(doc, ["S.No.", "Bidder Name", "Bidder Code", "EMD Submitted"], [
            ("1", "M/s. PumpTech Industries", "BID-001", "Yes"),
            ("2", "M/s. Flow Solutions Pvt. Ltd.", "BID-002", "Yes"),
            ("3", "M/s. Rotating Equipment Co.", "BID-003", "Yes"),
            ("4", "M/s. Supreme Pumps Ltd.", "BID-004", "Yes"),
        ])
        doc.add_heading("2. Price Comparative Statement", level=1)
        self._add_table(doc, ["Item", "Qty", "BID-001 (Rs)", "BID-002 (Rs)", "BID-003 (Rs)"], [
            ("Impeller (SS 316)", "2", "1,80,000", "1,65,000", "1,72,000"),
            ("Mechanical Seal Set", "4", "2,40,000", "2,55,000", "2,32,000"),
            ("Bearing Assembly", "6", "96,000", "90,000", "88,000"),
            ("Coupling Element", "2", "36,000", "34,000", "38,000"),
            ("Total (Ex-GST)", "-", "5,52,000", "5,44,000", "5,30,000"),
        ])
        doc.add_heading("3. Recommendation", level=1)
        doc.add_paragraph("Award to BID-003 (L1 - Lowest Price) subject to management approval.")
        doc.add_heading("4. Evaluation Committee", level=1)
        self._add_sign_table(doc, ["Chairperson (Procurement)", "Member (Maintenance)", "Member (Finance)", "Member (Quality)"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_purchase_order_docx(self, filename="Purchase_Order.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Purchase Order (PO)", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("PO No.", "PO/2026-XXXX"), ("PO Date", "DD/MM/YYYY"), ("Vendor", "M/s. Supplier Name"),
            ("Vendor Code", "VEN-XXXX"), ("NIT / Tender Ref.", "TENDER/2026-XXX"), ("Delivery Date", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Item Details", level=1)
        self._add_table(doc, ["S.No.", "Description", "HSN", "Qty", "Unit", "Rate (Rs)", "Amount (Rs)"], [
            ("1", "Impeller - SS 316", "8413", "2", "No.", "88,000", "1,76,000"),
            ("2", "Mechanical Seal Set", "8484", "4", "Set", "62,000", "2,48,000"),
            ("3", "Bearing Assembly SKF 6312", "8482", "6", "No.", "16,000", "96,000"),
            ("4", "Coupling Element", "8483", "2", "No.", "18,000", "36,000"),
        ])
        doc.add_heading("2. Commercial Terms", level=1)
        self._add_table(doc, ["Term", "Value"], [
            ("Total Amount (Ex-GST)", "Rs 5,56,000"), ("GST @ 18%", "Rs 1,00,080"),
            ("Grand Total", "Rs 6,56,080"), ("Freight", "Included (FOC at site)"),
            ("Payment Terms", "100% within 30 days of delivery + testing"), ("Validity", "90 days from PO date"),
        ])
        doc.add_heading("3. Delivery Instructions", level=1)
        doc.add_paragraph("Delivery to: Stores, Unit 1 Gate. Packaging: Export worthy / Wooden crate. Documents: Invoice, Packing List, MTC, Warranty Certificate.")
        doc.add_heading("4. Authorized Signatories", level=1)
        self._add_sign_table(doc, ["Purchase Officer", "Maintenance Head (Indenter)"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_contract_agreement_docx(self, filename="Contract_Agreement.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Contract Agreement", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("(General Conditions of Contract + Special Conditions of Contract)")
        self._add_meta_table(doc, [
            ("Contract No.", "CONTRACT/2026-XXX"), ("Work Order Ref.", "WO/2026-XXXX"),
            ("Contractor", "M/s. Contractor Name"), ("Contract Type", "Lump Sum / Rate Contract / AMC"),
            ("Contract Value", "Rs XX,XX,XXX"), ("Contract Period", "DD/MM/YYYY to DD/MM/YYYY"),
        ])
        doc.add_heading("PART A - General Conditions of Contract (GCC)", level=1)
        doc.add_heading("1. Definitions", level=2)
        doc.add_paragraph("Employer means the Purchaser / Owner company. Contractor means the firm awarded the contract.")
        doc.add_heading("2. Scope of Work", level=2)
        doc.add_paragraph("The Contractor shall execute the work as described in Scope of Work (Annexure-A).")
        doc.add_heading("3. Contract Price & Payment", level=2)
        doc.add_paragraph("3.1 Fixed price, inclusive of taxes (except GST). 3.2 Payment per milestone. 3.3 Retention: 5%.")
        doc.add_heading("4. Performance Security", level=2)
        doc.add_paragraph("BG for 5% of contract value within 30 days of LOA. Valid for 60 days beyond DLP.")
        doc.add_heading("5. Defect Liability Period", level=2)
        doc.add_paragraph("DLP: 12 months from commissioning or 18 months from delivery, whichever is earlier.")
        doc.add_heading("6. Liquidated Damages", level=2)
        doc.add_paragraph("LD @ 0.5% per week, max 10% of contract value.")
        doc.add_heading("PART B - Special Conditions of Contract (SCC)", level=1)
        doc.add_heading("1. Site Terms", level=2)
        doc.add_paragraph("Contractor shall comply with plant safety rules. Work hours: 06:00-18:00.")
        doc.add_heading("2. Testing & Commissioning", level=2)
        doc.add_paragraph("FAT at vendor premises. SAT after installation. 72-hour performance test.")
        doc.add_heading("3. Annexures", level=2)
        doc.add_paragraph("Annexure-A: Scope of Work. Annexure-B: BOQ. Annexure-C: Delivery Schedule. Annexure-D: Drawings. Annexure-E: BG Format.")
        doc.add_heading("4. Signatories", level=1)
        self._add_sign_table(doc, ["Employer - Authorized", "Contractor - Authorized", "Witness 1", "Witness 2"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_bank_guarantee_docx(self, filename="Bank_Guarantee.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Bank Guarantee (BG)", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("BG No.", "BG/BANK/2026-XXXX"), ("BG Date", "DD/MM/YYYY"),
            ("Contract / PO Ref.", "CONTRACT/2026-XXX"), ("BG Type", "Performance BG / Advance BG"),
            ("BG Value", "Rs XX,XX,XXX"),
        ])
        doc.add_heading("1. BG Details", level=1)
        self._add_table(doc, ["Parameter", "Value"], [
            ("Beneficiary", "M/s. Employer / Owner Company"), ("Applicant", "M/s. Contractor / Supplier"),
            ("Issuing Bank", "Bank Name, Branch"), ("BG Amount", "Rs XX,XX,XXX (Rupees XX Lakhs Only)"),
            ("BG Percentage", "5% of Contract Value"), ("Valid Until", "DD/MM/YYYY + 60 days"),
            ("Claim Period", "30 days from expiry"),
        ])
        doc.add_heading("2. BG Text", level=1)
        doc.add_paragraph("We, [Bank Name], at the request of [Applicant], do hereby irrevocably undertake to pay to [Beneficiary] an amount not exceeding [BG Amount] upon written demand.")
        doc.add_paragraph("2.1 Issued in connection with Contract No. [CONTRACT/2026-XXX]. 2.2 Valid until [Expiry Date]. 2.3 Subject to URDG 758. 2.4 Governed by laws of India.")
        doc.add_heading("3. Signatories", level=1)
        self._add_sign_table(doc, ["Bank Authorized Signatory", "Applicant - Authorized", "Beneficiary - Received By"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_delivery_challan_docx(self, filename="Delivery_Challan_GRN.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Delivery Challan / Goods Receipt Note (GRN)", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        self._add_meta_table(doc, [
            ("GRN No.", "GRN/2026-XXXX"), ("Date", "DD/MM/YYYY"), ("PO Reference", "PO/2026-XXXX"),
            ("Supplier", "M/s. Supplier Name"), ("Delivery Note", "DN/XXXX"), ("Vehicle / LR No.", "Vehicle No."),
        ])
        doc.add_heading("1. Goods Received", level=1)
        self._add_table(doc, ["S.No.", "Description", "PO Qty", "Received", "Accepted", "Rejected", "Remarks"], [
            ("1", "Impeller - SS 316", "2", "2", "2", "0", "As per spec"),
            ("2", "Mechanical Seal Set", "4", "4", "4", "0", "As per spec"),
            ("3", "Bearing Assembly", "6", "6", "5", "1", "Visible mark"),
            ("4", "Coupling Element", "2", "2", "2", "0", "As per spec"),
        ])
        doc.add_heading("2. Inspection Findings", level=1)
        doc.add_paragraph("Visual: Packaging intact. Dimensional: Within tolerance. Certificates: MTC, Warranty received.")
        doc.add_heading("3. GRN Summary", level=1)
        self._add_table(doc, ["Item", "Value"], [
            ("Total Items", "4"), ("Qty Received", "14"), ("Qty Accepted", "13"),
            ("Qty Rejected", "1"), ("Invoice Verified", "Yes - INV/XXXX"),
        ])
        doc.add_heading("4. Signatories", level=1)
        self._add_sign_table(doc, ["Store Keeper", "Quality Inspector", "Maintenance Indenter"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_gem_procurement_review_docx(self, filename="Gem_Procurement_Review.docx"):
        doc = self._init_doc()
        title = doc.add_heading("GeM Procurement Review", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("Government e-Marketplace (GeM) - Procurement Analysis & Review")
        self._add_meta_table(doc, [
            ("Review No.", "GEM/REVIEW/2026-XX"), ("Period", "Jan 2026 - Mar 2026 (Q4 FY 2025-26)"),
            ("Department", "Maintenance & Materials"), ("Prepared By", "Procurement Officer"), ("Date", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. GeM Order Summary", level=1)
        self._add_table(doc, ["S.No.", "Category", "Orders", "Value (Rs)"], [
            ("1", "Office Supplies", "45", "2,25,000"),
            ("2", "IT Equipment", "12", "8,50,000"),
            ("3", "Safety Equipment (PPE)", "8", "3,40,000"),
            ("4", "Tools & Hand Tools", "15", "1,80,000"),
            ("5", "Lubricants & Chemicals", "6", "4,20,000"),
            ("Total", "-", "86", "20,15,000"),
        ])
        doc.add_heading("2. Savings Analysis", level=1)
        doc.add_paragraph("GeM average savings vs market price: 24-27% across categories.")
        doc.add_heading("3. Recommendations", level=1)
        doc.add_paragraph("1. Restrict to 4.5+ star sellers. 2. Use bid/auction for specialized items. 3. Automated GRN matching. 4. Monthly order consolidation.")
        doc.add_heading("4. Compliance", level=1)
        doc.add_paragraph("All purchases comply with GFR 2017. MII preference applied. GeM BID/Auction for orders > Rs 50 lakhs.")
        doc.add_heading("5. Signatories", level=1)
        self._add_sign_table(doc, ["Procurement Officer", "Materials Manager"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_pqc_document_docx(self, filename="PQC_Document.docx"):
        doc = self._init_doc()
        title = doc.add_heading("Pre-Qualification Criterion (PQC) Document", level=0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph("For Empanelment of Vendors / Contractors")
        self._add_meta_table(doc, [
            ("PQC No.", "PQC/2026-XXX"), ("Category", "Mechanical Spares / Civil Contractor"),
            ("Validity", "2 Years from empanelment"), ("Issuing Authority", "Procurement & Contracts"),
            ("Closing Date", "DD/MM/YYYY"),
        ])
        doc.add_heading("1. Eligibility Criteria", level=1)
        self._add_table(doc, ["S.No.", "Criterion", "Requirement"], [
            ("1", "Company Registration", "Companies Act / Partnership Act"),
            ("2", "Minimum Turnover", "Rs XX Crores per annum"),
            ("3", "Years in Business", "Minimum 5 years"),
            ("4", "ISO Certification", "ISO 9001:2015"),
            ("5", "GST Registration", "Valid GSTIN"),
        ])
        doc.add_heading("2. Technical Qualification", level=1)
        self._add_table(doc, ["S.No.", "Criterion", "Points"], [
            ("1", "Similar work experience", "30"), ("2", "Equipment owned", "15"),
            ("3", "Technical manpower", "15"), ("4", "Quality management", "10"),
            ("5", "Safety record", "10"), ("6", "References", "10"),
        ])
        doc.add_heading("3. Minimum Qualifying Marks", level=1)
        doc.add_paragraph("Technical: 50/100. Commercial: 20/40. Combined: 65/140.")
        doc.add_heading("4. Evaluation Committee", level=1)
        self._add_sign_table(doc, ["Chairperson (Procurement)", "Member (Maintenance)", "Member (Finance)", "Member (Quality)"])
        path = os.path.join(OUTPUT_DIR, filename)
        doc.save(path)
        return path

    def generate_materials_review_pptx(self, filename="Materials_Review_Presentation.pptx"):
        prs = Presentation()
        prs.slide_width = PptxInches(13.333)
        prs.slide_height = PptxInches(7.5)
        slide_layout = prs.slide_layouts[5]
        slide1 = prs.slides.add_slide(slide_layout)
        title_shape = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(2))
        tf = title_shape.text_frame
        tf.text = "Materials Management Review"
        tf.paragraphs[0].font.size = PptxInches(1.2)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = 1
        sub = slide1.shapes.add_textbox(PptxInches(0.5), PptxInches(4), PptxInches(12), PptxInches(1))
        sub.text_frame.text = "Quarterly Review - Q4 FY 2025-26\nProcurement & Materials"
        sub.text_frame.paragraphs[0].font.size = PptxInches(0.6)
        sub.text_frame.paragraphs[0].alignment = 1
        slide2 = prs.slides.add_slide(slide_layout)
        tb = slide2.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
        tb.text_frame.text = "Procurement KPIs"
        tb.text_frame.paragraphs[0].font.size = PptxInches(0.8)
        tb.text_frame.paragraphs[0].font.bold = True
        kpi = slide2.shapes.add_table(6, 4, PptxInches(0.5), PptxInches(1.5), PptxInches(12), PptxInches(4.5)).table
        for i, h in enumerate(['KPI', 'Target', 'Actual', 'Status']):
            kpi.cell(0, i).text = h
        for r, d in enumerate([
            ('PO Cycle Time', '< 10 days', '8.5 days', 'On Track'), ('Vendor Payment', '< 30 days', '28 days', 'On Track'),
            ('Contract Compliance', '100%', '97%', '3 deviations'), ('GeM Order %', '> 80%', '85%', 'On Track'),
            ('Cost Savings', '> 5%', '6.2%', 'Exceeding'),
        ]):
            for c, v in enumerate(d):
                kpi.cell(r+1, c).text = v
        slide3 = prs.slides.add_slide(slide_layout)
        tb3 = slide3.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12), PptxInches(1))
        tb3.text_frame.text = "Action Items"
        tb3.text_frame.paragraphs[0].font.size = PptxInches(0.8)
        tb3.text_frame.paragraphs[0].font.bold = True
        ab = slide3.shapes.add_textbox(PptxInches(1), PptxInches(1.5), PptxInches(11), PptxInches(5))
        af = ab.text_frame
        for i, a in enumerate([
            "1. Follow up on 5 pending MRO deliveries",
            "2. Complete GeM seller rating review",
            "3. Update PQC database for 12 vendors",
            "4. Finalize annual rate contract for safety equipment",
            "5. Review slow-moving inventory",
            "6. Escalate delayed civil materials PO",
            "7. Implement vendor performance dashboard",
        ]):
            if i == 0:
                af.text = a
            else:
                p = af.add_paragraph()
                p.text = a
                p.space_after = PptxInches(0.2)
                p.font.size = PptxInches(0.45)
        path = os.path.join(OUTPUT_DIR, filename)
        prs.save(path)
        return path
